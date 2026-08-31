"""Audit generated PPT against ppt_content.json for completeness and formatting."""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from app.services.ppt_shape_utils import (
    get_highlights_shape,
    get_key_activities_shape,
    has_combined_hl_ka_table,
    is_contd_title,
    is_delivery_slide_title,
    paragraph_text,
    service_suffix_from_title,
    slide_title_text,
)


def _norm(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip().lower())


def hl_story_lines(slide) -> list[str]:
    try:
        hl = get_highlights_shape(slide)
        cell = hl.table.cell(2, 0)
        lines = []
        for p in cell.text_frame._txBody.findall(qn("a:p")):
            t = paragraph_text(p).strip()
            if not t:
                continue
            if t.startswith("Stories ") or t.startswith("Sprint") or "current week" in t.lower():
                continue
            if t.startswith("-") or t.startswith("•") or len(t) > 20:
                lines.append(t.lstrip("-• ").strip())
        return lines
    except Exception as exc:
        return [f"__ERROR__:{exc}"]


def ka_layout(slide) -> dict:
    ka = get_key_activities_shape(slide)
    hl = None
    try:
        hl = get_highlights_shape(slide)
    except ValueError:
        pass

    if ka and hl:
        ka_bottom = ka.top + ka.height
        hl_bottom = hl.top + hl.height
        gap_emu = ka.top - hl_bottom
        return {
            "mode": "standalone",
            "ka_top_in": round(ka.top / 914400, 2),
            "ka_left_in": round(ka.left / 914400, 2),
            "ka_h_in": round(ka.height / 914400, 2),
            "hl_bottom_in": round(hl_bottom / 914400, 2),
            "gap_in": round(gap_emu / 914400, 2),
            "ka_body_empty": not ka.table.cell(1, 0).text.strip(),
            "overlap": ka.top < hl_bottom,
        }
    if has_combined_hl_ka_table(slide) and hl:
        return {"mode": "embedded", "combined": True, "ka_body_empty": True}
    return {"mode": "none", "ka_body_empty": True}


def expected_stories(slide_data: dict) -> list[str]:
    stories = []
    for sec in slide_data["sections"]:
        for bucket in ("completed", "released", "inprogress"):
            stories.extend(sec.get(bucket) or [])
    return stories


def audit(ppt_path: Path, json_path: Path) -> int:
    content = json.loads(json_path.read_text(encoding="utf-8"))
    prs = Presentation(str(ppt_path))

    issues: list[str] = []
    passes: list[str] = []

    # Group deck slides by project
    deck: dict[str, dict] = {}
    for i, slide in enumerate(prs.slides):
        title = slide_title_text(slide)
        if not is_delivery_slide_title(title):
            continue
        proj = service_suffix_from_title(title)
        deck.setdefault(proj, {"main": None, "contd": []})
        entry = {"idx": i + 1, "title": title, "stories": hl_story_lines(slide), "ka": ka_layout(slide)}
        if is_contd_title(title):
            deck[proj]["contd"].append(entry)
        else:
            deck[proj]["main"] = entry

    print(f"Auditing: {ppt_path.name}")
    print(f"Against:  {json_path.name} ({len(content['slides'])} projects, 44 stories expected)\n")

    for slide_data in content["slides"]:
        title = slide_data["title"]
        expected = expected_stories(slide_data)
        proj_deck = deck.get(title)

        print(f"{'='*60}")
        print(f"PROJECT: {title}")
        print(f"  Expected: {len(slide_data['sections'])} sprints, {len(expected)} stories")

        if not proj_deck or not proj_deck["main"]:
            issues.append(f"{title}: main slide missing from deck")
            print("  FAIL: main slide not found")
            continue

        main = proj_deck["main"]
        contd = proj_deck["contd"]
        all_deck_stories = list(main["stories"])
        for c in contd:
            all_deck_stories.extend(c["stories"])

        missing = []
        for story in expected:
            n = _norm(story)
            if not any(n in _norm(d) or _norm(d) in n for d in all_deck_stories):
                missing.append(story)

        extra = []
        for story in all_deck_stories:
            if story.startswith("__ERROR__"):
                continue
            n = _norm(story)
            if not any(n in _norm(e) or _norm(e) in n for e in expected):
                extra.append(story)

        if missing:
            issues.append(f"{title}: {len(missing)} story(s) missing from PPT")
            print(f"  FAIL: {len(missing)} missing stories")
            for s in missing[:5]:
                print(f"    - {s[:80]}")
        else:
            passes.append(f"{title}: all {len(expected)} stories present")
            print(f"  OK: all {len(expected)} stories found across {1 + len(contd)} slide(s)")

        if extra:
            issues.append(f"{title}: {len(extra)} unexpected story line(s) in PPT")
            print(f"  WARN: {len(extra)} extra lines (template leftovers?)")

        # KA checks
        ka_main = main["ka"]
        print(f"  KA main slide: mode={ka_main.get('mode')}, empty={ka_main.get('ka_body_empty', True)}")
        if ka_main.get("mode") == "standalone":
            if ka_main.get("overlap"):
                issues.append(f"{title}: KA overlaps Highlights table")
                print("  FAIL: KA overlaps HL")
            elif ka_main.get("gap_in", 0) < 0:
                issues.append(f"{title}: KA positioned above HL bottom")
                print("  FAIL: KA gap negative")
            else:
                print(f"  OK: KA below HL (gap={ka_main.get('gap_in')} in)")
        for c in contd:
            ck = c["ka"]
            print(f"  KA contd slide {c['idx']}: mode={ck.get('mode')}, empty={ck.get('ka_body_empty', True)}")
            if ck.get("mode") == "standalone" and not ck.get("ka_body_empty", True):
                issues.append(f"{title} contd slide {c['idx']}: KA not empty")

        if contd:
            print(f"  Continuation: {len(contd)} slide(s) at positions {[c['idx'] for c in contd]}")
            # verify contd immediately follow main chain
            chain = [main["idx"]] + [c["idx"] for c in sorted(contd, key=lambda x: x["idx"])]
            for j in range(len(chain) - 1):
                if chain[j + 1] != chain[j] + 1:
                    issues.append(f"{title}: contd slide {chain[j+1]} not immediately after {chain[j]}")
                    print(f"  WARN: slide gap between {chain[j]} and {chain[j+1]}")

    print(f"\n{'='*60}")
    print(f"SUMMARY: {len(passes)} passed, {len(issues)} issue(s)")
    for p in passes:
        print(f"  PASS: {p}")
    for issue in issues:
        print(f"  ISSUE: {issue}")
    return len(issues)


if __name__ == "__main__":
    ppt = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "output/HEB_Delivery_Status_v2_fixed.pptx"
    js = Path(sys.argv[2]) if len(sys.argv) > 2 else ROOT / "output/ppt_content.json"
    raise SystemExit(audit(ppt, js))
