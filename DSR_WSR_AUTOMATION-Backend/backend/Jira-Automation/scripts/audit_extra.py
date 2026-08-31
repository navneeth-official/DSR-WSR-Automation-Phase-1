"""Supplemental audit: extra lines, bullets, KA placement."""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from app.services.ppt_shape_utils import (
    get_highlights_shape,
    get_key_activities_shape,
    has_combined_hl_ka_table,
    is_delivery_slide_title,
    paragraph_text,
    slide_title_text,
)


def main() -> None:
    prs = Presentation(str(ROOT / "output/HEB_Delivery_Status_v2_fixed.pptx"))
    content = json.loads((ROOT / "output/ppt_content.json").read_text(encoding="utf-8"))
    tpl = Presentation(str(ROOT / "templates/G10X H-E-B WSR Haskell Location Pharmacy GSS PAM 01 Aug 2025.pptx"))

    def hl_lines(slide):
        hl = get_highlights_shape(slide)
        return [
            paragraph_text(p).strip()
            for p in hl.table.cell(2, 0).text_frame._txBody.findall(qn("a:p"))
            if paragraph_text(p).strip()
        ]

    for proj in ("Supplier Core Service", "Location Core Service"):
        expected: set[str] = set()
        for s in content["slides"]:
            if s["title"] == proj:
                for sec in s["sections"]:
                    for b in ("completed", "released", "inprogress"):
                        expected.update(sec.get(b) or [])
        print(f"=== {proj} line check ===")
        for i, slide in enumerate(prs.slides):
            t = slide_title_text(slide)
            if proj.lower() not in t.lower() or not is_delivery_slide_title(t):
                continue
            print(f" slide {i + 1}: {t[:55]}")
            for line in hl_lines(slide):
                if line.startswith("Sprint") or line.startswith("Stories") or "current week" in line.lower():
                    continue
                ok = any(line == e or line in e or e in line for e in expected)
                print(("  OK  " if ok else "  EXTRA"), line[:90])

    missing_bullet = has_bullet = 0
    for slide in prs.slides:
        t = slide_title_text(slide)
        if not is_delivery_slide_title(t):
            continue
        try:
            hl = get_highlights_shape(slide)
            for p in hl.table.cell(2, 0).text_frame._txBody.findall(qn("a:p")):
                txt = paragraph_text(p).strip()
                if not txt or txt.startswith("Sprint") or txt.startswith("Stories") or "current week" in txt.lower():
                    continue
                p_pr = p.find(qn("a:pPr"))
                bu = p_pr.find(qn("a:buChar")) if p_pr is not None else None
                if bu is None:
                    missing_bullet += 1
                else:
                    has_bullet += 1
        except ValueError:
            pass
    print("\n=== BULLET FORMATTING ===")
    print(f"Story lines WITH bullet XML: {has_bullet}")
    print(f"Story lines WITHOUT bullet XML: {missing_bullet}")

    print("\n=== KA PLACEMENT template vs output (inches) ===")
    pairs = [
        ("Pricing", 7, 5),
        ("GSS", 12, 15),
        ("Wentworth", 8, 7),
        ("Location", 9, 10),
        ("Cost", 2, 2),
        ("Pharmacy", 10, 14),
    ]
    for name, tpl_i, out_i in pairs:
        for label, deck, idx in (("TPL", tpl, tpl_i), ("OUT", prs, out_i)):
            s = deck.slides[idx]
            ka = get_key_activities_shape(s)
            hl = get_highlights_shape(s)
            if ka:
                gap = (ka.top - hl.top - hl.height) / 914400
                print(
                    f"{name:10} {label} slide{idx+1:2d}: "
                    f"KA top={ka.top/914400:.2f} left={ka.left/914400:.2f} "
                    f"h={ka.height/914400:.2f} gap={gap:.2f}"
                )
            else:
                print(f"{name:10} {label} slide{idx+1:2d}: no standalone KA, combined={has_combined_hl_ka_table(s)}")

    print("\n=== CONTD SLIDE LAYOUT ===")
    for i, slide in enumerate(prs.slides):
        t = slide_title_text(slide)
        if "(contd" not in t.lower():
            continue
        ka = get_key_activities_shape(slide)
        comb = has_combined_hl_ka_table(slide)
        has_hl = True
        try:
            get_highlights_shape(slide)
        except ValueError:
            has_hl = False
        ka_desc = f"standalone KA (empty={not ka.table.cell(1,0).text.strip()})" if ka else (
            "embedded KA rows" if comb else "no KA"
        )
        print(f" slide {i+1}: {t[:48]} | HL={has_hl} | {ka_desc}")


if __name__ == "__main__":
    main()
