"""
Build HEB delivery-status deck from G10X template.

G10X sustainment format (Supplier Core Services reference — slide 4):
- Paragraph styles from Cost Core template: lvl 0 sprint/current-week (round bullet),
  lvl 7 category headers (solid right-pointing arrowhead: Wingdings Ø only), lvl 1 story items (dash bullet).
- Story order per sprint: completed, released, in-progress.
- No blank line between category header and first story; one blank line between sprints.
- Fill main slide to Supplier capacity (~20 paragraph slots) before HL (Contd...).
- Key Activities only after Highlights finish; ~2 body lines between HL and KA tabs when on same slide.
- ka_contd_only: dense HL on main, KA on KA-only (Contd...) when content exceeds KA zone.
"""

import argparse
import copy
import io
import json
import os
import re
import shutil
import sys
from datetime import date
from pathlib import Path
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

_REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_REPO_ROOT))

from app.services.sprint_display import sprint_dates_from_section
from app.services.ppt_cover_date import format_wsr_cover_date, sync_cover_slide_wsr_date
from app.services.ppt_layout_metrics import hl_ka_tab_gap_emu, rendered_text_bottom_emu
from app.services.ppt_hl_bullets import set_story_line_text
from app.services.ppt_logo_sync import sync_heb_logo_from_main

G10X = str(_REPO_ROOT / "templates" / "G10X H-E-B WSR Sustainment 05 June 2026 .pptx")
OUTPUT = str(_REPO_ROOT / "output" / "HEB_Delivery_Status.pptx")
PHARMACY_CONTD_INDEX = 10
HIGHLIGHTS_CONTD_TEMPLATE_INDEX = 4  # G10X Supplier contd — highlights-only layout
HIGHLIGHTS_KA_CONTD_TEMPLATE_INDEX = 5  # G10X Supplier contd — highlights + key activities
SAFE_KA_BOTTOM_IN = 4.85
MIN_KA_ROW_HEIGHT = 228600  # ~0.25 in — never collapse visible KA rows to 0
MIN_HL_KA_GAP = 45720  # ~0.05 in — fallback when no G10X KA reference exists
MIN_TEXT_KA_CLEARANCE = 137160  # ~0.15 in — minimum text-to-KA clearance (G10X LoCo min)
HL_KA_FILL_THRESHOLD = 0.85  # at or above this fraction of template capacity, use full HL box
HL_OVERFLOW_SAFETY = 0  # fill to G10X mark before creating (Contd...)
EMPTY_KA_CONTD_MAX_EMU = 320040  # ~0.35 in — header strip on KA-only (Contd...) slides
HL_CONTENT_BOTTOM_PAD = 91440  # ~0.1 in below last text line inside HL table
# G10X story-bucket order within each sprint (Supplier reference slide 4)
STORY_BUCKET_ORDER = ("completed", "released", "inprogress")
CANONICAL_STYLE_SLIDE_INDEX = 2   # Cost Core — sustainment bullet styles (lvl 0/7/1)
CANONICAL_FILL_SLIDE_INDEX = 3    # Supplier — fill-to-mark capacity reference
# G10X category headers ("Stories completed this week", etc.)
# buChar Ø + buFont Wingdings → solid right-pointing arrowhead ONLY (see app/constants/ppt_bullets.py)
CATEGORY_HEADER_LEVEL = 7
CATEGORY_HEADER_BULLET = "\u00d8"
CATEGORY_HEADER_BULLET_FONT = "Wingdings,Sans-Serif"
CATEGORY_HEADER_BULLET_COLOR = "000000"

# Key activities preserved from original BSA (LoCo) template — moved to contd slide when needed
LOCO_KA_ITEMS = [
    "Mainframe Retirement",
    "Coordinate with team on UAT Feedback",
    "Refine any new stories on a need basis",
    "Continue backlog Refinement for Q3 work",
    "Continue ramp-up and documentation on FAM app",
]

# 0-based slide index -> delivery status content
SLIDES = {
    2: {
        "title": "Cost Core Service",
        "sprint_name": "Q3.02 FY26 Orion",
        "sprint_dates": "Jul 6 \u2013 Jul 19",
        "released": [
            "Validate buyer funding calculation",
            "Update invoice reconciliation workflow",
        ],
        "inprogress": [
            "Implement COGS adjustment API",
            "Optimize cost calculation engine",
            "Enhance contract validation logic",
        ],
        "completed": [
            "Add support for break-pack cost",
            "Fix supplier cost synchronization",
            "Improve cost audit logging",
        ],
    },
    3: {
        "title": "Supplier Core Service",
        "sprint_name": "Q3.02 FY26 Orion",
        "sprint_dates": "Jul 6 \u2013 Jul 19",
        "released": [
            "Validate supplier onboarding request payload",
            "Enable bulk supplier import via CSV",
            "Implement supplier contact validation rules",
        ],
        "inprogress": [
            "Implement supplier profile update endpoint",
            "Add audit logging for supplier modifications",
            "Add pagination and filtering to supplier listing API",
            "Create supplier synchronization job with ERP",
        ],
        "completed": [
            "Expose Supplier Catalog API for downstream systems",
            "Optimize supplier search using indexed queries",
            "Resolve duplicate supplier ID generation issue",
        ],
    },
    4: {
        "title": "Pricing Core Service",
        "sprint_name": "Q3.02 FY26 Polaris",
        "sprint_dates": "Jul 6 \u2013 Jul 19",
        "released": [
            "Validate promotional pricing rules",
            "Add unit price override support",
            "Implement price history endpoint",
        ],
        "inprogress": [
            "Enhance dynamic pricing engine",
            "Optimize price lookup performance",
            "Build bulk pricing update service",
            "Add pricing audit events",
        ],
        "completed": [
            "Resolve incorrect discount calculation",
            "Implement pricing cache refresh",
        ],
    },
    5: {
        "title": "Wentworth",
        "sprint_name": "Q3.01 FY26 Atlas",
        "sprint_dates": "Jun 29 \u2013 Jul 12",
        "released": [
            "Improve Wentworth reporting dashboard",
            "Implement user access validation",
        ],
        "inprogress": [
            "Develop notification service",
            "Optimize dashboard loading performance",
            "Enhance report export functionality",
        ],
        "completed": [
            "Fix authentication issue",
            "Update reporting filters",
        ],
    },
    6: {
        "title": "Location Core Service",
        "sprint_name": "Q2.14 FY26 Fornax",
        "sprint_dates": "Jun 1 \u2013 Jun 14",
        "released": [
            "Status Column Displays X for Offsite Warehouse",
        ],
        "inprogress": [
            "Changes to Offsite Warehouse \u2013 Create API",
            "Warehouse \u2013 Edit Offsite Warehouses",
        ],
        "completed": [
            "Offsite Warehouse List sorted by Offsite Name",
        ],
    },
    7: {
        "title": "Pharmacy and Wellness",
        "sprint_name": "Q3.03 FY26 Orion",
        "sprint_dates": "Jul 20 \u2013 Aug 2",
        "released": [
            "Validate prescription upload workflow",
            "Improve wellness appointment scheduling",
        ],
        "inprogress": [
            "Implement pharmacy inventory alerts",
            "Develop prescription reminder service",
            "Optimize patient profile retrieval",
        ],
        "completed": [
            "Add medication search filters",
            "Resolve duplicate prescription issue",
            "Improve pharmacy dashboard performance",
        ],
    },
    8: {
        "title": "Global Sourcing Solution",
        "sprint_name": "Q3.02 FY26 Pegasus",
        "sprint_dates": "Jul 6 \u2013 Jul 19",
        "released": [
            "Validate supplier quotation workflow",
            "Add sourcing request approval process",
            "Implement purchase request notifications",
        ],
        "inprogress": [
            "Develop supplier comparison module",
            "Optimize RFQ processing",
            "Enhance sourcing analytics dashboard",
            "Implement contract lifecycle integration",
        ],
        "completed": [
            "Improve supplier search",
            "Fix sourcing workflow validation",
        ],
    },
    9: {
        "title": "LoCo",
        "sprint_name": "Q3.03 FY26 Vega",
        "sprint_dates": "Jul 20 \u2013 Aug 2",
        "released": [
            "Validate local inventory synchronization",
            "Add regional warehouse configuration",
        ],
        "inprogress": [
            "Implement location mapping service",
            "Optimize inventory allocation logic",
            "Enhance local order routing",
            "Develop warehouse capacity monitoring",
        ],
        "completed": [
            "Improve location search performance",
            "Resolve inventory synchronization issue",
        ],
    },
}

# 0-based indices of obsolete continuation slides to remove (highest first)
DELETE_SLIDES = []  # populated dynamically if Pharmacy (Contd...) exists
LAYOUT_HINTS: dict = {}


def load_layout_hints(path: str | Path | None) -> dict:
    """Load per-service layout hints JSON for rebuild/repair."""
    global LAYOUT_HINTS
    if not path:
        LAYOUT_HINTS = {}
        return LAYOUT_HINTS
    p = Path(path)
    if not p.is_file():
        LAYOUT_HINTS = {}
        return LAYOUT_HINTS
    with open(p, encoding="utf-8-sig") as f:
        LAYOUT_HINTS = json.load(f)
    return LAYOUT_HINTS


def build_content(raw, totals_raw=None):
    totals = totals_raw or raw
    if raw.get("sections"):
        return {
            "title": raw["title"],
            "sections": [_section_display_content(s) for s in raw["sections"]],
        }

    released = raw["released"]
    inprogress = raw["inprogress"]
    completed = raw["completed"]
    total = len(totals["released"]) + len(totals["inprogress"]) + len(totals["completed"])
    sprint_status = totals.get("sprint_status", "In-progress")
    return {
        "title": raw["title"],
        "sprint_bold": f"Sprint \u2013 {raw['sprint_name']}, {sprint_status} ",
        "sprint_light": (
            f"({raw['sprint_dates']}) Stories (Total \u2013 {total}, "
            f"Done \u2013 {len(totals['completed'])}, In-review \u2013 {len(totals['released'])}, "
            f"In-progress \u2013 {len(totals['inprogress'])})"
        ),
        "released_count": str(len(released)),
        "released_items": released,
        "inprogress_count": str(len(inprogress)),
        "inprogress_items": inprogress,
        "completed_count": str(len(completed)),
        "completed_items": completed,
    }


def _section_display_content(section):
    """Build populate_highlights_cell payload for one sprint section."""
    released = section["released"]
    inprogress = section["inprogress"]
    completed = section["completed"]
    total = len(released) + len(inprogress) + len(completed)
    sprint_status = section.get("sprint_status", "In-progress")
    return {
        "sprint_bold": f"Sprint \u2013 {section['sprint_name']}, {sprint_status} ",
        "sprint_light": (
            f"({section['sprint_dates']}) Stories (Total \u2013 {total}, "
            f"Done \u2013 {len(completed)}, In-review \u2013 {len(released)}, "
            f"In-progress \u2013 {len(inprogress)})"
        ),
        "released_count": str(len(released)),
        "released_items": released,
        "inprogress_count": str(len(inprogress)),
        "inprogress_items": inprogress,
        "completed_count": str(len(completed)),
        "completed_items": completed,
    }


def _flatten_story_raw(raw):
    """Flatten multi-sprint sections into one raw dict for split planning."""
    if not raw.get("sections"):
        return raw
    released, inprogress, completed = [], [], []
    for section in raw["sections"]:
        released.extend(section["released"])
        inprogress.extend(section["inprogress"])
        completed.extend(section["completed"])
    first = raw["sections"][0]
    return {
        "title": raw["title"],
        "sprint_name": first["sprint_name"],
        "sprint_dates": first["sprint_dates"],
        "sprint_status": first.get("sprint_status", "In-progress"),
        "released": released,
        "inprogress": inprogress,
        "completed": completed,
        "key_activities": raw.get("key_activities", []),
        "sections": raw["sections"],
    }


def _sprint_dates_from_payload(payload: dict) -> str:
    """Prefer ISO sprint bounds from JSON (DB source of truth)."""
    return sprint_dates_from_section(payload)


def _normalize_sprint_section(section: dict) -> dict:
    """Ensure sprint_dates always matches sprint_start_date / sprint_end_date."""
    normalized = dict(section)
    normalized["sprint_dates"] = _sprint_dates_from_payload(section)
    return normalized


def load_slides_from_json(path):
    """Load slide chunks from ppt_content.json into the SLIDES dict shape."""
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    chunks = data.get("slides", data)
    slides = {}
    for i, chunk in enumerate(chunks, start=2):
        entry = {
            "title": chunk["title"],
            "key_activities": chunk.get("key_activities", []),
        }
        if chunk.get("sections"):
            entry["sections"] = [
                _normalize_sprint_section(s) for s in chunk["sections"]
            ]
        else:
            entry.update({
                "sprint_name": chunk["sprint_name"],
                "sprint_start_date": chunk.get("sprint_start_date"),
                "sprint_end_date": chunk.get("sprint_end_date"),
                "sprint_dates": _sprint_dates_from_payload(chunk),
                "sprint_status": chunk.get("sprint_status", "In-progress"),
                "released": chunk.get("released", []),
                "inprogress": chunk.get("inprogress", []),
                "completed": chunk.get("completed", []),
            })
        slides[i] = entry
    return slides


def resolve_ka_items(raw):
    """Return KA bullet list; empty DB content → one placeholder row for manual entry."""
    items = raw.get("key_activities") or []
    if items:
        return items
    return [""]


def count_filled_ka_items(items_or_shape):
    """Non-empty KA bullets only — placeholders must not inflate layout height."""
    if isinstance(items_or_shape, list):
        return sum(1 for item in items_or_shape if str(item).strip())
    cell = items_or_shape.table.cell(1, 0)
    return sum(1 for p in cell.text_frame.paragraphs if p.text.strip())


def ka_layout_item_count(raw):
    """Item count used for KA height / fit geometry (0 when only placeholders)."""
    return count_filled_ka_items(resolve_ka_items(raw))


def finalize_key_activities(prs, service_title, main_idx, raw, g10x_layout, layout_mode):
    """Apply Key Activities content (or placeholders) on the correct slide."""
    items = resolve_ka_items(raw)
    g10x_ref = Presentation(G10X)
    profile = get_ka_layout_profile(g10x_ref, g10x_layout)
    contd_ka_ref = get_contd_ka_reference(g10x_ref)

    def apply_to_slide(slide, use_contd_position=False):
        ka = _ensure_ka_shape_on_slide(slide, g10x_ref, g10x_layout)
        if not ka:
            return
        set_ka_items(ka, items)
        pos_ref = contd_ka_ref if use_contd_position and contd_ka_ref else None
        fit_key_activities_table(ka, profile, position_ref=pos_ref)
        if use_contd_position and contd_ka_ref:
            ka.top = contd_ka_ref.top
            ka.left = contd_ka_ref.left
            ka.width = contd_ka_ref.width
        if use_contd_position and count_filled_ka_items(items) == 0:
            ka.height = min(ka.height, EMPTY_KA_CONTD_MAX_EMU)
            ka.table.rows[1].height = max(
                ka.height - ka.table.rows[0].height, MIN_KA_ROW_HEIGHT
            )

    if layout_mode in ("hl_ka_contd", "ka_contd_only", "supplier_contd"):
        contd_idx = _find_contd_slide_with_ka(prs, service_title)
        if contd_idx is not None:
            apply_to_slide(
                prs.slides[contd_idx],
                use_contd_position=(layout_mode == "ka_contd_only"),
            )
    elif layout_mode == "ka_on_main":
        apply_to_slide(prs.slides[main_idx])
    else:
        apply_to_slide(prs.slides[main_idx])


def parse_args():
    parser = argparse.ArgumentParser(description="Build HEB delivery-status deck.")
    parser.add_argument(
        "--content",
        default=os.environ.get("PPT_CONTENT_JSON", ""),
        help="Path to ppt_content.json from generate_ppt_content.py",
    )
    parser.add_argument(
        "--output",
        default=OUTPUT,
        help="Output .pptx path",
    )
    parser.add_argument(
        "--layout-hints",
        default="",
        help="Path to layout hints JSON (pack_all_sections_on_main per service)",
    )
    parser.add_argument(
        "--template",
        default=G10X,
        help="Path to the WSR .pptx template used to build the deck",
    )
    parser.add_argument(
        "--wsr-end-date",
        default="",
        help="Cover slide date (YYYY-MM-DD). Default: report_end_date from --content JSON",
    )
    return parser.parse_args()


def set_single_run_text(p_elem, text):
    runs = p_elem.findall(qn("a:r"))
    if not runs:
        return
    runs[0].find(qn("a:t")).text = text
    for extra in runs[1:]:
        p_elem.remove(extra)


def set_two_run_header(p_elem, prefix, count_suffix):
    runs = p_elem.findall(qn("a:r"))
    if len(runs) >= 2:
        runs[0].find(qn("a:t")).text = prefix
        runs[1].find(qn("a:t")).text = f"\u2013 {count_suffix}"
        for extra in runs[2:]:
            p_elem.remove(extra)
    else:
        set_single_run_text(p_elem, f"{prefix}\u2013 {count_suffix}")


def _ensure_category_header_bullet(p_elem):
    """
    G10X rule: category headers use ONLY the solid right-pointing arrowhead bullet.
    Requires buChar Ø (U+00D8) at list level 7 WITH buFont Wingdings — no •, dash, or ▶.
    """
    from pptx.oxml import parse_xml

    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    pPr = p_elem.find(qn("a:pPr"))
    if pPr is None:
        pPr = parse_xml(f'<a:pPr xmlns:a="{ns}"/>')
        p_elem.insert(0, pPr)
    pPr.set("lvl", str(CATEGORY_HEADER_LEVEL))
    for tag in ("a:buNone", "a:buAutoNum", "a:buBlip"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)

    bu_clr = pPr.find(qn("a:buClr"))
    if bu_clr is None:
        pPr.append(
            parse_xml(
                f'<a:buClr xmlns:a="{ns}">'
                f'<a:srgbClr val="{CATEGORY_HEADER_BULLET_COLOR}"/></a:buClr>'
            )
        )

    bu_font = pPr.find(qn("a:buFont"))
    if bu_font is None:
        pPr.append(
            parse_xml(
                f'<a:buFont xmlns:a="{ns}" typeface="{CATEGORY_HEADER_BULLET_FONT}"/>'
            )
        )
    else:
        bu_font.set("typeface", CATEGORY_HEADER_BULLET_FONT)

    bu = pPr.find(qn("a:buChar"))
    if bu is None:
        pPr.append(
            parse_xml(
                f'<a:buChar xmlns:a="{ns}" char="{CATEGORY_HEADER_BULLET}"/>'
            )
        )
    else:
        bu.set("char", CATEGORY_HEADER_BULLET)


def _apply_g10x_paragraph_spacing(p_elem, is_story=False):
    """G10X rule 15: category→story is tight (spcBef=0); no extra line gaps."""
    pPr = p_elem.find(qn("a:pPr"))
    if pPr is None:
        return
    if is_story:
        for tag in ("a:spcBef", "a:spcAft"):
            existing = pPr.find(qn(tag))
            if existing is not None:
                pPr.remove(existing)
        from pptx.oxml import parse_xml
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pPr.insert(0, parse_xml(f'<a:spcBef xmlns:a="{ns}"><a:spcPts val="0"/></a:spcBef>'))


def _make_sprint_gap_para(section_tmpl):
    """G10X rule: exactly one blank paragraph between sprint blocks."""
    gap = copy.deepcopy(section_tmpl["current_week"])
    set_single_run_text(gap, "")
    pPr = gap.find(qn("a:pPr"))
    if pPr is not None:
        for tag in ("a:buChar", "a:buAutoNum", "a:buBlip", "a:buFont"):
            el = pPr.find(qn(tag))
            if el is not None:
                pPr.remove(el)
        if pPr.find(qn("a:buNone")) is None:
            from pptx.oxml import parse_xml
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            pPr.append(parse_xml(f'<a:buNone xmlns:a="{ns}"/>'))
    _apply_g10x_paragraph_spacing(gap)
    return gap


def replace_bullet_block(txBody, first_index, last_index, items):
    all_p = txBody.findall(qn("a:p"))
    template_p = all_p[first_index]
    block = all_p[first_index : last_index + 1]
    anchor = block[-1]

    new_paragraphs = []
    for text in items:
        new_p = copy.deepcopy(template_p)
        set_story_line_text(new_p, text)
        _apply_g10x_paragraph_spacing(new_p, is_story=True)
        new_paragraphs.append(new_p)

    prev = anchor
    for new_p in new_paragraphs:
        prev.addnext(new_p)
        prev = new_p

    for p in block:
        txBody.remove(p)


def normalize_title_text(text):
    """Strip zero-width spaces and normalize Delivery status dash for matching."""
    cleaned = (text or "").replace("\u200b", "").replace("\xa0", " ").strip()
    en_dash = "\u2013"
    return re.sub(
        r"(Delivery status)\s+-\s+",
        rf"\1 {en_dash} ",
        cleaned,
        count=1,
        flags=re.I,
    )


def set_title_suffix(slide, suffix, ref_title_shape):
    clean_suffix = normalize_title_text(suffix)
    for s in slide.shapes:
        if s.shape_id == 2 and s.has_text_frame:
            s.top = ref_title_shape.top
            s.height = ref_title_shape.height
            s.width = ref_title_shape.width
            s.left = ref_title_shape.left

            ref_rPr = ref_title_shape.text_frame.paragraphs[0].runs[0]._r.find(qn("a:rPr"))
            p = s.text_frame.paragraphs[0]
            if not p.runs:
                p.add_run()
            run = p.runs[0]
            old_rPr = run._r.find(qn("a:rPr"))
            if old_rPr is not None:
                run._r.remove(old_rPr)
            if ref_rPr is not None:
                run._r.insert(0, copy.deepcopy(ref_rPr))
            run.text = normalize_title_text(f"Delivery status \u2013 {clean_suffix}")
            for extra in p.runs[1:]:
                p._p.remove(extra._r)
            return


def get_highlights_shape(slide):
    """Main 3x3 highlights table — shape 9 on most slides, shape 7 on Wentworth."""
    for sid in (9, 7):
        sh = next((s for s in slide.shapes if s.shape_id == sid and s.has_table), None)
        if sh and len(sh.table.rows) == 3 and len(sh.table.columns) >= 3:
            return sh
    raise ValueError("No highlights table found on slide")


def get_key_activities_shape(slide):
    for sid in (6, 8, 4):
        sh = next((s for s in slide.shapes if s.shape_id == sid and s.has_table), None)
        if sh and len(sh.table.rows) == 2 and len(sh.table.columns) == 1:
            return sh
    return None


G10X_LAYOUT_BY_TITLE = {
    "Cost Core Service": 2,
    "Supplier Core Services": 3,
    "Supplier Core Service": 3,  # alias — matches ppt_mapping.PPT_SLIDE_TITLES
    "Pricing Core Service": 6,
    "Pricing Core": 6,
    "Wentworth": 7,
    "Location Core Service": 8,
    "Pharmacy and Wellness": 9,
    "Global Sourcing Solution": 11,
    "LoCo": 12,
    "Patronage Travel": 2,
    "PATRV": 2,
}

# Approximate chars per visual line at G10X Highlights content width (12pt Manrope)
_CHARS_PER_VISUAL_LINE = {0: 92, 1: 72, 7: 80}
_DEFAULT_CHARS_PER_LINE = 85
_VISUAL_LINE_HEIGHT_BUFFER = 1.12  # extra factor for wrapped lines


def get_g10x_layout_slide(g10x_prs, title):
    """Pick the G10X slide whose highlights/KA spacing matches this service."""
    idx = G10X_LAYOUT_BY_TITLE.get(title, 2)  # default: Cost Core
    return g10x_prs.slides[idx]


def _ka_fields_from_ref(ref_ka, ref_hl):
    """Build Key Activities layout constants from a reference KA table shape."""
    ka_r0 = ref_ka.table.rows[0].height
    ka_ref_r1 = ref_ka.table.rows[1].height
    ka_ref_pad = ref_ka.height - ka_r0 - ka_ref_r1
    ka_ref_items = max(
        sum(
            1
            for p in ref_ka.table.cell(1, 0).text_frame.paragraphs
            if p.text.strip()
        ),
        1,
    )
    standard_gap = ref_ka.top - ref_hl.top - ref_hl.height
    return {
        "ref_ka": ref_ka,
        "ka_r0": ka_r0,
        "ka_ref_r1": ka_ref_r1,
        "ka_ref_pad": ka_ref_pad,
        "ka_ref_items": ka_ref_items,
        "ref_ka_height": ref_ka.height,
        "ref_ka_top": ref_ka.top,
        "ref_ka_bottom": ref_ka.top + ref_ka.height,
        "ka_header_cell_ref": ref_ka.table.cell(0, 0),
        "ka_content_cell_ref": ref_ka.table.cell(1, 0),
        "standard_gap": standard_gap,
    }


def get_canonical_ka_reference(g10x_prs):
    """Cost Core on-slide KA — canonical shape for services whose G10X slide has no KA."""
    _, ref_ka = get_g10x_highlights_and_ka(g10x_prs.slides[CANONICAL_STYLE_SLIDE_INDEX])
    return ref_ka


def get_ka_layout_profile(g10x_prs, g10x_layout_slide):
    """KA layout profile; uses Cost Core KA when the service G10X slide has none."""
    profile = build_layout_profile(g10x_layout_slide)
    ref_ka = profile.get("ref_ka") or get_canonical_ka_reference(g10x_prs)
    if not ref_ka:
        return profile
    return {**profile, **_ka_fields_from_ref(ref_ka, profile["ref_hl"]), "ref_ka": ref_ka}


def _ensure_ka_shape_on_slide(slide, g10x_prs, g10x_layout):
    """
    Insert a Key Activities table on the slide when missing.

    Clones from the service G10X reference or Cost Core canonical KA so every
    project slide can host the KA tab even when the deck template omits it.
    """
    ka = get_key_activities_shape(slide)
    if ka:
        return ka
    ka_profile = get_ka_layout_profile(g10x_prs, g10x_layout)
    ref_ka = ka_profile.get("ref_ka")
    if not ref_ka:
        return None
    newel = copy.deepcopy(ref_ka.element)
    slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    return get_key_activities_shape(slide)


def ensure_key_activities_tab(slide, g10x_prs, g10x_layout, profile=None):
    """
    Ensure the Key Activities table exists on ``slide`` and seat it below HL.

    Returns the KA shape, or None if the canonical reference is unavailable.
    """
    ka = _ensure_ka_shape_on_slide(slide, g10x_prs, g10x_layout)
    if not ka:
        return None
    if profile is None:
        profile = build_layout_profile(g10x_layout)
    try:
        hl = get_highlights_shape(slide)
    except ValueError:
        return ka
    ka_profile = get_ka_layout_profile(g10x_prs, g10x_layout)
    ka_h = fit_key_activities_table(ka, ka_profile)
    _position_ka_below_hl(hl, ka, profile, ka_h=ka_h)
    return ka


def g10x_main_has_on_slide_ka(g10x_layout_slide):
    """True when the G10X reference main slide keeps Key Activities on the same slide."""
    _, ref_ka = get_g10x_highlights_and_ka(g10x_layout_slide)
    return ref_ka is not None


def get_g10x_highlights_and_ka(g10x_slide):
    ref_hl = next(
        s
        for s in g10x_slide.shapes
        if s.has_table and s.table.cell(0, 0).text.strip() == "Highlights"
    )
    ref_ka = next(
        (
            s
            for s in g10x_slide.shapes
            if s.has_table and "Key activities" in s.table.cell(0, 0).text
        ),
        None,
    )
    return ref_hl, ref_ka


# Footer safe zone — G10X Cost Core KA bottom (~6.29 in) is the tallest single-slide layout.
MAX_KA_BOTTOM_EMU = 5747107  # 6.29 in


def build_layout_profile(g10x_slide):
    """Extract Highlights + Key Activities spacing constants from a G10X reference slide."""
    ref_hl, ref_ka = get_g10x_highlights_and_ka(g10x_slide)
    r0 = ref_hl.table.rows[0].height
    r1 = ref_hl.table.rows[1].height
    ref_r2 = ref_hl.table.rows[2].height
    ref_pad = ref_hl.height - r0 - r1 - ref_r2
    hl_paras = ref_hl.table.cell(2, 0).text_frame._txBody.findall(qn("a:p"))
    # Capacity = every paragraph slot in the G10X template cell (incl. empty placeholders).
    # The original deck fills all slots before using (Contd...).
    ref_para_count = max(len(hl_paras), 1)
    ref_filled_para_count = max(
        sum(1 for p in hl_paras if _paragraph_text(p).strip()),
        1,
    )
    standard_gap = None
    ka_profile = {}
    if ref_ka:
        ka_r0 = ref_ka.table.rows[0].height
        ka_ref_r1 = ref_ka.table.rows[1].height
        ka_ref_pad = ref_ka.height - ka_r0 - ka_ref_r1
        ka_ref_items = max(
            sum(
                1
                for p in ref_ka.table.cell(1, 0).text_frame.paragraphs
                if p.text.strip()
            ),
            1,
        )
        standard_gap = ref_ka.top - ref_hl.top - ref_hl.height
        expanded_hl_height = ref_ka.top + ref_ka.height - ref_hl.top
        expanded_content_h = expanded_hl_height - r0 - r1 - ref_pad
        story_slots = max(ref_para_count - 2, 1)
        per_story_h = ref_r2 / story_slots if story_slots else ref_r2
        expanded_story_slots = max(int(expanded_content_h / per_story_h), 1) if per_story_h else 1
        expanded_capacity = expanded_story_slots + 2
        ka_profile = {
            "ka_r0": ka_r0,
            "ka_ref_r1": ka_ref_r1,
            "ka_ref_pad": ka_ref_pad,
            "ka_ref_items": ka_ref_items,
            "ref_ka_height": ref_ka.height,
            "ref_ka_top": ref_ka.top,
            "ref_ka_bottom": ref_ka.top + ref_ka.height,
            "ka_header_cell_ref": ref_ka.table.cell(0, 0),
            "ka_content_cell_ref": ref_ka.table.cell(1, 0),
            "expanded_hl_height": expanded_hl_height,
            "expanded_content_h": max(expanded_content_h, ref_r2),
            "expanded_capacity": max(expanded_capacity, ref_para_count),
        }
    return {
        "ref_hl": ref_hl,
        "ref_ka": ref_ka,
        "r0": r0,
        "r1": r1,
        "ref_r2": ref_r2,
        "ref_pad": ref_pad,
        "ref_para_count": ref_para_count,
        "ref_filled_para_count": ref_filled_para_count,
        "ref_hl_height": ref_hl.height,
        "ref_hl_top": ref_hl.top,
        "standard_gap": standard_gap,
        "hl_content_cell_ref": ref_hl.table.cell(2, 0),
        **ka_profile,
    }


def apply_cell_format(cell, ref_cell, top_align=False):
    """Copy G10X cell margins/anchor so content spacing matches the template."""
    ref_tcPr = ref_cell._tc.find(qn("a:tcPr"))
    tcPr = cell._tc.get_or_add_tcPr()
    for attr in ("anchor", "marT", "marB", "marL", "marR"):
        if ref_tcPr is not None and ref_tcPr.get(attr) is not None:
            tcPr.set(attr, ref_tcPr.get(attr))
        elif attr in tcPr.attrib:
            del tcPr.attrib[attr]

    ref_tf = ref_cell.text_frame
    tf = cell.text_frame
    tf.margin_top = ref_tf.margin_top
    tf.margin_bottom = ref_tf.margin_bottom
    tf.margin_left = ref_tf.margin_left
    tf.margin_right = ref_tf.margin_right

    ref_bodyPr = ref_tf._txBody.find(qn("a:bodyPr"))
    body = tf._txBody
    bodyPr = body.find(qn("a:bodyPr"))
    if bodyPr is None:
        from pptx.oxml import parse_xml
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        bodyPr = parse_xml(f'<a:bodyPr xmlns:a="{ns}"/>')
        body.insert(0, bodyPr)
    if ref_bodyPr is not None:
        for attr in ("anchor", "tIns", "bIns", "lIns", "rIns"):
            val = ref_bodyPr.get(attr)
            if val is not None:
                bodyPr.set(attr, val)

    if top_align:
        _set_cell_top_align(cell)


def _set_cell_top_align(cell):
    """
    Top-align text in a content cell.
    G10X fills the Highlights row to full height but text always starts flush
    below the header bar — not vertically centred in the tall cell.
    """
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.set("anchor", "t")
    bodyPr = cell.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        from pptx.oxml import parse_xml
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        bodyPr = parse_xml(f'<a:bodyPr xmlns:a="{ns}" anchor="t"/>')
        cell.text_frame._txBody.insert(0, bodyPr)
    else:
        bodyPr.set("anchor", "t")


def _set_table_shape_height(table_shape, row_heights, target_h):
    """Assign row heights then restore shape height (pptx resets height on row change)."""
    for i, h in enumerate(row_heights):
        table_shape.table.rows[i].height = h
    table_shape.height = target_h


def _hl_story_slots(ref_para_count):
    """Story line slots in the Highlights content row (excludes 2 sprint header lines)."""
    return max(ref_para_count - 2, 1)


def _calc_content_height(para_count, ref_para_count, ref_content_h, min_h, *, story_slots=None):
    """Scale a content row height based on paragraph/item count."""
    denom = story_slots if story_slots is not None else ref_para_count
    denom = max(denom, 1)
    per_unit = ref_content_h / denom
    if para_count >= ref_para_count:
        return max(int(per_unit * para_count), min_h)
    if para_count >= ref_para_count * 0.85:
        return ref_content_h
    return max(int(per_unit * para_count), min_h)


def _calc_table_pad(count, ref_count, ref_pad, at_least=91440):
    """Scale bottom table padding; full pad when content is at or above reference."""
    if count >= ref_count:
        return ref_pad
    if count >= ref_count * HL_KA_FILL_THRESHOLD:
        return ref_pad
    return max(int(ref_pad * count / ref_count), at_least)


def _hl_ka_tab_gap(profile):
    """Gap between HL table bottom and KA table top (~2 canonical body lines)."""
    return hl_ka_tab_gap_emu(profile.get("canonical_line_height_emu"))


def _hl_ka_gap(profile):
    """Legacy G10X shape-border overlap — not used for KA placement (see _hl_ka_tab_gap)."""
    gap = profile.get("standard_gap")
    return MIN_HL_KA_GAP if gap is None else gap


def _max_hl_height_before_ka(hl, ka_h, profile):
    """Legacy footer cap from HL top — prefer ka_fits_below_highlights after content-sized HL."""
    return int(MAX_KA_BOTTOM_EMU - ka_h - hl.top - _hl_ka_tab_gap(profile))


def _normal_layout_cap(profile, budgets):
    """Supplier-scale slot count — avoids small per-slide G10X refs forcing expanded HL."""
    return profile.get("canonical_para_slot_count", budgets["template_cap"])


def _hl_text_bottom_for_ka(hl, profile):
    """Bottom of last HL text — same anchor as format evaluation (text_ka_clearance_in)."""
    return _hl_rendered_text_bottom(hl, profile)


def _ka_top_after_hl(hl, profile):
    """KA top = HL table bottom + tab gap (same rule for every service)."""
    return hl.top + int(hl.height or 0) + _hl_ka_tab_gap(profile)


def _effective_ka_top(hl, profile):
    """KA top: tab gap below HL table, but never above last HL text + clearance."""
    tab_top = _ka_top_after_hl(hl, profile)
    text_floor = _hl_rendered_text_bottom(hl, profile) + MIN_TEXT_KA_CLEARANCE
    return max(tab_top, text_floor)


def _ka_top_after_hl_text(hl, profile):
    """Legacy text-anchored KA top — prefer _ka_top_after_hl for tab-to-tab spacing."""
    return _hl_text_bottom_for_ka(hl, profile) + _hl_ka_tab_gap(profile)


def _hl_capacity_above_ka(profile, ka_item_count=0):
    """Paragraph slots that fit above on-slide KA within the footer safe zone."""
    return _ka_anchored_hl_capacity(profile, ka_item_count)


def _ensure_hl_table_fits_rendered_text(hl, profile, max_h=None):
    """Grow the HL table so wrap-aware text is not clipped (returns target height)."""
    min_h = _min_hl_height_for_text(None, profile, hl=hl)
    if max_h is not None:
        min_h = min(min_h, max_h)
    target = max(int(hl.height or 0), min_h)
    r0, r1 = profile["r0"], profile["r1"]
    min_pad = max(int(profile["ref_pad"] * 0.15), 45720)
    for _ in range(4):
        content_h = max(target - r0 - r1 - min_pad, profile["r0"])
        _set_table_shape_height(hl, [r0, r1, content_h], target)
        rendered_bottom = _hl_rendered_text_bottom(hl, profile)
        needed = int(rendered_bottom - hl.top + HL_CONTENT_BOTTOM_PAD)
        if needed <= target:
            return target
        target = needed
        if max_h is not None:
            target = min(target, max_h)
    return target


def _desired_ka_top(hl, profile):
    """KA top: ~2 body lines below the Highlights table border (not text-anchored)."""
    return _effective_ka_top(hl, profile)


def _clamp_ka_top_to_footer(ka_h: int, desired_top: int) -> int:
    """Keep KA bottom within the footer safe zone (above logo / confidential bar)."""
    return min(desired_top, max(0, MAX_KA_BOTTOM_EMU - ka_h))


def _position_ka_below_hl(hl, ka, profile, *, ka_h: int | None = None):
    """Place KA below Highlights with G10X gap; clamp so footer is not covered."""
    if ka_h is None:
        ka_h = int(ka.height or 0)
    ka.top = _clamp_ka_top_to_footer(ka_h, _desired_ka_top(hl, profile))
    return ka.top


def _cell_margins(cell):
    body_pr = cell._tc.find(qn("a:txBody")).find(qn("a:bodyPr"))
    return (
        int(body_pr.get("marT") or 45720),
        int(body_pr.get("marB") or 45720),
    )


def _inner_content_bottom(hl):
    """Bottom of the text area inside the Highlights content cell."""
    r0, r1, r2 = (hl.table.rows[i].height for i in range(3))
    _, mar_b = _cell_margins(hl.table.cell(2, 0))
    return hl.top + r0 + r1 + r2 - mar_b


def _paragraph_bullet_level(p_elem):
    """Return list level from paragraph XML (0=sprint, 7=category, 1=story dash)."""
    pPr = p_elem.find(qn("a:pPr"))
    if pPr is None or pPr.get("lvl") is None:
        return 0
    return int(pPr.get("lvl"))


def _estimate_para_visual_lines(p_elem):
    """Estimate rendered lines for a paragraph (long GPT titles wrap in PowerPoint)."""
    text = _paragraph_text(p_elem).strip()
    if not text:
        return 0
    lvl = _paragraph_bullet_level(p_elem)
    chars = _CHARS_PER_VISUAL_LINE.get(lvl, _DEFAULT_CHARS_PER_LINE)
    return max(1, (len(text) + chars - 1) // chars)


def _count_visual_lines_in_hl(hl):
    """Total visual lines in Highlights content cell (paragraph count + wrapping)."""
    tx_body = hl.table.cell(2, 0).text_frame._txBody
    total = 0
    for p in tx_body.findall(qn("a:p")):
        if _paragraph_text(p).strip():
            total += _estimate_para_visual_lines(p)
    return total


def _effective_hl_line_count(hl, profile):
    """Paragraph slots vs wrapped lines — use whichever needs more vertical space."""
    para_count = _count_hl_paragraphs(hl)
    visual = _count_visual_lines_in_hl(hl)
    return max(para_count, visual)


def _hl_per_line_height_emu(profile):
    """
    Per-line EMU for wrap-aware HL text estimates.

    Always anchor on the Supplier canonical story-line height when available.
    Per-project ref_r2 / slot-count can be too large (Wentworth) or too small
    (Location) relative to real wrapped rendering — both cause bad HL sizing.
    """
    ref_para_count = max(profile["ref_para_count"], 1)
    project_pll = profile["ref_r2"] / max(ref_para_count - 2, 1)
    canonical_pll = profile.get("canonical_line_height_emu")
    if canonical_pll is None:
        return project_pll
    return canonical_pll


def _estimated_text_bottom(hl, profile):
    """Where rendered text ends inside HL (row height alone can underestimate wraps)."""
    r0, r1 = hl.table.rows[0].height, hl.table.rows[1].height
    mar_t, mar_b = _cell_margins(hl.table.cell(2, 0))
    per_line = _hl_per_line_height_emu(profile)
    line_count = _effective_hl_line_count(hl, profile)
    content_h = int(per_line * line_count * _VISUAL_LINE_HEIGHT_BUFFER)
    return hl.top + r0 + r1 + mar_t + content_h


def _min_hl_height_for_text(para_count, profile, hl=None):
    """Minimum HL table height so all paragraphs remain visible above KA."""
    r0, r1 = profile["r0"], profile["r1"]
    ref_pad = profile["ref_pad"]
    effective = para_count
    if hl is not None:
        effective = _effective_hl_line_count(hl, profile)
    per_line = _hl_per_line_height_emu(profile)
    content_h = max(
        int(per_line * effective * _VISUAL_LINE_HEIGHT_BUFFER),
        profile["r0"],
    )
    min_pad = max(int(ref_pad * 0.3), 91440)
    return r0 + r1 + content_h + min_pad


def _hl_rendered_text_bottom(hl, profile):
    """Wrap-aware bottom of visible Highlights text (matches format evaluation)."""
    est = _estimated_text_bottom(hl, profile)
    inner = _inner_content_bottom(hl)
    return max(est, inner)


def _is_dense_hl(hl, profile):
    """True when HL content is at or above the G10X fill mark (Supplier reference)."""
    para_count = _count_hl_paragraphs(hl)
    ref_filled = profile.get("canonical_fill_para_count")
    if ref_filled is None:
        return False
    return para_count >= ref_filled * HL_KA_FILL_THRESHOLD


def _hl_content_only_dimensions(hl, profile, max_h=None):
    """
    Highlights table height from rendered text only (+ small bottom pad).
    Never shorter than the minimum height required to hold all paragraphs/lines.
    """
    r0, r1 = profile["r0"], profile["r1"]
    min_pad = max(int(profile["ref_pad"] * 0.15), 45720)
    para_count = _count_hl_paragraphs(hl)
    min_h = _min_hl_height_for_text(para_count, profile, hl=hl)
    est_text_bottom = _estimated_text_bottom(hl, profile)
    target_h = int(est_text_bottom - hl.top + HL_CONTENT_BOTTOM_PAD)
    floor_h = r0 + r1 + profile["r0"] + min_pad
    target_h = max(target_h, floor_h, min_h)
    if max_h is not None:
        target_h = min(target_h, max_h)
    content_h = max(target_h - r0 - r1 - min_pad, profile["r0"])
    return target_h, content_h, min_pad


def _fit_hl_content_only(hl, profile, max_h=None):
    """Apply content-only Highlights sizing (text + pad, no G10X budget stretch)."""
    ref_hl = profile["ref_hl"]
    hl.left = ref_hl.left
    hl.top = ref_hl.top
    hl.width = ref_hl.width
    apply_cell_format(
        hl.table.cell(2, 0),
        profile["hl_content_cell_ref"],
        top_align=True,
    )
    target_h, content_h, pad = _hl_content_only_dimensions(hl, profile, max_h=max_h)
    r0, r1 = profile["r0"], profile["r1"]
    _set_table_shape_height(hl, [r0, r1, content_h], target_h)
    return target_h


def _count_hl_paragraphs(hl):
    """Count non-empty paragraphs in the Highlights content cell."""
    tx_body = hl.table.cell(2, 0).text_frame._txBody
    return sum(
        1 for p in tx_body.findall(qn("a:p")) if _paragraph_text(p).strip()
    )


def estimate_paragraph_count(raw):
    """Paragraphs the highlights cell will contain for this story data."""
    if raw.get("sections"):
        return _sections_total_paragraph_count(raw["sections"])

    return _single_section_paragraph_count(raw)


def _bucket_story_count(section, bucket, display=False):
    if display:
        key = f"{bucket}_items"
        return len(section.get(key) or [])
    return len(section.get(bucket) or [])


def _single_section_paragraph_count(section, display=False):
    """
    Paragraph slots for one sprint block — must match populate_highlights_* output.

    Continued sections omit the sprint line and current-week row; omitted category
    headers are not rendered on (Contd…) slides.
    """
    continued = section.get("continued_section")
    omit = set(section.get("omit_category_headers") or [])
    n = 0 if continued else 2
    for bucket in STORY_BUCKET_ORDER:
        count = _bucket_story_count(section, bucket, display=display)
        if not count:
            continue
        if bucket not in omit:
            n += 1
        n += count
    return n


def _sections_total_paragraph_count(sections, display=False):
    """Total paragraphs including one blank line between sprint blocks (rule 5)."""
    n = 0
    for i, section in enumerate(sections):
        if i > 0:
            n += 1
        n += _single_section_paragraph_count(section, display=display)
    return n


def _main_slide_paragraph_count(released, inprogress, completed):
    """Paragraphs on main slide for given section prefixes (rule 16 order)."""
    return _single_section_paragraph_count({
        "released": released,
        "inprogress": inprogress,
        "completed": completed,
    })


def _capacity_for_content_row(profile, content_row_h):
    """Paragraph slots that fit in a content-row height (incl. 2 sprint headers)."""
    story_slots_ref = max(profile["ref_para_count"] - 2, 1)
    per_story = profile["ref_r2"] / story_slots_ref
    if per_story <= 0:
        return profile["ref_para_count"]
    story_slots = max(int(content_row_h / per_story), 1)
    return story_slots + 2


def get_hl_budgets(profile):
    """
    G10X vertical budgets for the Highlights table (see Supplier main vs Contd):

    - template  : original single-slide Highlights cell height
    - with_ka   : Highlights grows through the on-slide KA zone
    - full_main : KA moves to (Contd...) — Highlights fills to MAX_KA_BOTTOM
    """
    ref_hl = profile["ref_hl"]
    ref_ka = profile.get("ref_ka")
    r0, r1, pad = profile["r0"], profile["r1"], profile["ref_pad"]

    template_h = profile["ref_hl_height"]
    template_content = profile["ref_r2"]
    template_cap = profile["ref_para_count"]

    with_ka_h = profile.get("expanded_hl_height")
    if ref_ka and not with_ka_h:
        with_ka_h = ref_ka.top + ref_ka.height - ref_hl.top
    with_ka_content = max((with_ka_h or template_h) - r0 - r1 - pad, template_content)

    full_main_h = MAX_KA_BOTTOM_EMU - ref_hl.top - MIN_HL_KA_GAP
    full_main_content = max(full_main_h - r0 - r1 - pad, template_content)

    return {
        "template_h": template_h,
        "template_content": template_content,
        "template_cap": template_cap,
        "with_ka_h": with_ka_h or template_h,
        "with_ka_content": with_ka_content,
        "with_ka_cap": _capacity_for_content_row(profile, with_ka_content),
        "full_main_h": full_main_h,
        "full_main_content": full_main_content,
        "full_main_cap": _capacity_for_content_row(profile, full_main_content),
    }


def _valid_main_prefix(nr, ni, nc, raw):
    """
    Sections on main must appear in order (released → in-progress → completed)
    without skipping an earlier section that still has overflow items.
    """
    if nc > 0 and raw["inprogress"] and ni == 0:
        return False
    if (ni > 0 or nc > 0) and raw["released"] and nr == 0:
        return False
    return True


def _max_paras_fit_above_ka(profile, ka_item_count=0):
    """Paragraphs that fit above on-slide KA (HL table + tab gap + KA within footer)."""
    ka_h = _estimate_ka_table_height(profile, ka_item_count)
    gap = _hl_ka_tab_gap(profile)
    max_h = int(MAX_KA_BOTTOM_EMU - ka_h - gap - profile["ref_hl_top"])
    cap = profile.get("ref_filled_para_count", profile["ref_para_count"])
    for n in range(cap, 1, -1):
        if _min_hl_height_for_text(n, profile) <= max_h:
            return n
    return 2


def _ka_anchored_hl_capacity(profile, ka_item_count=0):
    """Paragraph slots that fit above on-slide KA within the footer safe zone."""
    ka_h = _estimate_ka_table_height(profile, ka_item_count)
    gap = _hl_ka_tab_gap(profile)
    max_h = int(MAX_KA_BOTTOM_EMU - ka_h - gap - profile["ref_hl_top"])
    r0, r1, pad = profile["r0"], profile["r1"], profile["ref_pad"]
    content_h = max(max_h - r0 - r1 - pad, profile["r0"])
    return _capacity_for_content_row(profile, content_h)


def _estimate_ka_table_height(profile, item_count=None):
    """Total KA table height for item_count bullets (header + content + pad)."""
    n = 0 if item_count is None else max(int(item_count), 0)
    ka_r0 = profile.get("ka_r0")
    if ka_r0 is None:
        ref_ka = profile.get("ref_ka")
        if ref_ka is not None:
            ka_r0 = ref_ka.table.rows[0].height
        else:
            return profile.get("ref_ka_height", 914400)
    if n == 0:
        return ka_r0 + MIN_KA_ROW_HEIGHT + min(profile.get("ka_ref_pad", 0), 45720)
    per_item = profile["ka_ref_r1"] / max(profile.get("ka_ref_items", 1), 1)
    content_h = max(int(per_item * n), MIN_KA_ROW_HEIGHT)
    return ka_r0 + content_h + profile.get("ka_ref_pad", 0)


def ka_fits_below_highlights(hl, profile, ka_item_count=None, *, ka_profile=None):
    """
    True when KA fits a fixed gap below the last HL text line within the footer zone.
    """
    kp = ka_profile or profile
    if kp.get("ka_r0") is None and kp.get("ref_ka") is not None:
        kp = {**kp, **_ka_fields_from_ref(kp["ref_ka"], profile["ref_hl"])}
    ka_h = _estimate_ka_table_height(kp, ka_item_count)
    return _effective_ka_top(hl, profile) + ka_h <= MAX_KA_BOTTOM_EMU


def apply_ka_on_main_slide(slide, profile, raw, g10x_prs, g10x_layout, budgets=None):
    """
    Size HL to content and place KA below the last HL text when footer geometry allows.
    Returns True only when KA fits on the main slide.
    """
    hl = get_highlights_shape(slide)
    items = resolve_ka_items(raw)
    item_count = count_filled_ka_items(items)
    budgets = budgets or get_hl_budgets(profile)
    ka_profile = get_ka_layout_profile(g10x_prs, g10x_layout)

    if _is_dense_hl(hl, profile):
        fit_highlights_table(hl, profile, layout_mode="expanded", budgets=budgets)
    else:
        _fit_hl_content_only(hl, profile)
    _ensure_hl_table_fits_rendered_text(hl, profile)

    if not ka_fits_below_highlights(hl, profile, item_count, ka_profile=ka_profile):
        return False

    ka = _ensure_ka_shape_on_slide(slide, g10x_prs, g10x_layout)
    if not ka:
        return False

    set_ka_items(ka, items)
    ref_ka = ka_profile.get("ref_ka")
    if ref_ka:
        ka.left = ref_ka.left
        ka.width = ref_ka.width

    _position_ka_below_hl(hl, ka, profile)
    ka_h = fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
    _position_ka_below_hl(hl, ka, profile, ka_h=ka_h)
    if ka.top + ka_h > MAX_KA_BOTTOM_EMU:
        return False
    return True


def _section_paragraph_count(section):
    """Paragraph slots one sprint section uses in the Highlights cell."""
    return _single_section_paragraph_count(section)


def _section_has_stories(section):
    return bool(
        section.get("released") or section.get("inprogress") or section.get("completed")
    )


def _split_section_to_cap(section, remaining_cap):
    """Split one sprint section so the prefix fits in remaining_cap paragraph slots."""
    continued = section.get("continued_section")
    base = 0 if continued else 2
    if remaining_cap < base + 2:
        return None, dict(section)

    budget = remaining_cap - base
    main = {
        "sprint_name": section["sprint_name"],
        "sprint_dates": section["sprint_dates"],
        "sprint_status": section.get("sprint_status", "In-progress"),
        "released": [],
        "inprogress": [],
        "completed": [],
    }
    if continued:
        main["continued_section"] = True
        main["omit_category_headers"] = list(section.get("omit_category_headers") or [])
    tail = {
        "sprint_name": section["sprint_name"],
        "sprint_dates": section["sprint_dates"],
        "sprint_status": section.get("sprint_status", "In-progress"),
        "released": list(section.get("released", [])),
        "inprogress": list(section.get("inprogress", [])),
        "completed": list(section.get("completed", [])),
    }

    for bucket in STORY_BUCKET_ORDER:
        items = tail[bucket]
        if not items:
            continue
        omit_hdr = bucket in set(section.get("omit_category_headers") or [])
        need_hdr = 0 if omit_hdr else 1
        if budget < need_hdr + 1:
            break
        take = min(len(items), budget - need_hdr)
        if take <= 0:
            break
        main[bucket] = items[:take]
        tail[bucket] = items[take:]
        budget -= need_hdr + take

    if not _section_has_stories(main):
        return None, dict(section)
    if not _section_has_stories(tail):
        return main, None
    tail = dict(tail)
    tail["continued_section"] = True
    return main, tail


def _peel_hl_ka_suffix(sections, hl_ka_cap):
    """
    Split ``sections`` into an HL-only prefix and an HL+KA suffix (maximum suffix).

    Prefer ``_minimum_hl_ka_suffix`` when packing contd chains — it leaves more
    content on HL-only slides before the final HL+KA slide.
    """
    if not sections:
        return [], []
    total = sum(_section_paragraph_count(s) for s in sections)
    if total <= hl_ka_cap:
        return [], [dict(s) for s in sections]

    ka_sections: list[dict] = []
    used = 0
    split_at = len(sections)

    for idx in range(len(sections) - 1, -1, -1):
        section = sections[idx]
        sec_count = _section_paragraph_count(section)
        gap = 1 if ka_sections else 0
        if used + gap + sec_count <= hl_ka_cap:
            ka_sections.insert(0, dict(section))
            used += gap + sec_count
            split_at = idx
            continue

        remaining = hl_ka_cap - used - (1 if ka_sections else 0)
        if remaining >= 4:
            tail_part, keep_part = _split_section_to_cap(section, remaining)
            if tail_part and _section_has_stories(tail_part):
                tail_part = dict(tail_part)
                tail_part["continued_section"] = True
                ka_sections.insert(0, tail_part)
            if keep_part and _section_has_stories(keep_part):
                only_sections = [dict(s) for s in sections[:idx]]
                only_sections.append(keep_part)
                return only_sections, ka_sections
        break

    only_sections = [dict(s) for s in sections[:split_at]]
    return only_sections, ka_sections


def _minimum_hl_ka_suffix(sections, hl_ka_cap):
    """
    Smallest non-empty HL+KA suffix (from the end) that fits ``hl_ka_cap``.

    Pushes the maximum sprint content onto prior HL-only continuation slides.
    """
    if not sections:
        return [], []
    total = sum(_section_paragraph_count(s) for s in sections)
    if total <= hl_ka_cap:
        return [], [dict(s) for s in sections]

    if len(sections) >= 2:
        last = sections[-1]
        if _section_paragraph_count(last) <= hl_ka_cap and _section_has_stories(last):
            return [dict(s) for s in sections[:-1]], [dict(last)]

    return _peel_hl_ka_suffix(sections, hl_ka_cap)


def _compact_contd_chain(chain, hl_only_cap, hl_ka_cap):
    """
    Drop sparse HL-only slides by merging into the following HL+KA slide.

    A separate HL-only slide is only kept when it carries at least
  ``HL_KA_FILL_THRESHOLD`` of the full-column capacity.
    """
    if len(chain) < 2:
        return chain

    merged: list[tuple[str, list]] = []
    idx = 0
    min_fill = max(int(hl_only_cap * HL_KA_FILL_THRESHOLD), 6)
    while idx < len(chain):
        mode, sections = chain[idx]
        if (
            mode == "hl_only"
            and idx + 1 < len(chain)
            and chain[idx + 1][0] == "hl_ka"
        ):
            only_count = sum(_section_paragraph_count(s) for s in sections)
            ka_secs = chain[idx + 1][1]
            ka_count = sum(_section_paragraph_count(s) for s in ka_secs)
            if (
                only_count < min_fill
                and only_count + ka_count <= hl_ka_cap
            ):
                merged.append(("hl_ka", sections + ka_secs))
                idx += 2
                continue
        merged.append((mode, list(sections)))
        idx += 1
    return merged


def plan_contd_slide_chain(contd_raw, profile, ka_item_count=0):
    """
    Split overflow highlights across (Contd…) slides.

    * HL-only slides are filled to ``full_main_cap`` before another is added.
    * Paragraph counts match rendered output (continued sections, omitted headers).
    * Sparse HL-only segments are merged into the final HL+KA slide.
    """
    budgets = get_hl_budgets(profile)
    if contd_raw.get("sections"):
        remaining = [dict(s) for s in contd_raw["sections"]]
    else:
        remaining = [dict(contd_raw)]

    canonical_slots = profile.get(
        "canonical_para_slot_count",
        profile.get("canonical_fill_para_count", budgets["full_main_cap"]),
    )
    hl_only_cap = min(budgets["full_main_cap"], canonical_slots)
    hl_ka_cap = _ka_anchored_hl_capacity(profile, ka_item_count)
    chain: list[tuple[str, list]] = []

    while remaining:
        total = sum(_section_paragraph_count(s) for s in remaining)
        if total <= hl_ka_cap:
            chain.append(("hl_ka", remaining))
            break

        packed, rest = _pack_sections_to_cap(remaining, hl_only_cap)
        if not packed:
            chain.append(("hl_ka", remaining))
            break

        if not rest:
            hl_only_part, hl_ka_part = _minimum_hl_ka_suffix(packed, hl_ka_cap)
            combined = hl_only_part + hl_ka_part
            combined_n = sum(_section_paragraph_count(s) for s in combined)
            only_n = sum(_section_paragraph_count(s) for s in hl_only_part)
            min_fill = max(int(hl_only_cap * HL_KA_FILL_THRESHOLD), 6)
            if hl_only_part and only_n < min_fill and combined_n <= hl_ka_cap:
                chain.append(("hl_ka", combined))
            elif hl_only_part:
                chain.append(("hl_only", hl_only_part))
                chain.append(("hl_ka", hl_ka_part))
            else:
                chain.append(("hl_ka", hl_ka_part or remaining))
            break

        chain.append(("hl_only", packed))
        remaining = rest

    return _compact_contd_chain(chain, hl_only_cap, hl_ka_cap)


def plan_sections_split(raw, profile, layout_hints=None):
    """
    Pack sprint sections onto the main slide up to the G10X full-column fill mark.

    When Key Activities move to (Contd…), the main slide uses ``full_main_cap``
    (entire Highlights column), not the smaller KA-anchored capacity. Overflow
    starts only after the main slide is filled to that limit.
    """
    budgets = get_hl_budgets(profile)
    sections = raw.get("sections", [])
    canonical_fill = profile.get("canonical_fill_para_count", budgets["full_main_cap"])
    canonical_slots = profile.get(
        "canonical_para_slot_count", canonical_fill
    )
    service_hints = (layout_hints or LAYOUT_HINTS or {}).get(raw.get("title"), {})
    normal_cap = _normal_layout_cap(profile, budgets)
    main_pack_cap = min(budgets["full_main_cap"], canonical_slots)

    total = estimate_paragraph_count(raw)

    if service_hints.get("pack_all_sections_on_main") or service_hints.get("suppress_hl_contd"):
        if total <= main_pack_cap:
            if total <= normal_cap:
                return raw, None, "normal", budgets
            return raw, None, "expanded", budgets

    if total <= main_pack_cap:
        if total <= normal_cap:
            return raw, None, "normal", budgets
        return raw, None, "expanded", budgets

    main_sections = []
    contd_sections = []
    used = 0

    for idx, section in enumerate(sections):
        sec_count = _section_paragraph_count(section)
        gap = 1 if main_sections else 0
        if used + gap + sec_count <= main_pack_cap:
            main_sections.append(dict(section))
            used += gap + sec_count
            continue

        remaining = main_pack_cap - used - (1 if main_sections else 0)
        main_part, tail_part = _split_section_to_cap(section, remaining)
        if main_part and _section_has_stories(main_part):
            if main_sections:
                used += 1
            main_sections.append(main_part)
            used += _section_paragraph_count(main_part)
        if tail_part and _section_has_stories(tail_part):
            tail_part = dict(tail_part)
            tail_part["continued_section"] = True
            omit = [
                bucket
                for bucket in STORY_BUCKET_ORDER
                if main_part and main_part.get(bucket) and tail_part.get(bucket)
            ]
            if omit:
                tail_part["omit_category_headers"] = omit
            contd_sections.append(tail_part)
        contd_sections.extend(dict(s) for s in sections[idx + 1 :])
        break

    if not contd_sections:
        if total <= normal_cap:
            return raw, None, "normal", budgets
        return raw, None, "expanded", budgets

    # HL-UTIL-01: skip a sparse contd only when overflow is tiny and main is nearly full.
    main_util = used / main_pack_cap if main_pack_cap else 1.0
    if main_util < HL_KA_FILL_THRESHOLD:
        contd_paras = sum(
            _section_paragraph_count(s) for s in contd_sections
        )
        if contd_paras <= 4 and total <= main_pack_cap + 2:
            if total <= normal_cap:
                return raw, None, "normal", budgets
            return raw, None, "expanded", budgets

    main_raw = {**raw, "sections": main_sections}
    contd_raw = {**raw, "sections": contd_sections}
    return main_raw, contd_raw, "hl_ka_contd", budgets


def plan_highlights_split(raw, profile):
    """
    Plan HL layout mode from paragraph count vs footer-safe capacity above KA.
    KA placement (main vs contd) is decided later by ka_fits_below_highlights().
    """
    budgets = get_hl_budgets(profile)
    total = estimate_paragraph_count(raw)
    normal_cap = _normal_layout_cap(profile, budgets)
    on_slide_cap = min(
        _max_paras_fit_above_ka(profile),
        profile.get("ref_filled_para_count", budgets["template_cap"]),
        budgets["full_main_cap"],
    )

    if total <= normal_cap:
        return raw, None, "normal", budgets
    if total <= on_slide_cap:
        return raw, None, "expanded", budgets
    safe_cap = max(on_slide_cap - HL_OVERFLOW_SAFETY, 2)
    main, contd = split_raw_for_layout(raw, safe_cap)
    if contd:
        return main, contd, "hl_ka_contd", budgets
    return raw, None, "expanded", budgets


def split_raw_for_layout(raw, ref_para_count):
    """Pack maximum story content on main (up to ref_para_count) in section order."""
    if estimate_paragraph_count(raw) <= ref_para_count:
        return raw, None

    n_rel = len(raw["released"])
    n_ip = len(raw["inprogress"])
    n_comp = len(raw["completed"])

    best = None
    best_key = (-1, -1, -1, -1)

    for nr in range(n_rel + 1):
        for ni in range(n_ip + 1):
            for nc in range(n_comp + 1):
                if not _valid_main_prefix(nr, ni, nc, raw):
                    continue
                count = _main_slide_paragraph_count(
                    raw["released"][:nr],
                    raw["inprogress"][:ni],
                    raw["completed"][:nc],
                )
                if count <= ref_para_count:
                    key = (count, nr, ni, nc)
                    if key > best_key:
                        best_key = key
                        best = (nr, ni, nc)

    if best is None:
        return raw, None

    nr, ni, nc = best
    main = {
        "title": raw["title"],
        "sprint_name": raw["sprint_name"],
        "sprint_dates": raw["sprint_dates"],
        "released": raw["released"][:nr],
        "inprogress": raw["inprogress"][:ni],
        "completed": raw["completed"][:nc],
    }
    contd = {
        "released": raw["released"][nr:],
        "inprogress": raw["inprogress"][ni:],
        "completed": raw["completed"][nc:],
    }

    if not (contd["released"] or contd["inprogress"] or contd["completed"]):
        return raw, None

    contd_raw = {
        "title": raw["title"],
        "sprint_name": raw["sprint_name"],
        "sprint_dates": raw["sprint_dates"],
        **contd,
    }
    return main, contd_raw


def fit_highlights_table(hl, profile, layout_mode="normal", budgets=None, max_height=None):
    """
    Size the Highlights table.

    Sparse (<85% fill): content-only — table bottom = last text line + small pad.
    Dense: G10X budget stretch (template / expanded / full_main).
    """
    budgets = budgets or get_hl_budgets(profile)
    ref_hl = profile["ref_hl"]
    hl.left = ref_hl.left
    hl.top = ref_hl.top
    hl.width = ref_hl.width

    content_cell = hl.table.cell(2, 0)
    apply_cell_format(content_cell, profile["hl_content_cell_ref"], top_align=True)

    r0, r1 = profile["r0"], profile["r1"]
    ref_pad = profile["ref_pad"]
    ref_para_count = profile["ref_para_count"]

    if layout_mode == "normal":
        layout_max = budgets["template_h"]
        max_content = budgets["template_content"]
    elif layout_mode == "expanded":
        layout_max = budgets["with_ka_h"]
        max_content = budgets["with_ka_content"]
    else:  # full_main, contd
        layout_max = budgets["full_main_h"]
        max_content = budgets["full_main_content"]

    max_h = layout_max
    if max_height is not None:
        max_h = min(max_h, max_height)

    if not _is_dense_hl(hl, profile):
        return _fit_hl_content_only(hl, profile, max_h=max_h)

    if layout_mode in ("normal", "expanded"):
        para_count = _count_hl_paragraphs(hl)
        effective = _effective_hl_line_count(hl, profile)
        min_h = _min_hl_height_for_text(para_count, profile, hl=hl)
        ref_filled = profile.get(
            "canonical_fill_para_count",
            profile.get("ref_filled_para_count", ref_para_count),
        )
        content_h = _calc_content_height(
            effective,
            ref_para_count,
            max_content,
            profile["r0"],
            story_slots=_hl_story_slots(ref_para_count),
        )
        content_h = min(content_h, max_content)
        pad = _calc_table_pad(para_count, ref_para_count, ref_pad)
        shrunk_h = r0 + r1 + content_h + pad
        target_h = min(max(min_h, shrunk_h), max_h)
        # At G10X reference density, stretch HL to the fill mark when text allows.
        if (
            para_count >= ref_filled * HL_KA_FILL_THRESHOLD
            and min_h <= max_h
        ):
            target_h = max_h
            content_h = max(max_h - r0 - r1 - pad, profile["r0"])
        elif target_h > shrunk_h:
            content_h = max(target_h - r0 - r1 - pad, profile["r0"])
    else:  # full_main, contd — rule 7: fill to Supplier mark when dense
        para_count = _count_hl_paragraphs(hl)
        effective = _effective_hl_line_count(hl, profile)
        ref_filled = profile.get(
            "canonical_fill_para_count",
            profile.get("ref_filled_para_count", ref_para_count),
        )
        min_h = _min_hl_height_for_text(para_count, profile, hl=hl)
        if para_count >= ref_filled * HL_KA_FILL_THRESHOLD:
            content_h = max(max_content, min_h - r0 - r1 - ref_pad)
            pad = ref_pad
            target_h = max(max_h, min_h)
        else:
            content_h = _calc_content_height(
                effective,
                ref_para_count,
                max_content,
                profile["r0"],
                story_slots=_hl_story_slots(ref_para_count),
            )
            content_h = min(content_h, max_content)
            pad = _calc_table_pad(para_count, ref_para_count, ref_pad)
            target_h = min(max(min_h, r0 + r1 + content_h + pad), max_h)

    _set_table_shape_height(hl, [r0, r1, content_h], target_h)
    return target_h


def fit_key_activities_table(ka, profile, position_ref=None):
    """
    G10X Key Activities tab rules:
    - Fixed left/width from reference.
    - Row 0 (header) keeps the slide's existing height unless it is 0.
    - Content row (row 1) scales with bullet count; table bottom pad from reference.
    - Content cell margins/anchor copied from G10X (marT≈54000, tf margin≈45720).
    - Returns target_h (total KA table height).
    """
    ref_ka = position_ref or profile.get("ref_ka")
    if ref_ka is not None:
        ka.left = ref_ka.left
        ka.width = ref_ka.width

    apply_cell_format(ka.table.cell(0, 0), profile["ka_header_cell_ref"])
    content_cell = ka.table.cell(1, 0)
    apply_cell_format(content_cell, profile["ka_content_cell_ref"], top_align=True)

    item_count = count_filled_ka_items(ka)
    ka_ref_r1 = profile["ka_ref_r1"]
    ka_ref_pad = profile["ka_ref_pad"]
    ka_ref_items = profile["ka_ref_items"]
    ref_ka_height = profile["ref_ka_height"]

    ka_r0 = ka.table.rows[0].height
    if ka_r0 == 0:
        ka_r0 = profile["ka_r0"]
    if ka_r0 == 0:
        ka_r0 = MIN_KA_ROW_HEIGHT

    min_content = MIN_KA_ROW_HEIGHT
    if item_count == 0:
        content_h = MIN_KA_ROW_HEIGHT if ka_r0 > 0 else max(MIN_KA_ROW_HEIGHT, 91440)
        pad = 22860 if ka_r0 == 0 else min(ka_ref_pad, 45720)
        target_h = ka_r0 + content_h + pad
        ka.table.rows[0].height = ka_r0
        ka.table.rows[1].height = content_h
        ka.height = target_h
        return target_h

    content_h = _calc_content_height(
        item_count, ka_ref_items, ka_ref_r1, min_content
    )
    if item_count > 0 and content_h == 0:
        content_h = max(ka_ref_r1, MIN_KA_ROW_HEIGHT)
    pad = _calc_table_pad(item_count, ka_ref_items, ka_ref_pad)

    if item_count >= ka_ref_items * HL_KA_FILL_THRESHOLD and item_count <= ka_ref_items:
        content_h = max(ka_ref_r1, MIN_KA_ROW_HEIGHT) if item_count else ka_ref_r1
        pad = ka_ref_pad
        target_h = ref_ka_height
        if target_h < ka_r0 + content_h + pad:
            target_h = ka_r0 + content_h + pad
    else:
        target_h = ka_r0 + content_h + pad

    ka.table.rows[0].height = ka_r0
    ka.table.rows[1].height = content_h
    ka.height = target_h
    return target_h


def fit_highlights_layout(slide, g10x_ref_slide, position_ka=True, layout_mode="normal", budgets=None, profile=None):
    """
    Size Highlights to content, then place KA a fixed gap below the last HL text line.
    """
    if profile is None:
        profile = build_layout_profile(g10x_ref_slide)
    if budgets is None:
        budgets = get_hl_budgets(profile)
    ref_ka = profile["ref_ka"]
    hl = get_highlights_shape(slide)
    ka = get_key_activities_shape(slide) if position_ka and ref_ka else None

    fit_highlights_table(
        hl,
        profile,
        layout_mode=layout_mode,
        budgets=budgets,
    )

    if not (position_ka and ref_ka and ka and layout_mode not in ("contd", "full_main")):
        return

    ka_h = fit_key_activities_table(ka, profile)
    _position_ka_below_hl(hl, ka, profile)
    fit_key_activities_table(ka, profile)


def fit_contd_slide_layout(contd_slide, g10x_layout, g10x_prs=None, budgets=None):
    """
    (Contd...) slide: size Highlights to content, then KA below last HL text + gap.
    """
    profile = build_layout_profile(g10x_layout)
    ka_profile = (
        get_ka_layout_profile(g10x_prs, g10x_layout)
        if g10x_prs is not None
        else profile
    )
    if budgets is None:
        budgets = get_hl_budgets(profile)
    ref_hl = profile["ref_hl"]
    ref_ka = ka_profile.get("ref_ka")
    hl = get_highlights_shape(contd_slide)
    ka = get_key_activities_shape(contd_slide)
    if not ka:
        fit_highlights_table(hl, profile, layout_mode="full_main", budgets=budgets)
        return

    hl.left = ref_hl.left
    hl.top = ref_hl.top
    hl.width = ref_hl.width
    apply_cell_format(
        hl.table.cell(2, 0),
        profile["hl_content_cell_ref"],
        top_align=True,
    )

    if ref_ka:
        ka.left = ref_ka.left
        ka.width = ref_ka.width

    ka_h = fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
    gap = _hl_ka_tab_gap(profile)
    max_hl_bottom = MAX_KA_BOTTOM_EMU - ka_h - gap
    max_h = int(max_hl_bottom - ref_hl.top)
    min_h = profile["r0"] + profile["r1"] + profile["r0"]
    max_h = max(max_h, min_h)

    if _is_dense_hl(hl, profile):
        fit_highlights_table(
            hl, profile, layout_mode="expanded", budgets=budgets, max_height=max_h
        )
    else:
        _ensure_hl_table_fits_rendered_text(hl, profile, max_h=max_h)

    ka_h = fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
    _position_ka_below_hl(hl, ka, profile, ka_h=ka_h)


def tighten_slide_hl_ka_spacing(slide, g10x_layout, g10x_prs, budgets=None, profile=None):
    """
    Final pass: finalize HL size, then seat KA ~2 body lines below the HL table border.
    KA-only contd slides keep the Pharmacy reference block geometry (no HL on slide).
    """
    if profile is None:
        profile = build_layout_profile(g10x_layout)
    if budgets is None:
        budgets = get_hl_budgets(profile)

    try:
        hl = get_highlights_shape(slide)
    except ValueError:
        hl = None
    ka = get_key_activities_shape(slide)

    if not hl and ka:
        ka_profile = get_ka_layout_profile(g10x_prs, g10x_layout)
        contd_ref = get_contd_ka_reference(g10x_prs)
        fit_key_activities_table(ka, ka_profile, position_ref=contd_ref)
        if contd_ref:
            ka.top = contd_ref.top
            ka.left = contd_ref.left
            ka.width = contd_ref.width
        if count_filled_ka_items(ka) == 0:
            ka.height = min(ka.height, EMPTY_KA_CONTD_MAX_EMU)
            if ka.table.rows[0].height > 0:
                ka.table.rows[1].height = max(
                    ka.height - ka.table.rows[0].height, MIN_KA_ROW_HEIGHT
                )
        return

    if not hl:
        return

    if ka and hl:
        ka_profile = get_ka_layout_profile(g10x_prs, g10x_layout)
        ref_ka = ka_profile.get("ref_ka")
        ka_h_est = _estimate_ka_table_height(ka_profile, count_filled_ka_items(ka))
        gap = _hl_ka_tab_gap(profile)
        max_hl_h = int(
            MAX_KA_BOTTOM_EMU - ka_h_est - gap - profile["ref_hl_top"]
        )
        if not _is_dense_hl(hl, profile):
            _ensure_hl_table_fits_rendered_text(hl, profile)
        else:
            fit_highlights_table(
                hl,
                profile,
                layout_mode="expanded",
                budgets=budgets,
                max_height=profile["ref_hl_top"] + max_hl_h,
            )
        ka_h = fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
        _position_ka_below_hl(hl, ka, profile, ka_h=ka_h)
    elif hl:
        if _is_dense_hl(hl, profile):
            fit_highlights_table(hl, profile, layout_mode="full_main", budgets=budgets)
        else:
            _fit_hl_content_only(hl, profile, max_h=budgets["full_main_h"])


def _paragraph_text(p_elem):
    return "".join(n.text or "" for n in p_elem.iter() if n.tag.endswith("}t"))


def _normalize_match_text(text):
    """Normalize text for section-header matching (nbsp, zero-width space)."""
    return (text or "").replace("\u200b", "").replace("\xa0", " ").lower()


def strip_trailing_empty_paragraphs(txBody):
    """Remove trailing blank bullet placeholders from a highlights cell."""
    all_p = txBody.findall(qn("a:p"))
    while len(all_p) > 2:
        if _paragraph_text(all_p[-1]).strip():
            break
        txBody.remove(all_p[-1])
        all_p = txBody.findall(qn("a:p"))


def _advance_offset_after_bullets(txBody, header_offset, items, has_category_header=True):
    """
    Replace story bullets immediately after category header (G10X rule 5/15).
    No blank line between header and first story.
    """
    if not items:
        all_p = txBody.findall(qn("a:p"))
        return min(header_offset + (1 if has_category_header else 0), len(all_p))
    bullet_index = header_offset + (1 if has_category_header else 0)
    replace_bullet_block(txBody, bullet_index, bullet_index, items)
    return bullet_index + len(items)


def _append_bucket_templates(new_paras, section_tmpl, section, omit=()):
    """Append paragraph templates in G10X bucket order (rule 16)."""
    for bucket in STORY_BUCKET_ORDER:
        items = section.get(bucket) or section.get(f"{bucket}_items")
        if not items:
            continue
        if bucket in omit:
            new_paras.append(copy.deepcopy(section_tmpl[f"{bucket}_bullet"]))
            continue
        new_paras.extend([
            copy.deepcopy(section_tmpl[f"{bucket}_hdr"]),
            copy.deepcopy(section_tmpl[f"{bucket}_bullet"]),
        ])


def _fill_story_buckets(hl_txBody, offset, section, omit=(), display=False):
    """Write category headers + story bullets in G10X order (rule 16)."""
    for bucket in STORY_BUCKET_ORDER:
        if display:
            items = section.get(f"{bucket}_items") or []
            count = section.get(f"{bucket}_count", str(len(items)))
        else:
            items = section.get(bucket) or []
            count = str(len(items))
        if not items:
            continue

        all_p = hl_txBody.findall(qn("a:p"))
        has_header = bucket not in omit
        if has_header:
            hdr_p = all_p[offset]
            if bucket == "completed":
                set_two_run_header(
                    hdr_p,
                    "Stories completed this week ",
                    f"{count} stories",
                )
            elif bucket == "released":
                set_two_run_header(
                    hdr_p,
                    "Stories released for partner review ",
                    f"{count} stories",
                )
            else:
                set_single_run_text(
                    hdr_p,
                    f"Stories in-progress this week \u2013 {count} stories",
                )
        offset = _advance_offset_after_bullets(
            hl_txBody, offset, items, has_category_header=has_header
        )
    return offset


def get_canonical_style_cell(g10x_prs):
    """Cost Core cell — canonical lvl 0 sprint / lvl 7 category / lvl 1 dash story bullets."""
    return get_highlights_template_cell(g10x_prs.slides[CANONICAL_STYLE_SLIDE_INDEX])


def get_canonical_fill_para_count(g10x_prs):
    """Supplier slide filled paragraph count — the fill-to-mark reference (~19)."""
    supplier_slide = g10x_prs.slides[CANONICAL_FILL_SLIDE_INDEX]
    ref_hl, _ = get_g10x_highlights_and_ka(supplier_slide)
    paras = ref_hl.table.cell(2, 0).text_frame._txBody.findall(qn("a:p"))
    return max(sum(1 for p in paras if _paragraph_text(p).strip()), 1)


def get_canonical_para_slot_count(g10x_prs):
    """Supplier slide total paragraph slots (incl. sprint gap placeholder) — ~20."""
    supplier_slide = g10x_prs.slides[CANONICAL_FILL_SLIDE_INDEX]
    ref_hl, _ = get_g10x_highlights_and_ka(supplier_slide)
    paras = ref_hl.table.cell(2, 0).text_frame._txBody.findall(qn("a:p"))
    return max(len(paras), 1)


def get_canonical_line_height_emu(g10x_prs):
    """Supplier story-line height — stable wrap estimate across all services."""
    supplier_slide = g10x_prs.slides[CANONICAL_FILL_SLIDE_INDEX]
    ref_hl, _ = get_g10x_highlights_and_ka(supplier_slide)
    r2 = ref_hl.table.rows[2].height
    para_slots = max(len(ref_hl.table.cell(2, 0).text_frame._txBody.findall(qn("a:p"))), 1)
    return int(r2 / max(para_slots - 2, 1))


def _is_sustainment_story_bullet(p_elem):
    """True when paragraph matches Supplier reference dash story bullet (lvl 1)."""
    if _paragraph_bullet_level(p_elem) != 1:
        return False
    pPr = p_elem.find(qn("a:pPr"))
    if pPr is None:
        return False
    bu = pPr.find(qn("a:buChar"))
    return bu is not None


def _find_canonical_story_bullet(paras):
    """First lvl-1 dash bullet in Cost Core template (Supplier story style)."""
    for p in paras:
        if _is_sustainment_story_bullet(p):
            return p
    comp_idx = _find_para_index(paras, "completed this week")
    if comp_idx is not None and comp_idx + 1 < len(paras):
        return paras[comp_idx + 1]
    return None


def get_highlights_template_cell(g10x_slide):
    hl_ref = next(
        s
        for s in g10x_slide.shapes
        if s.has_table and s.table.cell(0, 0).text.strip() == "Highlights"
    )
    return hl_ref.table.cell(2, 0)


def _find_para_index(paras, *needles, start=0):
    for i in range(start, len(paras)):
        text = _normalize_match_text(_paragraph_text(paras[i]))
        if any(n in text for n in needles):
            return i
    return None


def _find_sprint_para_index(paras):
    """Sprint summary line — not 'Current week sprint status' or 'SPUR' labels."""
    for i, p in enumerate(paras):
        text = _normalize_match_text(_paragraph_text(p)).strip()
        if not text or text == "spur":
            continue
        if "current week sprint status" in text:
            continue
        if text.startswith("sprint") or "sprint \u2013" in text or "sprint -" in text:
            return i
    return None


def _find_same_slide_bullet_para(paras):
    """Any story bullet paragraph from this slide (for style fallback)."""
    for p in paras:
        text = _paragraph_text(p).strip()
        if not text:
            continue
        tl = _normalize_match_text(text)
        if any(
            x in tl
            for x in (
                "released for partner",
                "in-progress this week",
                "in progress this week",
                "completed this week",
                "current week sprint status",
            )
        ):
            continue
        pPr = p.find(qn("a:pPr"))
        if pPr is not None and pPr.find(qn("a:buChar")) is not None:
            return p
    return None


def set_sprint_line(p_elem, sprint_bold, sprint_light):
    """Set sprint header text, preserving two-run bold/light when available."""
    runs = p_elem.findall(qn("a:r"))
    if len(runs) >= 2:
        runs[0].find(qn("a:t")).text = sprint_bold
        runs[1].find(qn("a:t")).text = sprint_light
        for extra in runs[2:]:
            p_elem.remove(extra)
    else:
        set_single_run_text(p_elem, sprint_bold + sprint_light)


def discover_section_templates(canonical_cell, global_completed=None):
    """
    Sustainment paragraph styles always come from Cost Core (Supplier reference).
    lvl 0: sprint / current-week (round bullet)
    lvl 7: category headers (arrow bullet, bold)
    lvl 1: story items (dash bullet) — aligned with current-week, not category
    """
    paras = canonical_cell.text_frame._txBody.findall(qn("a:p"))
    story_bullet = _find_canonical_story_bullet(paras)

    sprint_idx = _find_sprint_para_index(paras)
    cw_idx = _find_para_index(paras, "current week sprint status")
    comp_idx = _find_para_index(paras, "completed this week")
    rel_idx = _find_para_index(paras, "released for partner review")
    prog_idx = _find_para_index(paras, "in-progress this week", "in progress this week")

    def pick(idx):
        if idx is not None and idx < len(paras):
            return paras[idx]
        return None

    def pick_bullet(hdr_idx):
        if hdr_idx is not None and hdr_idx + 1 < len(paras):
            candidate = paras[hdr_idx + 1]
            if _is_sustainment_story_bullet(candidate):
                return candidate
        return story_bullet

    required = {
        "sprint": pick(sprint_idx) if pick(sprint_idx) is not None else paras[0],
        "current_week": pick(cw_idx) if pick(cw_idx) is not None else paras[1],
        "completed_hdr": pick(comp_idx),
        "completed_bullet": pick_bullet(comp_idx) if pick_bullet(comp_idx) is not None else story_bullet,
        "released_hdr": pick(rel_idx),
        "released_bullet": pick_bullet(rel_idx) if pick_bullet(rel_idx) is not None else story_bullet,
        "inprogress_hdr": pick(prog_idx),
        "inprogress_bullet": pick_bullet(prog_idx) if pick_bullet(prog_idx) is not None else story_bullet,
    }
    if global_completed:
        if required["completed_hdr"] is None:
            required["completed_hdr"] = global_completed[0]
        if required["completed_bullet"] is None:
            required["completed_bullet"] = global_completed[1]
    for key, val in list(required.items()):
        if val is None:
            if key.endswith("_hdr"):
                required[key] = required.get("completed_hdr")
            elif key.endswith("_bullet"):
                required[key] = story_bullet
    if any(v is None for v in required.values()):
        raise RuntimeError("Could not resolve highlights section templates from G10X")
    return {key: copy.deepcopy(val) for key, val in required.items()}


def populate_highlights_cell(hl_cell, section_tmpl, content):
    """Populate highlights — supports multiple sprint sections per slide."""
    hl_txBody = hl_cell.text_frame._txBody

    for p in list(hl_txBody.findall(qn("a:p"))):
        hl_txBody.remove(p)

    sections = content.get("sections")
    if not sections:
        sections = [content]

    new_paras = []
    for si, sec in enumerate(sections):
        if si > 0:
            new_paras.append(_make_sprint_gap_para(section_tmpl))
        new_paras.append(copy.deepcopy(section_tmpl["sprint"]))
        new_paras.append(copy.deepcopy(section_tmpl["current_week"]))
        _append_bucket_templates(new_paras, section_tmpl, sec)
    for p in new_paras:
        hl_txBody.append(p)

    offset = 0
    for si, sec in enumerate(sections):
        if si > 0:
            offset += 1
        all_p = hl_txBody.findall(qn("a:p"))
        set_sprint_line(
            all_p[offset],
            sec["sprint_bold"],
            sec["sprint_light"],
        )
        offset += 2
        offset = _fill_story_buckets(hl_txBody, offset, sec, display=True)

    strip_trailing_empty_paragraphs(hl_txBody)


def populate_highlights_contd_cell(hl_cell, section_tmpl, contd_raw):
    """Highlights overflow on (Contd...) — supports multi-sprint sections."""
    if contd_raw.get("sections"):
        _populate_highlights_contd_sections(
            hl_cell, section_tmpl, contd_raw["sections"]
        )
        return

    hl_txBody = hl_cell.text_frame._txBody

    for p in list(hl_txBody.findall(qn("a:p"))):
        hl_txBody.remove(p)

    omit = set(contd_raw.get("omit_category_headers", []))

    new_paras = []
    _append_bucket_templates(new_paras, section_tmpl, contd_raw, omit=omit)
    for p in new_paras:
        hl_txBody.append(p)

    offset = 0
    offset = _fill_story_buckets(hl_txBody, offset, contd_raw, omit=omit)

    strip_trailing_empty_paragraphs(hl_txBody)


def _populate_highlights_contd_sections(hl_cell, section_tmpl, sections):
    """Populate (Contd...) highlights from sprint section dicts."""
    hl_txBody = hl_cell.text_frame._txBody
    for p in list(hl_txBody.findall(qn("a:p"))):
        hl_txBody.remove(p)

    new_paras = []
    for si, section in enumerate(sections):
        omit = set(section.get("omit_category_headers", []))
        continued = section.get("continued_section")
        if si > 0 and not continued:
            new_paras.append(_make_sprint_gap_para(section_tmpl))
        if not continued:
            new_paras.append(copy.deepcopy(section_tmpl["sprint"]))
            new_paras.append(copy.deepcopy(section_tmpl["current_week"]))
        _append_bucket_templates(new_paras, section_tmpl, section, omit=omit)
    for p in new_paras:
        hl_txBody.append(p)

    offset = 0
    for si, section in enumerate(sections):
        omit = set(section.get("omit_category_headers", []))
        continued = section.get("continued_section")
        if si > 0 and not continued:
            offset += 1
        display = _section_display_content(section)
        all_p = hl_txBody.findall(qn("a:p"))
        if not continued:
            set_sprint_line(
                all_p[offset],
                display["sprint_bold"],
                display["sprint_light"],
            )
            offset += 2
        offset = _fill_story_buckets(
            hl_txBody, offset, display, omit=omit, display=True
        )

    strip_trailing_empty_paragraphs(hl_txBody)


def get_contd_ka_reference(g10x_prs):
    """G10X Pharmacy (Contd...) slide — canonical KA position for KA-only contd slides."""
    return get_key_activities_shape(g10x_prs.slides[PHARMACY_CONTD_INDEX])


def cleanup_ka_contd_orphans(contd_slide):
    """Remove empty Highlights placeholder tables left on KA-only contd slides."""
    ka = get_key_activities_shape(contd_slide)
    for shape in list(contd_slide.shapes):
        if not shape.has_table or shape is ka:
            continue
        try:
            hdr = shape.table.cell(0, 0).text.strip()
        except (IndexError, AttributeError):
            hdr = ""
        if "Key activities" not in hdr:
            delete_shape(shape)


def move_ka_from_main_to_contd(main_slide, contd_slide, g10x_layout, g10x_prs=None):
    """Key Activities appear only after highlights finish — on the (Contd...) slide."""
    if g10x_prs is not None:
        _ensure_ka_shape_on_slide(main_slide, g10x_prs, g10x_layout)
    main_ka = get_key_activities_shape(main_slide)
    if not main_ka:
        if g10x_prs is not None:
            _ensure_ka_shape_on_slide(contd_slide, g10x_prs, g10x_layout)
        return

    contd_ka = get_key_activities_shape(contd_slide)
    if contd_ka:
        delete_shape(contd_ka)

    newel = copy.deepcopy(main_ka.element)
    contd_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    delete_shape(main_ka)

    if g10x_prs is not None and not get_key_activities_shape(contd_slide):
        _ensure_ka_shape_on_slide(contd_slide, g10x_prs, g10x_layout)


def fit_hl_only_contd_layout(contd_slide, g10x_layout, budgets=None):
    """HL-only (Contd...) slide: size Highlights to content up to the footer fill mark."""
    profile = build_layout_profile(g10x_layout)
    budgets = budgets or get_hl_budgets(profile)
    hl = get_highlights_shape(contd_slide)
    ref_hl = profile["ref_hl"]
    hl.left = ref_hl.left
    hl.top = ref_hl.top
    hl.width = ref_hl.width
    apply_cell_format(
        hl.table.cell(2, 0),
        profile["hl_content_cell_ref"],
        top_align=True,
    )
    if _is_dense_hl(hl, profile):
        fit_highlights_table(hl, profile, layout_mode="full_main", budgets=budgets)
    else:
        _fit_hl_content_only(hl, profile, max_h=budgets["full_main_h"])


def ensure_hl_only_contd_slide(
    prs,
    g10x_prs,
    service_title,
    main_idx,
    contd_raw,
    ref_title,
    g10x_layout,
    section_tmpl,
):
    """
    (Contd...) with highlights overflow only — Key Activities stays on the main slide.
    G10X pattern when HL + KA fit together on main but highlights need one more slide.
    """
    contd_idx = find_contd_slide(prs, service_title)
    main_slide = prs.slides[main_idx]

    if contd_idx is None:
        template_slide = g10x_prs.slides[HIGHLIGHTS_CONTD_TEMPLATE_INDEX]
        new_slide = prs.slides.add_slide(main_slide.slide_layout)
        clear_slide_shapes(new_slide)
        copy_shapes_to_slide(template_slide, new_slide)
        move_slide_after(prs, len(prs.slides) - 1, main_idx)
        contd_idx = main_idx + 1

    contd_slide = prs.slides[contd_idx]
    set_title_suffix(contd_slide, f"{service_title}  (Contd\u2026)", ref_title)

    hl_shape = get_highlights_shape(contd_slide)
    populate_highlights_contd_cell(
        hl_shape.table.cell(2, 0),
        section_tmpl,
        contd_raw,
    )
    fit_hl_only_contd_layout(contd_slide, g10x_layout)

    sync_heb_logo_from_main(main_slide, contd_slide)

    return contd_idx


def ensure_highlights_contd_slide(
    prs,
    g10x_prs,
    service_title,
    main_idx,
    contd_raw,
    ref_title,
    g10x_layout,
    section_tmpl,
    profile=None,
    raw=None,
    layout_mode="hl_ka_contd",
):
    """
    (Contd...) slide(s): remaining highlights, then Key Activities on the last slide.

  When one (Contd…) slide is not enough, additional HL-only slides are inserted
  instead of overlapping Highlights and Key Activities.
    """
    if profile is None:
        profile = build_layout_profile(g10x_layout)
    if layout_mode == "supplier_contd":
        return ensure_hl_only_contd_slide(
            prs,
            g10x_prs,
            service_title,
            main_idx,
            contd_raw,
            ref_title,
            g10x_layout,
            section_tmpl,
        )
    for idx in sorted(find_contd_slides_for_service(prs, service_title), reverse=True):
        delete_slide(prs, idx)
    return ensure_contd_slide_chain(
        prs,
        g10x_prs,
        service_title,
        main_idx,
        contd_raw,
        ref_title,
        g10x_layout,
        section_tmpl,
        profile,
        raw if raw is not None else contd_raw,
    )


def ensure_ka_contd_slide(prs, g10x_prs, service_title, main_idx, ref_title, g10x_layout):
    """
    (Contd...) slide with Key Activities only — all story content stays on main.
    Uses the G10X Pharmacy (Contd) layout which has KA but no Highlights table.
    """
    contd_idx = find_contd_slide(prs, service_title)
    main_slide = prs.slides[main_idx]

    if contd_idx is None:
        template_slide = g10x_prs.slides[PHARMACY_CONTD_INDEX]
        new_slide = prs.slides.add_slide(main_slide.slide_layout)
        clear_slide_shapes(new_slide)
        copy_shapes_to_slide(template_slide, new_slide)
        move_slide_after(prs, len(prs.slides) - 1, main_idx)
        contd_idx = main_idx + 1

    contd_slide = prs.slides[contd_idx]
    set_title_suffix(contd_slide, f"{service_title}  (Contd\u2026)", ref_title)

    cleanup_ka_contd_orphans(contd_slide)
    move_ka_from_main_to_contd(main_slide, contd_slide, g10x_layout, g10x_prs)

    profile = get_ka_layout_profile(g10x_prs, g10x_layout)
    contd_ka_ref = get_contd_ka_reference(g10x_prs)
    ka = _ensure_ka_shape_on_slide(contd_slide, g10x_prs, g10x_layout)
    if ka and contd_ka_ref:
        fit_key_activities_table(ka, profile)
        ka.top = contd_ka_ref.top
        ka.left = contd_ka_ref.left
        ka.width = contd_ka_ref.width

    sync_heb_logo_from_main(main_slide, contd_slide)

    return contd_idx


def find_completed_templates():
    src = Presentation(G10X)
    for slide in src.slides:
        for sh in slide.shapes:
            if not sh.has_table:
                continue
            try:
                if sh.table.cell(0, 0).text.strip() != "Highlights":
                    continue
            except (IndexError, AttributeError):
                continue
            paras = sh.table.cell(2, 0).text_frame._txBody.findall(qn("a:p"))
            for i, p in enumerate(paras):
                text = _normalize_match_text(_paragraph_text(p))
                if "completed this week" in text:
                    bullet_idx = i + 1
                    if bullet_idx < len(paras):
                        return paras[i], paras[bullet_idx]
    raise RuntimeError("Could not find completed-section template paragraphs")


def delete_shape(shape):
    shape._element.getparent().remove(shape._element)


# Map service name -> title substring used to locate slide in deck
TITLE_SEARCH = {
    "Supplier Core Service": "Supplier Core",
    "LoCo": "LoCo",
    "Pricing Core": "Pricing Core",
    "Patronage Travel": "Patronage Travel",
    "PATRV": "Patronage Travel",
}

CANONICAL_DELIVERY_CLONE_TITLE = "Cost Core Service"


def _last_main_delivery_slide_index(prs) -> int | None:
    last: int | None = None
    for i, slide in enumerate(prs.slides):
        title = _delivery_slide_title_text(slide)
        if _is_delivery_status_slide_title(title) and "(contd" not in title.lower():
            last = i
    return last


def ensure_delivery_slide(
    prs,
    service_title: str,
    ref_title_shape,
) -> int:
    """
    Return the 0-based index of the delivery slide for ``service_title``.

    Clones the Cost Core delivery slide when the G10X deck has no matching slide
  (new team tracks introduced after template publish).
    """
    slide_idx = find_slide_by_title(prs, service_title)
    if slide_idx is not None:
        return slide_idx

    src_idx = find_slide_by_title(prs, CANONICAL_DELIVERY_CLONE_TITLE)
    if src_idx is None:
        raise RuntimeError(
            f"Cannot create slide for {service_title!r}: "
            f"{CANONICAL_DELIVERY_CLONE_TITLE} template slide missing"
        )

    src_slide = prs.slides[src_idx]
    insert_after = _last_main_delivery_slide_index(prs)
    if insert_after is None:
        insert_after = src_idx

    prs.slides.add_slide(src_slide.slide_layout)
    new_idx = len(prs.slides) - 1
    new_slide = prs.slides[new_idx]
    clear_slide_shapes(new_slide)
    copy_shapes_to_slide(src_slide, new_slide)
    move_slide_after(prs, new_idx, insert_after)
    new_idx = insert_after + 1

    set_title_suffix(new_slide, service_title, ref_title_shape)
    print(f"Created delivery slide for new track: {service_title}")
    return new_idx


def ensure_delivery_slides_for_titles(
    prs,
    service_titles: set[str],
    ref_title_shape,
) -> None:
    for title in sorted(service_titles):
        ensure_delivery_slide(prs, title, ref_title_shape)


def find_slide_by_title(prs, title_fragment, exclude_contd=True):
    search_terms = [TITLE_SEARCH.get(title_fragment, title_fragment)]
    if title_fragment not in search_terms:
        search_terms.append(title_fragment)
    for search in search_terms:
        for i, slide in enumerate(prs.slides):
            title = normalize_title_text(
                next((s.text_frame.text for s in slide.shapes if s.shape_id == 2), "")
            )
            if exclude_contd and "(Contd" in title:
                continue
            if search in title and "Delivery status" in title:
                return i
    return None


def find_contd_slide(prs, title_fragment):
    indices = find_contd_slides_for_service(prs, title_fragment)
    return indices[0] if indices else None


def _find_contd_slide_with_ka(prs, title_fragment):
    """Last (Contd…) slide for a service that carries Key Activities."""
    for idx in reversed(find_contd_slides_for_service(prs, title_fragment)):
        try:
            if get_key_activities_shape(prs.slides[idx]):
                return idx
        except ValueError:
            continue
    indices = find_contd_slides_for_service(prs, title_fragment)
    return indices[-1] if indices else None


def _title_search_terms(title_fragment: str) -> list[str]:
    primary = TITLE_SEARCH.get(title_fragment, title_fragment)
    terms = [primary]
    if title_fragment not in terms:
        terms.append(title_fragment)
    return terms


def find_contd_slides_for_service(prs, title_fragment):
    """All (Contd…) slides for a service in deck order."""
    out: list[int] = []
    for search in _title_search_terms(title_fragment):
        needle = search.lower()
        for i, slide in enumerate(prs.slides):
            if i in out:
                continue
            title = normalize_title_text(
                next((s.text_frame.text for s in slide.shapes if s.shape_id == 2), "")
            ).lower()
            if needle in title and "(contd" in title:
                out.append(i)
    return out


def _create_contd_slide_after(prs, g10x_prs, after_idx, template_index):
    """Insert a new (Contd…) slide from a G10X template immediately after ``after_idx``."""
    layout = prs.slides[after_idx].slide_layout
    template_slide = g10x_prs.slides[template_index]
    prs.slides.add_slide(layout)
    new_idx = len(prs.slides) - 1
    new_slide = prs.slides[new_idx]
    clear_slide_shapes(new_slide)
    copy_shapes_to_slide(template_slide, new_slide)
    move_slide_after(prs, new_idx, after_idx)
    return after_idx + 1


def _pack_sections_to_cap(sections, cap):
    """Pack sprint sections onto one slide up to ``cap`` paragraph slots."""
    packed: list[dict] = []
    used = 0
    idx = 0
    while idx < len(sections):
        section = sections[idx]
        sec_count = _section_paragraph_count(section)
        gap = 1 if packed else 0
        if used + gap + sec_count <= cap:
            packed.append(dict(section))
            used += gap + sec_count
            idx += 1
            continue

        remaining_slots = cap - used - (1 if packed else 0)
        if remaining_slots < 4 and packed:
            return packed, [dict(s) for s in sections[idx:]]

        main_part, tail_part = _split_section_to_cap(
            section, max(remaining_slots, 4)
        )
        if main_part and _section_has_stories(main_part):
            if packed:
                used += 1
            packed.append(main_part)
        rest: list[dict] = []
        if tail_part and _section_has_stories(tail_part):
            tail_part = dict(tail_part)
            tail_part["continued_section"] = True
            rest.append(tail_part)
        rest.extend(dict(s) for s in sections[idx + 1 :])
        return packed, rest
    return packed, []


def _chunk_raw_from_contd(contd_raw, sections):
    if contd_raw.get("sections"):
        return {**contd_raw, "sections": sections}
    if len(sections) == 1:
        return sections[0]
    return {**contd_raw, "sections": sections}


def _append_hl_only_contd_slide(
    prs,
    g10x_prs,
    service_title,
    after_idx,
    contd_raw,
    ref_title,
    g10x_layout,
    section_tmpl,
    main_slide,
):
    contd_idx = _create_contd_slide_after(
        prs, g10x_prs, after_idx, HIGHLIGHTS_CONTD_TEMPLATE_INDEX
    )
    contd_slide = prs.slides[contd_idx]
    set_title_suffix(contd_slide, f"{service_title}  (Contd\u2026)", ref_title)
    hl_shape = get_highlights_shape(contd_slide)
    populate_highlights_contd_cell(
        hl_shape.table.cell(2, 0), section_tmpl, contd_raw
    )
    fit_hl_only_contd_layout(contd_slide, g10x_layout)
    sync_heb_logo_from_main(main_slide, contd_slide)
    return contd_idx


def _append_hl_ka_contd_slide(
    prs,
    g10x_prs,
    service_title,
    after_idx,
    contd_raw,
    ref_title,
    g10x_layout,
    section_tmpl,
    main_slide,
):
    contd_idx = _create_contd_slide_after(
        prs, g10x_prs, after_idx, HIGHLIGHTS_KA_CONTD_TEMPLATE_INDEX
    )
    contd_slide = prs.slides[contd_idx]
    set_title_suffix(contd_slide, f"{service_title}  (Contd\u2026)", ref_title)
    hl_shape = get_highlights_shape(contd_slide)
    populate_highlights_contd_cell(
        hl_shape.table.cell(2, 0), section_tmpl, contd_raw
    )
    move_ka_from_main_to_contd(main_slide, contd_slide, g10x_layout, g10x_prs)
    fit_contd_slide_layout(contd_slide, g10x_layout, g10x_prs)
    ensure_key_activities_tab(contd_slide, g10x_prs, g10x_layout)
    sync_heb_logo_from_main(main_slide, contd_slide)
    return contd_idx


def ensure_contd_slide_chain(
    prs,
    g10x_prs,
    service_title,
    main_idx,
    contd_raw,
    ref_title,
    g10x_layout,
    section_tmpl,
    profile,
    raw,
):
    """
    Create one or more (Contd…) slides — HL-only continuations first, then HL + KA.

    When overflow does not fit on a single (Contd…) slide, additional HL-only
    slides are inserted instead of overlapping Highlights and Key Activities.
    """
    ka_item_count = ka_layout_item_count(raw)
    chain = plan_contd_slide_chain(contd_raw, profile, ka_item_count)
    main_slide = prs.slides[main_idx]
    after_idx = main_idx

    for mode, sections in chain:
        chunk_raw = _chunk_raw_from_contd(contd_raw, sections)
        if mode == "hl_only":
            after_idx = _append_hl_only_contd_slide(
                prs,
                g10x_prs,
                service_title,
                after_idx,
                chunk_raw,
                ref_title,
                g10x_layout,
                section_tmpl,
                main_slide,
            )
        else:
            after_idx = _append_hl_ka_contd_slide(
                prs,
                g10x_prs,
                service_title,
                after_idx,
                chunk_raw,
                ref_title,
                g10x_layout,
                section_tmpl,
                main_slide,
            )
    return after_idx


def prepare_deck_from_g10x(template_path: str | None = None):
    """Copy the WSR template and remove obsolete continuation slides."""
    source = template_path or G10X
    shutil.copy2(source, OUTPUT)
    prs = Presentation(OUTPUT)
    remove_obsolete_contd_slides(prs)
    remove_pharmacy_contd_if_present(prs)
    prs.save(OUTPUT)
    return OUTPUT


def remove_pharmacy_contd_if_present(prs):
    to_delete = []
    for i, slide in enumerate(prs.slides):
        title = next((s.text_frame.text for s in slide.shapes if s.shape_id == 2), "")
        if "Pharmacy" in title and "(Contd" in title:
            to_delete.append(i)
    for idx in sorted(to_delete, reverse=True):
        delete_slide(prs, idx)


def copy_shapes_to_slide(src_slide, dst_slide):
    for shape in src_slide.shapes:
        newel = copy.deepcopy(shape.element)
        dst_slide.shapes._spTree.insert_element_before(newel, "p:extLst")


def clear_slide_shapes(slide):
    for shape in list(slide.shapes):
        delete_shape(shape)


def move_slide_after(prs, slide_index, after_index):
    sld_id_lst = prs.slides._sldIdLst
    elements = list(sld_id_lst)
    el = elements[slide_index]
    elements.pop(slide_index)
    elements.insert(after_index + 1, el)
    for child in list(sld_id_lst):
        sld_id_lst.remove(child)
    for child in elements:
        sld_id_lst.append(child)


def set_ka_items(ka_shape, items):
    from app.services.ppt_ka_bullets import replace_ka_bullet_block

    txBody = ka_shape.table.cell(1, 0).text_frame._txBody
    all_p = txBody.findall(qn("a:p"))
    if all_p and items:
        last = min(2, len(all_p) - 1)
        replace_ka_bullet_block(txBody, 0, last, items)


def ka_would_overflow(slide):
    ka = get_key_activities_shape(slide)
    if not ka:
        return False
    bottom_in = (ka.top + ka.height) / 914400
    items = sum(
        1 for p in ka.table.cell(1, 0).text_frame.paragraphs if p.text.strip()
    )
    return bottom_in > SAFE_KA_BOTTOM_IN or items > 3


def ensure_loco_contd_slide(prs, g10x_prs, loco_idx):
    """LoCo sprint highlights on main slide; key activities on (Contd...) slide."""
    main_slide = prs.slides[loco_idx]
    cost_idx = find_slide_by_title(prs, "Cost Core Service")
    ref_title = next(s for s in prs.slides[cost_idx].shapes if s.shape_id == 2)
    main_ka = get_key_activities_shape(main_slide)
    ka_items = LOCO_KA_ITEMS
    if main_ka:
        existing = [
            p.text.strip()
            for p in main_ka.table.cell(1, 0).text_frame.paragraphs
            if p.text.strip()
        ]
        if existing:
            ka_items = existing
        delete_shape(main_ka)

    contd_idx = find_contd_slide(prs, "LoCo")

    if contd_idx is None:
        template_contd = g10x_prs.slides[PHARMACY_CONTD_INDEX]
        layout = main_slide.slide_layout
        new_slide = prs.slides.add_slide(layout)
        clear_slide_shapes(new_slide)
        copy_shapes_to_slide(template_contd, new_slide)
        move_slide_after(prs, len(prs.slides) - 1, loco_idx)
        contd_idx = loco_idx + 1

    contd_slide = prs.slides[contd_idx]
    set_title_suffix(contd_slide, "LoCo  (Contd\u2026)", ref_title)
    remove_duplicate_loco_contd_slides(prs, contd_idx)

    ka = get_key_activities_shape(contd_slide)
    if not ka:
        ka = _ensure_ka_shape_on_slide(contd_slide, g10x_prs, g10x_prs.slides[CANONICAL_STYLE_SLIDE_INDEX])
    if ka:
        set_ka_items(ka, ka_items)
        loco_profile = build_layout_profile(g10x_prs.slides[12])
        fit_key_activities_table(ka, loco_profile)
        pharm_contd_ka = next(
            s
            for s in g10x_prs.slides[PHARMACY_CONTD_INDEX].shapes
            if s.has_table and "Key activities" in s.table.cell(0, 0).text
        )
        ka.top = pharm_contd_ka.top
        ka.left = pharm_contd_ka.left
        ka.width = pharm_contd_ka.width

    sync_heb_logo_from_main(main_slide, contd_slide)

    loco_main_idx = find_slide_by_title(prs, "LoCo")
    contd_idx = find_contd_slide(prs, "LoCo")
    if (
        loco_main_idx is not None
        and contd_idx is not None
        and contd_idx != loco_main_idx + 1
    ):
        move_slide_after(prs, contd_idx, loco_main_idx)


def remove_obsolete_contd_slides(prs):
    """Remove Supplier continuation slides; keep Pharmacy contd for LoCo reuse."""
    to_delete = []
    for i, slide in enumerate(prs.slides):
        title = next((s.text_frame.text for s in slide.shapes if s.shape_id == 2), "")
        if "Supplier" in title and "(Contd" in title:
            to_delete.append(i)
    for idx in sorted(to_delete, reverse=True):
        delete_slide(prs, idx)


def remove_duplicate_loco_contd_slides(prs, keep_idx):
    """Remove extra LoCo (Contd...) slides beyond the one we keep."""
    to_delete = []
    for i, slide in enumerate(prs.slides):
        if i == keep_idx:
            continue
        title = next((s.text_frame.text for s in slide.shapes if s.shape_id == 2), "")
        if "LoCo" in title and "(Contd" in title:
            to_delete.append(i)
    for idx in sorted(to_delete, reverse=True):
        delete_slide(prs, idx)


def remove_redundant_ka_only_contd(prs, service_title, main_idx):
    """Drop KA-only (Contd...) when Key Activities already fit on the main slide."""
    if not get_key_activities_shape(prs.slides[main_idx]):
        return
    contd_idx = find_contd_slide(prs, service_title)
    if contd_idx is None:
        return
    contd_slide = prs.slides[contd_idx]
    for shape in contd_slide.shapes:
        if not shape.has_table:
            continue
        try:
            if shape.table.cell(0, 0).text.strip() != "Highlights":
                continue
            if any(
                _paragraph_text(p).strip()
                for p in shape.table.cell(2, 0).text_frame._txBody.findall(qn("a:p"))
            ):
                return
        except (IndexError, AttributeError):
            continue
    delete_slide(prs, contd_idx)


def cleanup_orphan_contd_slides(prs):
    """Remove empty (Contd...) slides left from the G10X template."""
    to_delete = []
    for i, slide in enumerate(prs.slides):
        title = next((s.text_frame.text for s in slide.shapes if s.shape_id == 2), "")
        if "(Contd" not in title:
            continue
        has_hl = any(
            s.has_table and s.table.cell(0, 0).text.strip() == "Highlights"
            for s in slide.shapes
            if s.has_table
        )
        has_ka = any(
            s.has_table and "Key activities" in s.table.cell(0, 0).text
            for s in slide.shapes
            if s.has_table
        )
        if not has_hl and not has_ka:
            to_delete.append(i)
    for idx in sorted(to_delete, reverse=True):
        delete_slide(prs, idx)


def _delivery_slide_title_text(slide) -> str:
    return normalize_title_text(
        next((s.text_frame.text for s in slide.shapes if s.shape_id == 2), "")
    )


def _is_delivery_status_slide_title(title: str) -> bool:
    lower = title.lower()
    return "delivery status" in lower


def _service_title_in_slide(title: str, service_title: str) -> bool:
    lower = title.lower()
    for search in _title_search_terms(service_title):
        if search.lower() in lower:
            return True
    return False


def remove_unpopulated_delivery_slides(prs, populated_titles: set[str]) -> int:
    """
    Remove delivery-status main slides (and their (Contd…) chain) that are not
    in this week's ``ppt_content.json``. Keeps Index, Matters of Attention, etc.
    """
    populated = {normalize_title_text(t) for t in populated_titles}
    to_delete: list[int] = []

    for i, slide in enumerate(prs.slides):
        title = _delivery_slide_title_text(slide)
        if not _is_delivery_status_slide_title(title):
            continue
        if "(contd" in title.lower():
            if not any(_service_title_in_slide(title, pop) for pop in populated):
                to_delete.append(i)
            continue
        if not any(_service_title_in_slide(title, pop) for pop in populated):
            to_delete.append(i)

    removed = 0
    for idx in sorted(set(to_delete), reverse=True):
        delete_slide(prs, idx)
        removed += 1
    return removed


def finalize_slide_order(prs, service_titles: list[str] | None = None):
    """Place each service's (Contd…) slide(s) immediately after its main delivery slide."""
    titles = service_titles or [raw["title"] for raw in SLIDES.values()]
    for title in reversed(titles):
        main_idx = find_slide_by_title(prs, title)
        if main_idx is None:
            continue
        contd_indices = find_contd_slides_for_service(prs, title)
        for offset in range(len(contd_indices)):
            main_idx = find_slide_by_title(prs, title)
            contd_indices = find_contd_slides_for_service(prs, title)
            if offset >= len(contd_indices):
                break
            contd_idx = contd_indices[offset]
            desired = main_idx + 1 + offset
            if contd_idx != desired:
                move_slide_after(prs, contd_idx, desired - 1)


def delete_slide(presentation, index):
    slide_id_list = presentation.slides._sldIdLst
    slide_ids = list(slide_id_list)
    presentation.part.drop_rel(slide_ids[index].rId)
    slide_id_list.remove(slide_ids[index])


# Index slide: map cell text fragments -> deck title search (delivery-status slides first).
INDEX_ENTRY_RULES: tuple[tuple[tuple[str, ...], str, bool], ...] = (
    (("cost core",), "Cost Core Service", True),
    (("supplier core",), "Supplier Core Service", True),
    (("pricing core",), "Pricing Core Service", True),
    (("wentworth",), "Wentworth", True),
    (("location core",), "Location Core Service", True),
    (("pharmacy", "wellness"), "Pharmacy and Wellness", True),
    (("global sourcing",), "Global Sourcing Solution", True),
    (("loco",), "LoCo", True),
    (("bsa",), "LoCo", True),
    (("matters", "attention"), "Matters Of Attention", False),
    (("team allocation",), "Team Allocation", False),
)

_INDEX_NUMBER_RE = re.compile(r"^\d{1,2}\u200b?$")


def _normalize_index_text(text: str) -> str:
    return normalize_title_text(text).lower()


def _cell_full_text(cell) -> str:
    return _normalize_index_text(cell.text_frame.text)


def _find_slide_by_plain_title(prs, title_fragment: str) -> int | None:
    needle = title_fragment.lower()
    for i, slide in enumerate(prs.slides):
        title = normalize_title_text(
            next((s.text_frame.text for s in slide.shapes if s.shape_id == 2), "")
        ).lower()
        if needle in title:
            return i
    return None


def _resolve_index_target_slide(prs, cell_text: str) -> int | None:
    """Map an index table cell to the 0-based slide index in the current deck."""
    for needles, search_title, delivery_only in INDEX_ENTRY_RULES:
        if not all(n in cell_text for n in needles):
            continue
        if delivery_only:
            idx = find_slide_by_title(prs, search_title)
        else:
            idx = _find_slide_by_plain_title(prs, search_title)
        if idx is not None:
            return idx
    return None


def _find_index_slide_index(prs) -> int | None:
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_id == 2 and shape.has_text_frame:
                if normalize_title_text(shape.text_frame.text) == "Index":
                    return i
    return None


def _find_index_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def _is_index_number_paragraph(paragraph) -> bool:
    text = normalize_title_text("".join(run.text for run in paragraph.runs))
    return bool(_INDEX_NUMBER_RE.match(text))


def _set_paragraph_slide_number(paragraph, slide_number: int) -> None:
    display = f"{slide_number:02d}"
    if paragraph.runs:
        paragraph.runs[0].text = display
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = display


def _set_element_hyperlink_to_slide(hlink_elem, index_slide, target_slide_idx, prs) -> None:
    target_part = prs.slides[target_slide_idx].part
    slide_rid = index_slide.part.relate_to(target_part, RT.SLIDE)
    hlink_elem.set(qn("r:id"), slide_rid)


def _update_cell_index_entry(cell, index_slide, target_slide_idx, prs) -> bool:
    """Update displayed slide number and internal hyperlinks for one index cell."""
    changed = False
    display_num = target_slide_idx + 1

    for paragraph in cell.text_frame.paragraphs:
        if _is_index_number_paragraph(paragraph):
            old = normalize_title_text("".join(run.text for run in paragraph.runs))
            new = f"{display_num:02d}"
            if old != new:
                _set_paragraph_slide_number(paragraph, display_num)
                changed = True

    tx_body = cell.text_frame._txBody
    for p_elem in tx_body.findall(qn("a:p")):
        for hlink in p_elem.iter(qn("a:hlinkClick")):
            _set_element_hyperlink_to_slide(hlink, index_slide, target_slide_idx, prs)
            changed = True

    return changed


def _index_table_cells_row_major(table):
    """Return table cells in reading order (left-to-right, top-to-bottom)."""
    cells = []
    for row in table.rows:
        cells.extend(row.cells)
    return cells


def _cell_has_index_content(cell) -> bool:
    """True when a table cell carries an index label (not just empty padding)."""
    text = normalize_title_text(cell.text_frame.text)
    if not text:
        return False
    cell_text = _cell_full_text(cell)
    for needles, _search_title, _delivery_only in INDEX_ENTRY_RULES:
        if all(n in cell_text for n in needles):
            return True
    for paragraph in cell.text_frame.paragraphs:
        if _is_index_number_paragraph(paragraph):
            return True
    return False


def _clear_index_cell_completely(cell) -> None:
    """Remove all visible text and slide hyperlinks from an index cell."""
    tx_body = cell.text_frame._txBody
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""
    for p_elem in tx_body.findall(qn("a:p")):
        for hlink in list(p_elem.findall(qn("a:hlinkClick"))):
            p_elem.remove(hlink)


def _clone_cell_text_body(dst_cell, src_tx_body) -> None:
    """Replace destination cell body with a deep copy of source ``a:txBody``."""
    dst_tc = dst_cell._tc
    old_tx = dst_tc.find(qn("a:txBody"))
    if old_tx is not None:
        dst_tc.remove(old_tx)
    dst_tc.insert(0, copy.deepcopy(src_tx_body))


def _default_index_tx_body(table) -> object | None:
    for cell in _index_table_cells_row_major(table):
        if "cost core" in _cell_full_text(cell):
            return copy.deepcopy(cell.text_frame._txBody)
    for cell in _index_table_cells_row_major(table):
        if _cell_has_index_content(cell):
            return copy.deepcopy(cell.text_frame._txBody)
    return None


def _set_index_label_text(tx_body, label: str) -> None:
    clean = normalize_title_text(label)
    for p_elem in tx_body.findall(qn("a:p")):
        texts = [t.text or "" for t in p_elem.iter(qn("a:t"))]
        joined = normalize_title_text("".join(texts))
        if joined and not _INDEX_NUMBER_RE.match(joined):
            set_single_run_text(p_elem, clean)
            return


def _find_index_cell_tx_body(table, needles: tuple[str, ...]) -> object | None:
    """Original index cell body for a label fragment (e.g. matters of attention)."""
    for cell in _index_table_cells_row_major(table):
        if not _cell_has_index_content(cell):
            continue
        cell_text = _cell_full_text(cell)
        if all(n in cell_text for n in needles):
            return copy.deepcopy(cell.text_frame._txBody)
    return None


def _non_delivery_slots_from_template(table) -> dict[tuple[str, ...], int]:
    """Map non-delivery INDEX_ENTRY_RULES needles to their G10X template cell index."""
    slots: dict[tuple[str, ...], int] = {}
    for slot_idx, cell in enumerate(_index_table_cells_row_major(table)):
        cell_text = _cell_full_text(cell)
        for needles, _search_title, delivery_only in INDEX_ENTRY_RULES:
            if delivery_only:
                continue
            if all(n in cell_text for n in needles):
                slots[tuple(needles)] = slot_idx
    return slots


def _find_non_delivery_slide_index(
    prs,
    needles: tuple[str, ...],
    search_title: str,
) -> int | None:
    target_idx = _find_slide_by_plain_title(prs, search_title)
    if target_idx is not None:
        return target_idx
    for fragment in needles:
        target_idx = _find_slide_by_plain_title(prs, fragment)
        if target_idx is not None:
            return target_idx
    return None


def _visible_non_delivery_slots(
    table,
    template_slots: dict[tuple[str, ...], int],
    delivery_count: int,
) -> dict[tuple[str, ...], int]:
    """
    When delivery tracks fit the first row, place MOA / Team Allocation on row 2
    (cells 3–4) so they appear directly below the delivery row instead of the
    deep template slots (8 / 10) that sit near the footer.
    """
    n_cols = len(table.columns)
    if delivery_count > n_cols:
        return template_slots

    placements: dict[tuple[str, ...], int] = {}
    next_slot = delivery_count
    for needles, _, delivery_only in INDEX_ENTRY_RULES:
        if delivery_only:
            continue
        key = tuple(needles)
        if key not in template_slots:
            continue
        while next_slot in template_slots.values():
            next_slot += 1
        placements[key] = next_slot
        next_slot += 1
    return placements


def _index_slot_assignments_for_populated_week(
    prs,
    table,
    populated_titles: set[str],
) -> list[tuple[int, int, object]]:
    """
    Index placements for delivery tracks (compacted top cells) plus Matters of
    Attention and Team Allocation (visible row when few tracks, else G10X slots).
    """
    template_tx = _default_index_tx_body(table)
    if template_tx is None:
        entries = _collect_active_index_entries(prs, table)
        return [
            (slot_idx, target_idx, tx_body)
            for slot_idx, (target_idx, tx_body) in enumerate(entries)
        ]

    non_delivery_slots = _non_delivery_slots_from_template(table)
    reserved_slots = set(non_delivery_slots.values())

    sorted_titles = sorted(
        populated_titles,
        key=lambda title: find_slide_by_title(prs, title) or 10_000,
    )
    assignments: list[tuple[int, int, object]] = []
    seen_targets: set[int] = set()
    next_slot = 0

    for title in sorted_titles:
        target_idx = find_slide_by_title(prs, title)
        if target_idx is None or target_idx in seen_targets:
            continue
        while next_slot in reserved_slots:
            next_slot += 1
        slot_idx = next_slot
        seen_targets.add(target_idx)
        next_slot += 1
        tx_body = copy.deepcopy(template_tx)
        _set_index_label_text(tx_body, title)
        assignments.append((slot_idx, target_idx, tx_body))

    visible_slots = _visible_non_delivery_slots(
        table, non_delivery_slots, len(assignments)
    )

    for needles, search_title, delivery_only in INDEX_ENTRY_RULES:
        if delivery_only:
            continue
        key = tuple(needles)
        slot_idx = visible_slots.get(key) or non_delivery_slots.get(key)
        if slot_idx is None:
            continue
        target_idx = _find_non_delivery_slide_index(prs, needles, search_title)
        if target_idx is None or target_idx in seen_targets:
            continue
        tx_body = _find_index_cell_tx_body(table, needles)
        if tx_body is None and template_tx is not None:
            tx_body = copy.deepcopy(template_tx)
            _set_index_label_text(tx_body, search_title)
        if tx_body is None:
            continue
        seen_targets.add(target_idx)
        assignments.append((slot_idx, target_idx, tx_body))

    return assignments


def _collect_active_index_entries(prs, table) -> list[tuple[int, object]]:
    """
    Index entries that still resolve to a slide in the deck, in template order.

    Returns ``(target_slide_idx, txBody_element)`` pairs. Skips removed projects
    (e.g. LoCo with no delivery slide) and duplicate targets.
    """
    entries: list[tuple[int, object]] = []
    seen_targets: set[int] = set()
    for cell in _index_table_cells_row_major(table):
        if not _cell_has_index_content(cell):
            continue
        cell_text = _cell_full_text(cell)
        target_idx = _resolve_index_target_slide(prs, cell_text)
        if target_idx is None or target_idx in seen_targets:
            continue
        seen_targets.add(target_idx)
        entries.append((target_idx, copy.deepcopy(cell.text_frame._txBody)))
    return entries


def reflow_index_slide(prs, populated_titles: set[str] | None = None) -> int:
    """
    Drop index rows for missing projects and compact remaining entries.

    Projects without a delivery slide (or other mapped slide) are removed entirely
    from the Index table — label and number — and surviving entries shift left/up
    with no empty placeholder cells.

    When ``populated_titles`` is set (WSR JSON mode), index rows are rebuilt for
    every populated delivery slide, including dynamically added team tracks.
    Matters of Attention and Team Allocation keep their G10X template cell slots.
    """
    index_idx = _find_index_slide_index(prs)
    if index_idx is None:
        return 0

    index_slide = prs.slides[index_idx]
    table = _find_index_table(index_slide)
    if table is None:
        return 0

    if populated_titles:
        assignments = _index_slot_assignments_for_populated_week(
            prs, table, populated_titles
        )
    else:
        active_entries = _collect_active_index_entries(prs, table)
        assignments = [
            (slot_idx, target_idx, tx_body)
            for slot_idx, (target_idx, tx_body) in enumerate(active_entries)
        ]
    slot_cells = _index_table_cells_row_major(table)

    for cell in slot_cells:
        if _cell_has_index_content(cell):
            _clear_index_cell_completely(cell)

    updated = 0
    for slot_idx, target_idx, tx_body in assignments:
        if slot_idx >= len(slot_cells):
            break
        cell = slot_cells[slot_idx]
        _clone_cell_text_body(cell, tx_body)
        if _update_cell_index_entry(cell, index_slide, target_idx, prs):
            updated += 1

    return updated


def sync_index_slide_numbers(prs, populated_titles: set[str] | None = None) -> int:
    """
    Reflow the Index slide: remove missing projects and sync numbers/hyperlinks.

    Call after ``finalize_slide_order()`` once unpopulated delivery slides are removed.
    """
    return reflow_index_slide(prs, populated_titles)


def main():
    args = parse_args()
    global OUTPUT
    OUTPUT = args.output
    load_layout_hints(args.layout_hints or None)

    slides_data = SLIDES
    using_json = False
    wsr_end_date = args.wsr_end_date.strip() or None
    if args.content and Path(args.content).is_file():
        with open(args.content, encoding="utf-8") as f:
            content_payload = json.load(f)
        slides_data = load_slides_from_json(args.content)
        if not wsr_end_date:
            wsr_end_date = content_payload.get("report_end_date")
        using_json = True
        print(f"Loaded {len(slides_data)} slide(s) from {args.content}")

    deck_path = prepare_deck_from_g10x(G10X)
    prs = Presentation(deck_path)
    g10x = Presentation(G10X)

    cost_idx = find_slide_by_title(prs, "Cost Core Service")
    if cost_idx is None:
        raise RuntimeError("Cost Core Service slide not found in template")
    ref_slide = prs.slides[cost_idx]
    ref_title = next(s for s in prs.slides[cost_idx].shapes if s.shape_id == 2)

    populated_titles = {raw["title"] for raw in slides_data.values()}
    if using_json and populated_titles:
        ensure_delivery_slides_for_titles(prs, populated_titles, ref_title)

    completed_hdr, completed_bullet = find_completed_templates()
    canonical_style_cell = get_canonical_style_cell(g10x)

    updated = 0
    canonical_fill = get_canonical_fill_para_count(g10x)
    canonical_slots = get_canonical_para_slot_count(g10x)
    canonical_line_height = get_canonical_line_height_emu(g10x)
    for _key, raw in slides_data.items():
        slide_idx = find_slide_by_title(prs, raw["title"])
        if slide_idx is None:
            print(f"Warning: slide not found for {raw['title']} (skipped)")
            continue
        g10x_layout = get_g10x_layout_slide(g10x, raw["title"])
        section_tmpl = discover_section_templates(
            canonical_style_cell, (completed_hdr, completed_bullet)
        )
        profile = build_layout_profile(g10x_layout)
        profile["canonical_fill_para_count"] = canonical_fill
        profile["canonical_para_slot_count"] = canonical_slots
        profile["canonical_line_height_emu"] = canonical_line_height
        budgets = get_hl_budgets(profile)
        story_raw = _flatten_story_raw(raw)
        if raw.get("sections"):
            main_raw, contd_raw, layout_mode, budgets = plan_sections_split(
                raw, profile, layout_hints=LAYOUT_HINTS
            )
        else:
            split_input = {
                "title": story_raw["title"],
                "sprint_name": story_raw["sprint_name"],
                "sprint_dates": story_raw["sprint_dates"],
                "sprint_status": story_raw.get("sprint_status", "In-progress"),
                "released": story_raw["released"],
                "inprogress": story_raw["inprogress"],
                "completed": story_raw["completed"],
            }
            main_raw, contd_raw, layout_mode, budgets = plan_highlights_split(
                split_input, profile
            )

        content = build_content(main_raw, totals_raw=story_raw)
        slide = prs.slides[slide_idx]
        set_title_suffix(slide, content["title"], ref_title)
        hl_shape = get_highlights_shape(slide)
        populate_highlights_cell(
            hl_shape.table.cell(2, 0),
            section_tmpl,
            content,
        )
        ka_item_count = ka_layout_item_count(raw)
        hl_overflow_contd = layout_mode in ("hl_ka_contd", "supplier_contd")

        if hl_overflow_contd:
            ensure_highlights_contd_slide(
                prs,
                g10x,
                raw["title"],
                slide_idx,
                contd_raw,
                ref_title,
                g10x_layout,
                section_tmpl,
                profile=profile,
                raw=raw,
                layout_mode=layout_mode,
            )
            main_ka = get_key_activities_shape(slide)
            if main_ka:
                delete_shape(main_ka)

        ka_on_main = False
        if contd_raw is None and not hl_overflow_contd:
            if apply_ka_on_main_slide(
                slide, profile, raw, g10x, g10x_layout, budgets
            ):
                ka_on_main = True
                layout_mode = "ka_on_main"
                remove_redundant_ka_only_contd(prs, raw["title"], slide_idx)
            else:
                main_ka = get_key_activities_shape(slide)
                if main_ka:
                    delete_shape(main_ka)

        if not ka_on_main:
            if layout_mode in ("normal", "expanded"):
                fit_mode = layout_mode
            elif layout_mode in ("hl_ka_contd", "supplier_contd", "ka_contd_only"):
                fit_mode = "full_main"
            else:
                fit_mode = "normal"
            fit_highlights_layout(
                slide,
                g10x_layout,
                position_ka=False,
                layout_mode=fit_mode,
                budgets=budgets,
                profile=profile,
            )

        if not ka_on_main and contd_raw is None and not hl_overflow_contd:
            hl_shape = get_highlights_shape(slide)
            ka_profile = get_ka_layout_profile(g10x, g10x_layout)
            _ensure_hl_table_fits_rendered_text(hl_shape, profile)
            if not ka_fits_below_highlights(
                hl_shape, profile, ka_item_count, ka_profile=ka_profile
            ):
                layout_mode = "ka_contd_only"
                ensure_ka_contd_slide(
                    prs, g10x, raw["title"], slide_idx, ref_title, g10x_layout
                )
                main_slide = prs.slides[slide_idx]
                main_ka = get_key_activities_shape(main_slide)
                if main_ka:
                    delete_shape(main_ka)
                hl_shape = get_highlights_shape(slide)
                if _is_dense_hl(hl_shape, profile):
                    fit_highlights_table(
                        hl_shape, profile, layout_mode="full_main", budgets=budgets
                    )
                else:
                    _fit_hl_content_only(
                        hl_shape, profile, max_h=budgets["full_main_h"]
                    )

        finalize_key_activities(
            prs, raw["title"], slide_idx, raw, g10x_layout, layout_mode
        )

        # Guarantee KA tab on main or (Contd...) for every populated project.
        ka_contd_idx = _find_contd_slide_with_ka(prs, raw["title"])
        if layout_mode in ("hl_ka_contd", "supplier_contd", "ka_contd_only"):
            if ka_contd_idx is not None:
                ensure_key_activities_tab(
                    prs.slides[ka_contd_idx], g10x, g10x_layout, profile
                )
        elif ka_contd_idx is None and layout_mode != "ka_on_main":
            ensure_key_activities_tab(slide, g10x, g10x_layout, profile)

        tighten_slide_hl_ka_spacing(
            slide, g10x_layout, g10x, budgets=budgets, profile=profile
        )
        for contd_idx in find_contd_slides_for_service(prs, raw["title"]):
            tighten_slide_hl_ka_spacing(
                prs.slides[contd_idx],
                g10x_layout,
                g10x,
                budgets=budgets,
                profile=profile,
            )
        updated += 1

    populated_titles = {raw["title"] for raw in slides_data.values()}
    if using_json:
        removed = remove_unpopulated_delivery_slides(prs, populated_titles)
        if removed:
            print(f"Removed {removed} unpopulated delivery-status slide(s)")

    service_titles = [raw["title"] for raw in slides_data.values()]
    finalize_slide_order(prs, service_titles)
    cleanup_orphan_contd_slides(prs)
    index_updates = sync_index_slide_numbers(
        prs,
        populated_titles if using_json else None,
    )
    if index_updates:
        print(
            f"Reflowed index slide: {index_updates} entr"
            f"{'y' if index_updates == 1 else 'ies'} (removed projects without content)"
        )
        index_idx = _find_index_slide_index(prs)
        if index_idx is not None:
            table = _find_index_table(prs.slides[index_idx])
            if table is not None:
                labels = []
                for j, cell in enumerate(_index_table_cells_row_major(table)):
                    t = _cell_full_text(cell)
                    if t:
                        labels.append(f"[{j}] {t.replace(chr(10), ' / ')}")
                if labels:
                    print(f"Index entries: {' | '.join(labels)}")

    if wsr_end_date:
        if sync_cover_slide_wsr_date(prs, wsr_end_date):
            print(
                f"Updated cover slide WSR date -> "
                f"{format_wsr_cover_date(date.fromisoformat(wsr_end_date))}"
            )
        else:
            print("Warning: cover slide date placeholder not found")

    out = deck_path
    try:
        prs.save(out)
    except PermissionError:
        out = deck_path.replace(".pptx", "_refreshed.pptx")
        prs.save(out)
    print(f"Built deck with {updated} delivery-status slides -> {out}")


if __name__ == "__main__":
    main()
