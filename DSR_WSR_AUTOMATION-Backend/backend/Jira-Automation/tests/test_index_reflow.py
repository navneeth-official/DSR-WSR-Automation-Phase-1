"""Unit tests for index-slide reflow helpers."""

from types import SimpleNamespace

from scripts.update_delivery_status import INDEX_ENTRY_RULES


def test_index_entry_rules_cover_all_delivery_services():
    delivery_titles = {
        rule[1]
        for rule in INDEX_ENTRY_RULES
        if rule[2]
    }
    expected = {
        "Cost Core Service",
        "Supplier Core Service",
        "Pricing Core Service",
        "Wentworth",
        "Location Core Service",
        "Pharmacy and Wellness",
        "Global Sourcing Solution",
        "LoCo",
    }
    assert expected <= delivery_titles


def test_non_delivery_index_slots_match_g10x_template():
    from pptx import Presentation

    from scripts.update_delivery_status import (
        G10X,
        _find_index_table,
        _index_table_cells_row_major,
        _non_delivery_slots_from_template,
        normalize_title_text,
    )

    prs = Presentation(G10X)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_id == 2 and normalize_title_text(shape.text_frame.text) == "Index":
                table = _find_index_table(slide)
                slots = _non_delivery_slots_from_template(table)
                assert slots.get(("matters", "attention")) == 8
                assert slots.get(("team allocation",)) == 10
                cells = _index_table_cells_row_major(table)
                assert len(cells) >= 11
                return
    raise AssertionError("Index slide not found in G10X template")


def test_visible_non_delivery_slots_when_few_delivery_tracks():
    from pptx import Presentation

    from scripts.update_delivery_status import (
        G10X,
        _find_index_table,
        _non_delivery_slots_from_template,
        _visible_non_delivery_slots,
        normalize_title_text,
    )

    prs = Presentation(G10X)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_id == 2 and normalize_title_text(shape.text_frame.text) == "Index":
                table = _find_index_table(slide)
                template_slots = _non_delivery_slots_from_template(table)
                visible = _visible_non_delivery_slots(table, template_slots, 3)
                assert visible.get(("matters", "attention")) == 3
                assert visible.get(("team allocation",)) == 4
                return
    raise AssertionError("Index slide not found in G10X template")


def test_cell_has_index_content_detects_loco_label():
    from scripts.update_delivery_status import _cell_has_index_content

    cell = SimpleNamespace(
        text_frame=SimpleNamespace(
            text="LoCo (BSA)",
            paragraphs=[],
        )
    )
    assert _cell_has_index_content(cell)
