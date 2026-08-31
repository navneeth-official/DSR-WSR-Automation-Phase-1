"""Tests for KA-only continuation slide layout."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from app.services.ppt_layout_metrics import EMU_PER_INCH
from app.services.ppt_shape_utils import get_key_activities_shape
from app.services.template_profile import scan_template_profile
from app.wsr_engine.continuation_engine import position_ka_at_contd_top
from app.wsr_engine.hl_ka_normalizer import find_ka_only_contd_template_index
from app.wsr_engine.models import TitleFormat

SUSTAINMENT = (
    Path(__file__).resolve().parents[2]
    / "G10X H-E-B WSR Sustainment 05 June 2026 .pptx"
)


@pytest.mark.skipif(not SUSTAINMENT.is_file(), reason="Sustainment template missing")
def test_find_ka_only_contd_template():
    prs = Presentation(str(SUSTAINMENT))
    idx = find_ka_only_contd_template_index(prs)
    assert idx is not None
    ka = get_key_activities_shape(prs.slides[idx])
    assert ka is not None
    assert ka.top / EMU_PER_INCH < 2.0


@pytest.mark.skipif(not SUSTAINMENT.is_file(), reason="Sustainment template missing")
def test_position_ka_at_contd_top_moves_ka_up():
    from app.wsr_engine.continuation_engine import ensure_ka_only_contd_slide

    prs = Presentation(str(SUSTAINMENT))
    template_prs = Presentation(str(SUSTAINMENT))
    profile = scan_template_profile(SUSTAINMENT)
    svc = profile.services["Wentworth"]
    main_idx = svc.main_slide_index
    main_ka = get_key_activities_shape(prs.slides[main_idx])
    assert main_ka is not None
    main_ka_top = main_ka.top

    title_format = TitleFormat(prefix="Delivery status", separator=" - ", contd_marker="(Contd..)")
    contd_indices, _ = ensure_ka_only_contd_slide(
        prs,
        template_prs,
        "Wentworth",
        main_idx,
        svc,
        main_idx,
        title_format,
    )
    contd_ka = get_key_activities_shape(prs.slides[contd_indices[0]])
    ref_idx = find_ka_only_contd_template_index(template_prs)
    ref_ka = get_key_activities_shape(template_prs.slides[ref_idx])

    assert contd_ka.top == ref_ka.top
    assert contd_ka.top < main_ka_top


@pytest.mark.skipif(not SUSTAINMENT.is_file(), reason="Sustainment template missing")
def test_ka_contd_only_main_shrinks_and_top_aligns_hl():
    """Main slide with KA on contd only: HL fits content and text is top-aligned."""
    from pptx.oxml.ns import qn

    from app.services.ppt_shape_utils import get_highlights_shape
    from app.wsr_engine.content_parser import load_content, section_display_content
    from app.wsr_engine.formatter import discover_section_templates, fill_highlights_on_slide, get_canonical_style_cell
    from app.wsr_engine.ka_layout import finalize_project_hl_ka
    from app.wsr_engine.overflow_engine import build_layout_profile, plan_overflow
    from app.wsr_engine.project_matcher import match_projects
    from app.wsr_engine.template_analyzer import analyze_template
    from app.wsr_engine.continuation_engine import ensure_ka_only_contd_slide

    content_path = Path(__file__).resolve().parents[2] / "output" / "ppt_content.json"
    if not content_path.is_file():
        pytest.skip("ppt_content.json missing")

    prs = Presentation(str(SUSTAINMENT))
    template_prs = Presentation(str(SUSTAINMENT))
    profile = scan_template_profile(SUSTAINMENT)
    svc = profile.services["Cost Core Service"]
    main_idx = svc.main_slide_index
    main_slide = prs.slides[main_idx]
    template_slide = template_prs.slides[main_idx]

    content = load_content(content_path)
    template = analyze_template(SUSTAINMENT)
    project = match_projects(template, content.projects, None)["Cost Core Service"]
    layout_profile = build_layout_profile(template_slide)
    overflow = plan_overflow(project, layout_profile)
    assert overflow.ka_contd_only

    tmpl = discover_section_templates(get_canonical_style_cell(template_prs, profile))
    fill_highlights_on_slide(main_slide, tmpl, overflow.main_sections)

    title_format = TitleFormat(prefix="Delivery status", separator=" - ", contd_marker="(Contd..)")
    contd_indices, _ = ensure_ka_only_contd_slide(
        prs,
        template_prs,
        "Cost Core Service",
        main_idx,
        svc,
        main_idx,
        title_format,
    )

    template_hl_h = get_highlights_shape(template_slide).height
    finalize_project_hl_ka(
        prs,
        main_idx,
        contd_indices,
        template_prs,
        svc,
        profile,
        layout_profile=layout_profile,
        ka_contd_only=True,
    )

    hl = get_highlights_shape(main_slide)
    cell = hl.table.cell(2, 0)
    tcPr = cell._tc.find(qn("a:tcPr"))
    assert tcPr is not None and tcPr.get("anchor") == "t"
    assert hl.height < template_hl_h
    assert hl.height / EMU_PER_INCH < template_hl_h / EMU_PER_INCH - 0.3
