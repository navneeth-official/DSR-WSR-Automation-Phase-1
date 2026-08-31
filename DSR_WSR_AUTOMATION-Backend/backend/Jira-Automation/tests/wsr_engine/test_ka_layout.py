"""Tests for KA tab insertion on ka_mode=none slides."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from app.services.ppt_shape_utils import get_highlights_shape, get_key_activities_shape, slide_title_text
from app.services.template_profile import scan_template_profile
from app.wsr_engine.content_parser import load_content, section_display_content
from app.wsr_engine.formatter import discover_section_templates, fill_highlights_on_slide, get_canonical_style_cell
from app.services.ppt_layout_metrics import apply_fixed_story_line_metrics, FOOTER_MAX_BOTTOM_EMU, hl_ka_tab_gap_emu
from app.wsr_engine.ka_layout import (
    effective_ka_top,
    ensure_key_activities_on_main_slide,
    finalize_project_hl_ka,
    fit_hl_only_on_slide,
    ka_overlaps_footer,
    ka_overlaps_hl_text,
    rendered_hl_text_bottom,
)
from app.wsr_engine.overflow_engine import build_layout_profile, plan_overflow, build_overflow_profile
from app.wsr_engine.project_matcher import match_projects
from app.wsr_engine.template_analyzer import analyze_template

HASKELL = (
    Path(__file__).resolve().parents[2]
    / "templates"
    / "G10X H-E-B WSR Haskell Location Pharmacy GSS PAM 01 Aug 2025.pptx"
)
CONTENT = Path(__file__).resolve().parents[2] / "output" / "ppt_content.json"


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_ensure_ka_inserts_tab_on_cost_core():
    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    svc = profile.services["Cost Core Service"]
    slide = prs.slides[svc.main_slide_index]
    assert get_key_activities_shape(slide) is None
    ensure_key_activities_on_main_slide(slide, prs, svc, profile)
    assert get_key_activities_shape(slide) is not None


@pytest.mark.skipif(
    not HASKELL.is_file() or not CONTENT.is_file(),
    reason="Haskell template or ppt_content.json missing",
)
def test_ensure_ka_sits_at_fixed_tab_gap_on_dense_supplier_slide():
    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    svc = profile.services["Supplier Core Service"]
    slide = prs.slides[svc.main_slide_index]

    content = load_content(CONTENT)
    template = analyze_template(HASKELL)
    supplier = match_projects(template, content.projects, None)["Supplier Core Service"]
    sections = [section_display_content(s) for s in supplier.sections[:1]]
    tmpl = discover_section_templates(get_canonical_style_cell(prs, profile))
    fill_highlights_on_slide(slide, tmpl, sections)
    ensure_key_activities_on_main_slide(slide, prs, svc, profile)

    hl = get_highlights_shape(slide)
    ka = get_key_activities_shape(slide)
    assert ka is not None
    template_slide = prs.slides[svc.main_slide_index]
    layout = apply_fixed_story_line_metrics(build_layout_profile(template_slide))
    layout.setdefault("canonical_para_count", layout["ref_para_count"])
    assert not ka_overlaps_hl_text(slide, layout)
    assert not ka_overlaps_footer(slide)
    gap = hl_ka_tab_gap_emu(layout.get("canonical_line_height_emu"))
    assert ka.top - rendered_hl_text_bottom(hl, layout) >= gap - 1000


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_effective_ka_top_uses_rendered_text_bottom():
    """KA placement follows rendered HL text + tab gap."""
    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    svc = profile.services["Supplier Core Service"]
    slide = prs.slides[svc.main_slide_index]
    hl = get_highlights_shape(slide)
    layout = apply_fixed_story_line_metrics(build_layout_profile(slide))
    ka_h = 900000
    gap = hl_ka_tab_gap_emu(layout.get("canonical_line_height_emu"))
    expected = rendered_hl_text_bottom(hl, layout) + gap
    assert effective_ka_top(hl, layout, ka_h) == min(expected, FOOTER_MAX_BOTTOM_EMU - ka_h)


@pytest.mark.skipif(
    not HASKELL.is_file() or not CONTENT.is_file(),
    reason="Haskell template or ppt_content.json missing",
)
def test_ka_only_on_last_slide_when_overflow():
    """Cost Core overflow: main slide has no KA; last contd slide has KA."""
    from app.wsr_engine.continuation_engine import ensure_continuation_slides
    from app.wsr_engine.overflow_engine import main_slide_capacity
    from app.wsr_engine.models import TitleFormat

    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    svc = profile.services["Cost Core Service"]
    main_idx = svc.main_slide_index
    main_slide = prs.slides[main_idx]
    template_slide = prs.slides[main_idx]

    content = load_content(CONTENT)
    template = analyze_template(HASKELL)
    cost = match_projects(template, content.projects, None)["Cost Core Service"]
    tmpl = discover_section_templates(get_canonical_style_cell(prs, profile))

    layout_profile = build_overflow_profile(prs, template_slide, svc, profile)
    overflow = plan_overflow(cost, layout_profile, main_cap=main_slide_capacity(layout_profile))
    fill_highlights_on_slide(main_slide, tmpl, overflow.main_sections)

    contd_indices: list[int] = []
    title_format = TitleFormat(prefix="Delivery Status", separator=" - ", contd_marker="(Contd..)")
    if overflow.continuation_chains:
        contd_indices, _ = ensure_continuation_slides(
            prs,
            prs,
            "Cost Core Service",
            main_idx,
            overflow.continuation_chains,
            svc,
            main_idx,
            title_format,
            tmpl,
        )

    finalize_project_hl_ka(
        prs,
        main_idx,
        contd_indices,
        prs,
        svc,
        profile,
        layout_profile=layout_profile,
        project_name="Cost Core Service",
        title_format=title_format,
    )

    assert overflow.continuation_chains, "Cost Core fixture should overflow to contd"
    assert get_key_activities_shape(main_slide) is None
    ka_contd_indices = [
        i
        for i in contd_indices
        if get_key_activities_shape(prs.slides[i]) is not None
    ]
    if not ka_contd_indices:
        ka_contd_indices = [
            i
            for i, slide in enumerate(prs.slides)
            if "(contd" in slide_title_text(slide).lower()
            and "cost core" in slide_title_text(slide).lower()
            and get_key_activities_shape(slide) is not None
        ]
    assert ka_contd_indices, "Cost Core should have KA on a (Contd..) slide"
    last_contd = prs.slides[contd_indices[-1]]
    if get_key_activities_shape(last_contd) is None:
        last_contd = prs.slides[ka_contd_indices[-1]]
    assert "(contd" in slide_title_text(last_contd).lower()

    hl = get_highlights_shape(last_contd)
    fit_layout = dict(layout_profile)
    fit_layout.setdefault("canonical_para_count", fit_layout["ref_para_count"])
    fit_layout.setdefault(
        "canonical_per_line_emu",
        int(fit_layout["ref_r2"] / max(fit_layout["ref_para_count"] - 2, 1)),
    )
    hl_h_in = hl.height / 914400
    assert hl_h_in < 3.5, f"contd HL should be content-sized, got {hl_h_in:.2f}in"
    assert not ka_overlaps_hl_text(last_contd, fit_layout)


@pytest.mark.skipif(
    not HASKELL.is_file() or not CONTENT.is_file(),
    reason="Haskell template or ppt_content.json missing",
)
def test_supplier_dense_overflow_splits_ka_to_dedicated_contd(tmp_path):
    """Full v2 build: dense Supplier overflow must split HL and KA contd slides."""
    from app.wsr_engine.main import WsrEngine

    out = tmp_path / "supplier_build.pptx"
    report = WsrEngine().run(HASKELL, CONTENT, out)
    assert not report.errors

    prs = Presentation(str(out))
    supplier_contd = [
        i
        for i, slide in enumerate(prs.slides)
        if "supplier" in slide_title_text(slide).lower()
        and "(contd" in slide_title_text(slide).lower()
    ]
    hl_only = []
    ka_only = []
    shared = []
    for idx in supplier_contd:
        slide = prs.slides[idx]
        ka = get_key_activities_shape(slide)
        try:
            hl = get_highlights_shape(slide)
        except ValueError:
            hl = None
        if hl and ka:
            shared.append(idx)
        elif hl:
            hl_only.append(idx)
        elif ka:
            ka_only.append(idx)

    assert hl_only, "Expected HL-only Supplier (Contd..) slide(s)"
    assert ka_only, "Expected dedicated KA-only Supplier (Contd..) slide"
    assert not shared, "Supplier should not keep HL+KA on the same contd slide when dense"

    for idx in hl_only:
        slide = prs.slides[idx]
        fit_layout = apply_fixed_story_line_metrics(build_layout_profile(slide))
        if get_key_activities_shape(slide) is not None:
            assert not ka_overlaps_hl_text(slide, fit_layout)


HL_CONTD_TEMPLATE_INDEX = 4


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_sparse_hl_contd_shrinks_to_content():
    """Sparse HL-only contd: table height follows line count, not template row-2 slack."""
    from app.services.ppt_layout_metrics import EMU_PER_INCH
    from app.wsr_engine.formatter import discover_section_templates, get_canonical_style_cell

    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    contd_slide = prs.slides[HL_CONTD_TEMPLATE_INDEX]
    hl = get_highlights_shape(contd_slide)
    template_h_in = hl.height / EMU_PER_INCH
    assert template_h_in > 4.0, "contd template should start with tall HL row"

    tmpl = discover_section_templates(get_canonical_style_cell(prs, profile))
    sparse_sections = [
        {"title": "Story A", "inprogress": [], "completed": [], "planned": []},
        {"title": "Story B", "inprogress": [], "completed": [], "planned": []},
        {"title": "Story C", "inprogress": [], "completed": [], "planned": []},
        {"title": "Story D", "inprogress": [], "completed": [], "planned": []},
    ]
    from app.wsr_engine.formatter import fill_highlights_on_slide

    fill_highlights_on_slide(contd_slide, tmpl, sparse_sections)

    main_slide = prs.slides[profile.services["Cost Core Service"].main_slide_index]
    fit_profile = apply_fixed_story_line_metrics(build_layout_profile(main_slide))
    fit_hl_only_on_slide(contd_slide, fit_profile)

    fitted_h_in = hl.height / EMU_PER_INCH
    assert fitted_h_in < template_h_in - 2.0, (
        f"sparse contd HL should shrink well below template, got {fitted_h_in:.2f}in"
    )
    assert fitted_h_in < 3.0, f"sparse 4-line contd HL should be ~2–3in, got {fitted_h_in:.2f}in"
