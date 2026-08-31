"""Tests for semantic slide detection."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from app.wsr_engine.project_detector import build_project_maps, classify_slide

TEMPLATE = Path(__file__).resolve().parents[2] / "templates" / "wsr_template.pptx"


@pytest.mark.skipif(not TEMPLATE.is_file(), reason="wsr_template.pptx not present")
def test_classify_delivery_main_slides():
    prs = Presentation(TEMPLATE)
    slides = [classify_slide(slide, i) for i, slide in enumerate(prs.slides)]
    mains = [s for s in slides if s.slide_type == "PROJECT_MAIN"]
    assert len(mains) >= 1


@pytest.mark.skipif(not TEMPLATE.is_file(), reason="wsr_template.pptx not present")
def test_build_project_maps():
    prs = Presentation(TEMPLATE)
    slides = [classify_slide(slide, i) for i, slide in enumerate(prs.slides)]
    projects = build_project_maps(slides)
    assert len(projects) >= 1
