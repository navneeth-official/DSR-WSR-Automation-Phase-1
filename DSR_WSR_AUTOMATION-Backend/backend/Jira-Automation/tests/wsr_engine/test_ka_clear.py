"""Tests for KA empty-state clearing."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from app.constants.ppt_bullets import KA_BULLET_CHAR
from app.services.ppt_ka_bullets import ka_bullet_char
from app.services.ppt_shape_utils import get_key_activities_shape
from app.wsr_engine.formatter import clear_ka_on_slide, set_ka_empty_items, set_ka_items

HASKELL = (
    Path(__file__).resolve().parents[2]
    / "templates"
    / "G10X H-E-B WSR Haskell Location Pharmacy GSS PAM 01 Aug 2025.pptx"
)


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_clear_ka_removes_template_sample_text():
    prs = Presentation(str(HASKELL))
    slide = prs.slides[7]  # Pricing Core Service — standalone KA with sample bullets
    ka = get_key_activities_shape(slide)
    assert ka is not None
    assert ka.table.cell(1, 0).text.strip()

    clear_ka_on_slide(slide)
    body = ka.table.cell(1, 0).text
    assert "pricing protection" not in body.lower()
    assert body.strip() == "" or all(not p.text.strip() for p in ka.table.cell(1, 0).text_frame.paragraphs)


def test_set_ka_empty_items_uses_placeholders():
    if not HASKELL.is_file():
        pytest.skip("Haskell sample template missing")
    prs = Presentation(str(HASKELL))
    ka = get_key_activities_shape(prs.slides[7])
    set_ka_empty_items(ka)
    paras = ka.table.cell(1, 0).text_frame.paragraphs
    assert len(paras) >= 1
    assert not any(p.text.strip() for p in paras)


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_set_ka_items_uses_round_bullet_not_dash():
    prs = Presentation(str(HASKELL))
    ka = get_key_activities_shape(prs.slides[7])
    set_ka_items(ka, ["First key activity", "Second key activity"])
    tx_body = ka.table.cell(1, 0).text_frame._txBody
    from pptx.oxml.ns import qn

    from app.services.ppt_shape_utils import paragraph_text

    filled = [p for p in tx_body.findall(qn("a:p")) if paragraph_text(p).strip()]
    assert len(filled) == 2
    for p in filled:
        assert ka_bullet_char(p) == KA_BULLET_CHAR
        p_pr = p.find(qn("a:pPr"))
        assert p_pr is not None and p_pr.get("lvl") == "0"


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_clear_ka_on_contd_template_with_many_bullets():
    """KA-only contd templates may have 5+ sample bullets — all must clear."""
    from app.wsr_engine.hl_ka_normalizer import find_ka_only_contd_template_index

    prs = Presentation(str(HASKELL))
    idx = find_ka_only_contd_template_index(prs)
    assert idx is not None
    slide = prs.slides[idx]
    ka = get_key_activities_shape(slide)
    body = ka.table.cell(1, 0).text
    assert "tobacco" in body.lower()

    clear_ka_on_slide(slide)
    body = ka.table.cell(1, 0).text
    assert "tobacco" not in body.lower()
    assert "ccr" not in body.lower()
    assert body.strip() == "" or all(
        not p.text.strip() for p in ka.table.cell(1, 0).text_frame.paragraphs
    )
