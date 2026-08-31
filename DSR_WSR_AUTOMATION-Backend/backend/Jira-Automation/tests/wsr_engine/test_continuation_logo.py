"""Tests for HEB logo on continuation slides."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.ppt_logo_sync import find_heb_logo_on_slide, sync_heb_logo_from_main
from app.services.ppt_shape_utils import slide_title_text
from app.services.template_profile import scan_template_profile
from app.wsr_engine.continuation_engine import create_contd_slide_after
from app.wsr_engine.hl_ka_normalizer import find_hl_only_contd_template_index

AUG = (
    Path(__file__).resolve().parents[2]
    / "templates"
    / "G10X H-E-B WSR Haskell Location Pharmacy GSS PAM 01 Aug 2025.pptx"
)


def _heb_logo_ok(slide) -> bool:
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            _ = sh.image.blob
            return True
        except (AttributeError, ValueError, KeyError):
            continue
    return False


@pytest.mark.skipif(not AUG.is_file(), reason="Aug Haskell template missing")
def test_copy_shapes_breaks_heb_logo_until_sync():
    prs = Presentation(str(AUG))
    profile = scan_template_profile(AUG)
    main_idx = profile.services["Cost Core Service"].main_slide_index
    template_idx = find_hl_only_contd_template_index(prs)
    assert template_idx is not None

    contd_idx = create_contd_slide_after(prs, prs, main_idx, template_idx)
    main_slide = prs.slides[main_idx]
    contd_slide = prs.slides[contd_idx]

    assert find_heb_logo_on_slide(main_slide) is not None
    assert not _heb_logo_ok(contd_slide)

    assert sync_heb_logo_from_main(main_slide, contd_slide) is True
    assert _heb_logo_ok(contd_slide)
    assert "(contd" not in slide_title_text(main_slide).lower()
