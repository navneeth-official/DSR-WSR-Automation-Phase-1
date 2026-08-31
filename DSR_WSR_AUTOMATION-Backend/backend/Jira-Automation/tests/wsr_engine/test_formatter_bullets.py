"""Tests for G10X bullet formatting and HL/KA layout preservation."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from app.services.ppt_shape_utils import (
    get_highlights_shape,
    get_key_activities_shape,
    has_combined_hl_ka_table,
    paragraph_text,
    slide_title_text,
)
from app.services.template_profile import scan_template_profile
from app.wsr_engine.formatter import (
    discover_section_templates,
    fill_highlights_on_slide,
    get_canonical_style_cell,
    section_templates_from_profile,
)
from app.wsr_engine.hl_ka_normalizer import (
    find_hl_only_contd_template_index,
    normalize_deck_hl_ka_layouts,
)

HASKELL_JULY = (
    Path(__file__).resolve().parents[2]
    / "templates"
    / "G10X H-E-B WSR Haskell Location Pharmacy GSS PAM 11 July 2025.pptx"
)
HASKELL_AUG = (
    Path(__file__).resolve().parents[2]
    / "templates"
    / "G10X H-E-B WSR Haskell Location Pharmacy GSS PAM 01 Aug 2025.pptx"
)
HASKELL = HASKELL_JULY


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_discover_section_templates_preserves_template_story_bullets():
    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    cell = get_canonical_style_cell(prs, profile)
    tmpl = discover_section_templates(cell)
    bullet = tmpl["completed_bullet"]
    p_pr = bullet.find(qn("a:pPr"))
    assert p_pr is not None
    assert p_pr.get("lvl") == "2"
    bu = p_pr.find(qn("a:buChar"))
    assert bu is not None
    assert bu.get("char") == "\uf0a7"
    buf = p_pr.find(qn("a:buFont"))
    assert buf is not None and "Wingdings" in buf.get("typeface", "")


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_discover_section_templates_preserves_category_header_bullets():
    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    cell = get_canonical_style_cell(prs, profile)
    tmpl = discover_section_templates(cell)
    hdr = tmpl["completed_hdr"]
    p_pr = hdr.find(qn("a:pPr"))
    assert p_pr is not None
    assert p_pr.get("lvl") == "1"
    bu = p_pr.find(qn("a:buChar"))
    assert bu is not None
    assert bu.get("char") == "o"


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_fill_highlights_applies_story_bullets():
    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    cell = get_canonical_style_cell(prs, profile)
    tmpl = discover_section_templates(cell)
    slide = prs.slides[profile.services["Cost Core Service"].main_slide_index]

    fill_highlights_on_slide(
        slide,
        tmpl,
        [
            {
                "sprint_bold": "Sprint 1 ",
                "sprint_light": "Jun 01 – Jun 14",
                "completed_items": ["Story A"],
                "completed_count": "1",
                "released_items": [],
                "released_count": "0",
                "inprogress_items": [],
                "inprogress_count": "0",
            }
        ],
    )

    hl = get_highlights_shape(slide)
    story_found = False
    for p in hl.table.cell(2, 0).text_frame._txBody.findall(qn("a:p")):
        txt = paragraph_text(p).strip()
        if txt == "Story A":
            story_found = True
            p_pr = p.find(qn("a:pPr"))
            assert p_pr is not None and p_pr.find(qn("a:buChar")) is not None
            assert p_pr.get("lvl") == "2"
            bu = p_pr.find(qn("a:buChar"))
            assert bu.get("char") == "\uf0a7"
            runs = p.findall(qn("a:r"))
            assert not any(
                r.find(qn("a:rPr")) is not None and r.find(qn("a:rPr")).get("b") == "1"
                for r in runs
            )
    assert story_found


@pytest.mark.skipif(not HASKELL_AUG.is_file(), reason="Aug Haskell template missing")
def test_fill_supplier_preserves_text_prefix_dash_bullet():
    """Aug template stores story dash as a text run (₋\\t), not buChar."""
    prs = Presentation(str(HASKELL_AUG))
    profile = scan_template_profile(HASKELL_AUG)
    svc = profile.services["Supplier Core Service"]
    slide = prs.slides[svc.main_slide_index]
    section_tmpl = section_templates_from_profile(slide, svc)

    fill_highlights_on_slide(
        slide,
        section_tmpl,
        [
            {
                "sprint_bold": "Sprint – Test ",
                "sprint_light": "(Jun 01 – Jun 15) Stories",
                "completed_items": [],
                "completed_count": "0",
                "released_items": [],
                "released_count": "0",
                "inprogress_items": ["Validate supplier onboarding workflow"],
                "inprogress_count": "1",
            }
        ],
    )

    hl = get_highlights_shape(slide)
    story_found = False
    for p in hl.table.cell(2, 0).text_frame._txBody.findall(qn("a:p")):
        txt = paragraph_text(p).strip()
        if txt.endswith("Validate supplier onboarding workflow"):
            story_found = True
            assert txt.startswith("\u208b"), f"expected text-prefix dash bullet, got {txt[:10]!r}"
            p_pr = p.find(qn("a:pPr"))
            assert p_pr is None or p_pr.find(qn("a:buChar")) is None
    assert story_found


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_normalize_preserves_embedded_ka_layouts():
    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    wentworth_idx = profile.services["Wentworth"].main_slide_index
    assert has_combined_hl_ka_table(prs.slides[wentworth_idx])

    normalize_deck_hl_ka_layouts(prs, profile)
    assert has_combined_hl_ka_table(prs.slides[wentworth_idx])
    assert get_key_activities_shape(prs.slides[wentworth_idx]) is None


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_contd_multi_sprint_does_not_leave_template_text():
    from app.wsr_engine.formatter import populate_highlights_contd_cell
    from app.wsr_engine.placeholder_locator import highlights_content_cell

    prs = Presentation(str(HASKELL))
    profile = scan_template_profile(HASKELL)
    cell = get_canonical_style_cell(prs, profile)
    tmpl = discover_section_templates(cell)
    slide = prs.slides[profile.services["Supplier Core Service"].contd_slide_indices[0]]
    hl_cell = highlights_content_cell(slide)

    populate_highlights_contd_cell(
        hl_cell,
        tmpl,
        {
            "sections": [
                {
                    "sprint_bold": "Sprint – Q3.01 FY26 Phoenix, inprogress ",
                    "sprint_light": "(Jun 01 – Jun 15) Stories",
                    "completed_items": ["Story A"],
                    "completed_count": "1",
                    "released_items": [],
                    "released_count": "0",
                    "inprogress_items": [],
                    "inprogress_count": "0",
                },
                {
                    "sprint_bold": "Sprint – Q3.02 FY26 Orion, In-progress ",
                    "sprint_light": "(Jun 10 – Jul 01) Stories",
                    "completed_items": ["Story B"],
                    "completed_count": "1",
                    "released_items": [],
                    "released_count": "0",
                    "inprogress_items": [],
                    "inprogress_count": "0",
                },
            ]
        },
    )

    lines = [
        paragraph_text(p).strip()
        for p in hl_cell.text_frame._txBody.findall(qn("a:p"))
        if paragraph_text(p).strip()
    ]
    assert not any("yugoslavic" in line.lower() for line in lines)
    assert not any("fireball island" in line.lower() for line in lines)
    assert "Story A" in lines
    assert "Story B" in lines
    assert any("Orion" in line for line in lines)
    orion_idx = next(i for i, line in enumerate(lines) if "Orion" in line)
    assert "Current week sprint status" in lines[orion_idx + 1]


@pytest.mark.skipif(not HASKELL.is_file(), reason="Haskell sample template missing")
def test_find_hl_only_contd_template():
    prs = Presentation(str(HASKELL))
    idx = find_hl_only_contd_template_index(prs)
    assert idx is not None
    slide = prs.slides[idx]
    assert "(contd" in slide_title_text(slide).lower()
    assert get_key_activities_shape(slide) is None
    assert not has_combined_hl_ka_table(slide)
