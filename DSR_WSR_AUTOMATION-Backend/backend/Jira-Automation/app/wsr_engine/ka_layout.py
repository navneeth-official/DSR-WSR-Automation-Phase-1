"""Wrap-aware Highlights / Key Activities layout — template-agnostic."""

from __future__ import annotations

import copy
import logging

from pptx import Presentation

from app.services.ppt_layout_metrics import (
    EMU_PER_INCH,
    FIXED_STORY_LINE_HEIGHT_EMU,
    FOOTER_MAX_BOTTOM_EMU,
    MIN_TEXT_KA_CLEARANCE_IN,
    apply_fixed_story_line_metrics,
    count_hl_paragraphs,
    count_visual_lines_in_hl,
    estimated_text_bottom_emu,
    hl_ka_tab_gap_emu,
    rendered_text_bottom_emu,
)
from app.services.ppt_shape_utils import get_highlights_shape, get_key_activities_shape, slide_title_text
from app.services.template_profile import ServiceProfile, TemplateProfile
from app.wsr_engine.hl_ka_normalizer import find_standalone_reference_slide
from app.wsr_engine.models import TitleFormat
from app.wsr_engine.overflow_engine import build_layout_profile

logger = logging.getLogger(__name__)

MAX_KA_BOTTOM_EMU = FOOTER_MAX_BOTTOM_EMU
MIN_HL_BOTTOM_PAD_EMU = int(0.1 * EMU_PER_INCH)
HL_CONTENT_BOTTOM_PAD_EMU = int(0.1 * EMU_PER_INCH)
_VISUAL_LINE_BUFFER = 1.12


def _canonical_per_line(profile: dict) -> int:
    return int(profile.get("canonical_per_line_emu") or FIXED_STORY_LINE_HEIGHT_EMU)


def _rendered_text_kw(profile: dict) -> dict:
    return {
        "ref_para_count": profile.get("canonical_para_count", profile["ref_para_count"]),
        "ref_r2": profile["ref_r2"],
        "per_line_emu": _canonical_per_line(profile),
    }


def rendered_hl_text_bottom(hl_shape, profile: dict) -> int:
    """Wrap-aware HL text end — never underestimate vs table row height."""
    return rendered_text_bottom_emu(hl_shape, **_rendered_text_kw(profile))


def hl_content_needs_ka_contd(hl_shape, profile: dict, ka_height_emu: int) -> bool:
    """True when rendered HL text + KA + tab gap cannot fit above the footer."""
    gap = hl_ka_tab_gap_emu(profile.get("canonical_line_height_emu"))
    text_bottom = rendered_hl_text_bottom(hl_shape, profile)
    return text_bottom + gap + ka_height_emu > MAX_KA_BOTTOM_EMU


def _set_table_shape_height(hl_shape, row_heights: list[int], target_h: int) -> None:
    table = hl_shape.table
    for i, h in enumerate(row_heights):
        table.rows[i].height = h
    hl_shape.height = target_h


def _min_hl_height_for_content(hl_shape, profile: dict) -> int:
    """Minimum HL table height so wrap-aware text is not clipped."""
    r0, r1 = profile["r0"], profile["r1"]
    ref_pad = profile["ref_pad"]
    per_line = _canonical_per_line(profile)
    effective_lines = max(count_hl_paragraphs(hl_shape), count_visual_lines_in_hl(hl_shape))
    content_h = int(per_line * effective_lines * _VISUAL_LINE_BUFFER)
    return int(r0 + r1 + content_h + MIN_HL_BOTTOM_PAD_EMU)


def fit_hl_to_content(hl_shape, profile: dict, max_h: int | None = None) -> int:
    """
    Resize Highlights table to wrap-aware text (+ small pad), never larger than required.

    Uses line-based text bottom only — not physical row-2 height from cloned contd
    templates, which would leave sparse HL boxes oversized.

    Returns the applied table height in EMU.
    """
    r0, r1 = profile["r0"], profile["r1"]
    per_line = _canonical_per_line(profile)

    text_bottom = hl_text_bottom_emu(hl_shape, profile)
    min_h = _min_hl_height_for_content(hl_shape, profile)
    target_h = max(int(text_bottom - hl_shape.top + HL_CONTENT_BOTTOM_PAD_EMU), min_h)
    if max_h is not None:
        target_h = min(target_h, max_h)

    min_pad = MIN_HL_BOTTOM_PAD_EMU
    min_content_h = max(int(per_line * 1), 45720)
    content_h = max(target_h - r0 - r1 - min_pad, min_content_h)
    _set_table_shape_height(hl_shape, [r0, r1, content_h], target_h)
    return target_h


def hl_table_needs_ka_contd(hl_shape, profile: dict, ka_height_emu: int) -> bool:
    """True when minimum HL table height cannot share the slide with KA."""
    gap = hl_ka_tab_gap_emu()
    min_h = _min_hl_height_for_content(hl_shape, profile)
    allowed = int(MAX_KA_BOTTOM_EMU - ka_height_emu - gap - hl_shape.top)
    return min_h > allowed


def hl_text_bottom_emu(hl_shape, profile: dict) -> int:
    """Wrap-aware text end using template line metrics (not inflated cell row height)."""
    return estimated_text_bottom_emu(
        hl_shape,
        ref_para_count=profile.get("canonical_para_count", profile["ref_para_count"]),
        ref_r2=profile["ref_r2"],
        per_line_emu=_canonical_per_line(profile),
    )


def effective_ka_top(hl_shape, profile: dict, ka_height: int) -> int:
    """KA top = rendered HL text bottom + tab gap, footer-clamped."""
    gap = hl_ka_tab_gap_emu(profile.get("canonical_line_height_emu"))
    text_bottom = rendered_hl_text_bottom(hl_shape, profile)
    tab_top = int(text_bottom + gap)
    max_top = MAX_KA_BOTTOM_EMU - ka_height
    return min(tab_top, max_top)


def ka_bottom_emu(slide) -> int | None:
    ka = get_key_activities_shape(slide)
    if ka is None:
        return None
    return int(ka.top + ka.height)


def ka_overlaps_footer(slide, max_bottom_emu: int = MAX_KA_BOTTOM_EMU) -> bool:
    bottom = ka_bottom_emu(slide)
    return bottom is not None and bottom > max_bottom_emu


def ka_overlaps_hl_text(slide, profile: dict) -> bool:
    """True when KA sits above rendered HL text + standard tab gap."""
    ka = get_key_activities_shape(slide)
    if ka is None:
        return False
    try:
        hl = get_highlights_shape(slide)
    except ValueError:
        return False
    gap = hl_ka_tab_gap_emu(profile.get("canonical_line_height_emu"))
    min_ka_top = rendered_hl_text_bottom(hl, profile) + gap
    return ka.top < min_ka_top


def _merge_slide_geometry(layout_profile: dict, slide_profile: dict) -> dict:
    """Use row metrics from the slide being fitted (contd templates differ from main)."""
    merged = dict(layout_profile)
    for key in ("r0", "r1", "ref_r2", "ref_pad", "ref_para_count", "ref_hl_height", "ref_hl_top", "ref_hl"):
        if key in slide_profile:
            merged[key] = slide_profile[key]
    return merged


def fit_main_slide_hl_ka(
    slide,
    template_prs: Presentation,
    profile: TemplateProfile,
    layout_profile: dict | None = None,
) -> bool:
    """
    Size HL to content and place KA below rendered HL text + tab gap.

    Returns False when HL and KA cannot share the slide (caller should split KA
    to a dedicated continuation slide).
    """
    ka = get_key_activities_shape(slide)
    if ka is None:
        return False
    try:
        hl = get_highlights_shape(slide)
    except ValueError:
        return False

    ref_idx = find_standalone_reference_slide(template_prs, profile)
    if ref_idx is None:
        return False

    ref_slide = template_prs.slides[ref_idx]
    ref_ka = get_key_activities_shape(ref_slide)
    if ref_ka is None:
        return False

    slide_profile = build_layout_profile(slide)
    layout_profile = apply_fixed_story_line_metrics(
        _merge_slide_geometry(dict(layout_profile or slide_profile), slide_profile)
    )

    ka_h = int(ka.height or ref_ka.height)
    gap = hl_ka_tab_gap_emu()
    max_hl_h = int(MAX_KA_BOTTOM_EMU - ka_h - gap - hl.top)

    if hl_table_needs_ka_contd(hl, layout_profile, ka_h):
        return False

    fit_hl_to_content(hl, layout_profile, max_h=max_hl_h)

    if hl_content_needs_ka_contd(hl, layout_profile, ka_h):
        return False

    ka.left = ref_ka.left
    ka.width = ref_ka.width
    ka.height = ka_h
    ka.top = effective_ka_top(hl, layout_profile, ka_h)

    clearance_in = (ka.top - rendered_hl_text_bottom(hl, layout_profile)) / EMU_PER_INCH
    if clearance_in + 1e-6 < MIN_TEXT_KA_CLEARANCE_IN:
        logger.warning("KA clearance below minimum after fit (%.3f in)", clearance_in)
        return False
    if ka_overlaps_hl_text(slide, layout_profile):
        logger.warning("KA overlaps rendered HL text after fit")
        return False
    if ka_overlaps_footer(slide):
        logger.warning("KA overlaps footer after fit")
        return False
    return True


def fit_hl_only_on_slide(slide, layout_profile: dict | None = None) -> int:
    """Size Highlights to content only, up to the full HL-only slide budget."""
    try:
        hl = get_highlights_shape(slide)
    except ValueError:
        return 0
    slide_profile = build_layout_profile(slide)
    if layout_profile is None:
        layout_profile = slide_profile
    else:
        layout_profile = _merge_slide_geometry(dict(layout_profile), slide_profile)
    layout_profile = apply_fixed_story_line_metrics(layout_profile)
    layout_profile.setdefault("canonical_para_count", layout_profile["ref_para_count"])
    max_h = int(MAX_KA_BOTTOM_EMU - layout_profile["ref_hl"].top - MIN_HL_BOTTOM_PAD_EMU)
    height = fit_hl_to_content(hl, layout_profile, max_h=max_h)
    from app.wsr_engine.formatter import top_align_highlights_on_slide

    top_align_highlights_on_slide(slide)
    return height


def ensure_key_activities_on_slide(
    slide,
    template_prs: Presentation,
    profile: TemplateProfile,
) -> object | None:
    """Clone empty KA when missing on any project slide."""
    ka = get_key_activities_shape(slide)
    if ka is not None:
        return ka

    ref_idx = find_standalone_reference_slide(template_prs, profile)
    if ref_idx is None:
        logger.warning("No standalone KA reference slide in template")
        return None
    ref_ka = get_key_activities_shape(template_prs.slides[ref_idx])
    if ref_ka is None:
        return None
    new_el = copy.deepcopy(ref_ka.element)
    slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    ka = get_key_activities_shape(slide)
    if ka is None:
        return None
    logger.debug("Inserted KA tab on slide")
    return ka


def ensure_key_activities_on_main_slide(
    slide,
    template_prs: Presentation,
    service: ServiceProfile,
    profile: TemplateProfile,
):
    """Clone empty KA when missing; fit HL+KA geometry on main slides."""
    if service.ka_mode == "embedded":
        return get_key_activities_shape(slide)

    ka = ensure_key_activities_on_slide(slide, template_prs, profile)
    if ka is None:
        return None

    from app.wsr_engine.formatter import populate_ka_on_slide

    populate_ka_on_slide(slide, [], service)
    fit_main_slide_hl_ka(slide, template_prs, profile)
    return ka


def _slide_fit_profile(slide, base_profile: dict) -> dict:
    slide_profile = build_layout_profile(slide)
    merged = apply_fixed_story_line_metrics(
        _merge_slide_geometry(dict(base_profile), slide_profile)
    )
    merged.setdefault("canonical_para_count", merged["ref_para_count"])
    return merged


def _split_ka_to_contd_slide(
    prs: Presentation,
    template_prs: Presentation,
    *,
    project_name: str,
    main_idx: int,
    hl_slide_idx: int,
    service: ServiceProfile,
    title_format: TitleFormat,
    ka_items: list[str],
    fit_profile: dict,
) -> int:
    """Move KA off ``hl_slide_idx`` onto a new KA-only continuation slide."""
    from app.wsr_engine.continuation_engine import append_ka_only_contd_slide
    from app.wsr_engine.formatter import populate_ka_on_slide, remove_ka_from_slide

    hl_slide = prs.slides[hl_slide_idx]
    remove_ka_from_slide(hl_slide)
    fit_hl_only_on_slide(hl_slide, fit_profile)

    ka_idx = append_ka_only_contd_slide(
        prs,
        template_prs,
        project_name,
        main_idx,
        hl_slide_idx,
        service,
        title_format,
    )
    populate_ka_on_slide(prs.slides[ka_idx], ka_items, service)
    return ka_idx


def _apply_hl_ka_or_split(
    prs: Presentation,
    template_prs: Presentation,
    *,
    slide_idx: int,
    main_idx: int,
    project_name: str,
    service: ServiceProfile,
    profile: TemplateProfile,
    title_format: TitleFormat | None,
    ka_items: list[str],
    fit_profile: dict,
) -> None:
    """Populate KA on ``slide_idx`` and fit HL+KA, splitting to KA contd when needed."""
    from app.wsr_engine.formatter import populate_ka_on_slide, remove_ka_from_slide

    slide = prs.slides[slide_idx]
    remove_ka_from_slide(slide)
    ensure_key_activities_on_slide(slide, template_prs, profile)
    populate_ka_on_slide(slide, ka_items, service)

    slide_fit = _slide_fit_profile(slide, fit_profile)
    if fit_main_slide_hl_ka(slide, template_prs, profile, layout_profile=slide_fit):
        return

    if title_format is None or not project_name:
        logger.warning("HL+KA overflow on slide %s but no title_format — leaving cramped layout", slide_idx)
        return

    logger.info(
        "HL+KA overflow on %s slide %s — moving Key Activities to dedicated (Contd..) slide",
        project_name,
        slide_idx + 1,
    )
    _split_ka_to_contd_slide(
        prs,
        template_prs,
        project_name=project_name,
        main_idx=main_idx,
        hl_slide_idx=slide_idx,
        service=service,
        title_format=title_format,
        ka_items=ka_items,
        fit_profile=fit_profile,
    )


def finalize_project_hl_ka(
    prs: Presentation,
    main_idx: int,
    contd_indices: list[int],
    template_prs: Presentation,
    service: ServiceProfile,
    profile: TemplateProfile,
    layout_profile: dict | None = None,
    ka_on_main: bool = False,
    ka_contd_only: bool = False,
    key_activities: list[str] | None = None,
    project_name: str = "",
    title_format: TitleFormat | None = None,
) -> None:
    """
    Place Key Activities only on the last slide of a project; tighten HL on every slide.

    When overflow continues on (Contd..) slides, earlier slides use HL-only layout
    (full vertical budget). The final slide holds HL + KA when they fit; otherwise
    KA moves to an additional KA-only (Contd..) slide.
    """
    from app.wsr_engine.formatter import populate_ka_on_slide, remove_ka_from_slide

    ka_items = list(key_activities or [])
    if title_format is None:
        title_format = TitleFormat(
            prefix="Delivery Status", separator=" - ", contd_marker="(Contd..)"
        )

    template_slide = template_prs.slides[service.main_slide_index]
    fit_profile = apply_fixed_story_line_metrics(
        dict(layout_profile or build_layout_profile(template_slide))
    )
    fit_profile.setdefault("canonical_para_count", fit_profile["ref_para_count"])

    if service.ka_mode == "embedded":
        for idx in contd_indices:
            fit_hl_only_on_slide(prs.slides[idx], fit_profile)
        populate_ka_on_slide(prs.slides[main_idx], ka_items, service)
        return

    if ka_on_main:
        _apply_hl_ka_or_split(
            prs,
            template_prs,
            slide_idx=main_idx,
            main_idx=main_idx,
            project_name=project_name,
            service=service,
            profile=profile,
            title_format=title_format,
            ka_items=ka_items,
            fit_profile=fit_profile,
        )
        return

    if ka_contd_only:
        main_slide = prs.slides[main_idx]
        remove_ka_from_slide(main_slide)
        fit_hl_only_on_slide(main_slide, fit_profile)
        if contd_indices:
            populate_ka_on_slide(prs.slides[contd_indices[-1]], ka_items, service)
        return

    last_idx = contd_indices[-1] if contd_indices else main_idx
    hl_only_indices = [main_idx] + [i for i in contd_indices if i != last_idx]

    for idx in hl_only_indices:
        slide = prs.slides[idx]
        remove_ka_from_slide(slide)
        fit_hl_only_on_slide(slide, fit_profile)

    pname = project_name
    if not pname:
        raw = slide_title_text(prs.slides[main_idx])
        if "-" in raw:
            pname = raw.split("-", 1)[-1].strip()
            if "(contd" in pname.lower():
                pname = pname[: pname.lower().index("(contd")].strip()

    _apply_hl_ka_or_split(
        prs,
        template_prs,
        slide_idx=last_idx,
        main_idx=main_idx,
        project_name=pname,
        service=service,
        profile=profile,
        title_format=title_format,
        ka_items=ka_items,
        fit_profile=fit_profile,
    )
