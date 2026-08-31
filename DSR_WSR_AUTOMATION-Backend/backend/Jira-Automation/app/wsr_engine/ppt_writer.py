"""Deck preparation and content application."""

from __future__ import annotations

import logging
import shutil
from datetime import date
from pathlib import Path

from pptx import Presentation

from app.wsr_engine.continuation_engine import ensure_continuation_slides, ensure_ka_only_contd_slide
from app.wsr_engine.content_parser import section_display_content
from app.wsr_engine.formatter import (
    clear_ka_on_project_slides,
    fill_highlights_on_slide,
    section_templates_from_profile,
)
from app.wsr_engine.hl_ka_normalizer import normalize_deck_hl_ka_layouts
from app.wsr_engine.ka_layout import finalize_project_hl_ka
from app.wsr_engine.index_updater import rebuild_index_from_projects, reflow_index_slide
from app.wsr_engine.models import BuildReport, ProjectContent, ProjectMap, TemplateModel
from app.wsr_engine.overflow_engine import (
    build_overflow_profile,
    build_layout_profile,
    main_slide_capacity,
    plan_overflow,
)
from app.wsr_engine.project_deletion import find_main_slide_index
from app.wsr_engine.slide_order import remove_project_contd_slides
from app.services.template_profile import TemplateProfile
from app.wsr_engine.slide_ops import save_presentation_clean

logger = logging.getLogger(__name__)


class PptWriter:
    def __init__(self, template_path: Path | str, working_path: Path | str | None = None):
        self.template_path = Path(template_path).resolve()
        self.working_path = Path(working_path) if working_path else None
        self.prs: Presentation | None = None
        self.template_prs: Presentation | None = None
        self.template_model: TemplateModel | None = None

    def prepare(self, template_model: TemplateModel) -> Presentation:
        self.template_model = template_model
        if self.working_path:
            shutil.copy2(self.template_path, self.working_path)
            self.prs = Presentation(str(self.working_path))
        else:
            self.prs = Presentation(str(self.template_path))

        self.template_prs = Presentation(str(self.template_path))
        normalize_deck_hl_ka_layouts(self.prs, template_model.profile)
        return self.prs

    def apply_project(
        self,
        project_map: ProjectMap,
        content: ProjectContent,
        profile: TemplateProfile,
        report: BuildReport,
    ) -> None:
        assert self.prs is not None and self.template_prs is not None
        assert self.template_model is not None

        service = profile.services.get(project_map.project_name)
        if service is None:
            service = profile.service_for_title(project_map.project_name)
        if service is None:
            report.warnings.append(f"No service profile for {project_map.project_name}; skipping")
            return

        main_idx = find_main_slide_index(self.prs, project_map.project_name)
        if main_idx is None or main_idx >= len(self.prs.slides):
            report.errors.append(f"Main slide not found for {project_map.project_name}")
            return

        main_slide = self.prs.slides[main_idx]
        template_slide = self.template_prs.slides[service.main_slide_index]
        section_tmpl = section_templates_from_profile(template_slide, service)

        layout_profile = build_overflow_profile(
            self.template_prs, template_slide, service, profile
        )
        overflow = plan_overflow(
            content, layout_profile, main_cap=main_slide_capacity(layout_profile)
        )

        try:
            fill_highlights_on_slide(
                main_slide,
                section_tmpl,
                overflow.main_sections,
            )
            report.inserted_projects.append(project_map.project_name)
        except ValueError as exc:
            report.warnings.append(f"Highlights fill failed for {project_map.project_name}: {exc}")
            clear_ka_on_project_slides(self.prs, project_map.project_name, template_slide, service)
            remove_project_contd_slides(self.prs, project_map.project_name)
            return

        contd_indices: list[int] = []
        if overflow.continuation_chains:
            title_format = self.template_model.title_format
            if title_format is None:
                from app.wsr_engine.models import TitleFormat
                title_format = TitleFormat(prefix="Delivery Status", separator=" - ", contd_marker="(Contd..)")

            contd_indices, created = ensure_continuation_slides(
                self.prs,
                self.template_prs,
                project_map.project_name,
                main_idx,
                overflow.continuation_chains,
                service,
                service.main_slide_index,
                title_format,
                section_tmpl,
            )
            report.continuations_created += created
        elif overflow.ka_contd_only:
            title_format = self.template_model.title_format
            if title_format is None:
                from app.wsr_engine.models import TitleFormat
                title_format = TitleFormat(prefix="Delivery Status", separator=" - ", contd_marker="(Contd..)")
            contd_indices, created = ensure_ka_only_contd_slide(
                self.prs,
                self.template_prs,
                project_map.project_name,
                main_idx,
                service,
                service.main_slide_index,
                title_format,
            )
            report.continuations_created += created
        else:
            remove_project_contd_slides(self.prs, project_map.project_name)

        finalize_project_hl_ka(
            self.prs,
            main_idx,
            contd_indices,
            self.template_prs,
            service,
            profile,
            layout_profile=layout_profile,
            ka_on_main=overflow.ka_on_main,
            ka_contd_only=overflow.ka_contd_only,
            key_activities=content.key_activities,
            project_name=project_map.project_name,
            title_format=self.template_model.title_format,
        )

    def update_index(
        self,
        projects: list[ProjectMap],
        *,
        skeleton_mode: bool = False,
        content_labels: dict[str, str] | None = None,
    ) -> int:
        assert self.prs is not None
        if skeleton_mode:
            return rebuild_index_from_projects(
                self.prs,
                projects,
                content_labels,
                self.template_model.index_layout,
            )
        return reflow_index_slide(self.prs, projects)

    def save(self, output_path: Path | str) -> Path:
        assert self.prs is not None
        return save_presentation_clean(self.prs, output_path)
