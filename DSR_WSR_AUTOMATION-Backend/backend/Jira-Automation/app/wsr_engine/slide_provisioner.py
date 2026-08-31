"""Provision delivery slides from a skeleton template prototype."""

from __future__ import annotations

import copy
import logging

from pptx import Presentation
from pptx.oxml.ns import qn

from app.services.ppt_logo_sync import sync_header_pictures_from_reference
from app.services.ppt_shape_utils import find_title_shape
from app.wsr_engine.models import ProjectContent, ProjectMap, TitleFormat
from app.wsr_engine.slide_ops import clear_slide_shapes, copy_shapes_to_slide, move_slide_after

logger = logging.getLogger(__name__)


def resolve_display_name(content: ProjectContent, aliases: dict[str, str]) -> str:
    """Canonical delivery slide title suffix for a content chunk."""
    if content.project_name and content.project_name in aliases:
        return aliases[content.project_name]
    if content.title in aliases.values():
        return content.title
    return content.title or content.project_name


def set_delivery_title(slide, display_name: str, title_format: TitleFormat, ref_title_shape) -> None:
    prefix = title_format.prefix
    sep = title_format.separator
    text = f"{prefix}{sep}{display_name}"

    try:
        title_shape = find_title_shape(slide)
    except RuntimeError:
        new_el = copy.deepcopy(ref_title_shape.element)
        slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        title_shape = find_title_shape(slide)

    ref_rPr = None
    if ref_title_shape.has_text_frame and ref_title_shape.text_frame.paragraphs:
        runs = ref_title_shape.text_frame.paragraphs[0].runs
        if runs:
            ref_rPr = runs[0]._r.find(qn("a:rPr"))

    paragraph = title_shape.text_frame.paragraphs[0]
    if not paragraph.runs:
        paragraph.add_run()
    run = paragraph.runs[0]
    old_rPr = run._r.find(qn("a:rPr"))
    if old_rPr is not None:
        run._r.remove(old_rPr)
    if ref_rPr is not None:
        run._r.insert(0, copy.deepcopy(ref_rPr))
    run.text = text
    for extra in paragraph.runs[1:]:
        paragraph._p.remove(extra._r)


def provision_project_slides(
    prs: Presentation,
    prototype_idx: int,
    content_projects: list[ProjectContent],
    title_format: TitleFormat,
    aliases: dict[str, str] | None = None,
) -> list[ProjectMap]:
    """
    Rename the prototype slide for the first project and clone it for the rest.

    New slides are inserted immediately before trailing static slides (Matters, Team
    Allocation, Thank you).
    """
    if not content_projects:
        return []

    aliases = dict(aliases or {})
    prototype_slide = prs.slides[prototype_idx]
    ref_title = find_title_shape(prototype_slide)
    project_maps: list[ProjectMap] = []

    for idx, content in enumerate(content_projects):
        display_name = resolve_display_name(content, aliases)

        if idx == 0:
            slide_idx = prototype_idx
            slide = prototype_slide
        else:
            layout = prototype_slide.slide_layout
            prs.slides.add_slide(layout)
            new_idx = len(prs.slides) - 1
            slide = prs.slides[new_idx]
            clear_slide_shapes(slide)
            copy_shapes_to_slide(prototype_slide, slide)
            sync_header_pictures_from_reference(prototype_slide, slide)
            insert_after = prototype_idx + idx - 1
            move_slide_after(prs, new_idx, insert_after)
            slide_idx = insert_after + 1
            slide = prs.slides[slide_idx]

        set_delivery_title(slide, display_name, title_format, ref_title)
        project_maps.append(ProjectMap(project_name=display_name, main_slide_index=slide_idx))
        logger.info("Provisioned delivery slide %d -> %s", slide_idx, display_name)

    return project_maps
