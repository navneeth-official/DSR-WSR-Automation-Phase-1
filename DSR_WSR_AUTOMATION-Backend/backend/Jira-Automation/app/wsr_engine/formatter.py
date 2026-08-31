"""Formatting-preserving Highlights population — clone template XML only."""

from __future__ import annotations

import copy
import re
from typing import Any

from pptx.oxml.ns import qn

from app.constants.ppt_bullets import (
    CATEGORY_HEADER_BULLET,
    CATEGORY_HEADER_BULLET_COLOR,
    CATEGORY_HEADER_BULLET_FONT,
    CATEGORY_HEADER_LEVEL,
)
from app.services.ppt_shape_utils import (
    get_highlights_shape,
    get_key_activities_shape,
    paragraph_text,
)
from app.services.ppt_hl_bullets import (
    clear_story_line_text,
    is_text_prefix_story_bullet,
    set_single_run_text,
    set_story_line_text,
)
from app.services.ppt_ka_bullets import replace_ka_bullet_block
from app.services.template_profile import ParagraphRoles, ServiceProfile, TemplateProfile
from app.wsr_engine.content_parser import section_display_content
from app.wsr_engine.models import SprintSection
from app.wsr_engine.placeholder_locator import highlights_content_cell

STORY_BUCKET_ORDER = ("completed", "released", "inprogress")
STORY_BULLET_CHAR = "\u208b"  # ₋ dash — G10X story line bullet (Sustainment reference)


def _normalize_match_text(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").replace("\u200b", "").replace("\xa0", " ").strip().lower())


def _paragraph_bullet_level(p_elem) -> int:
    pPr = p_elem.find(qn("a:pPr"))
    if pPr is None or pPr.get("lvl") is None:
        return 0
    return int(pPr.get("lvl"))


def _looks_like_category(text: str) -> bool:
    t = _normalize_match_text(text)
    return any(
        needle in t
        for needles in (
            ("completed this week",),
            ("released for partner", "released for partner review"),
            ("in-progress this week", "in progress this week"),
        )
        for needle in needles
    ) or ("stories " in t and "week" in t)


def _paragraph_is_bold(p_elem) -> bool:
    for run in p_elem.findall(qn("a:r")):
        r_pr = run.find(qn("a:rPr"))
        if r_pr is not None and r_pr.get("b") == "1":
            return True
    return False


def _story_after_header(paras, hdr_idx: int | None):
    if hdr_idx is None:
        return None
    for i in range(hdr_idx + 1, len(paras)):
        text = paragraph_text(paras[i]).strip()
        if not text:
            continue
        if _looks_like_category(text):
            continue
        if "current week sprint status" in _normalize_match_text(text):
            continue
        return paras[i]
    return None


def _is_story_bullet(p_elem) -> bool:
    text = paragraph_text(p_elem).strip()
    if not text or _looks_like_category(text):
        return False
    if "current week sprint status" in _normalize_match_text(text):
        return False
    if _paragraph_is_bold(p_elem) and _paragraph_bullet_level(p_elem) == 1:
        return False
    p_pr = p_elem.find(qn("a:pPr"))
    if p_pr is not None and p_pr.find(qn("a:buChar")) is not None:
        return True
    if is_text_prefix_story_bullet(p_elem):
        return True
    return text.startswith(STORY_BULLET_CHAR) or text.startswith("-") or text.startswith("\u2013")


def _find_canonical_story_bullet(paras):
    for needle in (
        "completed this week",
        "released for partner",
        "in-progress this week",
        "in progress this week",
    ):
        story = _story_after_header(paras, _find_para_index(paras, needle))
        if story is not None:
            return story
    for p in paras:
        if _is_story_bullet(p):
            return p
    return None


def _normalize_story_bullet_template(p_elem) -> None:
    """Force lvl-1 dash bullet, non-bold — matches G10X story line style."""
    from pptx.oxml import parse_xml

    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for run in p_elem.findall(qn("a:r")):
        r_pr = run.find(qn("a:rPr"))
        if r_pr is not None and r_pr.get("b") == "1":
            r_pr.set("b", "0")

    p_pr = p_elem.find(qn("a:pPr"))
    if p_pr is None:
        p_pr = parse_xml(f'<a:pPr xmlns:a="{ns}"/>')
        p_elem.insert(0, p_pr)
    p_pr.set("lvl", "1")
    for tag in ("a:buNone", "a:buAutoNum", "a:buBlip"):
        el = p_pr.find(qn(tag))
        if el is not None:
            p_pr.remove(el)

    bu = p_pr.find(qn("a:buChar"))
    if bu is None:
        p_pr.append(parse_xml(f'<a:buChar xmlns:a="{ns}" char="{STORY_BULLET_CHAR}"/>'))
    else:
        bu.set("char", STORY_BULLET_CHAR)
    _apply_g10x_paragraph_spacing(p_elem, is_story=True)


def _find_sprint_para_index(paras) -> int | None:
    for i, p in enumerate(paras):
        text = _normalize_match_text(paragraph_text(p)).strip()
        if not text or text == "spur":
            continue
        if "current week sprint status" in text:
            continue
        if text.startswith("sprint") or "sprint \u2013" in text or "sprint -" in text:
            return i
    return None


def _find_para_index(paras, *needles: str, start: int = 0) -> int | None:
    for i in range(start, len(paras)):
        text = _normalize_match_text(paragraph_text(paras[i]))
        if any(n in text for n in needles):
            return i
    return None


def _ensure_category_header_bullet(p_elem) -> None:
    from pptx.oxml import parse_xml

    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    pPr = p_elem.find(qn("a:pPr"))
    if pPr is None:
        pPr = parse_xml(f'<a:pPr xmlns:a="{ns}"/>')
        p_elem.insert(0, pPr)
    pPr.set("lvl", str(CATEGORY_HEADER_LEVEL))
    for tag in ("a:buNone", "a:buAutoNum", "a:buBlip"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)

    bu_clr = pPr.find(qn("a:buClr"))
    if bu_clr is None:
        pPr.append(
            parse_xml(
                f'<a:buClr xmlns:a="{ns}">'
                f'<a:srgbClr val="{CATEGORY_HEADER_BULLET_COLOR}"/></a:buClr>'
            )
        )

    bu_font = pPr.find(qn("a:buFont"))
    if bu_font is None:
        pPr.append(
            parse_xml(
                f'<a:buFont xmlns:a="{ns}" typeface="{CATEGORY_HEADER_BULLET_FONT}"/>'
            )
        )
    else:
        bu_font.set("typeface", CATEGORY_HEADER_BULLET_FONT)

    bu = pPr.find(qn("a:buChar"))
    if bu is None:
        pPr.append(
            parse_xml(
                f'<a:buChar xmlns:a="{ns}" char="{CATEGORY_HEADER_BULLET}"/>'
            )
        )
    else:
        bu.set("char", CATEGORY_HEADER_BULLET)


def _apply_g10x_paragraph_spacing(p_elem, is_story: bool = False) -> None:
    pPr = p_elem.find(qn("a:pPr"))
    if pPr is None:
        return
    if is_story:
        for tag in ("a:spcBef", "a:spcAft"):
            existing = pPr.find(qn(tag))
            if existing is not None:
                pPr.remove(existing)
        from pptx.oxml import parse_xml

        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pPr.insert(0, parse_xml(f'<a:spcBef xmlns:a="{ns}"><a:spcPts val="0"/></a:spcBef>'))


def find_canonical_style_slide_index(template_prs, profile: TemplateProfile | None = None) -> int:
    if profile is not None:
        for name, svc in profile.services.items():
            if "cost core" in name.lower():
                return svc.main_slide_index
    return 2 if len(template_prs.slides) > 2 else 0


def get_canonical_style_cell(template_prs, profile: TemplateProfile | None = None):
    idx = find_canonical_style_slide_index(template_prs, profile)
    hl = get_highlights_shape(template_prs.slides[idx])
    return hl.table.cell(2, 0)


def discover_section_templates(canonical_cell, global_completed=None) -> dict[str, Any]:
    """Clone sprint / category / story paragraph styles from the active template as-is."""
    paras = canonical_cell.text_frame._txBody.findall(qn("a:p"))
    story_bullet = _find_canonical_story_bullet(paras)

    sprint_idx = _find_sprint_para_index(paras)
    cw_idx = _find_para_index(paras, "current week sprint status")
    comp_idx = _find_para_index(paras, "completed this week")
    rel_idx = _find_para_index(paras, "released for partner review")
    prog_idx = _find_para_index(paras, "in-progress this week", "in progress this week")

    def pick(idx):
        if idx is not None and idx < len(paras):
            return paras[idx]
        return None

    def pick_bullet(hdr_idx):
        story = _story_after_header(paras, hdr_idx)
        if story is not None:
            return story
        return story_bullet

    required = {
        "sprint": pick(sprint_idx) if pick(sprint_idx) is not None else paras[0],
        "current_week": pick(cw_idx) if pick(cw_idx) is not None else paras[1],
        "completed_hdr": pick(comp_idx),
        "completed_bullet": pick_bullet(comp_idx) if pick_bullet(comp_idx) is not None else story_bullet,
        "released_hdr": pick(rel_idx),
        "released_bullet": pick_bullet(rel_idx) if pick_bullet(rel_idx) is not None else story_bullet,
        "inprogress_hdr": pick(prog_idx),
        "inprogress_bullet": pick_bullet(prog_idx) if pick_bullet(prog_idx) is not None else story_bullet,
    }
    if global_completed:
        if required["completed_hdr"] is None:
            required["completed_hdr"] = global_completed[0]
        if required["completed_bullet"] is None:
            required["completed_bullet"] = global_completed[1]

    fallback_hdr = (
        required["completed_hdr"]
        if required["completed_hdr"] is not None
        else required["released_hdr"]
        if required["released_hdr"] is not None
        else required["inprogress_hdr"]
    )
    fallback_bullet = (
        story_bullet
        if story_bullet is not None
        else required["released_bullet"]
        if required["released_bullet"] is not None
        else required["inprogress_bullet"]
    )

    for key, val in list(required.items()):
        if val is None:
            if key.endswith("_hdr"):
                required[key] = fallback_hdr
            elif key.endswith("_bullet"):
                required[key] = fallback_bullet
    if any(v is None for v in required.values()):
        raise RuntimeError("Could not resolve highlights section templates from template")

    result: dict[str, Any] = {}
    for key, val in required.items():
        result[key] = copy.deepcopy(val)
    return result


def _pick_para(paras: list, idx: int | None, fallback: int = 0):
    if idx is not None and 0 <= idx < len(paras):
        return paras[idx]
    return paras[min(fallback, len(paras) - 1)]


def section_templates_from_profile(layout_slide, service: ServiceProfile) -> dict[str, Any]:
    """Deep-copy paragraph templates from the service reference slide."""
    hl = get_highlights_shape(layout_slide)
    cell = hl.table.cell(service.hl_content_row, service.hl_content_col)
    paras = cell.text_frame._txBody.findall(qn("a:p"))
    roles = service.paragraph_roles

    story_para = _find_canonical_story_bullet(paras)
    if story_para is None:
        story_para = _pick_para(paras, roles.completed_bullet, 0)

    def hdr(key: str) -> Any:
        idx = getattr(roles, f"{key}_hdr")
        return copy.deepcopy(_pick_para(paras, idx, roles.completed_hdr or 2))

    def bullet(key: str) -> Any:
        idx = getattr(roles, f"{key}_bullet")
        hdr_idx = getattr(roles, f"{key}_hdr")
        story = _story_after_header(paras, hdr_idx)
        if story is not None:
            return copy.deepcopy(story)
        if idx is not None and 0 <= idx < len(paras):
            candidate = paras[idx]
            if _is_story_bullet(candidate):
                return copy.deepcopy(candidate)
        return copy.deepcopy(story_para)

    return {
        "sprint": copy.deepcopy(_pick_para(paras, roles.sprint, 0)),
        "current_week": copy.deepcopy(_pick_para(paras, roles.current_week, 1)),
        "completed_hdr": hdr("completed"),
        "completed_bullet": bullet("completed"),
        "released_hdr": hdr("released"),
        "released_bullet": bullet("released"),
        "inprogress_hdr": hdr("inprogress"),
        "inprogress_bullet": bullet("inprogress"),
    }


def set_two_run_header(p_elem, prefix: str, count_suffix: str) -> None:
    runs = p_elem.findall(qn("a:r"))
    if len(runs) >= 2:
        runs[0].find(qn("a:t")).text = prefix
        runs[1].find(qn("a:t")).text = f"\u2013 {count_suffix}"
        for extra in runs[2:]:
            p_elem.remove(extra)
    else:
        set_single_run_text(p_elem, f"{prefix}\u2013 {count_suffix}")


def _cleared_story_bullet_template(section_tmpl: dict, bucket: str):
    bullet = copy.deepcopy(section_tmpl[f"{bucket}_bullet"])
    clear_story_line_text(bullet)
    return bullet


def _make_sprint_gap_para(section_tmpl: dict) -> Any:
    gap = copy.deepcopy(section_tmpl["current_week"])
    set_single_run_text(gap, "")
    pPr = gap.find(qn("a:pPr"))
    if pPr is not None:
        for tag in ("a:buChar", "a:buAutoNum", "a:buBlip", "a:buFont"):
            el = pPr.find(qn(tag))
            if el is not None:
                pPr.remove(el)
        if pPr.find(qn("a:buNone")) is None:
            from pptx.oxml import parse_xml

            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            pPr.append(parse_xml(f'<a:buNone xmlns:a="{ns}"/>'))
    _apply_g10x_paragraph_spacing(gap)
    return gap


def replace_bullet_block(txBody, first_index: int, last_index: int, items: list[str]) -> None:
    all_p = txBody.findall(qn("a:p"))
    template_p = all_p[first_index]
    block = all_p[first_index : last_index + 1]
    anchor = block[-1]

    new_paragraphs = []
    for text in items:
        new_p = copy.deepcopy(template_p)
        set_story_line_text(new_p, text)
        new_paragraphs.append(new_p)

    prev = anchor
    for new_p in new_paragraphs:
        prev.addnext(new_p)
        prev = new_p

    for p in block:
        txBody.remove(p)


def strip_trailing_empty_paragraphs(txBody) -> None:
    all_p = txBody.findall(qn("a:p"))
    while len(all_p) > 2:
        if paragraph_text(all_p[-1]).strip():
            break
        txBody.remove(all_p[-1])
        all_p = txBody.findall(qn("a:p"))


def set_sprint_line(p_elem, sprint_bold: str, sprint_light: str) -> None:
    runs = p_elem.findall(qn("a:r"))
    if len(runs) >= 2:
        runs[0].find(qn("a:t")).text = sprint_bold
        runs[1].find(qn("a:t")).text = sprint_light
        for extra in runs[2:]:
            p_elem.remove(extra)
    elif runs:
        set_single_run_text(p_elem, sprint_bold + sprint_light)


def _advance_offset_after_bullets(txBody, header_offset: int, items: list, has_category_header: bool = True) -> int:
    if not items:
        all_p = txBody.findall(qn("a:p"))
        return min(header_offset + (1 if has_category_header else 0), len(all_p))
    bullet_index = header_offset + (1 if has_category_header else 0)
    replace_bullet_block(txBody, bullet_index, bullet_index, items)
    return bullet_index + len(items)


def _append_bucket_templates(new_paras, section_tmpl, section, omit=()) -> None:
    for bucket in STORY_BUCKET_ORDER:
        items = section.get(bucket) or section.get(f"{bucket}_items")
        if not items:
            continue
        if bucket in omit:
            new_paras.append(_cleared_story_bullet_template(section_tmpl, bucket))
            continue
        new_paras.extend([
            copy.deepcopy(section_tmpl[f"{bucket}_hdr"]),
            _cleared_story_bullet_template(section_tmpl, bucket),
        ])


def _fill_story_buckets(hl_txBody, offset: int, section, omit=(), display: bool = False) -> int:
    for bucket in STORY_BUCKET_ORDER:
        if display:
            items = section.get(f"{bucket}_items") or []
            count = section.get(f"{bucket}_count", str(len(items)))
        else:
            items = section.get(bucket) or []
            count = str(len(items))
        if not items:
            continue

        all_p = hl_txBody.findall(qn("a:p"))
        has_header = bucket not in omit
        if has_header:
            hdr_p = all_p[offset]
            if bucket == "completed":
                set_two_run_header(hdr_p, "Stories completed this week ", f"{count} stories")
            elif bucket == "released":
                set_two_run_header(hdr_p, "Stories released for partner review ", f"{count} stories")
            else:
                set_single_run_text(hdr_p, f"Stories in-progress this week \u2013 {count} stories")
        offset = _advance_offset_after_bullets(
            hl_txBody, offset, items, has_category_header=has_header
        )
    return offset


def set_cell_top_align(cell) -> None:
    """Top-align text in a table cell (G10X HL content always starts flush below header)."""
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.set("anchor", "t")
    bodyPr = cell.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        from pptx.oxml import parse_xml

        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        bodyPr = parse_xml(f'<a:bodyPr xmlns:a="{ns}" anchor="t"/>')
        cell.text_frame._txBody.insert(0, bodyPr)
    else:
        bodyPr.set("anchor", "t")


def top_align_highlights_on_slide(slide) -> None:
    set_cell_top_align(highlights_content_cell(slide))


def clear_highlights_content_cell(slide) -> None:
    """Remove all paragraph content from the Highlights cell."""
    cell = highlights_content_cell(slide)
    hl_txBody = cell.text_frame._txBody
    for p in list(hl_txBody.findall(qn("a:p"))):
        hl_txBody.remove(p)


def _delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def remove_ka_from_slide(slide) -> None:
    """Remove standalone KA shape or clear embedded KA rows — for HL-only contd slides."""
    ka = get_key_activities_shape(slide)
    if ka is not None:
        _delete_shape(ka)
        return
    clear_embedded_ka_in_hl(slide)


def remove_hl_from_slide(slide) -> None:
    """Remove the Highlights table — for KA-only continuation slides."""
    try:
        hl = get_highlights_shape(slide)
    except ValueError:
        return
    _delete_shape(hl)


def prepare_contd_slide_for_highlights(slide) -> None:
    """Strip template leftovers before filling HL-only continuation content."""
    remove_ka_from_slide(slide)
    try:
        clear_highlights_content_cell(slide)
    except ValueError:
        pass


def populate_highlights_cell(hl_cell, section_tmpl: dict, content: dict) -> None:
    hl_txBody = hl_cell.text_frame._txBody
    for p in list(hl_txBody.findall(qn("a:p"))):
        hl_txBody.remove(p)

    sections = content.get("sections")
    if sections is None:
        sections = [content]
    if not sections:
        return
    new_paras = []
    for si, sec in enumerate(sections):
        if si > 0:
            new_paras.append(_make_sprint_gap_para(section_tmpl))
        new_paras.append(copy.deepcopy(section_tmpl["sprint"]))
        new_paras.append(copy.deepcopy(section_tmpl["current_week"]))
        _append_bucket_templates(new_paras, section_tmpl, sec)

    for p in new_paras:
        hl_txBody.append(p)

    offset = 0
    for si, sec in enumerate(sections):
        if si > 0:
            offset += 1
        all_p = hl_txBody.findall(qn("a:p"))
        set_sprint_line(all_p[offset], sec["sprint_bold"], sec["sprint_light"])
        offset += 2
        offset = _fill_story_buckets(hl_txBody, offset, sec, display=True)

    strip_trailing_empty_paragraphs(hl_txBody)


def _populate_contd_sections(hl_cell, section_tmpl: dict, sections: list) -> None:
    hl_txBody = hl_cell.text_frame._txBody
    for p in list(hl_txBody.findall(qn("a:p"))):
        hl_txBody.remove(p)

    new_paras = []
    for si, section in enumerate(sections):
        omit = set(section.get("omit_category_headers", []))
        continued = section.get("continued_section")
        if si > 0 and not continued:
            new_paras.append(_make_sprint_gap_para(section_tmpl))
        if not continued:
            new_paras.append(copy.deepcopy(section_tmpl["sprint"]))
            new_paras.append(copy.deepcopy(section_tmpl["current_week"]))
        _append_bucket_templates(new_paras, section_tmpl, section, omit=omit)

    for p in new_paras:
        hl_txBody.append(p)

    offset = 0
    for si, section in enumerate(sections):
        omit = set(section.get("omit_category_headers", []))
        continued = section.get("continued_section")
        if si > 0 and not continued:
            offset += 1
        if not continued:
            all_p = hl_txBody.findall(qn("a:p"))
            set_sprint_line(all_p[offset], section["sprint_bold"], section["sprint_light"])
            offset += 2
        offset = _fill_story_buckets(hl_txBody, offset, section, omit=omit, display=True)

    strip_trailing_empty_paragraphs(hl_txBody)


def populate_highlights_contd_cell(hl_cell, section_tmpl: dict, contd_raw: dict) -> None:
    if contd_raw.get("sections"):
        _populate_contd_sections(hl_cell, section_tmpl, contd_raw["sections"])
        return

    hl_txBody = hl_cell.text_frame._txBody
    for p in list(hl_txBody.findall(qn("a:p"))):
        hl_txBody.remove(p)

    omit = set(contd_raw.get("omit_category_headers", []))
    new_paras = []
    _append_bucket_templates(new_paras, section_tmpl, contd_raw, omit=omit)
    for p in new_paras:
        hl_txBody.append(p)

    _fill_story_buckets(hl_txBody, 0, contd_raw, omit=omit)
    strip_trailing_empty_paragraphs(hl_txBody)


def set_ka_empty_items(ka_shape, placeholder_rows: int | None = None) -> None:
    """Leave KA with blank placeholder bullets for manual BSA entry (matches v1)."""
    from app.constants.ppt_mapping import KA_PLACEHOLDER_ROWS

    count = placeholder_rows if placeholder_rows is not None else KA_PLACEHOLDER_ROWS
    items = [""] * max(count, 1)
    tx_body = ka_shape.table.cell(1, 0).text_frame._txBody
    all_p = tx_body.findall(qn("a:p"))
    if not all_p:
        return
    replace_ka_bullet_block(tx_body, 0, len(all_p) - 1, items)


def resolve_ka_items(key_activities: list[str]) -> list[str]:
    """Return KA bullets; empty content -> placeholder rows for manual entry."""
    items = [str(item).strip() for item in key_activities if str(item).strip()]
    if items:
        return items
    from app.constants.ppt_mapping import KA_PLACEHOLDER_ROWS

    return [""] * max(KA_PLACEHOLDER_ROWS, 1)


def set_ka_items(ka_shape, items: list[str]) -> None:
    tx_body = ka_shape.table.cell(1, 0).text_frame._txBody
    all_p = tx_body.findall(qn("a:p"))
    if not all_p:
        return
    replace_ka_bullet_block(tx_body, 0, len(all_p) - 1, items)


def populate_embedded_ka_in_hl(
    slide,
    items: list[str],
    ka_header_row: int | None = None,
) -> None:
    """Fill embedded KA rows inside a combined Highlights table."""
    hl = get_highlights_shape(slide)
    table = hl.table
    if ka_header_row is None:
        for ri in range(len(table.rows)):
            try:
                label = table.cell(ri, 0).text.strip().lower()
            except (IndexError, AttributeError):
                continue
            if "key activit" in label:
                ka_header_row = ri
                break
    if ka_header_row is None:
        return

    content_start = ka_header_row + 1
    if content_start >= len(table.rows):
        return

    cell = table.cell(content_start, 0)
    tx_body = cell.text_frame._txBody
    all_p = tx_body.findall(qn("a:p"))
    if all_p:
        replace_ka_bullet_block(tx_body, 0, len(all_p) - 1, items)
    elif items:
        cell.text = items[0]

    for ri in range(content_start + 1, len(table.rows)):
        table.cell(ri, 0).text = ""


def populate_ka_on_slide(
    slide,
    key_activities: list[str],
    service: ServiceProfile | None = None,
) -> None:
    """Write KA bullets on standalone or embedded layouts."""
    items = resolve_ka_items(key_activities)
    ka = get_key_activities_shape(slide)
    if ka is not None:
        set_ka_items(ka, items)
        return
    if service is not None and service.ka_mode == "embedded":
        populate_embedded_ka_in_hl(slide, items, service.ka_embedded_start_row)


def clear_embedded_ka_in_hl(slide) -> None:
    """Clear KA rows embedded inside a combined Highlights table."""
    try:
        hl = get_highlights_shape(slide)
    except ValueError:
        return
    table = hl.table
    ka_header_row: int | None = None
    for ri in range(len(table.rows)):
        try:
            label = table.cell(ri, 0).text.strip().lower()
        except (IndexError, AttributeError):
            continue
        if "key activit" in label:
            ka_header_row = ri
            break
    if ka_header_row is None:
        return
    for ri in range(ka_header_row + 1, len(table.rows)):
        table.cell(ri, 0).text = ""


def clear_ka_on_slide(slide, template_slide=None, service: ServiceProfile | None = None) -> None:
    """Clear KA to blank placeholders on any slide layout (standalone or embedded)."""
    ka = get_key_activities_shape(slide)
    if ka is not None:
        set_ka_empty_items(ka)
        return
    if service is not None and service.ka_mode == "embedded":
        clear_embedded_ka_in_hl(slide)
        return
    clear_embedded_ka_in_hl(slide)


def clear_ka_on_main_slide(prs, project_name: str, template_slide=None, service=None) -> None:
    """Clear KA placeholders on the main slide only — contd slides are HL-only."""
    from app.wsr_engine.project_deletion import find_main_slide_index

    main_idx = find_main_slide_index(prs, project_name)
    if main_idx is not None and main_idx < len(prs.slides):
        clear_ka_on_slide(prs.slides[main_idx], template_slide, service)


def remove_ka_from_project_contd_slides(prs, project_name: str) -> None:
    """Strip KA from every (Contd..) slide for a project."""
    from app.services.ppt_shape_utils import slide_title_text
    from app.wsr_engine.project_deletion import find_main_slide_index, find_project_slide_indices

    main_idx = find_main_slide_index(prs, project_name)
    for idx in find_project_slide_indices(prs, project_name):
        if idx == main_idx:
            continue
        if "(contd" in slide_title_text(prs.slides[idx]).lower():
            remove_ka_from_slide(prs.slides[idx])


def clear_ka_on_project_slides(prs, project_name: str, template_slide=None, service=None) -> None:
    """Clear KA on main slide; remove KA from contd slides."""
    clear_ka_on_main_slide(prs, project_name, template_slide, service)
    remove_ka_from_project_contd_slides(prs, project_name)


def fill_highlights_on_slide(
    slide,
    section_tmpl: dict,
    sections: list[dict],
) -> None:
    cell = highlights_content_cell(slide)
    populate_highlights_cell(cell, section_tmpl, {"sections": sections})


def fill_contd_highlights_on_slide(
    slide,
    section_tmpl: dict,
    sections: list[dict],
) -> None:
    prepare_contd_slide_for_highlights(slide)
    cell = highlights_content_cell(slide)
    populate_highlights_contd_cell(cell, section_tmpl, {"sections": sections})
