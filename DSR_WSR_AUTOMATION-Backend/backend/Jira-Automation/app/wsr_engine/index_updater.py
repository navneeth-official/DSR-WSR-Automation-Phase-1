"""Generic index slide reflow and slide-number updates."""

from __future__ import annotations

import copy
import logging
import re

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from rapidfuzz import fuzz

from app.services.ppt_shape_utils import normalize_title_text, slide_title_text, iter_all_shapes
from app.wsr_engine.index_layout import INDEX_ENTRY_RULES, IndexLayout, discover_index_layout
from app.wsr_engine.models import ProjectMap

logger = logging.getLogger(__name__)

_INDEX_NUMBER_RE = re.compile(r"^\d{1,2}\u200b?$")

_INDEX_LABEL_KEYWORDS: tuple[str, ...] = (
    "cost core",
    "supplier core",
    "pricing core",
    "wentworth",
    "location core",
    "pharmacy",
    "wellness",
    "global sourcing",
    "product attribute",
    "loco",
    "bsa",
    "matters of attention",
    "team allocation",
)

# Keyword rules re-exported from index_layout for backwards compatibility.

_DELIVERY_INDEX_KEYWORDS = tuple(
    kw for kw in _INDEX_LABEL_KEYWORDS if kw not in ("matters of attention", "team allocation")
)


def _normalize_index_text(text: str) -> str:
    return normalize_title_text(text).lower()


def _find_index_slide_index(prs: Presentation) -> int | None:
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if _normalize_index_text(shape.text_frame.text) == "index":
                return i
    return None


def _find_index_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def _index_table_cells_row_major(table):
    cells = []
    for row in table.rows:
        cells.extend(row.cells)
    return cells


def _cell_has_index_content(cell) -> bool:
    text = _normalize_index_text(cell.text_frame.text)
    if not text:
        return False
    if any(kw in text for kw in _INDEX_LABEL_KEYWORDS):
        return True
    for paragraph in cell.text_frame.paragraphs:
        if _is_index_number_paragraph(paragraph):
            return True
    return False


def _is_index_number_paragraph(paragraph) -> bool:
    text = normalize_title_text("".join(run.text for run in paragraph.runs))
    return bool(_INDEX_NUMBER_RE.match(text))


def _set_paragraph_slide_number(paragraph, slide_number: int) -> None:
    display = f"{slide_number:02d}"
    if paragraph.runs:
        paragraph.runs[0].text = display
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = display


def _set_element_hyperlink_to_slide(hlink_elem, index_slide, target_slide_idx, prs) -> None:
    target_part = prs.slides[target_slide_idx].part
    slide_rid = index_slide.part.relate_to(target_part, RT.SLIDE)
    hlink_elem.set(qn("r:id"), slide_rid)


def _update_cell_index_entry(cell, index_slide, target_slide_idx, prs) -> bool:
    changed = False
    display_num = target_slide_idx + 1

    for paragraph in cell.text_frame.paragraphs:
        if _is_index_number_paragraph(paragraph):
            old = normalize_title_text("".join(run.text for run in paragraph.runs))
            new = f"{display_num:02d}"
            if old != new:
                _set_paragraph_slide_number(paragraph, display_num)
                changed = True

    tx_body = cell.text_frame._txBody
    for p_elem in tx_body.findall(qn("a:p")):
        for hlink in p_elem.iter(qn("a:hlinkClick")):
            _set_element_hyperlink_to_slide(hlink, index_slide, target_slide_idx, prs)
            changed = True

    return changed


def _clear_index_cell_completely(cell) -> None:
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""
    tx_body = cell.text_frame._txBody
    for p_elem in tx_body.findall(qn("a:p")):
        for hlink in list(p_elem.findall(qn("a:hlinkClick"))):
            p_elem.remove(hlink)


def _clone_cell_text_body(dst_cell, src_tx_body) -> None:
    dst_tc = dst_cell._tc
    old_tx = dst_tc.find(qn("a:txBody"))
    if old_tx is not None:
        dst_tc.remove(old_tx)
    dst_tc.insert(0, copy.deepcopy(src_tx_body))


def _resolve_index_target(
    prs: Presentation,
    cell_text: str,
    projects: list[ProjectMap],
) -> int | None:
    label = _normalize_index_text(cell_text)
    if not label:
        return None

    # Skip non-delivery index rows (team allocation, RAID, etc.)
    if "team allocation" in label or "raid status" in label:
        return None

    best_score = 0
    best_idx: int | None = None

    for proj in projects:
        score = fuzz.partial_ratio(label, proj.project_name.lower())
        if score > best_score and score >= 70:
            best_score = score
            best_idx = proj.main_slide_index

    if best_idx is not None and best_idx < len(prs.slides):
        return best_idx

    for i, slide in enumerate(prs.slides):
        title = normalize_title_text(slide_title_text(slide)).lower()
        if not title or "(contd" in title:
            continue
        score = fuzz.partial_ratio(label, title)
        if score > best_score and score >= 70:
            best_score = score
            best_idx = i

    return best_idx


def _collect_active_index_entries(prs: Presentation, table, projects: list[ProjectMap]) -> list[tuple[int, object]]:
    entries: list[tuple[int, object]] = []
    seen_targets: set[int] = set()

    for cell in _index_table_cells_row_major(table):
        if not _cell_has_index_content(cell):
            continue
        cell_text = _normalize_index_text(cell.text_frame.text)
        target_idx = _resolve_index_target(prs, cell_text, projects)
        if target_idx is None or target_idx in seen_targets:
            continue
        seen_targets.add(target_idx)
        entries.append((target_idx, copy.deepcopy(cell.text_frame._txBody)))

    return entries


def _is_static_index_cell(cell) -> bool:
    text = _normalize_index_text(cell.text_frame.text)
    return "matters of attention" in text or "team allocation" in text


def _cell_has_delivery_keyword(cell) -> bool:
    text = _normalize_index_text(cell.text_frame.text)
    return any(kw in text for kw in _DELIVERY_INDEX_KEYWORDS)


def _is_delivery_index_cell(cell) -> bool:
    if _is_static_index_cell(cell):
        return False
    text = _normalize_index_text(cell.text_frame.text)
    if not text:
        return False
    if any(kw in text for kw in _DELIVERY_INDEX_KEYWORDS):
        return True
    return _cell_has_index_content(cell)


def _find_index_cell_by_needles(table, needles: tuple[str, ...]):
    for cell in _index_table_cells_row_major(table):
        text = _normalize_index_text(cell.text_frame.text)
        if all(needle in text for needle in needles):
            return cell
    return None


def _index_cell_at(table, row: int, col: int):
    if row < 0 or col < 0 or row >= len(table.rows):
        return None
    if col >= len(table.columns):
        return None
    return table.cell(row, col)


def _write_index_cell(
    cell,
    index_slide,
    prs: Presentation,
    target_idx: int,
    label: str,
    touched: set[int],
) -> bool:
    _set_index_cell_content(cell, target_idx + 1, label)
    changed = _update_cell_index_entry(cell, index_slide, target_idx, prs)
    touched.add(_cell_identity(cell))
    return changed


def _write_static_index_cell(
    table,
    index_slide,
    prs: Presentation,
    row: int,
    col: int,
    title_fragment: str,
    label: str,
    touched: set[int],
) -> bool:
    cell = _index_cell_at(table, row, col)
    if cell is None:
        return False
    target_idx = _find_static_slide_index(prs, title_fragment)
    if target_idx is None:
        return False
    return _write_index_cell(cell, index_slide, prs, target_idx, label, touched)


def _match_project_for_index(projects: list[ProjectMap], search_title: str) -> ProjectMap | None:
    needle = search_title.lower()
    if needle == "loco":
        for proj in projects:
            name = proj.project_name.lower()
            if "loco" in name and "location core" not in name:
                return proj
        return None

    best: ProjectMap | None = None
    best_score = 0
    for proj in projects:
        name = proj.project_name.lower()
        if needle == "location core service" and "loco" in name and "location core" not in name:
            continue
        score = fuzz.partial_ratio(needle, name)
        if score > best_score and score >= 85:
            best_score = score
            best = proj
    return best


def _find_static_slide_index(prs: Presentation, title_fragment: str) -> int | None:
    needle = title_fragment.lower()
    index_idx = _find_index_slide_index(prs)
    for i, slide in enumerate(prs.slides):
        if index_idx is not None and i == index_idx:
            continue
        title = normalize_title_text(slide_title_text(slide)).lower()
        if title and needle in title and "delivery status" not in title:
            return i
        for shape in iter_all_shapes(slide.shapes):
            if shape.has_text_frame:
                body = normalize_title_text(shape.text_frame.text).lower()
                if needle in body and "delivery status" not in body:
                    return i
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        body = normalize_title_text(cell.text_frame.text).lower()
                        if needle in body and "delivery status" not in body:
                            return i
    return None


def _extract_index_label(cell) -> str:
    for paragraph in cell.text_frame.paragraphs:
        if _is_index_number_paragraph(paragraph):
            continue
        label = normalize_title_text("".join(run.text for run in paragraph.runs))
        if label:
            return label
    return ""


def _set_index_cell_content(cell, slide_number: int, label: str) -> None:
    has_number = any(
        _is_index_number_paragraph(paragraph)
        for paragraph in cell.text_frame.paragraphs
    )
    if not has_number:
        number_paragraph = cell.text_frame.paragraphs[0]
        _set_paragraph_slide_number(number_paragraph, slide_number)

    label_written = False
    for paragraph in cell.text_frame.paragraphs:
        if _is_index_number_paragraph(paragraph):
            _set_paragraph_slide_number(paragraph, slide_number)
            continue
        if not label_written:
            if paragraph.runs:
                paragraph.runs[0].text = label
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.add_run().text = label
            label_written = True
        else:
            for run in paragraph.runs:
                run.text = ""


def _cell_identity(cell) -> int:
    """Stable id for merged table cells (same tc may appear at multiple indices)."""
    return id(cell._tc)


def _collect_rebuild_entries(
    prs: Presentation,
    projects: list[ProjectMap],
    layout: IndexLayout,
    content_labels: dict[str, str],
) -> list[tuple[int, str]]:
    """Ordered (target_slide_idx, label) pairs following template slot order."""
    entries: list[tuple[int, str]] = []
    for slot in layout.slots:
        if slot.is_delivery:
            proj = _match_project_for_index(projects, slot.search_title)
            if proj is None:
                continue
            label = content_labels.get(proj.project_name, proj.project_name)
            entries.append((proj.main_slide_index, label))
            continue
        target_idx = _find_static_slide_index(prs, slot.search_title)
        if target_idx is None:
            continue
        entries.append((target_idx, slot.search_title))
    return entries


def rebuild_index_from_projects(
    prs: Presentation,
    projects: list[ProjectMap],
    content_labels: dict[str, str] | None = None,
    layout: IndexLayout | None = None,
) -> int:
    """Rebuild index by compacting entries row-major with column-aligned formatting."""
    index_idx = _find_index_slide_index(prs)
    if index_idx is None:
        logger.warning("Index slide not found")
        return 0

    index_slide = prs.slides[index_idx]
    table = _find_index_table(index_slide)
    if table is None:
        logger.warning("Index table not found")
        return 0

    slot_cells = _index_table_cells_row_major(table)
    delivery_slots = [cell for cell in slot_cells if _is_delivery_index_cell(cell)]
    if not delivery_slots and not any(_is_static_index_cell(c) for c in slot_cells):
        logger.warning("No index slots found; falling back to reflow")
        return reflow_index_slide(prs, projects)

    if layout is None or not layout.slots:
        layout = discover_index_layout(table)

    content_labels = content_labels or {}
    entries = _collect_rebuild_entries(prs, projects, layout, content_labels)
    if not entries:
        return 0

    rows = len(table.rows)
    cols = len(table.columns)
    positions = layout.grid_positions(rows, cols)[: len(entries)]

    for cell in slot_cells:
        if (
            _cell_has_index_content(cell)
            or _cell_has_delivery_keyword(cell)
            or _is_static_index_cell(cell)
        ):
            _clear_index_cell_completely(cell)

    updated = 0
    for (target_idx, label), (row, col) in zip(entries, positions):
        cell = _index_cell_at(table, row, col)
        if cell is None:
            continue
        proto = layout.column_prototypes.get(col)
        if proto is not None:
            _clone_cell_text_body(cell, proto)
        _set_index_cell_content(cell, target_idx + 1, label)
        if _update_cell_index_entry(cell, index_slide, target_idx, prs):
            updated += 1

    logger.info("Index rebuild: %d entries updated (compact grid)", updated)
    return updated


def reflow_index_slide(prs: Presentation, projects: list[ProjectMap]) -> int:
    index_idx = _find_index_slide_index(prs)
    if index_idx is None:
        logger.warning("Index slide not found")
        return 0

    index_slide = prs.slides[index_idx]
    table = _find_index_table(index_slide)
    if table is None:
        logger.warning("Index table not found")
        return 0

    active_entries = _collect_active_index_entries(prs, table, projects)
    slot_cells = _index_table_cells_row_major(table)

    for cell in slot_cells:
        if _cell_has_index_content(cell):
            _clear_index_cell_completely(cell)

    updated = 0
    for slot_idx, (target_idx, tx_body) in enumerate(active_entries):
        if slot_idx >= len(slot_cells):
            break
        cell = slot_cells[slot_idx]
        _clone_cell_text_body(cell, tx_body)
        if _update_cell_index_entry(cell, index_slide, target_idx, prs):
            updated += 1

    logger.info("Index reflow: %d entries updated", updated)
    return updated
