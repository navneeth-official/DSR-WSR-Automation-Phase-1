"""Reorder delivery slides and remove orphan continuation slides."""

from __future__ import annotations

import logging

from pptx import Presentation

from app.services.ppt_shape_utils import (
    get_highlights_shape,
    get_key_activities_shape,
    is_delivery_slide_title,
    paragraph_text,
    service_suffix_from_title,
    slide_title_text,
)
from app.wsr_engine.continuation_engine import find_contd_slides_for_project
from app.wsr_engine.project_deletion import find_main_slide_index
from app.wsr_engine.slide_ops import delete_slide, move_slide_after
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)


def remove_project_contd_slides(prs: Presentation, project_name: str) -> int:
    """Delete all (Contd..) slides for a project."""
    indices = find_contd_slides_for_project(prs, project_name)
    for idx in sorted(indices, reverse=True):
        delete_slide(prs, idx)
    return len(indices)


def finalize_slide_order(prs: Presentation, project_names: list[str]) -> None:
    """Place each project's (Contd..) slide(s) immediately after its main slide."""
    for name in reversed(project_names):
        for offset in range(32):
            main_idx = find_main_slide_index(prs, name)
            if main_idx is None:
                break
            contd_indices = find_contd_slides_for_project(prs, name)
            if offset >= len(contd_indices):
                break
            contd_idx = contd_indices[offset]
            desired = main_idx + 1 + offset
            if contd_idx != desired:
                move_slide_after(prs, contd_idx, desired - 1)


def _slide_has_hl_content(slide) -> bool:
    try:
        hl = get_highlights_shape(slide)
        cell = hl.table.cell(2, 0)
        for p in cell.text_frame._txBody.findall(qn("a:p")):
            if paragraph_text(p).strip():
                return True
    except (ValueError, IndexError, AttributeError):
        pass
    return False


def cleanup_orphan_contd_slides(prs: Presentation) -> int:
    """Remove empty (Contd..) slides with no highlights or KA body."""
    to_delete: list[int] = []
    for i, slide in enumerate(prs.slides):
        title = slide_title_text(slide)
        if "(contd" not in title.lower():
            continue
        has_hl = False
        has_ka = get_key_activities_shape(slide) is not None
        try:
            get_highlights_shape(slide)
            has_hl = True
        except ValueError:
            pass
        if not has_hl and not has_ka:
            to_delete.append(i)
            continue
        if has_hl and not _slide_has_hl_content(slide) and not has_ka:
            to_delete.append(i)
    removed = 0
    for idx in sorted(set(to_delete), reverse=True):
        delete_slide(prs, idx)
        removed += 1
    if removed:
        logger.info("Removed %d orphan contd slide(s)", removed)
    return removed


def delete_unmatched_delivery_slides(prs: Presentation, matched_names: set[str]) -> int:
    """Remove delivery slides whose project is not in matched_names (by title)."""
    to_delete: list[int] = []
    matched_lower = {n.lower() for n in matched_names}

    for i, slide in enumerate(prs.slides):
        title = slide_title_text(slide)
        if not is_delivery_slide_title(title):
            continue
        suffix = service_suffix_from_title(title).lower()
        if any(_names_match(suffix, name) for name in matched_lower):
            continue
        to_delete.append(i)

    removed = 0
    for idx in sorted(set(to_delete), reverse=True):
        delete_slide(prs, idx)
        removed += 1
    return removed


def _names_match(suffix: str, matched: str) -> bool:
    return suffix == matched or suffix in matched or matched in suffix
