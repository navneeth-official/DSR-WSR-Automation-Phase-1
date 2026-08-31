"""Remove template projects that have no matching content."""

from __future__ import annotations

import logging

from pptx import Presentation

from app.services.ppt_shape_utils import (
    is_delivery_slide_title,
    normalize_title_text,
    service_suffix_from_title,
    slide_title_text,
)
from app.wsr_engine.models import ProjectMap, TemplateModel
from app.wsr_engine.slide_ops import delete_slide

logger = logging.getLogger(__name__)


def _service_in_title(title: str, project_name: str) -> bool:
    suffix = service_suffix_from_title(title).lower()
    needle = project_name.lower()
    return needle in suffix or suffix in needle


def collect_slides_to_delete(
    template: TemplateModel,
    matched_template_names: set[str],
) -> list[int]:
    to_delete: set[int] = set()
    for proj in template.projects:
        if proj.project_name in matched_template_names:
            continue
        for idx in proj.all_slide_indices:
            to_delete.add(idx)
    return sorted(to_delete, reverse=True)


def delete_unmatched_projects(
    prs: Presentation,
    template: TemplateModel,
    matched_template_names: set[str],
) -> int:
    indices = collect_slides_to_delete(template, matched_template_names)
    removed = 0
    for idx in indices:
        if idx < len(prs.slides):
            delete_slide(prs, idx)
            removed += 1
            logger.info("Deleted slide %d for unmatched project", idx)
    return removed


def find_project_slide_indices(prs: Presentation, project_name: str) -> list[int]:
    indices: list[int] = []
    for i, slide in enumerate(prs.slides):
        title = slide_title_text(slide)
        if not is_delivery_slide_title(title):
            continue
        if _service_in_title(title, project_name):
            indices.append(i)
    return indices


def find_main_slide_index(prs: Presentation, project_name: str) -> int | None:
    """Return the main (non-contd) delivery slide index for a project."""
    for idx in find_project_slide_indices(prs, project_name):
        title = slide_title_text(prs.slides[idx]).lower()
        if "(contd" not in title:
            return idx
    indices = find_project_slide_indices(prs, project_name)
    return indices[0] if indices else None


def refresh_project_maps_after_deletion(
    prs: Presentation,
    project_names: list[str],
) -> list[ProjectMap]:
    projects: list[ProjectMap] = []

    for name in project_names:
        indices = find_project_slide_indices(prs, name)
        if not indices:
            continue

        main_idx = find_main_slide_index(prs, name)
        contd: list[int] = []
        if main_idx is not None:
            for idx in indices:
                if idx == main_idx:
                    continue
                title = slide_title_text(prs.slides[idx]).lower()
                if "(contd" in title:
                    contd.append(idx)

        if main_idx is None:
            main_idx = indices[0]

        projects.append(
            ProjectMap(
                project_name=name,
                main_slide_index=main_idx,
                continuation_indices=sorted(contd),
            )
        )

    return projects
