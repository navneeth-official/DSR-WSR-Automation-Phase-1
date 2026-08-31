"""Height-based overflow planning for Highlights content."""

from __future__ import annotations

from typing import Any

from pptx.oxml.ns import qn
from pptx.slide import Slide

from pptx import Presentation

from app.services.ppt_layout_metrics import hl_ka_fits_on_main_slide
from app.services.ppt_shape_utils import get_highlights_shape, get_key_activities_shape, iter_all_shapes, paragraph_text
from app.services.template_profile import ServiceProfile, TemplateProfile
from app.wsr_engine.hl_ka_normalizer import find_standalone_reference_slide
from app.wsr_engine.models import OverflowPlan, ProjectContent
from app.wsr_engine.content_parser import section_display_content
from app.wsr_engine.placeholder_locator import count_template_paragraph_slots

STORY_BUCKET_ORDER = ("completed", "released", "inprogress")
MAX_KA_BOTTOM_EMU = 5747107
# Universal HL-only line budget (main + contd slides without KA tab).
UNIVERSAL_HL_ONLY_CAP = 26
# When HL line count is below this, HL and KA share the main slide.
HL_KA_SAME_SLIDE_MAX = 20
_CHARS_PER_LINE = {0: 92, 1: 72, 7: 80}
_DEFAULT_CHARS = 85


def _text_visual_lines(text: str, level: int = 1) -> int:
    t = (text or "").strip()
    if not t:
        return 0
    chars = _CHARS_PER_LINE.get(level, _DEFAULT_CHARS)
    return max(1, (len(t) + chars - 1) // chars)


def _section_paragraph_count(section: dict, display: bool = False) -> int:
    """Wrap-aware slot estimate — counts visual lines, not just paragraph slots."""
    continued = section.get("continued_section")
    omit = set(section.get("omit_category_headers") or [])
    n = 0 if continued else 2
    if not continued:
        if display:
            sprint = (section.get("sprint_bold") or "") + (section.get("sprint_light") or "")
        else:
            sprint = ""
        if sprint.strip():
            n += _text_visual_lines(sprint, level=0)
        n += 1  # current week line
    for bucket in STORY_BUCKET_ORDER:
        if display:
            items = section.get(f"{bucket}_items") or []
        else:
            items = section.get(bucket) or []
        if not items:
            continue
        if bucket not in omit:
            n += 1
        for item in items:
            n += _text_visual_lines(item, level=1)
    return n


def build_layout_profile(slide: Slide) -> dict[str, Any]:
    ref_hl = None
    ref_ka = None
    for shape in iter_all_shapes(slide.shapes):
        if not shape.has_table:
            continue
        try:
            header = shape.table.cell(0, 0).text.strip().lower()
        except (IndexError, AttributeError):
            continue
        if ref_hl is None and "highlights" in header:
            ref_hl = shape
        if ref_ka is None and "key activit" in header:
            ref_ka = shape

    if ref_hl is None:
        ref_hl = get_highlights_shape(slide)

    r0 = ref_hl.table.rows[0].height
    r1 = ref_hl.table.rows[1].height
    ref_r2 = ref_hl.table.rows[2].height
    ref_pad = ref_hl.height - r0 - r1 - ref_r2
    cell = ref_hl.table.cell(2, 0)
    hl_paras = cell.text_frame._txBody.findall(qn("a:p"))
    ref_para_count = max(len(hl_paras), 1)

    profile: dict[str, Any] = {
        "ref_hl": ref_hl,
        "ref_ka": ref_ka,
        "r0": r0,
        "r1": r1,
        "ref_r2": ref_r2,
        "ref_pad": ref_pad,
        "ref_para_count": ref_para_count,
        "canonical_para_count": ref_para_count,
        "canonical_per_line_emu": int(ref_r2 / max(ref_para_count - 2, 1)),
        "ref_hl_height": ref_hl.height,
        "ref_hl_top": ref_hl.top,
    }

    if ref_ka:
        expanded_hl_height = ref_ka.top + ref_ka.height - ref_hl.top
        profile["expanded_hl_height"] = expanded_hl_height
        profile["ref_ka_height"] = int(ref_ka.height)

    return profile


def _capacity_for_content_row(profile: dict, content_row_h: int) -> int:
    story_slots_ref = max(profile["ref_para_count"] - 2, 1)
    per_story = profile["ref_r2"] / story_slots_ref
    if per_story <= 0:
        return profile["ref_para_count"]
    story_slots = max(int(content_row_h / per_story), 1)
    return story_slots + 2


def get_hl_budgets(profile: dict) -> dict[str, Any]:
    ref_hl = profile.get("ref_hl")
    ref_ka = profile.get("ref_ka")
    r0, r1, pad = profile["r0"], profile["r1"], profile["ref_pad"]

    template_h = profile["ref_hl_height"]
    template_content = profile["ref_r2"]
    template_cap = profile["ref_para_count"]

    ref_hl_top = profile.get("ref_hl_top")
    if ref_hl_top is None and ref_hl is not None:
        ref_hl_top = ref_hl.top
    if ref_hl_top is None:
        ref_hl_top = 0

    with_ka_h = profile.get("expanded_hl_height")
    if ref_ka and not with_ka_h and hasattr(ref_ka, "top"):
        with_ka_h = ref_ka.top + ref_ka.height - ref_hl_top
    with_ka_content = max((with_ka_h or template_h) - r0 - r1 - pad, template_content)

    full_main_h = MAX_KA_BOTTOM_EMU - ref_hl_top - 91440
    full_main_content = max(full_main_h - r0 - r1 - pad, template_content)

    return {
        "template_cap": template_cap,
        "with_ka_cap": _capacity_for_content_row(profile, with_ka_content),
        "full_main_cap": _capacity_for_content_row(profile, full_main_content),
    }


def build_overflow_profile(
    template_prs: Presentation,
    template_slide: Slide,
    service: ServiceProfile,
    profile: TemplateProfile,
) -> dict[str, Any]:
    """
    Layout profile for overflow planning.

    When a service slide has no KA tab but one will be inserted (ka_mode=none),
    borrow geometry from a standalone HL+KA reference slide so packing reserves
    space for both tables.
    """
    layout_profile = build_layout_profile(template_slide)
    if service.ka_mode == "embedded":
        return layout_profile
    if layout_profile.get("ref_ka"):
        return layout_profile

    ref_idx = find_standalone_reference_slide(template_prs, profile)
    if ref_idx is None:
        return layout_profile

    ref_slide = template_prs.slides[ref_idx]
    ref_ka = get_key_activities_shape(ref_slide)
    ref_hl = get_highlights_shape(ref_slide)
    if ref_ka is None or ref_hl is None:
        return layout_profile

    merged = dict(layout_profile)
    merged["ref_ka"] = ref_ka
    merged["expanded_hl_height"] = ref_ka.top - ref_hl.top
    return merged


def main_slide_capacity(layout_profile: dict) -> int:
    """HL+KA single-slide capacity (used only when all content fits on one slide)."""
    budgets = get_hl_budgets(layout_profile)
    if layout_profile.get("ref_ka"):
        return budgets["with_ka_cap"]
    return budgets["full_main_cap"]


def hl_only_slide_capacity(layout_profile: dict) -> int:
    """HL-only capacity for main/contd slides when overflow continues on later slides."""
    return UNIVERSAL_HL_ONLY_CAP


def _total_display_lines(sections: list[dict]) -> int:
    """Wrap-aware HL line count including one-line gaps between sprints."""
    total = 0
    for idx, section in enumerate(sections):
        gap = 1 if idx > 0 else 0
        total += gap + _section_paragraph_count(section, display=True)
    return total


def _pack_sections_to_cap(sections: list[dict], cap: int) -> tuple[list[dict], list[dict]]:
    packed: list[dict] = []
    used = 0
    idx = 0
    while idx < len(sections):
        section = sections[idx]
        sec_count = _section_paragraph_count(section, display=True)
        gap = 1 if packed else 0
        if used + gap + sec_count <= cap:
            packed.append(dict(section))
            used += gap + sec_count
            idx += 1
            continue
        if not packed:
            return [dict(section)], [dict(s) for s in sections[idx + 1 :]]
        return packed, [dict(s) for s in sections[idx:]]
    return packed, []


def plan_overflow(
    project: ProjectContent,
    profile: dict,
    main_cap: int | None = None,
    hl_only_cap: int | None = None,
) -> OverflowPlan:
    budgets = get_hl_budgets(profile)
    if main_cap is not None:
        hl_ka_cap = main_cap
    else:
        hl_ka_cap = (
            budgets["with_ka_cap"] if profile.get("ref_ka") else UNIVERSAL_HL_ONLY_CAP
        )
    if hl_only_cap is None:
        hl_only_cap = UNIVERSAL_HL_ONLY_CAP

    display_sections = [section_display_content(s) for s in project.sections]
    total = _total_display_lines(display_sections)
    if total <= 0:
        total = sum(
            1 + len(s.completed) + len(s.released) + len(s.inprogress)
            for s in project.sections
        )

    if profile.get("ref_ka"):
        if total < HL_KA_SAME_SLIDE_MAX:
            if hl_ka_fits_on_main_slide(profile, total):
                return OverflowPlan(
                    main_sections=display_sections,
                    continuation_chains=[],
                    ka_on_main=True,
                )
            return OverflowPlan(
                main_sections=display_sections,
                continuation_chains=[],
                ka_contd_only=True,
            )
        if HL_KA_SAME_SLIDE_MAX < total < UNIVERSAL_HL_ONLY_CAP:
            return OverflowPlan(
                main_sections=display_sections,
                continuation_chains=[],
                ka_contd_only=True,
            )

    if total <= hl_ka_cap:
        return OverflowPlan(main_sections=display_sections, continuation_chains=[])

    main_sections: list[dict] = []
    contd_sections: list[dict] = []
    used = 0

    for idx, section in enumerate(display_sections):
        sec_count = _section_paragraph_count(section, display=True)
        gap = 1 if main_sections else 0
        if used + gap + sec_count <= hl_only_cap:
            main_sections.append(dict(section))
            used += gap + sec_count
            continue
        contd_sections.extend(dict(s) for s in display_sections[idx:])
        break

    if not contd_sections:
        return OverflowPlan(main_sections=display_sections, continuation_chains=[])

    continuation_chains: list[list[dict]] = []
    remaining = contd_sections
    while remaining:
        packed, rest = _pack_sections_to_cap(remaining, hl_only_cap)
        if packed:
            continuation_chains.append(packed)
        elif rest:
            break
        if not rest:
            break
        remaining = rest

    return OverflowPlan(main_sections=main_sections, continuation_chains=continuation_chains)
