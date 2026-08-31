"""Low-level PowerPoint slide operations."""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


def delete_slide(presentation, index: int) -> None:
    slide_id_list = presentation.slides._sldIdLst
    slide_ids = list(slide_id_list)
    slide_id = slide_ids[index]
    r_id = slide_id.rId
    slide_part = presentation.part.related_part(r_id)
    presentation.part.drop_rel(r_id)
    slide_id_list.remove(slide_id)
    scrub_slide_incoming_links(presentation, slide_part)


def scrub_slide_incoming_links(presentation, slide_part) -> None:
    """Remove hyperlink relationships from any slide that still points at *slide_part*."""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tx_body = shape.text_frame._txBody
            for p_elem in tx_body.findall(qn("a:p")):
                for hlink in list(p_elem.findall(qn("a:hlinkClick"))):
                    rid = hlink.get(qn("r:id"))
                    if not rid:
                        continue
                    try:
                        target = slide.part.related_part(rid)
                    except KeyError:
                        continue
                    if target is slide_part:
                        p_elem.remove(hlink)


def move_slide_after(prs, slide_index: int, after_index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    elements = list(sld_id_lst)
    el = elements[slide_index]
    elements.pop(slide_index)
    elements.insert(after_index + 1, el)
    for child in list(sld_id_lst):
        sld_id_lst.remove(child)
    for child in elements:
        sld_id_lst.append(child)


def delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def clear_slide_shapes(slide) -> None:
    for shape in list(slide.shapes):
        delete_shape(shape)


def copy_shapes_to_slide(src_slide, dst_slide) -> None:
    for shape in src_slide.shapes:
        newel = copy.deepcopy(shape.element)
        dst_slide.shapes._spTree.insert_element_before(newel, "p:extLst")


def _active_slide_parts(prs) -> set:
    presentation = prs.part
    return {presentation.related_part(sld.rId) for sld in prs.slides._sldIdLst}


def purge_orphan_slide_parts(prs) -> int:
    """Rename unreachable slide parts so they cannot collide on save."""
    from pptx.opc.packuri import PackURI

    package = prs.part.package
    active = _active_slide_parts(prs)
    renamed = 0
    for part in package.iter_parts():
        pname = str(part.partname)
        if not pname.startswith("/ppt/slides/slide"):
            continue
        if part in active:
            continue
        part.partname = PackURI(f"/ppt/slides/slide_orphan_{renamed}.xml")
        renamed += 1
    return renamed


def normalize_slide_partnames(prs) -> None:
    purge_orphan_slide_parts(prs)
    rids = [sld_id.rId for sld_id in prs.slides._sldIdLst]
    if rids:
        prs.part.rename_slide_parts(rids)


def save_presentation_clean(prs, path: str | Path) -> Path:
    normalize_slide_partnames(prs)
    out = Path(path)
    prs.save(str(out))
    return out.resolve()
