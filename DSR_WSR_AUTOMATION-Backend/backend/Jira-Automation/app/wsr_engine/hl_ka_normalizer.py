"""Split combined 5×3 HL+KA tables into separate shapes before fill."""

from __future__ import annotations

import copy
import logging

from pptx import Presentation

from app.services.ppt_shape_utils import (
    get_highlights_shape,
    get_key_activities_shape,
    has_combined_hl_ka_table,
    is_contd_title,
    is_delivery_slide_title,
    iter_all_shapes,
    service_suffix_from_title,
    slide_title_text,
)
from app.services.template_profile import TemplateProfile

logger = logging.getLogger(__name__)


def _delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def find_standalone_reference_slide(prs: Presentation, profile: TemplateProfile | None = None) -> int | None:
    """First delivery slide with separate HL and KA tables (not combined)."""
    if profile is not None:
        for svc in profile.services.values():
            if svc.ka_mode != "standalone":
                continue
            slide = prs.slides[svc.main_slide_index]
            if has_combined_hl_ka_table(slide):
                continue
            try:
                hl = get_highlights_shape(slide)
                ka = get_key_activities_shape(slide)
            except ValueError:
                continue
            if hl is not None and ka is not None:
                return svc.main_slide_index

    for idx, slide in enumerate(prs.slides):
        title = slide_title_text(slide)
        if not is_delivery_slide_title(title) or is_contd_title(title):
            continue
        if has_combined_hl_ka_table(slide):
            continue
        try:
            hl = get_highlights_shape(slide)
            ka = get_key_activities_shape(slide)
        except ValueError:
            continue
        if hl is not None and ka is not None:
            return idx
    return None


def find_hl_only_contd_template_index(template_prs: Presentation) -> int | None:
    """Template slide with Highlights only — no standalone or embedded KA."""
    for idx, slide in enumerate(template_prs.slides):
        title = slide_title_text(slide)
        if not is_delivery_slide_title(title) or not is_contd_title(title):
            continue
        if get_key_activities_shape(slide) is not None:
            continue
        if has_combined_hl_ka_table(slide):
            continue
        try:
            get_highlights_shape(slide)
            return idx
        except ValueError:
            continue
    return None


def find_ka_only_contd_template_index(template_prs: Presentation) -> int | None:
    """Template (Contd..) slide with KA only — no Highlights table."""
    for idx, slide in enumerate(template_prs.slides):
        title = slide_title_text(slide)
        if not is_delivery_slide_title(title) or not is_contd_title(title):
            continue
        if get_key_activities_shape(slide) is None:
            continue
        if has_combined_hl_ka_table(slide):
            continue
        try:
            get_highlights_shape(slide)
            continue
        except ValueError:
            return idx
    return None


def _find_combined_table(slide):
    for shape in iter_all_shapes(slide.shapes):
        if not shape.has_table or len(shape.table.columns) < 3:
            continue
        if len(shape.table.rows) < 5:
            continue
        try:
            h0 = shape.table.cell(0, 0).text.strip().lower()
            h3 = shape.table.cell(3, 0).text.strip().lower()
        except (IndexError, AttributeError):
            continue
        if "highlights" in h0 and "key activit" in h3:
            return shape
    return None


def normalize_slide_hl_ka_layout(slide, ref_slide) -> bool:
    """
    Replace a combined 5-row HL+KA table with separate HL and KA shapes
    cloned from a reference slide.
    """
    combined = _find_combined_table(slide)
    if combined is None:
        return False

    ref_hl = get_highlights_shape(ref_slide)
    ref_ka = get_key_activities_shape(ref_slide)
    if ref_ka is None:
        logger.warning("Reference slide has no standalone KA table; skip normalize")
        return False

    left = combined.left
    top = combined.top
    width = combined.width

    hl_el = copy.deepcopy(ref_hl.element)
    ka_el = copy.deepcopy(ref_ka.element)
    sp_tree = slide.shapes._spTree
    sp_tree.insert_element_before(hl_el, "p:extLst")
    sp_tree.insert_element_before(ka_el, "p:extLst")

    new_hl = slide.shapes[-2]
    new_ka = slide.shapes[-1]
    new_hl.left = left
    new_hl.top = top
    new_hl.width = width
    new_ka.left = ref_ka.left if hasattr(ref_ka, "left") else left
    new_ka.top = ref_ka.top if hasattr(ref_ka, "top") else top + ref_hl.height
    new_ka.width = ref_ka.width if hasattr(ref_ka, "width") else width

    _delete_shape(combined)
    logger.debug("Normalized combined HL+KA table on slide")
    return True


def _should_normalize_slide(slide, profile: TemplateProfile | None) -> bool:
    title = slide_title_text(slide)
    if not is_delivery_slide_title(title):
        return False
    if is_contd_title(title):
        return False
    if not has_combined_hl_ka_table(slide):
        return False
    if profile is None:
        return True
    suffix = service_suffix_from_title(title)
    svc = profile.services.get(suffix)
    if svc is None:
        return True
    return svc.ka_mode == "standalone"


def normalize_deck_hl_ka_layouts(prs: Presentation, profile: TemplateProfile | None = None) -> int:
    """
    Normalize delivery main slides that use combined HL+KA tables when the
    service profile expects standalone KA. Preserves embedded KA layouts.
    """
    ref_idx = find_standalone_reference_slide(prs, profile)
    if ref_idx is None:
        logger.info("No standalone HL+KA reference slide found; skip normalization")
        return 0

    ref_slide = prs.slides[ref_idx]
    changed = 0
    for slide in prs.slides:
        if not _should_normalize_slide(slide, profile):
            continue
        if normalize_slide_hl_ka_layout(slide, ref_slide):
            changed += 1

    logger.info("Normalized %d slide(s) with combined HL+KA tables", changed)
    return changed
