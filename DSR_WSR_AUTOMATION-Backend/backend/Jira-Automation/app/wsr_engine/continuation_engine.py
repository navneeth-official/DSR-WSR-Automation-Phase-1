"""Dynamic continuation slide creation and population."""

from __future__ import annotations

import copy
import logging

from pptx import Presentation
from pptx.oxml.ns import qn

from app.services.ppt_logo_sync import sync_header_pictures_from_reference
from app.services.ppt_shape_utils import find_title_shape, normalize_title_text, slide_title_text
from app.services.template_profile import ServiceProfile
from app.wsr_engine.formatter import (
    clear_ka_on_slide,
    fill_contd_highlights_on_slide,
    remove_hl_from_slide,
    remove_ka_from_slide,
)
from app.wsr_engine.hl_ka_normalizer import find_hl_only_contd_template_index, find_ka_only_contd_template_index
from app.wsr_engine.models import TitleFormat
from app.wsr_engine.slide_ops import (
    clear_slide_shapes,
    copy_shapes_to_slide,
    delete_slide,
    move_slide_after,
)

logger = logging.getLogger(__name__)


def set_contd_title(slide, project_name: str, ref_title_shape, title_format: TitleFormat) -> None:
    suffix = f"{project_name}  {title_format.contd_marker}"
    prefix = title_format.prefix
    sep = title_format.separator

    try:
        title_shape = find_title_shape(slide)
    except RuntimeError:
        new_el = copy.deepcopy(ref_title_shape.element)
        slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        title_shape = find_title_shape(slide)

    ref_rPr = None
    if ref_title_shape.has_text_frame and ref_title_shape.text_frame.paragraphs:
        runs = ref_title_shape.text_frame.paragraphs[0].runs
        if runs:
            ref_rPr = runs[0]._r.find(qn("a:rPr"))

    p = title_shape.text_frame.paragraphs[0]
    if not p.runs:
        p.add_run()
    run = p.runs[0]
    old_rPr = run._r.find(qn("a:rPr"))
    if old_rPr is not None:
        run._r.remove(old_rPr)
    if ref_rPr is not None:
        run._r.insert(0, copy.deepcopy(ref_rPr))
    run.text = f"{prefix}{sep}{suffix}"
    for extra in p.runs[1:]:
        p._p.remove(extra._r)


def find_contd_slides_for_project(prs: Presentation, project_name: str) -> list[int]:
    needle = project_name.lower()
    indices: list[int] = []
    for i, slide in enumerate(prs.slides):
        title = normalize_title_text(slide_title_text(slide)).lower()
        if needle in title and "(contd" in title:
            indices.append(i)
    return indices


def _resolve_hl_only_contd_template(
    template_prs: Presentation,
    service: ServiceProfile,
    fallback_idx: int,
) -> int:
    """Prefer an HL-only contd template; fall back to the project main slide."""
    global_idx = find_hl_only_contd_template_index(template_prs)
    if global_idx is not None:
        return global_idx

    for idx in service.contd_slide_indices:
        slide = template_prs.slides[idx]
        from app.services.ppt_shape_utils import get_key_activities_shape, has_combined_hl_ka_table

        if get_key_activities_shape(slide) is None and not has_combined_hl_ka_table(slide):
            try:
                from app.services.ppt_shape_utils import get_highlights_shape

                get_highlights_shape(slide)
                return idx
            except ValueError:
                continue

    logger.debug(
        "No HL-only contd template; using main slide %d as contd source",
        fallback_idx,
    )
    return fallback_idx


def create_contd_slide_after(
    prs: Presentation,
    template_prs: Presentation,
    after_idx: int,
    template_slide_idx: int,
) -> int:
    layout = prs.slides[after_idx].slide_layout
    template_slide = template_prs.slides[template_slide_idx]
    prs.slides.add_slide(layout)
    new_idx = len(prs.slides) - 1
    new_slide = prs.slides[new_idx]
    clear_slide_shapes(new_slide)
    copy_shapes_to_slide(template_slide, new_slide)
    sync_header_pictures_from_reference(template_slide, new_slide)
    move_slide_after(prs, new_idx, after_idx)
    return after_idx + 1


def ensure_continuation_slides(
    prs: Presentation,
    template_prs: Presentation,
    project_name: str,
    main_idx: int,
    continuation_chains: list[list[dict]],
    service: ServiceProfile,
    template_slide_idx: int,
    title_format: TitleFormat,
    section_tmpl: dict,
) -> tuple[list[int], int]:
    if not continuation_chains:
        return [], 0

    existing = find_contd_slides_for_project(prs, project_name)
    needed = len(continuation_chains)
    created = 0

    template_idx = _resolve_hl_only_contd_template(
        template_prs,
        service,
        service.contd_slide_indices[0] if service.contd_slide_indices else template_slide_idx,
    )

    insert_after = main_idx
    contd_indices: list[int] = []

    for chain_idx in range(needed):
        if chain_idx < len(existing):
            slide_idx = existing[chain_idx]
            contd_slide = prs.slides[slide_idx]
            clear_slide_shapes(contd_slide)
            copy_shapes_to_slide(template_prs.slides[template_idx], contd_slide)
            sync_header_pictures_from_reference(
                template_prs.slides[template_slide_idx], contd_slide
            )
        else:
            slide_idx = create_contd_slide_after(
                prs, template_prs, insert_after, template_idx
            )
            created += 1
            existing = find_contd_slides_for_project(prs, project_name)

        main_slide = prs.slides[main_idx]
        contd_slide = prs.slides[slide_idx]
        ref_title = find_title_shape(main_slide)
        set_contd_title(contd_slide, project_name, ref_title, title_format)

        fill_contd_highlights_on_slide(contd_slide, section_tmpl, continuation_chains[chain_idx])
        remove_ka_from_slide(contd_slide)
        sync_header_pictures_from_reference(
            template_prs.slides[template_slide_idx], contd_slide
        )

        contd_indices.append(slide_idx)
        insert_after = slide_idx

    keep = set(contd_indices)
    for idx in sorted(find_contd_slides_for_project(prs, project_name), reverse=True):
        if idx not in keep:
            delete_slide(prs, idx)

    return contd_indices, created


def _resolve_ka_only_contd_template(
    template_prs: Presentation,
    fallback_idx: int,
) -> int:
    """Prefer a KA-only contd template; fall back to the project main slide."""
    idx = find_ka_only_contd_template_index(template_prs)
    if idx is not None:
        return idx
    logger.debug(
        "No KA-only contd template; using main slide %d as contd source",
        fallback_idx,
    )
    return fallback_idx


def position_ka_at_contd_top(
    slide,
    template_prs: Presentation,
    *,
    layout_ref_slide=None,
) -> None:
    """Place KA at the top content coordinate used on KA-only contd template slides."""
    from app.services.ppt_shape_utils import get_highlights_shape, get_key_activities_shape

    ka = get_key_activities_shape(slide)
    if ka is None:
        return

    ref_idx = find_ka_only_contd_template_index(template_prs)
    if ref_idx is not None:
        ref_ka = get_key_activities_shape(template_prs.slides[ref_idx])
        ka.left = ref_ka.left
        ka.width = ref_ka.width
        ka.top = ref_ka.top
        return

    ref_slide = layout_ref_slide
    if ref_slide is not None:
        try:
            ref_hl = get_highlights_shape(ref_slide)
            ref_ka = get_key_activities_shape(ref_slide)
            ka.top = ref_hl.top
            ka.left = ref_ka.left if ref_ka is not None else ref_hl.left
            ka.width = ref_ka.width if ref_ka is not None else ref_hl.width
        except ValueError:
            pass


def ensure_ka_only_contd_slide(
    prs: Presentation,
    template_prs: Presentation,
    project_name: str,
    main_idx: int,
    service: ServiceProfile,
    template_slide_idx: int,
    title_format: TitleFormat,
) -> tuple[list[int], int]:
    """Create one continuation slide that holds only the KA tab (HL stays on main)."""
    existing = find_contd_slides_for_project(prs, project_name)
    template_slide = template_prs.slides[template_slide_idx]
    ka_contd_template_idx = _resolve_ka_only_contd_template(template_prs, template_slide_idx)
    ka_contd_template = template_prs.slides[ka_contd_template_idx]
    created = 0

    if existing:
        slide_idx = existing[0]
        contd_slide = prs.slides[slide_idx]
        clear_slide_shapes(contd_slide)
        copy_shapes_to_slide(ka_contd_template, contd_slide)
        sync_header_pictures_from_reference(template_slide, contd_slide)
    else:
        slide_idx = create_contd_slide_after(
            prs, template_prs, main_idx, ka_contd_template_idx
        )
        created = 1
        contd_slide = prs.slides[slide_idx]

    main_slide = prs.slides[main_idx]
    ref_title = find_title_shape(main_slide)
    set_contd_title(contd_slide, project_name, ref_title, title_format)
    remove_hl_from_slide(contd_slide)
    position_ka_at_contd_top(contd_slide, template_prs, layout_ref_slide=template_slide)
    clear_ka_on_slide(contd_slide, template_slide, service)
    sync_header_pictures_from_reference(template_slide, contd_slide)

    keep = {slide_idx}
    for idx in sorted(find_contd_slides_for_project(prs, project_name), reverse=True):
        if idx not in keep:
            delete_slide(prs, idx)

    return [slide_idx], created


def append_ka_only_contd_slide(
    prs: Presentation,
    template_prs: Presentation,
    project_name: str,
    main_idx: int,
    after_idx: int,
    service: ServiceProfile,
    title_format: TitleFormat,
) -> int:
    """Insert a KA-only (Contd..) slide immediately after ``after_idx``."""
    template_slide = template_prs.slides[service.main_slide_index]
    ka_contd_template_idx = _resolve_ka_only_contd_template(template_prs, service.main_slide_index)
    slide_idx = create_contd_slide_after(
        prs, template_prs, after_idx, ka_contd_template_idx
    )
    contd_slide = prs.slides[slide_idx]
    main_slide = prs.slides[main_idx]
    ref_title = find_title_shape(main_slide)
    set_contd_title(contd_slide, project_name, ref_title, title_format)
    remove_hl_from_slide(contd_slide)
    position_ka_at_contd_top(contd_slide, template_prs, layout_ref_slide=template_slide)
    sync_header_pictures_from_reference(template_slide, contd_slide)
    return slide_idx
