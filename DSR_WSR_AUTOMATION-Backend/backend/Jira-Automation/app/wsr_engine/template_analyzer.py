"""Analyze WSR template structure and build internal model."""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation

from app.services.template_profile import scan_template_profile
from app.wsr_engine.index_layout import discover_index_layout
from app.wsr_engine.index_updater import _find_index_table
from app.wsr_engine.models import TemplateModel
from app.wsr_engine.project_detector import (
    build_project_maps,
    classify_slide,
    detect_title_format,
)

logger = logging.getLogger(__name__)


def analyze_template(template_path: Path | str) -> TemplateModel:
    path = Path(template_path).resolve()
    if not path.is_file():
        raise FileNotFoundError(f"Template not found: {path}")

    prs = Presentation(str(path))
    slides = [classify_slide(slide, idx) for idx, slide in enumerate(prs.slides)]
    projects = build_project_maps(slides)
    index_idx = next((s.index for s in slides if s.slide_type == "INDEX"), None)
    title_format = detect_title_format(slides)
    profile = scan_template_profile(path)

    index_layout = None
    if index_idx is not None:
        table = _find_index_table(prs.slides[index_idx])
        if table is not None:
            index_layout = discover_index_layout(table)

    model = TemplateModel(
        template_path=str(path),
        slides=slides,
        projects=projects,
        index_slide_index=index_idx,
        title_format=title_format,
        profile=profile,
        index_layout=index_layout,
    )

    logger.info(
        "Analyzed %s: %d slides, %d projects, index=%s",
        path.name,
        len(slides),
        len(projects),
        index_idx,
    )
    return model
