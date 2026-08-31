"""Apply single coherent geometry corrections to a delivery-status deck."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

from app.geometry.metrics import (
    live_slide_metrics,
    metric_delta_summary,
    metrics_changed,
    slide_metric_snapshot,
)
from app.geometry.planner import plan_slide_repair
from app.geometry.types import (
    GeometryCorrectionResult,
    GeometryReport,
    RepairMode,
    SlideGeometryDelta,
    SlideRepairPlan,
)
from app.services.ppt_format_repair import (
    apply_layout_repair,
    compact_vertical_layout,
    fix_ka_footer_overflow,
    shrink_hl_only,
    shrink_ka_table,
    tighten_hl_and_position_ka,
)

from app.paths import G10X_TEMPLATE, SCRIPTS_DIR


def _uds_helpers():
    scripts = str(SCRIPTS_DIR)
    if scripts not in sys.path:
        sys.path.insert(0, scripts)
    import update_delivery_status as uds  # noqa: WPS433

    from pptx import Presentation as Prs

    g10x_prs = Prs(str(G10X_TEMPLATE))
    return uds, g10x_prs


class GeometryCorrector:
    """
    Deterministic layout correction from geometry inspection results.

    Applies at most one coherent repair per slide per iteration. All EMU
    deltas are computed inside ``ppt_format_repair`` / ``update_delivery_status``
    helpers — never from vision pixel measurements.

    Verifies measurable geometry change (EMU metrics) after each slide repair.
    """

    def correct(
        self,
        ppt_path: Path | str,
        report: GeometryReport,
        *,
        slide_plans: dict[int, SlideRepairPlan] | None = None,
        save: bool = True,
    ) -> GeometryCorrectionResult:
        path = Path(ppt_path).resolve()
        uds, g10x_prs = _uds_helpers()
        prs = Presentation(str(path))

        actions: list[str] = []
        modified_slides: list[int] = []
        unchanged_slides: list[int] = []
        slide_deltas: list[SlideGeometryDelta] = []
        failures: list[str] = []

        for slide_report in report.slides:
            idx = slide_report.slide_index
            plan = (
                slide_plans.get(idx)
                if slide_plans
                else plan_slide_repair(slide_report)
            )
            if plan is None or not plan.applies_change:
                continue

            if idx < 1 or idx > len(prs.slides):
                failures.append(f"Invalid slide index {idx}")
                continue

            slide = prs.slides[idx - 1]
            metrics_before = live_slide_metrics(slide, uds, g10x_prs)

            try:
                repair_changed = self._apply_plan(slide, plan, uds=uds, g10x_prs=g10x_prs)
            except Exception as exc:  # noqa: BLE001
                failures.append(f"Slide {idx}: {exc}")
                continue

            metrics_after = live_slide_metrics(slide, uds, g10x_prs)
            geo_changed = repair_changed and metrics_changed(metrics_before, metrics_after)

            delta = SlideGeometryDelta(
                slide_index=idx,
                changed=geo_changed,
                plan_mode=plan.mode.value,
                reason=plan.reason,
                metrics_before=slide_metric_snapshot(metrics_before),
                metrics_after=slide_metric_snapshot(metrics_after),
                metric_deltas=metric_delta_summary(metrics_before, metrics_after),
            )
            slide_deltas.append(delta)

            if geo_changed:
                modified_slides.append(idx)
                actions.append(f"slide_{idx}:{plan.mode.value}:{plan.reason}")
            else:
                unchanged_slides.append(idx)
                failures.append(
                    f"Slide {idx}: plan {plan.mode.value} produced no measurable geometry change"
                )

        if save and modified_slides:
            prs.save(str(path))

        return GeometryCorrectionResult(
            modified=bool(modified_slides),
            ppt_path=str(path),
            actions_applied=actions,
            slides_modified=modified_slides,
            slides_unchanged=unchanged_slides,
            slide_deltas=slide_deltas,
            failures=failures,
            message=(
                f"Applied {len(actions)} geometry correction(s)"
                if modified_slides
                else "No geometry changes applied"
            ),
        )

    @staticmethod
    def _apply_plan(slide, plan: SlideRepairPlan, *, uds, g10x_prs) -> bool:
        if plan.mode == RepairMode.TIGHTEN_AND_POSITION:
            return tighten_hl_and_position_ka(slide, uds, g10x_prs)

        if plan.mode == RepairMode.SHRINK_KA:
            return shrink_ka_table(slide, uds, g10x_prs)

        if plan.mode == RepairMode.SHRINK_HL:
            return shrink_hl_only(slide, uds, g10x_prs)

        if plan.mode == RepairMode.FIX_FOOTER_OVERFLOW:
            return fix_ka_footer_overflow(slide, uds, g10x_prs)

        if plan.mode == RepairMode.COMPACT_VERTICAL:
            return compact_vertical_layout(slide, uds, g10x_prs)

        if plan.mode == RepairMode.ENSURE_CLEARANCE:
            return apply_layout_repair(
                slide,
                uds,
                g10x_prs,
                raw=None,
                params={
                    "layout_mode": plan.layout_mode,
                    "expand_for_wrap": plan.expand_for_wrap,
                },
            )

        if plan.mode == RepairMode.EXPAND_AND_REFLOW:
            return apply_layout_repair(
                slide,
                uds,
                g10x_prs,
                raw=None,
                params={
                    "layout_mode": "expanded",
                    "expand_for_wrap": True,
                },
            )

        return False
