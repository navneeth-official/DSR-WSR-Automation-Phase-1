"""
End-to-end validation: compare vision measurements before and after layout correction.

Produces per-iteration reports with expected vs actual metric changes, issue
resolution status, and root-cause hints when corrections fail to help.
"""

from __future__ import annotations

import re
import shutil
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal

from app.layout import LayoutCorrector
from app.layout.shape_ops import PixelScale
from app.pipeline.dependencies import PipelineDependencies
from app.pipeline.layout_corrector import vision_report_to_inputs
from app.pipeline.types import PipelineConfig

MetricDirection = Literal["improved", "same", "worse", "unknown"]
OutcomeStatus = Literal[
    "resolved",
    "improved",
    "unchanged",
    "worsened",
    "not_applicable",
    "no_measurement",
]

TRACKED_METRICS = (
    "gap_between_sections",
    "unused_space_inside_highlight",
    "last_highlight_text_bottom",
    "keyactivities_title_top",
    "highlight_box_bottom",
)

ISSUE_METRIC_KEYS = (
    "gap_pixels",
    "unused_space_pixels",
    "overlap_pixels",
    "overflow_pixels",
)


@dataclass
class MetricChange:
    metric: str
    before: float | None
    after: float | None
    delta: float | None
    direction: MetricDirection
    ideal_direction: str  # "decrease" | "increase" | "toward_target" | "informational"

    def to_dict(self) -> dict[str, Any]:
        return {
            "metric": self.metric,
            "before": self.before,
            "after": self.after,
            "delta": self.delta,
            "direction": self.direction,
            "ideal_direction": self.ideal_direction,
        }


@dataclass
class CorrectionOutcome:
    slide_number: int | None
    action_summary: str
    action_type: str
    target: str
    delta_emu: int | None
    expected_metric: str | None
    expected_change_px: float | None
    actual_change_px: float | None
    expected_direction: str
    actual_direction: MetricDirection
    status: OutcomeStatus
    notes: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_number": self.slide_number,
            "action": self.action_summary,
            "action_type": self.action_type,
            "target": self.target,
            "delta_emu": self.delta_emu,
            "expected_metric": self.expected_metric,
            "expected_change_px": self.expected_change_px,
            "actual_change_px": self.actual_change_px,
            "expected_direction": self.expected_direction,
            "actual_direction": self.actual_direction,
            "status": self.status,
            "notes": self.notes,
        }


@dataclass
class IssueResolution:
    issue_id: str
    slide_number: int | None
    affected_object: str
    recommended_action: str
    before_measurement: dict[str, float]
    after_measurement: dict[str, float] | None
    status: Literal["resolved", "improved", "unchanged", "worsened", "new"]
    reason: str

    def to_dict(self) -> dict[str, Any]:
        return {
            "issue_id": self.issue_id,
            "slide_number": self.slide_number,
            "affected_object": self.affected_object,
            "recommended_action": self.recommended_action,
            "before_measurement": self.before_measurement,
            "after_measurement": self.after_measurement,
            "status": self.status,
            "reason": self.reason,
        }


@dataclass
class SlideIterationReport:
    slide_number: int | None
    title: str
    measurements_before: dict[str, float]
    measurements_after: dict[str, float]
    metric_changes: list[MetricChange] = field(default_factory=list)
    corrections_applied: list[dict[str, Any]] = field(default_factory=list)
    correction_outcomes: list[CorrectionOutcome] = field(default_factory=list)
    issues_before: list[dict[str, Any]] = field(default_factory=list)
    issues_after: list[dict[str, Any]] = field(default_factory=list)
    issue_resolutions: list[IssueResolution] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_number": self.slide_number,
            "title": self.title,
            "measurements_before": self.measurements_before,
            "measurements_after": self.measurements_after,
            "metric_changes": [m.to_dict() for m in self.metric_changes],
            "corrections_applied": self.corrections_applied,
            "correction_outcomes": [c.to_dict() for c in self.correction_outcomes],
            "issues_before": self.issues_before,
            "issues_after": self.issues_after,
            "issue_resolutions": [i.to_dict() for i in self.issue_resolutions],
        }


@dataclass
class IterationValidationReport:
    iteration: int
    deck_score_before: int
    deck_score_after: int
    deck_score_delta: int
    deck_pass_before: bool
    deck_pass_after: bool
    correction_modified: bool
    correction_message: str
    change_magnitude_emu: int
    slides: list[SlideIterationReport] = field(default_factory=list)
    no_improvement_corrections: list[CorrectionOutcome] = field(default_factory=list)
    regression_corrections: list[CorrectionOutcome] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "iteration": self.iteration,
            "deck_score_before": self.deck_score_before,
            "deck_score_after": self.deck_score_after,
            "deck_score_delta": self.deck_score_delta,
            "deck_pass_before": self.deck_pass_before,
            "deck_pass_after": self.deck_pass_after,
            "correction_modified": self.correction_modified,
            "correction_message": self.correction_message,
            "change_magnitude_emu": self.change_magnitude_emu,
            "slides": [s.to_dict() for s in self.slides],
            "no_improvement_corrections": [c.to_dict() for c in self.no_improvement_corrections],
            "regression_corrections": [c.to_dict() for c in self.regression_corrections],
        }


@dataclass
class CorrectionValidationReport:
    source_ppt: Path
    working_ppt: Path
    max_iterations: int
    stopped_reason: str
    iterations: list[IterationValidationReport] = field(default_factory=list)
    summary: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return {
            "source_ppt": str(self.source_ppt),
            "working_ppt": str(self.working_ppt),
            "max_iterations": self.max_iterations,
            "stopped_reason": self.stopped_reason,
            "iterations": [it.to_dict() for it in self.iterations],
            "summary": self.summary,
        }


def _num(value: Any) -> float | None:
    if isinstance(value, (int, float)):
        return float(value)
    return None


def _slide_measurements(slide_data: dict[str, Any]) -> dict[str, float]:
    raw = slide_data.get("measurements") or {}
    out: dict[str, float] = {}
    if isinstance(raw, dict):
        for key, val in raw.items():
            n = _num(val)
            if n is not None:
                out[str(key)] = n
    return out


def _metric_direction(
    before: float | None,
    after: float | None,
    *,
    ideal: str,
    tolerance_px: float = 2.0,
) -> MetricDirection:
    if before is None or after is None:
        return "unknown"
    delta = after - before
    if abs(delta) <= tolerance_px:
        return "same"
    if ideal == "decrease":
        return "improved" if delta < 0 else "worse"
    if ideal == "increase":
        return "improved" if delta > 0 else "worse"
    return "unknown"


def _ideal_for_metric(metric: str, issues: list[dict[str, Any]]) -> str:
    actions = {str(i.get("recommended_action") or "") for i in issues}
    if metric == "unused_space_inside_highlight":
        return "decrease"
    if metric == "gap_between_sections":
        if "move_section_up" in actions:
            return "decrease"
        if "move_section_down" in actions:
            return "increase"
    return "informational"


def _issue_key(issue: dict[str, Any]) -> tuple[str, str, str]:
    return (
        str(issue.get("issue_id") or ""),
        str(issue.get("affected_object") or ""),
        str(issue.get("recommended_action") or ""),
    )


def _primary_issue_metric(issue: dict[str, Any]) -> tuple[str, float] | None:
    measurement = issue.get("measurement") or {}
    if not isinstance(measurement, dict):
        return None
    action = str(issue.get("recommended_action") or "")
    if action in ("move_section_up", "move_section_down"):
        for key in ("gap_pixels", "overlap_pixels", "gap_between_sections"):
            val = _num(measurement.get(key))
            if val is not None:
                return key, val
    if action in ("reduce_unused_space", "decrease_textbox_height"):
        for key in ("unused_space_pixels", "unused_pixels", "waste_pixels"):
            val = _num(measurement.get(key))
            if val is not None:
                return key, val
    if action in ("increase_textbox_height", "expand_placeholder", "overflow_detected"):
        for key in ("overflow_pixels", "missing_pixels", "overlap_pixels"):
            val = _num(measurement.get(key))
            if val is not None:
                return key, val
    for key, val in measurement.items():
        n = _num(val)
        if n is not None:
            return str(key), n
    return None


def _resolve_issue_status(
    before_issue: dict[str, Any],
    after_issues: list[dict[str, Any]],
    metric_changes: dict[str, MetricChange],
) -> IssueResolution:
    issue_id = str(before_issue.get("issue_id") or "UNKNOWN")
    slide_number = before_issue.get("slide_number")
    affected = str(before_issue.get("affected_object") or "")
    action = str(before_issue.get("recommended_action") or "")
    before_meas = dict(before_issue.get("measurement") or {})

    matching = [
        i for i in after_issues
        if _issue_key(i) == _issue_key(before_issue)
        or (
            str(i.get("recommended_action") or "") == action
            and str(i.get("affected_object") or "") == affected
        )
    ]

    if not matching:
        return IssueResolution(
            issue_id=issue_id,
            slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
            affected_object=affected,
            recommended_action=action,
            before_measurement={k: float(v) for k, v in before_meas.items() if _num(v) is not None},
            after_measurement=None,
            status="resolved",
            reason="Issue no longer reported by vision after correction.",
        )

    after_issue = matching[0]
    after_meas = dict(after_issue.get("measurement") or {})

    primary = _primary_issue_metric(before_issue)
    if primary:
        metric_key, before_val = primary
        after_val = _num(after_meas.get(metric_key))
        if after_val is not None:
            if action in ("move_section_up", "reduce_unused_space", "decrease_textbox_height"):
                if after_val < before_val - 2:
                    return IssueResolution(
                        issue_id=issue_id,
                        slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
                        affected_object=affected,
                        recommended_action=action,
                        before_measurement={metric_key: before_val},
                        after_measurement={metric_key: after_val},
                        status="improved",
                        reason=f"{metric_key} decreased from {before_val} to {after_val} but issue persists.",
                    )
                if after_val > before_val + 2:
                    return IssueResolution(
                        issue_id=issue_id,
                        slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
                        affected_object=affected,
                        recommended_action=action,
                        before_measurement={metric_key: before_val},
                        after_measurement={metric_key: after_val},
                        status="worsened",
                        reason=f"{metric_key} increased from {before_val} to {after_val}.",
                    )

    if action == "move_section_up" and "gap_between_sections" in metric_changes:
        mc = metric_changes["gap_between_sections"]
        if mc.direction == "improved":
            return IssueResolution(
                issue_id=issue_id,
                slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
                affected_object=affected,
                recommended_action=action,
                before_measurement={k: float(v) for k, v in before_meas.items() if _num(v) is not None},
                after_measurement={k: float(v) for k, v in after_meas.items() if _num(v) is not None},
                status="improved",
                reason="Gap metric improved but vision still flags excessive spacing.",
            )

    if action in ("reduce_unused_space", "decrease_textbox_height"):
        mc = metric_changes.get("unused_space_inside_highlight")
        if mc and mc.direction == "improved":
            return IssueResolution(
                issue_id=issue_id,
                slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
                affected_object=affected,
                recommended_action=action,
                before_measurement={k: float(v) for k, v in before_meas.items() if _num(v) is not None},
                after_measurement={k: float(v) for k, v in after_meas.items() if _num(v) is not None},
                status="improved",
                reason="Unused space metric improved but vision still reports waste.",
            )

    return IssueResolution(
        issue_id=issue_id,
        slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
        affected_object=affected,
        recommended_action=action,
        before_measurement={k: float(v) for k, v in before_meas.items() if _num(v) is not None},
        after_measurement={k: float(v) for k, v in after_meas.items() if _num(v) is not None},
        status="unchanged",
        reason="Issue still present with similar measurements after correction.",
    )


def _expected_metric_for_action(action_type: str, target: str, detail: str) -> tuple[str | None, str, float | None]:
    detail_lower = detail.lower()
    if "gap" in detail_lower and "key activities" in detail_lower:
        if "down" in detail_lower:
            return "gap_between_sections", "increase", _extract_px(detail, "overlap")
        if "up" in detail_lower:
            return "gap_between_sections", "decrease", _extract_px(detail, "gap")
    if target == "highlights" and action_type == "resize_shape":
        if "reduced" in detail_lower or "shrink" in detail_lower:
            return "unused_space_inside_highlight", "decrease", _extract_px(detail, "unused")
        if "expand" in detail_lower:
            return "highlight_box_bottom", "increase", None
    if action_type == "maintain_gap":
        return "gap_between_sections", "increase", _extract_px(detail, "gap")
    return None, "informational", None


def _extract_px(detail: str, keyword: str) -> float | None:
    patterns = [
        rf"{keyword}\s+([0-9.]+)px",
        rf"{keyword}\s+([0-9.]+)\s*px",
        rf"\({keyword}\s+([0-9.]+)px\)",
        rf"unused space\s+([0-9.]+)px",
        rf"gap\s+([0-9.]+)px",
        rf"overlap\s+([0-9.]+)px",
    ]
    for pattern in patterns:
        match = re.search(pattern, detail, re.IGNORECASE)
        if match:
            return float(match.group(1))
    generic = re.search(r"([0-9.]+)px", detail)
    if generic:
        return float(generic.group(1))
    return None


def _emu_to_px(delta_emu: int, image_height_px: int, slide_height_emu: int) -> float:
    if slide_height_emu <= 0:
        return 0.0
    return abs(delta_emu) * image_height_px / slide_height_emu


def _outcome_status(
    expected_dir: str,
    actual_dir: MetricDirection,
    expected_px: float | None,
    actual_px: float | None,
) -> OutcomeStatus:
    if expected_dir == "informational":
        return "not_applicable"
    if actual_dir == "unknown":
        return "no_measurement"
    if actual_dir == "same":
        return "unchanged"
    if actual_dir == "improved":
        return "improved"
    if actual_dir == "worse":
        return "worsened"
    return "unchanged"


def _root_cause_hint(
    outcome: CorrectionOutcome,
    slide_report: SlideIterationReport,
) -> str:
    if outcome.status not in ("unchanged", "worsened", "no_measurement"):
        return ""

    hints: list[str] = []
    if outcome.status == "no_measurement":
        hints.append("Vision did not return the target metric after correction.")

    if outcome.actual_direction == "same" and outcome.delta_emu:
        hints.append(
            "Rendered slide may not reflect EMU change at vision resolution, "
            "or the correction magnitude is below vision estimation noise (~2px)."
        )

    if outcome.actual_direction == "worse":
        other_actions = [
            c for c in slide_report.correction_outcomes
            if c is not outcome and c.slide_number == outcome.slide_number
        ]
        if len(other_actions) >= 1:
            hints.append(
                "Conflicting corrections on the same slide may have counteracted "
                "this change (e.g. shrink Highlights then move Key Activities)."
            )

    if outcome.expected_metric == "gap_between_sections" and outcome.expected_direction == "decrease":
        hints.append(
            "Gap may still exceed minimum clearance threshold; vision re-flags "
            "until gap is at or below target."
        )

    if outcome.expected_metric == "unused_space_inside_highlight":
        hints.append(
            "Shrink amount may be less than reported unused space, or vision "
            "re-estimates box bounds differently after resize."
        )

    return " ".join(hints)


def _build_slide_report(
    before_slide: dict[str, Any],
    after_slide: dict[str, Any] | None,
    correction_slide: dict[str, Any] | None,
    *,
    image_height_px: int,
    slide_height_emu: int,
) -> SlideIterationReport:
    slide_number = before_slide.get("slide_number") or before_slide.get("slide_index")
    title = str(before_slide.get("title") or "")
    meas_before = _slide_measurements(before_slide)
    meas_after = _slide_measurements(after_slide) if after_slide else {}

    issues_before = list(before_slide.get("issues") or [])
    issues_after = list((after_slide or {}).get("issues") or [])

    ideal_gap = _ideal_for_metric("gap_between_sections", issues_before)
    ideal_waste = _ideal_for_metric("unused_space_inside_highlight", issues_before)

    metric_changes: list[MetricChange] = []
    metric_map: dict[str, MetricChange] = {}
    for metric in TRACKED_METRICS:
        before = meas_before.get(metric)
        after = meas_after.get(metric)
        ideal = ideal_waste if metric == "unused_space_inside_highlight" else (
            ideal_gap if metric == "gap_between_sections" else "informational"
        )
        direction = _metric_direction(before, after, ideal=ideal)
        change = MetricChange(
            metric=metric,
            before=before,
            after=after,
            delta=(after - before) if before is not None and after is not None else None,
            direction=direction,
            ideal_direction=ideal,
        )
        metric_changes.append(change)
        metric_map[metric] = change

    corrections_applied: list[dict[str, Any]] = []
    correction_outcomes: list[CorrectionOutcome] = []

    if correction_slide:
        for action_detail, delta in zip(
            correction_slide.get("actions") or [],
            correction_slide.get("deltas_emu") or [],
        ):
            corrections_applied.append(
                {"detail": action_detail, "delta_emu": delta}
            )

        action_types = correction_slide.get("actions") or []
        deltas = correction_slide.get("deltas_emu") or []
        # Re-parse from layout corrector structured data if available
        structured = correction_slide.get("structured_actions") or []
        if structured:
            for entry in structured:
                action_type = str(entry.get("action_type") or "")
                target = str(entry.get("target") or "")
                detail = str(entry.get("detail") or "")
                delta_emu = entry.get("delta_emu")
                expected_metric, expected_dir, expected_px = _expected_metric_for_action(
                    action_type, target, detail
                )
                if expected_px is None and isinstance(delta_emu, int):
                    expected_px = _emu_to_px(delta_emu, image_height_px, slide_height_emu)

                actual_px = None
                actual_dir: MetricDirection = "unknown"
                if expected_metric and expected_metric in metric_map:
                    mc = metric_map[expected_metric]
                    actual_px = mc.delta
                    actual_dir = mc.direction

                status = _outcome_status(expected_dir, actual_dir, expected_px, actual_px)
                outcome = CorrectionOutcome(
                    slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
                    action_summary=detail,
                    action_type=action_type,
                    target=target,
                    delta_emu=delta_emu,
                    expected_metric=expected_metric,
                    expected_change_px=expected_px,
                    actual_change_px=actual_px,
                    expected_direction=expected_dir,
                    actual_direction=actual_dir,
                    status=status,
                )
                outcome.notes = _root_cause_hint(outcome, SlideIterationReport(
                    slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
                    title=title,
                    measurements_before=meas_before,
                    measurements_after=meas_after,
                    metric_changes=metric_changes,
                    correction_outcomes=correction_outcomes,
                ))
                correction_outcomes.append(outcome)
        else:
            for detail, delta in zip(action_types, deltas):
                expected_metric, expected_dir, expected_px = _expected_metric_for_action(
                    "unknown", "", str(detail)
                )
                if expected_px is None and isinstance(delta, int):
                    expected_px = _emu_to_px(delta, image_height_px, slide_height_emu)
                actual_px = None
                actual_dir = "unknown"
                if expected_metric and expected_metric in metric_map:
                    mc = metric_map[expected_metric]
                    actual_px = mc.delta
                    actual_dir = mc.direction
                outcome = CorrectionOutcome(
                    slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
                    action_summary=str(detail),
                    action_type="unknown",
                    target="",
                    delta_emu=delta,
                    expected_metric=expected_metric,
                    expected_change_px=expected_px,
                    actual_change_px=actual_px,
                    expected_direction=expected_dir,
                    actual_direction=actual_dir,
                    status=_outcome_status(expected_dir, actual_dir, expected_px, actual_px),
                )
                correction_outcomes.append(outcome)

    issue_resolutions = [
        _resolve_issue_status(
            {**issue, "slide_number": slide_number},
            issues_after,
            metric_map,
        )
        for issue in issues_before
    ]

    # New issues after correction
    before_keys = {_issue_key(i) for i in issues_before}
    for issue in issues_after:
        if _issue_key(issue) not in before_keys:
            issue_resolutions.append(
                IssueResolution(
                    issue_id=str(issue.get("issue_id") or "NEW"),
                    slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
                    affected_object=str(issue.get("affected_object") or ""),
                    recommended_action=str(issue.get("recommended_action") or ""),
                    before_measurement={},
                    after_measurement={
                        k: float(v) for k, v in (issue.get("measurement") or {}).items()
                        if _num(v) is not None
                    },
                    status="new",
                    reason="New issue reported after correction (possible regression).",
                )
            )

    return SlideIterationReport(
        slide_number=int(slide_number) if isinstance(slide_number, (int, float)) else None,
        title=title,
        measurements_before=meas_before,
        measurements_after=meas_after,
        metric_changes=metric_changes,
        corrections_applied=corrections_applied,
        correction_outcomes=correction_outcomes,
        issues_before=issues_before,
        issues_after=issues_after,
        issue_resolutions=issue_resolutions,
    )


def run_correction_validation(
    ppt_path: Path,
    *,
    max_iterations: int = 3,
    output_dir: Path | None = None,
    working_copy: bool = True,
) -> CorrectionValidationReport:
    """
    Run before/after vision evaluation around each correction pass.

    Uses a working copy of the presentation so the source file is preserved.
    """
    source = ppt_path.resolve()
    if not source.is_file():
        raise FileNotFoundError(f"Presentation not found: {source}")

    if working_copy:
        work_path = source.with_name(f"{source.stem}_validation{source.suffix}")
        shutil.copy2(source, work_path)
    else:
        work_path = source

    deps = PipelineDependencies.create_default()
    cfg = PipelineConfig(
        max_iterations=max_iterations,
        keep_render_images=True,
        render_output_dir=output_dir,
    )

    render_dir = cfg.render_output_dir
    if render_dir is None:
        render_dir = Path(tempfile.mkdtemp(prefix="ppt_correction_validation_"))
        cfg = PipelineConfig(
            max_iterations=max_iterations,
            keep_render_images=True,
            render_output_dir=render_dir,
        )

    from pptx import Presentation

    prs = Presentation(str(work_path))
    slide_height_emu = prs.slide_height

    iterations: list[IterationValidationReport] = []
    stopped_reason = "iteration_limit"

    for round_num in range(1, max_iterations + 1):
        batch_before = deps.renderer.render_deck(
            work_path,
            output_dir=render_dir / f"iter_{round_num}_before",
            keep_images=True,
        )
        report_before = deps.vision_client.evaluate(batch_before)

        if deps.vision_client.passes(report_before):
            stopped_reason = "no_issues_before_correction"
            break

        image_paths = {
            s.slide_index: s.image_path
            for s in batch_before.slides
            if s.image_path.is_file()
        }
        inputs = vision_report_to_inputs(report_before, image_paths=image_paths)
        layout_corrector = LayoutCorrector()
        layout_result = layout_corrector.correct(work_path, inputs, save=True)
        correction = layout_result.to_correction_result()

        correction_details_by_slide: dict[int, dict[str, Any]] = {}
        for slide_entry in layout_result.slides:
            sn = slide_entry.slide_number
            if sn is None:
                continue
            correction_details_by_slide[int(sn)] = {
                "slide_number": sn,
                "modified": slide_entry.modified,
                "actions": [a.detail for a in slide_entry.actions],
                "deltas_emu": [
                    a.delta_emu for a in slide_entry.actions if a.delta_emu is not None
                ],
                "structured_actions": [
                    {
                        "action_type": a.action_type.value,
                        "target": a.target,
                        "detail": a.detail,
                        "delta_emu": a.delta_emu,
                    }
                    for a in slide_entry.actions
                ],
                "failures": slide_entry.failures,
            }

        batch_after = deps.renderer.render_deck(
            work_path,
            output_dir=render_dir / f"iter_{round_num}_after",
            keep_images=True,
        )
        report_after = deps.vision_client.evaluate(batch_after)

        image_heights = [
            PixelScale.from_paths(
                slide_width_emu=prs.slide_width,
                slide_height_emu=slide_height_emu,
                image_path=s.image_path,
            ).image_height_px
            for s in batch_after.slides
            if s.image_path.is_file()
        ]
        default_height = image_heights[0] if image_heights else 1080

        slides_report: list[SlideIterationReport] = []
        after_by_index = {
            int(s.get("slide_number") or s.get("slide_index") or 0): s
            for s in report_after.slides
        }

        for before_slide in report_before.slides:
            sn = int(before_slide.get("slide_number") or before_slide.get("slide_index") or 0)
            after_slide = after_by_index.get(sn)
            corr_slide = correction_details_by_slide.get(sn)
            slide_report = _build_slide_report(
                before_slide,
                after_slide,
                corr_slide,
                image_height_px=default_height,
                slide_height_emu=slide_height_emu,
            )
            slides_report.append(slide_report)

        all_outcomes = [o for s in slides_report for o in s.correction_outcomes]
        no_improvement = [
            o for o in all_outcomes
            if o.status in ("unchanged", "no_measurement")
        ]
        regressions = [o for o in all_outcomes if o.status == "worsened"]

        iter_report = IterationValidationReport(
            iteration=round_num,
            deck_score_before=report_before.deck_score,
            deck_score_after=report_after.deck_score,
            deck_score_delta=report_after.deck_score - report_before.deck_score,
            deck_pass_before=report_before.deck_pass,
            deck_pass_after=report_after.deck_pass,
            correction_modified=correction.modified,
            correction_message=correction.message,
            change_magnitude_emu=correction.change_magnitude_emu,
            slides=slides_report,
            no_improvement_corrections=no_improvement,
            regression_corrections=regressions,
        )
        iterations.append(iter_report)

        if deps.vision_client.passes(report_after):
            stopped_reason = "no_issues_after_correction"
            break

        if not correction.modified:
            stopped_reason = "corrector_noop"
            break

        if round_num >= max_iterations:
            stopped_reason = "iteration_limit"

    summary = _build_summary(iterations)
    return CorrectionValidationReport(
        source_ppt=source,
        working_ppt=work_path,
        max_iterations=max_iterations,
        stopped_reason=stopped_reason,
        iterations=iterations,
        summary=summary,
    )


def _rebuild_structured_actions(
    slide_entry: dict[str, Any],
    evaluation,
    mapper,
) -> list[dict[str, Any]]:
    """Deprecated — structured actions are captured directly from LayoutCorrector."""
    del slide_entry, evaluation, mapper
    return []


def _build_summary(iterations: list[IterationValidationReport]) -> dict[str, Any]:
    if not iterations:
        return {"message": "No iterations completed."}

    first = iterations[0]
    last = iterations[-1]
    total_resolved = 0
    total_improved = 0
    total_unchanged = 0
    total_worsened = 0
    total_new = 0
    failed_corrections: list[dict[str, Any]] = []

    for it in iterations:
        for slide in it.slides:
            for res in slide.issue_resolutions:
                if res.status == "resolved":
                    total_resolved += 1
                elif res.status == "improved":
                    total_improved += 1
                elif res.status == "unchanged":
                    total_unchanged += 1
                elif res.status == "worsened":
                    total_worsened += 1
                elif res.status == "new":
                    total_new += 1
        for outcome in it.no_improvement_corrections + it.regression_corrections:
            failed_corrections.append(outcome.to_dict())

    return {
        "iterations_run": len(iterations),
        "deck_score_first_before": first.deck_score_before,
        "deck_score_last_after": last.deck_score_after,
        "deck_score_net_change": last.deck_score_after - first.deck_score_before,
        "issues_resolved": total_resolved,
        "issues_improved_but_persistent": total_improved,
        "issues_unchanged": total_unchanged,
        "issues_worsened": total_worsened,
        "issues_new_after_correction": total_new,
        "corrections_without_measurable_improvement": len(
            [it for it in iterations for o in it.no_improvement_corrections]
        ),
        "corrections_causing_regression": len(
            [it for it in iterations for o in it.regression_corrections]
        ),
        "failed_corrections": failed_corrections,
    }


def format_validation_report_text(report: CorrectionValidationReport) -> str:
    """Human-readable summary for console / log file."""
    lines: list[str] = [
        "=" * 72,
        "VISION CORRECTION VALIDATION REPORT",
        "=" * 72,
        f"Source: {report.source_ppt}",
        f"Working copy: {report.working_ppt}",
        f"Stopped: {report.stopped_reason}",
        "",
    ]

    for it in report.iterations:
        lines.extend([
            f"--- Iteration {it.iteration} ---",
            f"Deck score: {it.deck_score_before} → {it.deck_score_after} "
            f"({'+' if it.deck_score_delta >= 0 else ''}{it.deck_score_delta})",
            f"Correction modified deck: {it.correction_modified} "
            f"(total |delta| EMU: {it.change_magnitude_emu})",
            "",
        ])
        for slide in it.slides:
            if not slide.corrections_applied and not slide.issues_before:
                continue
            lines.append(f"Slide {slide.slide_number}: {slide.title}")
            if slide.measurements_before:
                lines.append(f"  Measurements BEFORE: {slide.measurements_before}")
            if slide.measurements_after:
                lines.append(f"  Measurements AFTER:  {slide.measurements_after}")
            for mc in slide.metric_changes:
                if mc.before is not None or mc.after is not None:
                    lines.append(
                        f"  Metric {mc.metric}: {mc.before} → {mc.after} "
                        f"({mc.direction}, ideal={mc.ideal_direction})"
                    )
            for corr in slide.corrections_applied:
                lines.append(f"  Applied: {corr['detail']} (delta_emu={corr.get('delta_emu')})")
            for outcome in slide.correction_outcomes:
                lines.append(
                    f"  Outcome [{outcome.status}]: {outcome.action_summary} | "
                    f"expected {outcome.expected_direction} "
                    f"{outcome.expected_metric} ~{outcome.expected_change_px}px, "
                    f"actual Δ={outcome.actual_change_px}px ({outcome.actual_direction})"
                )
                if outcome.notes:
                    lines.append(f"    Note: {outcome.notes}")
            for res in slide.issue_resolutions:
                lines.append(
                    f"  Issue {res.issue_id} ({res.recommended_action}): {res.status} — {res.reason}"
                )
            lines.append("")

        if it.no_improvement_corrections:
            lines.append("  Corrections with NO measurable improvement:")
            for o in it.no_improvement_corrections:
                lines.append(f"    - Slide {o.slide_number}: {o.action_summary}")
        if it.regression_corrections:
            lines.append("  Corrections that WORSENED a metric:")
            for o in it.regression_corrections:
                lines.append(f"    - Slide {o.slide_number}: {o.action_summary}")
        lines.append("")

    lines.extend([
        "--- Summary ---",
        *[f"{k}: {v}" for k, v in report.summary.items()],
        "=" * 72,
    ])
    return "\n".join(lines)
