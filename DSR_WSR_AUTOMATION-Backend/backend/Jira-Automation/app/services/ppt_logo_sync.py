"""Re-insert header logos after shape XML copy (template-agnostic)."""

from __future__ import annotations

import io
import logging

from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.ppt_shape_utils import (
    is_delivery_slide_title,
    iter_all_shapes,
    slide_title_text,
)

logger = logging.getLogger(__name__)

HEB_LOGO_SHAPE_ID = 3
_HEADER_TOP_FRACTION = 0.28
_POSITION_TOLERANCE_EMU = 120000


def _delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _picture_has_blob(shape) -> bool:
    try:
        _ = shape.image.blob
        return True
    except (AttributeError, ValueError, KeyError):
        return False


def _slide_height(slide) -> int:
    try:
        return int(slide.part.presentation.slide_height)
    except (AttributeError, TypeError, ValueError):
        return 6858000


def _header_top_limit(slide) -> int:
    return int(_slide_height(slide) * _HEADER_TOP_FRACTION)


def _iter_pictures(shapes):
    for shape in iter_all_shapes(shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def _header_pictures(slide) -> list:
    """Return valid header-band pictures, left-to-right."""
    limit = _header_top_limit(slide)
    pics = [sh for sh in _iter_pictures(slide.shapes) if _picture_has_blob(sh) and sh.top <= limit]
    if pics:
        return sorted(pics, key=lambda sh: (sh.top, sh.left))

    by_id = next(
        (sh for sh in _iter_pictures(slide.shapes) if getattr(sh, "shape_id", None) == HEB_LOGO_SHAPE_ID),
        None,
    )
    if by_id is not None and _picture_has_blob(by_id):
        return [by_id]

    fallback = [sh for sh in _iter_pictures(slide.shapes) if _picture_has_blob(sh)]
    return sorted(fallback, key=lambda sh: (sh.top, sh.left))


def find_heb_logo_on_slide(slide):
    """Return the primary H-E-B header logo picture when present."""
    pics = _header_pictures(slide)
    if not pics:
        return None
    by_id = next((sh for sh in pics if getattr(sh, "shape_id", None) == HEB_LOGO_SHAPE_ID), None)
    return by_id or pics[0]


def _position_close(a, b) -> bool:
    return (
        abs(a.left - b.left) <= _POSITION_TOLERANCE_EMU
        and abs(a.top - b.top) <= _POSITION_TOLERANCE_EMU
    )


def _remove_broken_header_pictures(target_slide, reference_picture) -> None:
    for shape in list(_iter_pictures(target_slide.shapes)):
        same_id = getattr(shape, "shape_id", None) == getattr(reference_picture, "shape_id", None)
        if same_id or _position_close(shape, reference_picture) or not _picture_has_blob(shape):
            if shape.top <= _header_top_limit(target_slide) or same_id:
                _delete_shape(shape)


def sync_header_pictures_from_reference(reference_slide, target_slide) -> int:
    """
    Copy header pictures from reference onto target.

    ``copy_shapes_to_slide`` deep-copies shape XML but not image relationships,
    which leaves broken placeholders. Re-insert from reference image bytes.
    """
    ref_pictures = _header_pictures(reference_slide)
    if not ref_pictures:
        logger.warning("No header pictures found on reference slide")
        return 0

    synced = 0
    for ref in ref_pictures:
        existing = next(
            (
                sh
                for sh in _iter_pictures(target_slide.shapes)
                if getattr(sh, "shape_id", None) == getattr(ref, "shape_id", None)
            ),
            None,
        )
        if existing is not None and _picture_has_blob(existing):
            continue

        _remove_broken_header_pictures(target_slide, ref)
        target_slide.shapes.add_picture(
            io.BytesIO(ref.image.blob),
            ref.left,
            ref.top,
            width=ref.width,
            height=ref.height,
        )
        synced += 1

    return synced


def sync_heb_logo_from_main(main_slide, contd_slide) -> bool:
    """Backward-compatible wrapper."""
    return sync_header_pictures_from_reference(main_slide, contd_slide) > 0


def sync_all_delivery_slide_logos(prs, reference_slide) -> int:
    """Ensure every delivery slide has header logos from the template reference."""
    total = 0
    for slide in prs.slides:
        title = slide_title_text(slide)
        if not is_delivery_slide_title(title):
            continue
        total += sync_header_pictures_from_reference(reference_slide, slide)
    return total


def resolve_logo_reference_slide(template_prs, projects) -> object | None:
    """Pick a template slide that still has valid header pictures."""
    for project in projects:
        idx = project.main_slide_index
        if 0 <= idx < len(template_prs.slides):
            slide = template_prs.slides[idx]
            if _header_pictures(slide):
                return slide
    for slide in template_prs.slides:
        if _header_pictures(slide):
            return slide
    return None
