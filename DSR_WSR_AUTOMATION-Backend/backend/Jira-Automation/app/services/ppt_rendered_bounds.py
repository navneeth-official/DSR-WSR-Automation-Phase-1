"""Rendered text bounds for HL metrics — COM first, image second, estimate last."""

from __future__ import annotations

import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Literal

from app.rendering.com_text_bounds import ComHlTextBounds, ComTextBoundsMeasurer
from app.services.ppt_layout_metrics import (
    EMU_PER_INCH,
    hl_waste_below_text_in_estimated,
    rendered_text_bottom_emu,
)

MeasurementMethod = Literal["com", "image", "estimated"]


@dataclass(frozen=True)
class HlTextMeasurement:
    text_bottom_in: float
    hl_bottom_in: float
    waste_in: float
    method: MeasurementMethod

    def to_dict(self) -> dict[str, Any]:
        return {
            "text_bottom_in": self.text_bottom_in,
            "hl_bottom_in": self.hl_bottom_in,
            "waste_in": self.waste_in,
            "method": self.method,
        }


def _hl_content_region_emu(
    hl_shape,
    *,
    com_bounds: ComHlTextBounds | None = None,
) -> tuple[int, int, int, int]:
    """
    Scan region for image ink detection: left, top, right, bottom in slide EMU.

    Prefer COM content-row bounds (rendered layout). python-pptx row heights often
    under-report the content area and clip the lower bullets from the scan ROI.
    """
    if com_bounds is not None:
        return (
            int(com_bounds.content_left_in * EMU_PER_INCH),
            int(com_bounds.content_top_in * EMU_PER_INCH),
            int(com_bounds.content_right_in * EMU_PER_INCH),
            int(com_bounds.content_bottom_in * EMU_PER_INCH),
        )

    r0 = hl_shape.table.rows[0].height
    r1 = hl_shape.table.rows[1].height
    r2 = hl_shape.table.rows[2].height
    left = int(hl_shape.left)
    top = int(hl_shape.top + r0 + r1)
    right = int(hl_shape.left + hl_shape.width)
    bottom = int(top + r2)
    return left, top, right, bottom


def _emu_rect_to_px(
    left: int,
    top: int,
    right: int,
    bottom: int,
    *,
    slide_width_emu: int,
    slide_height_emu: int,
    image_width_px: int,
    image_height_px: int,
) -> tuple[int, int, int, int]:
    def conv_x(emu: int) -> int:
        return int(round(emu / slide_width_emu * image_width_px))

    def conv_y(emu: int) -> int:
        return int(round(emu / slide_height_emu * image_height_px))

    return (
        conv_x(left),
        conv_y(top),
        max(conv_x(right), conv_x(left) + 1),
        max(conv_y(bottom), conv_y(top) + 1),
    )


def measure_hl_text_bottom_from_image(
    image_path: Path | str,
    hl_shape,
    *,
    slide_width_emu: int | None = None,
    slide_height_emu: int | None = None,
    text_luminance_threshold: int = 200,
    hl_bottom_in: float | None = None,
    com_bounds: ComHlTextBounds | None = None,
) -> HlTextMeasurement | None:
    """
    Detect bottom-most dark text pixel inside the HL content cell on a rendered PNG.
    """
    try:
        from PIL import Image
    except ImportError:
        return None

    path = Path(image_path)
    if not path.is_file():
        return None

    slide_width_emu = slide_width_emu or 12192000  # 13.333 in default widescreen
    slide_height_emu = slide_height_emu or 6858000  # 7.5 in

    left, top, right, bottom = _hl_content_region_emu(hl_shape, com_bounds=com_bounds)
    if hl_bottom_in is None:
        hl_bottom_in = round((hl_shape.top + hl_shape.height) / EMU_PER_INCH, 4)

    with Image.open(path) as img:
        rgb = img.convert("RGB")
        w_px, h_px = rgb.size
        x0, y0, x1, y1 = _emu_rect_to_px(
            left,
            top,
            right,
            bottom,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            image_width_px=w_px,
            image_height_px=h_px,
        )
        x0 = max(0, min(x0, w_px - 1))
        x1 = max(x0 + 1, min(x1, w_px))
        y0 = max(0, min(y0, h_px - 1))
        y1 = max(y0 + 1, min(y1, h_px))

        crop = rgb.crop((x0, y0, x1, y1))
        pixels = crop.load()
        cw, ch = crop.size

        text_bottom_px: int | None = None
        for y in range(ch - 1, -1, -1):
            row_has_text = False
            for x in range(cw):
                r, g, b = pixels[x, y]
                if r < text_luminance_threshold and g < text_luminance_threshold and b < text_luminance_threshold:
                    row_has_text = True
                    break
            if row_has_text:
                text_bottom_px = y0 + y
                break

    if text_bottom_px is None:
        return None

    text_bottom_in = round(text_bottom_px / h_px * (slide_height_emu / EMU_PER_INCH), 4)
    waste_in = round(max(hl_bottom_in - text_bottom_in, 0.0), 4)
    return HlTextMeasurement(
        text_bottom_in=text_bottom_in,
        hl_bottom_in=hl_bottom_in,
        waste_in=waste_in,
        method="image",
    )


def measure_hl_text_estimated(hl_shape) -> HlTextMeasurement:
    """Fallback: canonical line-height model (legacy)."""
    ref_r2 = hl_shape.table.rows[2].height
    text_bottom_emu = rendered_text_bottom_emu(hl_shape, ref_r2=ref_r2)
    hl_bottom_emu = hl_shape.top + hl_shape.height
    text_bottom_in = round(text_bottom_emu / EMU_PER_INCH, 4)
    hl_bottom_in = round(hl_bottom_emu / EMU_PER_INCH, 4)
    waste_in = hl_waste_below_text_in_estimated(hl_shape, ref_r2=ref_r2)
    return HlTextMeasurement(
        text_bottom_in=text_bottom_in,
        hl_bottom_in=hl_bottom_in,
        waste_in=waste_in,
        method="estimated",
    )


def _from_com(bounds: ComHlTextBounds) -> HlTextMeasurement:
    return HlTextMeasurement(
        text_bottom_in=bounds.text_bottom_in,
        hl_bottom_in=bounds.hl_bottom_in,
        waste_in=bounds.waste_in,
        method="com",
    )



def _refine_com_with_image(
    com: HlTextMeasurement,
    image: HlTextMeasurement,
) -> HlTextMeasurement:
    """
    Combine COM table geometry with image ink detection.

    COM hl_bottom is authoritative for the container; image text_bottom is
    authoritative for visible ink on the rendered PNG.
    """
    return HlTextMeasurement(
        text_bottom_in=image.text_bottom_in,
        hl_bottom_in=com.hl_bottom_in,
        waste_in=round(max(com.hl_bottom_in - image.text_bottom_in, 0.0), 4),
        method="image",
    )


def measure_hl_text(
    hl_shape,
    *,
    com_bounds: ComHlTextBounds | None = None,
    image_path: Path | str | None = None,
    slide_width_emu: int | None = None,
    slide_height_emu: int | None = None,
) -> HlTextMeasurement:
    """COM → image refinement → image-only → estimated fallback chain."""
    com_result: HlTextMeasurement | None = None
    if com_bounds is not None:
        com_result = _from_com(com_bounds)

    image_result: HlTextMeasurement | None = None
    if image_path is not None:
        image_result = measure_hl_text_bottom_from_image(
            image_path,
            hl_shape,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            hl_bottom_in=com_result.hl_bottom_in if com_result else None,
            com_bounds=com_bounds,
        )

    if com_result is not None and image_result is not None:
        return _refine_com_with_image(com_result, image_result)
    if com_result is not None:
        return com_result
    if image_result is not None:
        return image_result
    return measure_hl_text_estimated(hl_shape)


def enrich_deck_rendered_hl_metrics(
    deck_data: dict[str, Any],
    ppt_path: str | Path,
    *,
    export_images_for_fallback: bool = True,
) -> dict[str, int]:
    """
    Attach rendered HL text metrics to each slide in deck_data.

    Sets rendered_text_bottom_in, hl_waste_below_text_in, hl_text_bounds_method,
    and updates estimated_text_bottom_in / clearance fields when measured.
    """
    ppt_path = Path(ppt_path).resolve()
    slides: list[dict[str, Any]] = deck_data.get("slides") or []
    stats = {"com": 0, "image": 0, "estimated": 0}

    com_by_slide: dict[int, ComHlTextBounds] = {}
    slide_width_emu = 12192000
    slide_height_emu = 6858000
    try:
        with ComTextBoundsMeasurer(ppt_path) as measurer:
            slide_width_emu = int(measurer.slide_width_in * EMU_PER_INCH)
            slide_height_emu = int(measurer.slide_height_in * EMU_PER_INCH)
            for slide_data in slides:
                if not slide_data.get("highlights"):
                    continue
                idx = int(slide_data["slide_index"])
                bounds = measurer.measure_slide(idx)
                if bounds is not None:
                    com_by_slide[idx] = bounds
    except Exception:
        pass

    image_by_slide: dict[int, Path] = {}
    slides_with_hl = [
        int(s["slide_index"]) for s in slides if s.get("highlights")
    ]
    if export_images_for_fallback and slides_with_hl:
        try:
            from app.services.ppt_slide_images import export_slides_to_png

            out_dir = Path(tempfile.mkdtemp(prefix="ppt_hl_bounds_", dir=ppt_path.parent))
            exported = export_slides_to_png(
                ppt_path, out_dir, slide_indices=slides_with_hl
            )
            for item in exported:
                image_by_slide[int(item["slide_index"])] = Path(item["image_path"])
        except Exception:
            image_by_slide = {}

    from pptx import Presentation

    prs = Presentation(str(ppt_path))
    slide_by_index = {i + 1: s for i, s in enumerate(prs.slides)}

    for slide_data in slides:
        if not slide_data.get("highlights"):
            continue
        idx = int(slide_data["slide_index"])
        pptx_slide = slide_by_index.get(idx)
        if pptx_slide is None:
            continue

        hl_shape = _get_highlights_shape_pptx(pptx_slide)
        if hl_shape is None:
            continue

        measurement = measure_hl_text(
            hl_shape,
            com_bounds=com_by_slide.get(idx),
            image_path=image_by_slide.get(idx),
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
        )

        method = measurement.method
        stats[method] = stats.get(method, 0) + 1

        slide_data["rendered_text_bottom_in"] = measurement.text_bottom_in
        slide_data["hl_bottom_measured_in"] = measurement.hl_bottom_in
        slide_data["hl_waste_below_text_in"] = measurement.waste_in
        slide_data["hl_text_bounds_method"] = method
        slide_data["estimated_text_bottom_in"] = measurement.text_bottom_in
        slide_data["hl_text_bottom_for_fit_in"] = measurement.text_bottom_in

        hl_bottom_emu = int(measurement.hl_bottom_in * EMU_PER_INCH)
        text_bottom_emu = int(measurement.text_bottom_in * EMU_PER_INCH)
        slide_data["hl_text_overflow_in"] = round(
            max((text_bottom_emu - hl_bottom_emu) / EMU_PER_INCH, 0), 4
        )

        ka = slide_data.get("key_activities")
        if ka and ka.get("position_in", {}).get("top") is not None:
            ka_top_in = float(ka["position_in"]["top"])
            slide_data["text_ka_clearance_in"] = round(
                ka_top_in - measurement.text_bottom_in, 4
            )

    deck_data["hl_bounds_measurement_stats"] = stats
    return stats


def _get_highlights_shape_pptx(slide):
    for sid in (9, 7):
        for shape in slide.shapes:
            if shape.shape_id == sid and shape.has_table and len(shape.table.rows) == 3:
                return shape
    return None
