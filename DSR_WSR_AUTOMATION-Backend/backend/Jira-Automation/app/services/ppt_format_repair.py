"""Apply FixPlan actions to a delivery-status PPTX (in-place and rebuild)."""

from __future__ import annotations

import json
import re
import subprocess
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.oxml.ns import qn

from app.services.ppt_format_fix_planner import ALLOWED_ACTIONS
from app.services.ppt_layout_metrics import EMU_PER_INCH, hl_ka_tab_gap_emu

from app.paths import G10X_TEMPLATE, PPT_BUILDER, REPO_ROOT, SCRIPTS_DIR
SPARSE_WASTE_FORCE_TIGHT_IN = 0.5
_HL_BOTTOM_TOLERANCE_EMU = 45720  # ~0.05 in


def _ensure_text_clearance(hl, ka, profile, uds) -> bool:
    """Expand HL height so rendered text fits above KA (uses UDS geometry helpers)."""
    text_bottom = uds._hl_rendered_text_bottom(hl, profile)
    hl_bottom = hl.top + hl.height
    if text_bottom <= hl_bottom - _HL_BOTTOM_TOLERANCE_EMU:
        return False

    max_h = None
    if ka is not None:
        ka_h = ka.height or uds._estimate_ka_table_height(profile, 0)
        gap = hl_ka_tab_gap_emu(profile.get("canonical_line_height_emu"))
        max_h = int(uds.MAX_KA_BOTTOM_EMU - ka_h - gap - hl.top)

    uds._fit_hl_content_only(hl, profile, max_h=max_h)
    return True


def _import_uds():
    scripts = str(SCRIPTS_DIR)
    if scripts not in sys.path:
        sys.path.insert(0, scripts)
    import update_delivery_status as uds  # noqa: WPS433

    return uds


def _paragraph_text(p_elem) -> str:
    return "".join(n.text or "" for n in p_elem.iter() if n.tag.endswith("}t"))


def _classify_role(text: str, level: int) -> str:
    if not text.strip():
        return "blank"
    low = text.lower()
    if low == "current week sprint status":
        return "current_week"
    if text.startswith("Sprint") or text.startswith("sprint"):
        return "sprint_line"
    if "stories completed this week" in low:
        return "category_completed"
    if "released for partner review" in low:
        return "category_released"
    if "in-progress this week" in low or "in progress this week" in low:
        return "category_inprogress"
    if level == 7:
        return "category_other"
    if level == 1:
        return "story_item"
    return "other"


def _para_level(p_elem) -> int:
    p_pr = p_elem.find(qn("a:pPr"))
    if p_pr is None or p_pr.get("lvl") is None:
        return 0
    return int(p_pr.get("lvl"))


def _service_title(title: str) -> str:
    base = re.sub(r"^Delivery status\s*[–-]\s*", "", title, flags=re.I)
    return re.sub(r"\s*\(Contd.*\)\s*$", "", base, flags=re.I).strip()


def _find_slide_by_index(prs, slide_index: int):
    if slide_index < 1 or slide_index > len(prs.slides):
        return None
    return prs.slides[slide_index - 1]


def _slide_title(slide) -> str:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_id == 2:
            return shape.text_frame.text.strip()
    return ""


def _find_slide_for_service(prs, service_title: str, slide_index: int | None = None):
    if slide_index is not None:
        slide = _find_slide_by_index(prs, slide_index)
        if slide:
            return slide
    target = service_title.lower()
    for slide in prs.slides:
        title = _slide_title(slide)
        if target in _service_title(title).lower():
            return slide
    return None


def _g10x_ref_for_slide(uds, g10x_prs, slide):
    title = _service_title(_slide_title(slide))
    layout = uds.get_g10x_layout_slide(g10x_prs, title)
    return layout


def apply_reflow_hl_ka(prs, slide, uds, g10x_prs, params: dict[str, Any]) -> bool:
    g10x_layout = _g10x_ref_for_slide(uds, g10x_prs, slide)
    layout_mode = params.get("layout_mode", "expanded")
    budgets = uds.get_hl_budgets(uds.build_layout_profile(g10x_layout))
    uds.fit_highlights_layout(
        slide,
        g10x_layout,
        position_ka=True,
        layout_mode=layout_mode,
        budgets=budgets,
    )
    tighten_hl_and_position_ka(slide, uds, g10x_prs)
    return True


def apply_fix_category_bullets(slide, uds) -> bool:
    hl = uds.get_highlights_shape(slide)
    if not hl:
        return False
    tx_body = hl.table.cell(2, 0).text_frame._txBody
    changed = False
    for p in tx_body.findall(qn("a:p")):
        text = _paragraph_text(p)
        role = _classify_role(text, _para_level(p))
        if role.startswith("category_"):
            uds._ensure_category_header_bullet(p)
            changed = True
    return changed


def apply_remove_extra_sprint_blanks(slide, uds) -> bool:
    hl = uds.get_highlights_shape(slide)
    if not hl:
        return False
    tx_body = hl.table.cell(2, 0).text_frame._txBody
    paras = tx_body.findall(qn("a:p"))
    story_roles = {
        "story_item", "category_inprogress", "category_released", "category_completed",
    }
    removed = 0
    i = 0
    while i < len(paras):
        if not _paragraph_text(paras[i]).strip():
            j = i + 1
            while j < len(paras) and not _paragraph_text(paras[j]).strip():
                j += 1
            prev_role = (
                _classify_role(_paragraph_text(paras[i - 1]), _para_level(paras[i - 1]))
                if i > 0 else ""
            )
            next_role = (
                _classify_role(_paragraph_text(paras[j]), _para_level(paras[j]))
                if j < len(paras) else ""
            )
            if prev_role in story_roles and next_role == "sprint_line" and j - i > 1:
                for k in range(j - 1, i, -1):
                    tx_body.remove(paras[k])
                    removed += 1
                paras = tx_body.findall(qn("a:p"))
                i = j
                continue
        i += 1
    return removed > 0


def apply_fix_title_en_dash(slide) -> bool:
    changed = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if "Delivery status - " not in tf.text:
            continue
        for para in tf.paragraphs:
            for run in para.runs:
                if "Delivery status - " in run.text:
                    run.text = run.text.replace("Delivery status - ", "Delivery status – ")
                    changed = True
        if "Delivery status - " in tf.text:
            tf.text = tf.text.replace("Delivery status - ", "Delivery status – ")
            changed = True
    return changed


def apply_remove_category_story_blanks(slide, uds) -> bool:
    """Remove blank line between category header and first story (HL-SPC-01)."""
    hl = uds.get_highlights_shape(slide)
    if not hl:
        return False
    tx_body = hl.table.cell(2, 0).text_frame._txBody
    paras = tx_body.findall(qn("a:p"))
    removed = 0
    i = 0
    while i < len(paras) - 1:
        role = _classify_role(_paragraph_text(paras[i]), _para_level(paras[i]))
        if role.startswith("category_"):
            j = i + 1
            while j < len(paras) and not _paragraph_text(paras[j]).strip():
                j += 1
            if j > i + 1 and j < len(paras):
                next_role = _classify_role(_paragraph_text(paras[j]), _para_level(paras[j]))
                if next_role == "story_item":
                    for k in range(j - 1, i, -1):
                        tx_body.remove(paras[k])
                        removed += 1
                    paras = tx_body.findall(qn("a:p"))
                    continue
        i += 1
    return removed > 0


def tighten_hl_and_position_ka(slide, uds, g10x_prs) -> bool:
    """Shrink HL/KA tables to content and enforce G10X text-to-KA clearance."""
    g10x_layout = _g10x_ref_for_slide(uds, g10x_prs, slide)
    profile = uds.build_layout_profile(g10x_layout)
    hl = uds.get_highlights_shape(slide)
    ka = uds.get_key_activities_shape(slide)
    if not hl and not ka:
        return False

    changed = False
    if not hl and ka:
        ka_profile = uds.get_ka_layout_profile(g10x_prs, g10x_layout)
        ref_ka = ka_profile.get("ref_ka")
        ka_h = uds.fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
        if ka.top + ka_h > uds.MAX_KA_BOTTOM_EMU:
            ka.top = max(0, uds.MAX_KA_BOTTOM_EMU - ka_h)
            changed = True
        return changed or True

    if not hl:
        return False

    text_bottom = uds._hl_rendered_text_bottom(hl, profile)
    hl_bottom = hl.top + hl.height
    if ka and text_bottom > hl_bottom - _HL_BOTTOM_TOLERANCE_EMU:
        if _ensure_text_clearance(hl, ka, profile, uds):
            hl_bottom = hl.top + hl.height
            text_bottom = uds._hl_rendered_text_bottom(hl, profile)
            changed = True

    pad = uds.HL_CONTENT_BOTTOM_PAD
    target_bottom = text_bottom + pad
    current_bottom = hl.top + hl.height
    waste_in = (current_bottom - text_bottom) / EMU_PER_INCH

    if current_bottom > target_bottom + 45720:
        new_h = target_bottom - hl.top
        if waste_in <= SPARSE_WASTE_FORCE_TIGHT_IN:
            min_h = uds._min_hl_height_for_text(
                uds._count_hl_paragraphs(hl), profile, hl=hl
            )
            new_h = max(new_h, min_h)
        r0, r1 = profile["r0"], profile["r1"]
        ref_pad = profile["ref_pad"]
        min_pad = max(int(ref_pad * 0.3), 91440)
        content_h = max(new_h - r0 - r1 - min_pad, profile["r0"])
        uds._set_table_shape_height(hl, [r0, r1, content_h], new_h)
        changed = True

    if ka:
        prev_top = ka.top
        if not uds._is_dense_hl(hl, profile):
            uds._fit_hl_content_only(hl, profile)
        ka_profile = uds.get_ka_layout_profile(g10x_prs, g10x_layout)
        ref_ka = ka_profile.get("ref_ka")
        ka_h = uds.fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
        uds._position_ka_below_hl(hl, ka, profile, ka_h=ka_h)
        if abs(ka.top - prev_top) > 4572:
            changed = True

    return changed


def shrink_ka_table(slide, uds, g10x_prs) -> bool:
    """Fit KA table height to bullet items only."""
    g10x_layout = _g10x_ref_for_slide(uds, g10x_prs, slide)
    ka = uds.get_key_activities_shape(slide)
    if not ka:
        return False
    ka_profile = uds.get_ka_layout_profile(g10x_prs, g10x_layout)
    ref_ka = ka_profile.get("ref_ka")
    prev_h = ka.height
    uds.fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
    return abs(ka.height - prev_h) > 4572


def fix_ka_footer_overflow(slide, uds, g10x_prs) -> bool:
    """Raise KA so its bottom sits within the footer safe zone."""
    g10x_layout = _g10x_ref_for_slide(uds, g10x_prs, slide)
    profile = uds.build_layout_profile(g10x_layout)
    hl = uds.get_highlights_shape(slide)
    ka = uds.get_key_activities_shape(slide)
    if not ka:
        return False

    changed = False
    ka_profile = uds.get_ka_layout_profile(g10x_prs, g10x_layout)
    ref_ka = ka_profile.get("ref_ka")
    ka_h = uds.fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
    max_top = max(0, uds.MAX_KA_BOTTOM_EMU - ka_h)
    if hl:
        gap = hl_ka_tab_gap_emu(profile.get("canonical_line_height_emu"))
        min_top = hl.top + hl.height + gap
        new_top = max(min_top, min(ka.top, max_top))
    else:
        new_top = max_top
    if ka.top != new_top:
        ka.top = new_top
        changed = True
    if ka.top + ka.height > uds.MAX_KA_BOTTOM_EMU:
        ka.top = max(0, uds.MAX_KA_BOTTOM_EMU - ka.height)
        changed = True
    return changed


def compact_vertical_layout(slide, uds, g10x_prs) -> bool:
    """
    Tighten sparse slides: HL sized to last text line + small pad (G10X), then
    seat KA directly below the HL table with standard gap — not footer-anchored.
    """
    g10x_layout = _g10x_ref_for_slide(uds, g10x_prs, slide)
    profile = uds.build_layout_profile(g10x_layout)
    hl = uds.get_highlights_shape(slide)
    ka = uds.get_key_activities_shape(slide)
    if not ka:
        return tighten_hl_and_position_ka(slide, uds, g10x_prs)

    changed = False
    ka_profile = uds.get_ka_layout_profile(g10x_prs, g10x_layout)
    ref_ka = ka_profile.get("ref_ka")

    if hl:
        prev_h = hl.height
        uds._fit_hl_content_only(hl, profile)
        if abs(hl.height - prev_h) > 4572:
            changed = True

    prev_top = ka.top
    ka_h = uds.fit_key_activities_table(ka, ka_profile, position_ref=ref_ka)
    if hl:
        uds._position_ka_below_hl(hl, ka, profile, ka_h=ka_h)
    else:
        ka.top = uds._clamp_ka_top_to_footer(ka_h, ka.top)
    if ka.top != prev_top:
        changed = True

    return changed


def shrink_hl_only(slide, uds, g10x_prs) -> bool:
    """Shrink Highlights to content without moving a footer-anchored KA."""
    g10x_layout = _g10x_ref_for_slide(uds, g10x_prs, slide)
    profile = uds.build_layout_profile(g10x_layout)
    budgets = uds.get_hl_budgets(profile)
    hl = uds.get_highlights_shape(slide)
    if not hl:
        return False
    prev_h = hl.height
    uds.fit_highlights_table(hl, profile, layout_mode="normal", budgets=budgets)
    return abs(hl.height - prev_h) > 4572


def _load_raw_by_service(content_json: Path) -> dict[str, dict[str, Any]]:
    with open(content_json, encoding="utf-8") as f:
        data = json.load(f)
    out: dict[str, dict[str, Any]] = {}
    for slide in data.get("slides", []):
        title = slide.get("title", "")
        if title:
            out[title] = slide
        pname = slide.get("project_name", "")
        if pname:
            out[pname] = slide
    return out


def apply_layout_repair(
    slide,
    uds,
    g10x_prs,
    raw: dict[str, Any] | None,
    params: dict[str, Any],
) -> bool:
    """Full geometry repair: shrink sparse HL, place KA on main, ensure clearance."""
    g10x_layout = _g10x_ref_for_slide(uds, g10x_prs, slide)
    profile = uds.build_layout_profile(g10x_layout)
    budgets = uds.get_hl_budgets(profile)
    layout_mode = params.get("layout_mode", "normal")
    changed = False

    hl = uds.get_highlights_shape(slide)
    if hl:
        uds.fit_highlights_table(hl, profile, layout_mode=layout_mode, budgets=budgets)
        changed = True

    if raw:
        ka = uds.get_key_activities_shape(slide)
        if ka and hl:
            if _ensure_text_clearance(hl, ka, profile, uds):
                changed = True
        elif uds.apply_ka_on_main_slide(
            slide, profile, raw, g10x_prs, g10x_layout, budgets
        ):
            changed = True
        else:
            uds.fit_highlights_layout(
                slide,
                g10x_layout,
                position_ka=True,
                layout_mode=params.get("layout_mode", "expanded"),
                budgets=budgets,
            )
            changed = True
    else:
        ka = uds.get_key_activities_shape(slide)
        if hl and ka:
            if _ensure_text_clearance(hl, ka, profile, uds):
                changed = True
        elif hl:
            uds.fit_highlights_layout(
                slide,
                g10x_layout,
                position_ka=bool(ka),
                layout_mode=layout_mode,
                budgets=budgets,
            )
            changed = True

    if tighten_hl_and_position_ka(slide, uds, g10x_prs):
        changed = True
    return changed


def apply_shrink_sparse_hl(
    slide,
    uds,
    g10x_prs,
    params: dict[str, Any],
    *,
    raw: dict[str, Any] | None = None,
) -> bool:
    if raw:
        return apply_layout_repair(slide, uds, g10x_prs, raw, params)
    hl = uds.get_highlights_shape(slide)
    if not hl:
        return False
    g10x_layout = _g10x_ref_for_slide(uds, g10x_prs, slide)
    profile = uds.build_layout_profile(g10x_layout)
    budgets = uds.get_hl_budgets(profile)
    layout_mode = params.get("layout_mode", "normal")
    uds.fit_highlights_table(hl, profile, layout_mode=layout_mode, budgets=budgets)
    ka = uds.get_key_activities_shape(slide)
    if ka:
        _ensure_text_clearance(hl, ka, profile, uds)
    return tighten_hl_and_position_ka(slide, uds, g10x_prs) or True


def apply_rebuild_with_hints(
    content_json: Path,
    ppt_path: Path,
    service_title: str,
    params: dict[str, Any],
    hints_accumulator: dict[str, Any],
) -> bool:
    entry = dict(params)
    entry.setdefault("pack_all_sections_on_main", True)
    hints_accumulator[service_title] = entry
    return True


def run_rebuild(
    content_json: Path,
    ppt_path: Path,
    layout_hints: dict[str, Any] | None = None,
) -> None:
    hints_path = ppt_path.with_suffix(".layout_hints.json")
    cmd = [
        sys.executable,
        str(PPT_BUILDER),
        "--content",
        str(content_json.resolve()),
        "--output",
        str(ppt_path.resolve()),
    ]
    if layout_hints:
        with open(hints_path, "w", encoding="utf-8") as f:
            json.dump(layout_hints, f, indent=2)
        cmd.extend(["--layout-hints", str(hints_path.resolve())])
    subprocess.run(cmd, cwd=str(REPO_ROOT), check=True)


def run_full_rebuild(content_json: Path, ppt_path: Path) -> None:
    """Rebuild entire deck from content JSON (builder applies latest layout rules)."""
    run_rebuild(content_json, ppt_path, layout_hints=None)


def apply_fix_plan(
    ppt_path: Path,
    fix_plan: dict[str, Any],
    *,
    content_json: Path | None = None,
    g10x_path: Path | None = None,
) -> dict[str, Any]:
    """
    Apply all fixes from a FixPlan. In-place fixes mutate ppt_path;
    rebuild_with_hints are collected and run once at the end if content_json given.
    """
    uds = _import_uds()
    g10x_file = g10x_path or G10X_TEMPLATE
    g10x_prs = Presentation(str(g10x_file))
    prs = Presentation(str(ppt_path))

    applied: list[dict[str, Any]] = []
    skipped: list[dict[str, Any]] = []
    rebuild_hints: dict[str, Any] = {}
    raw_by_service: dict[str, dict[str, Any]] = {}
    if content_json and content_json.is_file():
        raw_by_service = _load_raw_by_service(content_json)

    for fix in fix_plan.get("fixes", []):
        action = fix.get("action")
        if action not in ALLOWED_ACTIONS:
            skipped.append({**fix, "error": "unknown action"})
            continue

        slide_index = fix.get("slide_index")
        service_title = fix.get("service_title", "")
        params = fix.get("params") or {}

        if action == "rebuild_with_hints":
            if not content_json:
                skipped.append({**fix, "error": "content_json required for rebuild"})
                continue
            ok = apply_rebuild_with_hints(
                content_json, ppt_path, service_title, params, rebuild_hints
            )
            applied.append({**fix, "ok": ok, "deferred": True})
            continue

        slide = _find_slide_for_service(prs, service_title, slide_index)
        if slide is None:
            skipped.append({**fix, "error": "slide not found"})
            continue

        ok = False
        raw = raw_by_service.get(service_title)
        if action == "reflow_hl_ka":
            ok = apply_reflow_hl_ka(prs, slide, uds, g10x_prs, params)
        elif action == "fix_category_bullets":
            ok = apply_fix_category_bullets(slide, uds)
        elif action == "remove_extra_sprint_blanks":
            ok = apply_remove_extra_sprint_blanks(slide, uds)
        elif action == "remove_category_story_blanks":
            ok = apply_remove_category_story_blanks(slide, uds)
        elif action == "fix_title_en_dash":
            ok = apply_fix_title_en_dash(slide)
        elif action == "layout_repair":
            ok = apply_layout_repair(slide, uds, g10x_prs, raw, params)
        elif action == "shrink_sparse_hl":
            ok = apply_shrink_sparse_hl(slide, uds, g10x_prs, params, raw=raw)

        applied.append({**fix, "ok": ok})

    prs.save(str(ppt_path))

    rebuilt = False
    if rebuild_hints and content_json:
        run_rebuild(content_json, ppt_path, rebuild_hints)
        rebuilt = True

    return {
        "applied": applied,
        "skipped": skipped,
        "rebuilt": rebuilt,
        "layout_hints": rebuild_hints,
    }
