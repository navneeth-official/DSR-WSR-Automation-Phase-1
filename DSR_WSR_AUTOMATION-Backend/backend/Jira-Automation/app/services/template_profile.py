"""Scan a WSR template and build a profile for content insertion."""

from __future__ import annotations

import json
import re
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.oxml.ns import qn

from app.services.ppt_shape_utils import (
    get_highlights_shape,
    get_key_activities_shape,
    has_highlights_table,
    is_contd_title,
    is_delivery_slide_title,
    normalize_title_text,
    paragraph_text,
    service_suffix_from_title,
    slide_title_text,
)

CATEGORY_PATTERNS: dict[str, tuple[str, ...]] = {
    "completed": ("completed this week",),
    "released": ("released for partner", "released for partner review"),
    "inprogress": ("in-progress this week", "in progress this week"),
}


@dataclass
class ParagraphRoles:
    sprint: int = 0
    current_week: int = 1
    completed_hdr: int | None = None
    completed_bullet: int | None = None
    released_hdr: int | None = None
    released_bullet: int | None = None
    inprogress_hdr: int | None = None
    inprogress_bullet: int | None = None
    sprint_gap: int | None = None


@dataclass
class ServiceProfile:
    service_title: str
    main_slide_index: int
    contd_slide_indices: list[int] = field(default_factory=list)
    hl_content_row: int = 2
    hl_content_col: int = 0
    paragraph_roles: ParagraphRoles = field(default_factory=ParagraphRoles)
    ka_mode: str = "none"  # none | standalone | embedded
    ka_embedded_start_row: int | None = None
    hl_row_count: int = 3


@dataclass
class TemplateProfile:
    template_file: str
    services: dict[str, ServiceProfile] = field(default_factory=dict)

    def service_for_title(self, title_fragment: str) -> ServiceProfile | None:
        needle = title_fragment.strip().lower()
        for key, svc in self.services.items():
            if key.lower() in needle or needle in key.lower():
                return svc
        for svc in self.services.values():
            if needle in svc.service_title.lower():
                return svc
        return None

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)

    @classmethod
    def from_dict(cls, data: dict[str, Any]) -> TemplateProfile:
        services = {}
        for key, raw in data.get("services", {}).items():
            roles = ParagraphRoles(**raw.get("paragraph_roles", {}))
            services[key] = ServiceProfile(
                service_title=raw["service_title"],
                main_slide_index=raw["main_slide_index"],
                contd_slide_indices=list(raw.get("contd_slide_indices", [])),
                hl_content_row=raw.get("hl_content_row", 2),
                hl_content_col=raw.get("hl_content_col", 0),
                paragraph_roles=roles,
                ka_mode=raw.get("ka_mode", "none"),
                ka_embedded_start_row=raw.get("ka_embedded_start_row"),
                hl_row_count=raw.get("hl_row_count", 3),
            )
        return cls(
            template_file=data.get("template_file", ""),
            services=services,
        )


def _normalize_match(text: str) -> str:
    return re.sub(r"\s+", " ", normalize_title_text(text).lower())


def _find_para_index(paras, *needles: str, start: int = 0) -> int | None:
    for i in range(start, len(paras)):
        text = _normalize_match(paragraph_text(paras[i]))
        if any(n in text for n in needles):
            return i
    return None


def _find_sprint_para_index(paras) -> int | None:
    for i, p in enumerate(paras):
        text = _normalize_match(paragraph_text(p)).strip()
        if not text or text == "spur":
            continue
        if "current week sprint status" in text:
            continue
        if text.startswith("sprint") or "sprint -" in text or "sprint –" in text:
            return i
    return None


def _looks_like_category(text: str) -> bool:
    t = _normalize_match(text)
    return any(
        p in t
        for patterns in CATEGORY_PATTERNS.values()
        for p in patterns
    ) or ("stories " in t and "week" in t)


def _find_story_after_header(paras, hdr_idx: int | None) -> int | None:
    if hdr_idx is None:
        return None
    for i in range(hdr_idx + 1, len(paras)):
        text = paragraph_text(paras[i]).strip()
        if not text:
            continue
        if _looks_like_category(text):
            continue
        if "current week sprint status" in _normalize_match(text):
            continue
        return i
    return None


def _detect_ka_mode(slide, hl_shape) -> tuple[str, int | None]:
    if get_key_activities_shape(slide) is not None:
        return "standalone", None
    table = hl_shape.table
    for ri in range(len(table.rows)):
        try:
            label = table.cell(ri, 0).text.strip().lower()
        except (IndexError, AttributeError):
            continue
        if "key activit" in label:
            return "embedded", ri
    return "none", None


def _scan_paragraph_roles(paras) -> ParagraphRoles:
    sprint = _find_sprint_para_index(paras) or 0
    current_week = _find_para_index(paras, "current week sprint status") or 1
    comp_hdr = _find_para_index(paras, *CATEGORY_PATTERNS["completed"])
    rel_hdr = _find_para_index(paras, *CATEGORY_PATTERNS["released"])
    prog_hdr = _find_para_index(paras, *CATEGORY_PATTERNS["inprogress"])

    return ParagraphRoles(
        sprint=sprint,
        current_week=current_week,
        completed_hdr=comp_hdr,
        completed_bullet=_find_story_after_header(paras, comp_hdr),
        released_hdr=rel_hdr,
        released_bullet=_find_story_after_header(paras, rel_hdr),
        inprogress_hdr=prog_hdr,
        inprogress_bullet=_find_story_after_header(paras, prog_hdr),
    )


def _is_generic_delivery_suffix(suffix: str) -> bool:
    """True when the title has no project/service name after 'Delivery Status'."""
    norm = normalize_title_text(suffix).lower()
    return norm in ("delivery status", "delivery", "")


def scan_template_profile(template_path: Path | str) -> TemplateProfile:
    """Build a template profile by reading slide structure and paragraph XML roles."""
    path = Path(template_path).resolve()
    prs = Presentation(str(path))
    profile = TemplateProfile(template_file=path.name)

    mains: dict[str, int] = {}
    contds: dict[str, list[int]] = {}

    for idx, slide in enumerate(prs.slides):
        title = slide_title_text(slide)
        if not is_delivery_slide_title(title):
            continue
        suffix = service_suffix_from_title(title)
        is_contd = is_contd_title(title)
        hl_ok = False
        try:
            get_highlights_shape(slide)
            hl_ok = True
        except ValueError:
            pass

        if _is_generic_delivery_suffix(suffix):
            continue

        if is_contd:
            has_ka = get_key_activities_shape(slide) is not None
            if hl_ok or has_ka:
                contds.setdefault(suffix, []).append(idx)
        elif has_highlights_table(slide):
            mains[suffix] = idx

    for suffix, main_idx in mains.items():
        slide = prs.slides[main_idx]
        hl = get_highlights_shape(slide)
        cell = hl.table.cell(2, 0)
        paras = cell.text_frame._txBody.findall(qn("a:p"))
        ka_mode, ka_row = _detect_ka_mode(slide, hl)
        profile.services[suffix] = ServiceProfile(
            service_title=suffix,
            main_slide_index=main_idx,
            contd_slide_indices=sorted(contds.get(suffix, [])),
            paragraph_roles=_scan_paragraph_roles(paras),
            ka_mode=ka_mode,
            ka_embedded_start_row=ka_row,
            hl_row_count=len(hl.table.rows),
        )

    return profile


def clone_prototype_service_profile(
    profile: TemplateProfile,
    prototype_name: str,
    project_maps: list,
) -> TemplateProfile:
    """Copy the prototype service scan onto every provisioned project slide."""
    proto = profile.services.get(prototype_name)
    if proto is None:
        proto = profile.service_for_title(prototype_name)
    if proto is None and profile.services:
        proto = next(iter(profile.services.values()))
    if proto is None:
        return profile

    updated = TemplateProfile(template_file=profile.template_file, services=dict(profile.services))
    for proj in project_maps:
        cloned = ServiceProfile(
            service_title=proj.project_name,
            main_slide_index=proto.main_slide_index,
            contd_slide_indices=list(proto.contd_slide_indices),
            hl_content_row=proto.hl_content_row,
            hl_content_col=proto.hl_content_col,
            paragraph_roles=ParagraphRoles(
                sprint=proto.paragraph_roles.sprint,
                current_week=proto.paragraph_roles.current_week,
                completed_hdr=proto.paragraph_roles.completed_hdr,
                completed_bullet=proto.paragraph_roles.completed_bullet,
                released_hdr=proto.paragraph_roles.released_hdr,
                released_bullet=proto.paragraph_roles.released_bullet,
                inprogress_hdr=proto.paragraph_roles.inprogress_hdr,
                inprogress_bullet=proto.paragraph_roles.inprogress_bullet,
                sprint_gap=proto.paragraph_roles.sprint_gap,
            ),
            ka_mode=proto.ka_mode,
            ka_embedded_start_row=proto.ka_embedded_start_row,
            hl_row_count=proto.hl_row_count,
        )
        updated.services[proj.project_name] = cloned

    from app.wsr_engine.template_mode import is_placeholder_name

    if is_placeholder_name(prototype_name):
        updated.services.pop(prototype_name, None)

    return updated


def save_template_profile(profile: TemplateProfile, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(profile.to_dict(), indent=2), encoding="utf-8")
    return path


def load_template_profile(path: Path) -> TemplateProfile:
    return TemplateProfile.from_dict(json.loads(path.read_text(encoding="utf-8")))
