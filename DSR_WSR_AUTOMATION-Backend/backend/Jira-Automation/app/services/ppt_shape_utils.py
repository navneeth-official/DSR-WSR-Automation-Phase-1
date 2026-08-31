"""Shared PowerPoint shape traversal helpers for WSR template scanning and fill."""

from __future__ import annotations

from pptx.oxml.ns import qn


def iter_all_shapes(shapes):
    """Yield shapes on a slide, including nested group members."""
    for shape in shapes:
        yield shape
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            yield from iter_all_shapes(child_shapes)


def normalize_title_text(text: str) -> str:
    t = (text or "").replace("\u200b", "").replace("\xa0", " ")
    t = t.replace("\u2013", "-").replace("\u2014", "-")
    return t.strip()


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text or ""


def slide_title_text(slide) -> str:
    for shape in iter_all_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = normalize_title_text(shape_text(shape))
        if text.lower() == "index":
            return "Index"
    for shape in iter_all_shapes(slide.shapes):
        if shape.shape_id == 2 and shape.has_text_frame:
            return shape_text(shape)
    for shape in iter_all_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = normalize_title_text(shape_text(shape))
        if "delivery status" in text.lower():
            return shape_text(shape)
    return ""


def paragraph_text(p_elem) -> str:
    parts: list[str] = []
    for node in p_elem.iter():
        if node.tag == qn("a:t") and node.text:
            parts.append(node.text)
    return "".join(parts)


def is_delivery_slide_title(title: str) -> bool:
    return "delivery status" in normalize_title_text(title).lower()


def is_contd_title(title: str) -> bool:
    return "(contd" in normalize_title_text(title).lower()


def service_suffix_from_title(title: str) -> str:
    """Extract service name from 'Delivery Status - Foo (Contd..)'."""
    t = normalize_title_text(title)
    lower = t.lower()
    if "delivery status" in lower:
        t = t.split("-", 1)[-1].strip()
    if "(contd" in t.lower():
        t = t[: t.lower().index("(contd")].strip()
    return t.strip()


def is_highlights_table(shape) -> bool:
    if not shape.has_table:
        return False
    try:
        header = shape.table.cell(0, 0).text.strip().lower()
        rows = len(shape.table.rows)
        cols = len(shape.table.columns)
    except (IndexError, AttributeError):
        return False
    return "highlights" in header and rows >= 3 and cols >= 3


def has_highlights_table(slide) -> bool:
    """True when the slide contains a Highlights table (not bare title match)."""
    return any(is_highlights_table(shape) for shape in iter_all_shapes(slide.shapes))


def is_ka_table(shape) -> bool:
    if not shape.has_table:
        return False
    try:
        header = shape.table.cell(0, 0).text.strip().lower()
        rows = len(shape.table.rows)
        cols = len(shape.table.columns)
    except (IndexError, AttributeError):
        return False
    return "key activit" in header and rows >= 2 and cols >= 1


def has_combined_hl_ka_table(slide) -> bool:
    """True when HL and KA share a single 5+ row table."""
    for shape in iter_all_shapes(slide.shapes):
        if not shape.has_table or len(shape.table.columns) < 3:
            continue
        if len(shape.table.rows) < 5:
            continue
        try:
            h0 = shape.table.cell(0, 0).text.strip().lower()
            h3 = shape.table.cell(3, 0).text.strip().lower()
        except (IndexError, AttributeError):
            continue
        if "highlights" in h0 and "key activit" in h3:
            return True
    return False


def get_highlights_shape(slide):
    """Find Highlights table by header text (content-based, not shape ID)."""
    for shape in iter_all_shapes(slide.shapes):
        if is_highlights_table(shape):
            return shape
    for sid in (9, 7, 4):
        sh = next(
            (s for s in iter_all_shapes(slide.shapes) if s.shape_id == sid and s.has_table),
            None,
        )
        if sh and is_highlights_table(sh):
            return sh
    raise ValueError("No highlights table found on slide")


def get_key_activities_shape(slide):
    """Find standalone KA table by header text."""
    for shape in iter_all_shapes(slide.shapes):
        if is_ka_table(shape):
            return shape
    for sid in (6, 8, 4, 3, 5):
        sh = next(
            (s for s in iter_all_shapes(slide.shapes) if s.shape_id == sid and s.has_table),
            None,
        )
        if sh and is_ka_table(sh):
            return sh
    return None


def find_title_shape(slide):
    for shape in iter_all_shapes(slide.shapes):
        if shape.shape_id == 2 and shape.has_text_frame:
            return shape
    for shape in iter_all_shapes(slide.shapes):
        if shape.has_text_frame:
            text = normalize_title_text(shape_text(shape))
            if "delivery status" in text.lower():
                return shape
    raise RuntimeError("Title shape not found on slide")
