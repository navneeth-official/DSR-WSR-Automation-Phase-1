"""Detect and update the WSR cover/index slide date from template formatting."""

from __future__ import annotations

import re
from dataclasses import dataclass
from datetime import date

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Pt

from app.services.ppt_layout_metrics import EMU_PER_INCH

_ENGLISH_MONTHS = (
    "January",
    "February",
    "March",
    "April",
    "May",
    "June",
    "July",
    "August",
    "September",
    "October",
    "November",
    "December",
)
_ENGLISH_MONTHS_UPPER = tuple(m.upper() for m in _ENGLISH_MONTHS)
_MONTH_BY_NAME = {name.lower(): i + 1 for i, name in enumerate(_ENGLISH_MONTHS)}
_MONTH_BY_ABBR = {name[:3].lower(): i + 1 for i, name in enumerate(_ENGLISH_MONTHS)}

_COVER_DATE_RE = re.compile(
    r"^\s*(?P<day>\d{1,2})\s+(?P<month>[A-Za-z]{3,9}),?\s*(?P<year>\d{4})\s*$"
)
COVER_DATE_SLIDE_MARGIN_EMU = int(0.06 * EMU_PER_INCH)
MIN_COVER_DATE_FONT_PT = 14.0
_DEFAULT_COVER_DATE_FONT_PT = 24.0
# Average glyph width / font size for Calibri-like proportional fonts.
_CHAR_WIDTH_FACTOR = 0.52
_LINE_HEIGHT_FACTOR = 1.2

_SKIP_TEXT_MARKERS = (
    "weekly",
    "report",
    "hebxxx",
    "proprietary",
    "customer",
    "obsession",
    "sustainment",
    "cost,",
    "pharmacy",
    "wentworth",
    "location",
    "pricing",
    "supplier",
    "pam",
)


@dataclass(frozen=True)
class _ShapeGeometry:
    top: int
    left: int
    width: int
    height: int


@dataclass(frozen=True)
class CoverDateFormat:
    day_zero_pad: bool
    month_style: str  # title | upper_full | upper_abbrev
    comma_after_month: bool


def _month_token(month_num: int, profile: CoverDateFormat) -> str:
    if profile.month_style == "upper_abbrev":
        return _ENGLISH_MONTHS[month_num - 1][:3].upper()
    if profile.month_style == "upper_full":
        return _ENGLISH_MONTHS_UPPER[month_num - 1]
    return _ENGLISH_MONTHS[month_num - 1]


def _day_token(day: int, profile: CoverDateFormat) -> str:
    return f"{day:02d}" if profile.day_zero_pad else str(day)


def format_cover_date(value: date, profile: CoverDateFormat) -> str:
    day = _day_token(value.day, profile)
    month = _month_token(value.month, profile)
    comma = "," if profile.comma_after_month else ""
    return f"{day} {month}{comma} {value.year}"


def parse_cover_date_text(text: str) -> tuple[date, CoverDateFormat] | None:
    match = _COVER_DATE_RE.match(text.strip())
    if not match:
        return None

    day_s = match.group("day")
    month_s = match.group("month")
    year_s = match.group("year")
    month_key = month_s.lower()
    month_num = _MONTH_BY_NAME.get(month_key) or _MONTH_BY_ABBR.get(month_key[:3])
    if month_num is None:
        return None

    if len(month_s) == 3 and month_s.isupper():
        month_style = "upper_abbrev"
    elif month_s.isupper():
        month_style = "upper_full"
    else:
        month_style = "title"

    profile = CoverDateFormat(
        day_zero_pad=len(day_s) == 2 and day_s.startswith("0"),
        month_style=month_style,
        comma_after_month="," in text,
    )
    return date(int(year_s), month_num, int(day_s)), profile


def _is_cover_date_candidate(text: str, shape_name: str) -> bool:
    stripped = text.strip()
    if not stripped or "\n" in stripped:
        return False
    lower = stripped.lower()
    if any(marker in lower for marker in _SKIP_TEXT_MARKERS):
        return False
    if parse_cover_date_text(stripped) is not None:
        return True
    return "date" in shape_name.lower() and "weekly" not in lower


def find_cover_date_shape(cover_slide):
    parsed_match = None
    named_match = None
    for shape in cover_slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if parse_cover_date_text(text) is not None:
            parsed_match = (shape, parse_cover_date_text(text)[1])
            break
        if named_match is None and _is_cover_date_candidate(text, shape.name):
            named_match = shape
    if parsed_match is not None:
        return parsed_match
    if named_match is not None:
        return named_match, CoverDateFormat(
            day_zero_pad=True, month_style="title", comma_after_month=False
        )
    return None, None


def _update_shape_runs(shape, value: date, profile: CoverDateFormat) -> bool:
    if not shape.text_frame.paragraphs:
        return False
    runs = shape.text_frame.paragraphs[0].runs
    if len(runs) < 2:
        return False

    day = _day_token(value.day, profile)
    month = _month_token(value.month, profile)
    month_with_comma = f"{month}," if profile.comma_after_month else month
    year = str(value.year)

    if len(runs) == 5 and runs[2].text.rstrip().endswith(","):
        runs[0].text = day
        runs[2].text = month_with_comma
        runs[4].text = year
        return True

    if len(runs) >= 4:
        runs[0].text = day
        runs[2].text = month
        year_idx = 3 if len(runs) == 4 else 4
        year_run = runs[year_idx].text
        runs[year_idx].text = f" {year}" if year_run.startswith(" ") else year
        return True

    if len(runs) == 3:
        runs[0].text = day
        runs[1].text = f" {month_with_comma} " if profile.comma_after_month else f" {month} "
        runs[2].text = year
        return True

    return False


def apply_cover_date(shape, value: date, profile: CoverDateFormat) -> None:
    if _update_shape_runs(shape, value, profile):
        return
    shape.text_frame.text = format_cover_date(value, profile)


def _snapshot_geometry(shape) -> _ShapeGeometry:
    return _ShapeGeometry(shape.top, shape.left, shape.width, shape.height)


def _restore_geometry(shape, geometry: _ShapeGeometry) -> None:
    shape.top = geometry.top
    shape.left = geometry.left
    shape.width = geometry.width
    shape.height = geometry.height


def _text_frame_body_pr(text_frame):
    return text_frame._txBody.find(qn("a:bodyPr"))


def _ensure_norm_autofit(text_frame, *, single_line: bool) -> None:
    """Use normAutofit so text shrinks inside a fixed box instead of spAutoFit growing it."""
    body_pr = _text_frame_body_pr(text_frame)
    if body_pr is None:
        return

    for tag in ("a:spAutoFit", "a:noAutofit", "a:normAutofit"):
        existing = body_pr.find(qn(tag))
        if existing is not None:
            body_pr.remove(existing)

    ns = body_pr.nsmap.get("a", "http://schemas.openxmlformats.org/drawingml/2006/main")
    body_pr.append(parse_xml(f'<a:normAutofit xmlns:a="{ns}"/>'))
    if single_line:
        body_pr.set("wrap", "none")


def _text_inset_width_emu(shape) -> int:
    body_pr = _text_frame_body_pr(shape.text_frame)
    if body_pr is None:
        return shape.width

    left_inset = int(body_pr.get("lIns", 0))
    right_inset = int(body_pr.get("rIns", 0))
    return max(shape.width - left_inset - right_inset, int(0.25 * EMU_PER_INCH))


def _dominant_font_pt(shape) -> float:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return run.font.size.pt
    return _DEFAULT_COVER_DATE_FONT_PT


def _set_all_run_font_sizes(shape, font_pt: float) -> None:
    size = Pt(font_pt)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = size


def _estimate_paragraph_width_in(paragraph, font_pt: float) -> float:
    width_pt = 0.0
    for run in paragraph.runs:
        text = run.text or ""
        if not text:
            continue
        run_pt = run.font.size.pt if run.font.size is not None else font_pt
        width_pt += len(text) * run_pt * _CHAR_WIDTH_FACTOR
    return width_pt / 72.0


def _non_empty_paragraph_count(text_frame) -> int:
    return sum(1 for paragraph in text_frame.paragraphs if (paragraph.text or "").strip())


def _fit_cover_date_text_to_box(shape) -> None:
    available_width_in = _text_inset_width_emu(shape) / EMU_PER_INCH
    box_height_in = shape.height / EMU_PER_INCH
    line_count = max(_non_empty_paragraph_count(shape.text_frame), 1)
    font_pt = _dominant_font_pt(shape)

    while font_pt >= MIN_COVER_DATE_FONT_PT:
        max_width_in = max(
            (
                _estimate_paragraph_width_in(paragraph, font_pt)
                for paragraph in shape.text_frame.paragraphs
                if (paragraph.text or "").strip()
            ),
            default=0.0,
        )
        line_height_in = font_pt / 72.0 * _LINE_HEIGHT_FACTOR
        needed_height_in = line_count * line_height_in
        if max_width_in <= available_width_in + 0.015 and needed_height_in <= box_height_in + 0.015:
            return
        font_pt -= 1.0
        _set_all_run_font_sizes(shape, font_pt)


def _clamp_cover_date_to_slide(shape, geometry: _ShapeGeometry, slide_height_emu: int) -> None:
    _restore_geometry(shape, geometry)
    max_bottom = slide_height_emu - COVER_DATE_SLIDE_MARGIN_EMU
    bottom = shape.top + shape.height
    if bottom <= max_bottom:
        return

    overflow = bottom - max_bottom
    shape.top = max(0, shape.top - overflow)
    if shape.top + shape.height > max_bottom:
        shape.height = max(int(0.22 * EMU_PER_INCH), max_bottom - shape.top)


def guard_cover_date_layout(shape, slide_height_emu: int) -> None:
    """Keep cover date text inside its template slot and within slide bounds."""
    geometry = _snapshot_geometry(shape)
    single_line = _non_empty_paragraph_count(shape.text_frame) <= 1
    _ensure_norm_autofit(shape.text_frame, single_line=single_line)
    _fit_cover_date_text_to_box(shape)
    _clamp_cover_date_to_slide(shape, geometry, slide_height_emu)
    _fit_cover_date_text_to_box(shape)


def sync_cover_slide_wsr_date(prs: Presentation, wsr_date: date | str) -> bool:
    """Update slide-1 cover date using the template's existing date format."""
    if isinstance(wsr_date, str):
        wsr_date = date.fromisoformat(wsr_date)

    cover = prs.slides[0]
    shape, profile = find_cover_date_shape(cover)
    if shape is None or profile is None:
        return False

    apply_cover_date(shape, wsr_date, profile)
    guard_cover_date_layout(shape, prs.slide_height)
    return True


# Backward-compatible helpers used by older call sites/tests.
def detect_cover_date_format(cover_slide) -> str:
    _, profile = find_cover_date_shape(cover_slide)
    if profile is None:
        return "title"
    if profile.month_style.startswith("upper"):
        return "upper"
    return "title"


def format_wsr_cover_date(wsr_date: date, fmt: str = "title") -> str:
    profile = CoverDateFormat(
        day_zero_pad=True,
        month_style="upper_full" if fmt == "upper" else "title",
        comma_after_month=False,
    )
    return format_cover_date(wsr_date, profile)
