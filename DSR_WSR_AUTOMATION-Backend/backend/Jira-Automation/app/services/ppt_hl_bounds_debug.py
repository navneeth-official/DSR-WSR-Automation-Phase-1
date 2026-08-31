"""Debug visualization for HL rendered text bounds measurement."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.rendering.com_text_bounds import (
    ComHlTextBounds,
    POINTS_PER_INCH,
    ComTextBoundsMeasurer,
    measure_hl_text_bounds_com,
)
from app.services.ppt_layout_metrics import EMU_PER_INCH
from app.services.ppt_rendered_bounds import (
    _emu_rect_to_px,
    _get_highlights_shape_pptx,
    _hl_content_region_emu,
    measure_hl_text,
)
from app.services.ppt_slide_images import export_slides_to_png


@dataclass
class InkScanResult:
    """Visible text ink detected in the HL content cell from a rendered PNG."""

    text_bottom_px: int
    ink_bbox_px: tuple[int, int, int, int]  # x0, y0, x1, y1 on full slide image
    content_bbox_px: tuple[int, int, int, int]


@dataclass(frozen=True)
class HlBoundaryInches:
    """Axis-aligned rectangle in slide inches (top/bottom are Y; left/right are X)."""

    top: float
    bottom: float
    left: float
    right: float
    source: str
    object_name: str = ""

    @property
    def height(self) -> float:
        return round(self.bottom - self.top, 4)

    @property
    def width(self) -> float:
        return round(self.right - self.left, 4)


@dataclass
class HlGeometryAudit:
    """All HL reference boundaries on one slide for comparison."""

    slide_index: int
    title: str
    entire_table_com: HlBoundaryInches
    entire_table_pptx: HlBoundaryInches
    content_row_com: HlBoundaryInches
    content_row_pptx: HlBoundaryInches
    content_cell_com: HlBoundaryInches
    content_cell_pptx: HlBoundaryInches
    com_text_bottom_in: float | None
    image_text_bottom_in: float | None
    measured_hl_bottom_in: float
    measured_waste_in: float
    measurement_method: str


@dataclass
class HlBoundsDebugSnapshot:
    slide_index: int
    title: str
    measurement_method: str
    text_bottom_in: float
    hl_bottom_in: float
    hl_top_in: float
    waste_in: float
    com_text_bottom_in: float | None
    hl_table_px: tuple[int, int, int, int]
    ink_bbox_px: tuple[int, int, int, int]
    text_bottom_px: int
    hl_bottom_px: int
    slide_width_in: float
    slide_height_in: float
    geometry: HlGeometryAudit | None = None
    content_row_px: tuple[int, int, int, int] | None = None
    content_cell_px: tuple[int, int, int, int] | None = None
    pptx_table_bottom_px: int | None = None
    scan_roi_px: tuple[int, int, int, int] | None = None
    scan_roi_in: tuple[float, float, float, float] | None = None
    pptx_scan_roi_in: tuple[float, float, float, float] | None = None


def _boundary_to_px(
    boundary: HlBoundaryInches,
    *,
    slide_width_in: float,
    slide_height_in: float,
    image_width_px: int,
    image_height_px: int,
) -> tuple[int, int, int, int]:
    x0 = _in_to_px_x(boundary.left, slide_width_in=slide_width_in, image_width_px=image_width_px)
    y0 = _in_to_px_y(boundary.top, slide_height_in=slide_height_in, image_height_px=image_height_px)
    x1 = _in_to_px_x(boundary.right, slide_width_in=slide_width_in, image_width_px=image_width_px)
    y1 = _in_to_px_y(boundary.bottom, slide_height_in=slide_height_in, image_height_px=image_height_px)
    return x0, y0, x1, y1


def _pptx_hl_boundaries(hl_shape) -> dict[str, HlBoundaryInches]:
    """Boundaries derived from python-pptx shape XML (may differ from rendered COM layout)."""
    r0 = hl_shape.table.rows[0].height
    r1 = hl_shape.table.rows[1].height
    r2 = hl_shape.table.rows[2].height
    left = hl_shape.left / EMU_PER_INCH
    right = (hl_shape.left + hl_shape.width) / EMU_PER_INCH
    table_top = hl_shape.top / EMU_PER_INCH
    table_bottom = (hl_shape.top + hl_shape.height) / EMU_PER_INCH
    row_top = (hl_shape.top + r0 + r1) / EMU_PER_INCH
    row_bottom = (hl_shape.top + r0 + r1 + r2) / EMU_PER_INCH

    def _b(top: float, bottom: float, *, name: str) -> HlBoundaryInches:
        return HlBoundaryInches(
            top=round(top, 4),
            bottom=round(bottom, 4),
            left=round(left, 4),
            right=round(right, 4),
            source="python-pptx",
            object_name=name,
        )

    return {
        "entire_table": _b(table_top, table_bottom, name=f"shape id={hl_shape.shape_id}"),
        "content_row": _b(row_top, row_bottom, name="row 3 (sum of row heights)"),
        "content_cell": _b(row_top, row_bottom, name="cell (3,1) via row heights"),
    }


def _collect_com_hl_boundaries(
    ppt_path: Path,
    slide_index: int,
) -> tuple[dict[str, HlBoundaryInches], float | None, str]:
    """Boundaries from PowerPoint COM (rendered layout). Returns (boundaries, text_bottom, shape_name)."""
    import win32com.client  # type: ignore[import-untyped]

    from app.rendering.com_text_bounds import measure_hl_text_bounds_com

    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = 1
    prs = app.Presentations.Open(str(ppt_path), WithWindow=False)
    try:
        slide = prs.Slides(slide_index)
        hl = None
        for sid in (9, 7):
            for i in range(1, int(slide.Shapes.Count) + 1):
                sh = slide.Shapes(i)
                if int(sh.Id) == sid and sh.HasTable and int(sh.Table.Rows.Count) == 3:
                    hl = sh
                    break
            if hl:
                break
        if hl is None:
            raise RuntimeError("Highlights table not found via COM")

        def pt2in(v: float) -> float:
            return round(v / POINTS_PER_INCH, 4)

        table_left = pt2in(float(hl.Left))
        table_right = pt2in(float(hl.Left) + float(hl.Width))
        table_top = pt2in(float(hl.Top))
        table_bottom = pt2in(float(hl.Top) + float(hl.Height))

        content_cell = hl.Table.Cell(3, 1).Shape
        cell_top = pt2in(float(content_cell.Top))
        cell_bottom = pt2in(float(content_cell.Top) + float(content_cell.Height))
        cell_left = pt2in(float(content_cell.Left))
        cell_right = pt2in(float(content_cell.Left) + float(content_cell.Width))

        row_tops: list[float] = []
        row_bottoms: list[float] = []
        for ci in range(1, 4):
            cs = hl.Table.Cell(3, ci).Shape
            row_tops.append(pt2in(float(cs.Top)))
            row_bottoms.append(pt2in(float(cs.Top) + float(cs.Height)))

        com_bounds = measure_hl_text_bounds_com(slide)
        text_bottom = com_bounds.text_bottom_in if com_bounds else None

        def _b(
            top: float,
            bottom: float,
            left: float,
            right: float,
            *,
            name: str,
        ) -> HlBoundaryInches:
            return HlBoundaryInches(
                top=top,
                bottom=bottom,
                left=left,
                right=right,
                source="COM",
                object_name=name,
            )

        boundaries = {
            "entire_table": _b(
                table_top,
                table_bottom,
                table_left,
                table_right,
                name=f"{hl.Name} (Id={hl.Id})",
            ),
            "content_row": _b(
                min(row_tops),
                max(row_bottoms),
                table_left,
                table_right,
                name="Table row 3 (all columns)",
            ),
            "content_cell": _b(
                cell_top,
                cell_bottom,
                cell_left,
                cell_right,
                name="Table.Cell(3,1).Shape",
            ),
        }
        return boundaries, text_bottom, str(hl.Name)
    finally:
        prs.Close()
        app.Quit()


def collect_hl_geometry_audit(
    ppt_path: str | Path,
    slide_index: int,
    *,
    measurement: Any | None = None,
    com_text_bottom_in: float | None = None,
    image_text_bottom_in: float | None = None,
) -> HlGeometryAudit:
    """Collect every HL reference boundary on one slide."""
    ppt_path = Path(ppt_path).resolve()
    prs = Presentation(str(ppt_path))
    pptx_slide = prs.slides[slide_index - 1]
    hl_shape = _get_highlights_shape_pptx(pptx_slide)
    if hl_shape is None:
        raise RuntimeError(f"No Highlights table on slide {slide_index}")

    title = ""
    for sh in pptx_slide.shapes:
        if getattr(sh, "shape_id", None) == 2 and sh.has_text_frame:
            title = sh.text_frame.text.strip()
            break

    pptx_bounds = _pptx_hl_boundaries(hl_shape)
    com_bounds, com_text, _ = _collect_com_hl_boundaries(ppt_path, slide_index)

    if measurement is None:
        measurement = measure_hl_text(
            hl_shape,
            com_bounds=None,
            image_path=None,
        )

    return HlGeometryAudit(
        slide_index=slide_index,
        title=title,
        entire_table_com=com_bounds["entire_table"],
        entire_table_pptx=pptx_bounds["entire_table"],
        content_row_com=com_bounds["content_row"],
        content_row_pptx=pptx_bounds["content_row"],
        content_cell_com=com_bounds["content_cell"],
        content_cell_pptx=pptx_bounds["content_cell"],
        com_text_bottom_in=com_text_bottom_in if com_text_bottom_in is not None else com_text,
        image_text_bottom_in=image_text_bottom_in,
        measured_hl_bottom_in=measurement.hl_bottom_in,
        measured_waste_in=measurement.waste_in,
        measurement_method=measurement.method,
    )


def format_geometry_report(audit: HlGeometryAudit) -> str:
    """Human-readable inch report for all HL boundaries."""
    lines = [
        f"Slide {audit.slide_index}: {audit.title}",
        f"Measurement method: {audit.measurement_method}",
        "",
        "Entire HL table (COM — rendered geometry)",
        f"  Object: {audit.entire_table_com.object_name}",
        f"  Top:    {audit.entire_table_com.top:.4f} in",
        f"  Bottom: {audit.entire_table_com.bottom:.4f} in",
        f"  Height: {audit.entire_table_com.height:.4f} in",
        "",
        "Entire HL table (python-pptx — XML shape bbox)",
        f"  Object: {audit.entire_table_pptx.object_name}",
        f"  Top:    {audit.entire_table_pptx.top:.4f} in",
        f"  Bottom: {audit.entire_table_pptx.bottom:.4f} in",
        f"  Height: {audit.entire_table_pptx.height:.4f} in",
        f"  NOTE: pptx bottom is {audit.entire_table_pptx.bottom - audit.entire_table_com.bottom:.4f} in "
        "below COM/rendered bottom (phantom shape extent in XML)",
        "",
        "Content Row (COM)",
        f"  Object: {audit.content_row_com.object_name}",
        f"  Top:    {audit.content_row_com.top:.4f} in",
        f"  Bottom: {audit.content_row_com.bottom:.4f} in",
        f"  Height: {audit.content_row_com.height:.4f} in",
        "",
        "Content Row (python-pptx row-height sum)",
        f"  Object: {audit.content_row_pptx.object_name}",
        f"  Top:    {audit.content_row_pptx.top:.4f} in",
        f"  Bottom: {audit.content_row_pptx.bottom:.4f} in",
        f"  Height: {audit.content_row_pptx.height:.4f} in",
        "",
        "Content Cell (COM Cell 3,1)",
        f"  Object: {audit.content_cell_com.object_name}",
        f"  Top:    {audit.content_cell_com.top:.4f} in",
        f"  Bottom: {audit.content_cell_com.bottom:.4f} in",
        f"  Height: {audit.content_cell_com.height:.4f} in",
        "",
        "Content Cell (python-pptx)",
        f"  Top:    {audit.content_cell_pptx.top:.4f} in",
        f"  Bottom: {audit.content_cell_pptx.bottom:.4f} in",
        f"  Height: {audit.content_cell_pptx.height:.4f} in",
        "",
        f"Image Text Bottom:  {audit.image_text_bottom_in:.4f} in"
        if audit.image_text_bottom_in is not None
        else "Image Text Bottom:  n/a",
        f"COM Text Bottom:    {audit.com_text_bottom_in:.4f} in"
        if audit.com_text_bottom_in is not None
        else "COM Text Bottom:    n/a",
        "",
        f"hl_bottom used for waste: {audit.measured_hl_bottom_in:.4f} in "
        f"(= COM entire table bottom)",
        f"Measured Waste:           {audit.measured_waste_in:.4f} in",
    ]
    return "\n".join(lines)


def format_scan_roi_report(
    *,
    scan_roi_in: tuple[float, float, float, float],
    pptx_roi_in: tuple[float, float, float, float] | None = None,
) -> str:
    """Report the image ink-detection scan region."""
    left, top, right, bottom = scan_roi_in
    lines = [
        "Image scan ROI (region searched for dark pixels)",
        f"  Left:   {left:.4f} in",
        f"  Top:    {top:.4f} in",
        f"  Right:  {right:.4f} in",
        f"  Bottom: {bottom:.4f} in",
        f"  Height: {bottom - top:.4f} in",
    ]
    if pptx_roi_in is not None:
        _, pptx_top, _, pptx_bottom = pptx_roi_in
        omitted = bottom - pptx_bottom
        lines.extend(
            [
                "",
                "python-pptx ROI (legacy fallback when COM unavailable)",
                f"  Top:    {pptx_top:.4f} in",
                f"  Bottom: {pptx_bottom:.4f} in",
            ]
        )
        if omitted > 0.01:
            lines.append(
                f"  Would omit {omitted:.4f} in of rendered content below y={pptx_bottom:.4f} in"
            )
    return "\n".join(lines)


def _in_to_px_y(value_in: float, *, slide_height_in: float, image_height_px: int) -> int:
    return int(round(value_in / slide_height_in * image_height_px))


def _in_to_px_x(value_in: float, *, slide_width_in: float, image_width_px: int) -> int:
    return int(round(value_in / slide_width_in * image_width_px))


def scan_text_ink_in_content_region(
    image_path: Path | str,
    hl_shape,
    *,
    slide_width_emu: int,
    slide_height_emu: int,
    text_luminance_threshold: int = 200,
    com_bounds: ComHlTextBounds | None = None,
) -> InkScanResult | None:
    """Find visible text ink bbox and bottom-most row in the HL content scan ROI."""
    from PIL import Image

    path = Path(image_path)
    left, top, right, bottom = _hl_content_region_emu(hl_shape, com_bounds=com_bounds)

    with Image.open(path) as img:
        rgb = img.convert("RGB")
        w_px, h_px = rgb.size
        x0, y0, x1, y1 = _emu_rect_to_px(
            left,
            top,
            right,
            bottom,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            image_width_px=w_px,
            image_height_px=h_px,
        )
        x0 = max(0, min(x0, w_px - 1))
        x1 = max(x0 + 1, min(x1, w_px))
        y0 = max(0, min(y0, h_px - 1))
        y1 = max(y0 + 1, min(y1, h_px))

        crop = rgb.crop((x0, y0, x1, y1))
        pixels = crop.load()
        cw, ch = crop.size

        ink_min_x: int | None = None
        ink_min_y: int | None = None
        ink_max_x: int | None = None
        ink_max_y: int | None = None
        text_bottom_px: int | None = None

        for y in range(ch):
            for x in range(cw):
                r, g, b = pixels[x, y]
                if (
                    r < text_luminance_threshold
                    and g < text_luminance_threshold
                    and b < text_luminance_threshold
                ):
                    abs_x = x0 + x
                    abs_y = y0 + y
                    ink_min_x = abs_x if ink_min_x is None else min(ink_min_x, abs_x)
                    ink_min_y = abs_y if ink_min_y is None else min(ink_min_y, abs_y)
                    ink_max_x = abs_x if ink_max_x is None else max(ink_max_x, abs_x)
                    ink_max_y = abs_y if ink_max_y is None else max(ink_max_y, abs_y)
                    if text_bottom_px is None or abs_y > text_bottom_px:
                        text_bottom_px = abs_y

        if (
            text_bottom_px is None
            or ink_min_x is None
            or ink_min_y is None
            or ink_max_x is None
            or ink_max_y is None
        ):
            return None

        return InkScanResult(
            text_bottom_px=text_bottom_px,
            ink_bbox_px=(ink_min_x, ink_min_y, ink_max_x, ink_max_y),
            content_bbox_px=(x0, y0, x1, y1),
        )


def _com_hl_table_px(
    ppt_path: Path,
    slide_index: int,
    *,
    image_width_px: int,
    image_height_px: int,
    slide_width_in: float,
    slide_height_in: float,
) -> tuple[int, int, int, int]:
    """HL table rectangle in slide-image pixels from COM."""
    import win32com.client  # type: ignore[import-untyped]

    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = 1
    prs = app.Presentations.Open(str(ppt_path), WithWindow=False)
    try:
        slide = prs.Slides(slide_index)
        hl = None
        for sid in (9, 7):
            for i in range(1, int(slide.Shapes.Count) + 1):
                sh = slide.Shapes(i)
                if int(sh.Id) == sid and sh.HasTable and int(sh.Table.Rows.Count) == 3:
                    hl = sh
                    break
            if hl:
                break
        if hl is None:
            raise RuntimeError("Highlights table not found")

        left_in = float(hl.Left) / POINTS_PER_INCH
        top_in = float(hl.Top) / POINTS_PER_INCH
        right_in = (float(hl.Left) + float(hl.Width)) / POINTS_PER_INCH
        bottom_in = (float(hl.Top) + float(hl.Height)) / POINTS_PER_INCH
    finally:
        prs.Close()
        app.Quit()

    x0 = _in_to_px_x(left_in, slide_width_in=slide_width_in, image_width_px=image_width_px)
    y0 = _in_to_px_y(top_in, slide_height_in=slide_height_in, image_height_px=image_height_px)
    x1 = _in_to_px_x(right_in, slide_width_in=slide_width_in, image_width_px=image_width_px)
    y1 = _in_to_px_y(bottom_in, slide_height_in=slide_height_in, image_height_px=image_height_px)
    return x0, y0, x1, y1


def build_debug_snapshot(
    ppt_path: str | Path,
    slide_index: int,
    *,
    image_path: Path | str | None = None,
) -> HlBoundsDebugSnapshot:
    """Collect all geometry used for hl_waste_below_text_in on one slide."""
    ppt_path = Path(ppt_path).resolve()
    prs = Presentation(str(ppt_path))
    pptx_slide = prs.slides[slide_index - 1]
    hl_shape = _get_highlights_shape_pptx(pptx_slide)
    if hl_shape is None:
        raise RuntimeError(f"No Highlights table on slide {slide_index}")

    title = ""
    for sh in pptx_slide.shapes:
        if getattr(sh, "shape_id", None) == 2 and sh.has_text_frame:
            title = sh.text_frame.text.strip()
            break

    if image_path is None:
        out_dir = ppt_path.parent / "debug"
        out_dir.mkdir(parents=True, exist_ok=True)
        exported = export_slides_to_png(ppt_path, out_dir, slide_indices=[slide_index])
        image_path = Path(exported[0]["image_path"])

    image_path = Path(image_path)
    from PIL import Image

    with Image.open(image_path) as img:
        image_width_px, image_height_px = img.size

    com_bounds = None
    com_text_bottom_in: float | None = None
    slide_width_emu = 12192000
    slide_height_emu = 6858000
    slide_width_in = 13.333
    slide_height_in = 7.5

    try:
        with ComTextBoundsMeasurer(ppt_path) as measurer:
            slide_width_emu = int(measurer.slide_width_in * EMU_PER_INCH)
            slide_height_emu = int(measurer.slide_height_in * EMU_PER_INCH)
            slide_width_in = measurer.slide_width_in
            slide_height_in = measurer.slide_height_in
            com_bounds = measurer.measure_slide(slide_index)
            if com_bounds is not None:
                com_text_bottom_in = com_bounds.text_bottom_in
    except Exception:
        com_bounds = None

    measurement = measure_hl_text(
        hl_shape,
        com_bounds=com_bounds,
        image_path=image_path,
        slide_width_emu=slide_width_emu,
        slide_height_emu=slide_height_emu,
    )

    ink = scan_text_ink_in_content_region(
        image_path,
        hl_shape,
        slide_width_emu=slide_width_emu,
        slide_height_emu=slide_height_emu,
        com_bounds=com_bounds,
    )
    if ink is None:
        raise RuntimeError("No visible text ink detected in HL content region")

    scan_left, scan_top, scan_right, scan_bottom = _hl_content_region_emu(
        hl_shape, com_bounds=com_bounds
    )
    scan_roi_px = _emu_rect_to_px(
        scan_left,
        scan_top,
        scan_right,
        scan_bottom,
        slide_width_emu=slide_width_emu,
        slide_height_emu=slide_height_emu,
        image_width_px=image_width_px,
        image_height_px=image_height_px,
    )
    scan_roi_in = (
        round(scan_left / EMU_PER_INCH, 4),
        round(scan_top / EMU_PER_INCH, 4),
        round(scan_right / EMU_PER_INCH, 4),
        round(scan_bottom / EMU_PER_INCH, 4),
    )
    pptx_left, pptx_top, pptx_right, pptx_bottom = _hl_content_region_emu(hl_shape)
    pptx_scan_roi_in = (
        round(pptx_left / EMU_PER_INCH, 4),
        round(pptx_top / EMU_PER_INCH, 4),
        round(pptx_right / EMU_PER_INCH, 4),
        round(pptx_bottom / EMU_PER_INCH, 4),
    )

    hl_table_px = _com_hl_table_px(
        ppt_path,
        slide_index,
        image_width_px=image_width_px,
        image_height_px=image_height_px,
        slide_width_in=slide_width_in,
        slide_height_in=slide_height_in,
    )

    hl_bottom_px = _in_to_px_y(
        measurement.hl_bottom_in,
        slide_height_in=slide_height_in,
        image_height_px=image_height_px,
    )

    geometry = collect_hl_geometry_audit(
        ppt_path,
        slide_index,
        measurement=measurement,
        com_text_bottom_in=com_text_bottom_in,
        image_text_bottom_in=measurement.text_bottom_in,
    )

    conv_kw = dict(
        slide_width_in=slide_width_in,
        slide_height_in=slide_height_in,
        image_width_px=image_width_px,
        image_height_px=image_height_px,
    )
    content_row_px = _boundary_to_px(geometry.content_row_com, **conv_kw)
    content_cell_px = _boundary_to_px(geometry.content_cell_com, **conv_kw)
    pptx_table_bottom_px = _in_to_px_y(
        geometry.entire_table_pptx.bottom,
        slide_height_in=slide_height_in,
        image_height_px=image_height_px,
    )

    return HlBoundsDebugSnapshot(
        slide_index=slide_index,
        title=title,
        measurement_method=measurement.method,
        text_bottom_in=measurement.text_bottom_in,
        hl_bottom_in=measurement.hl_bottom_in,
        hl_top_in=com_bounds.hl_top_in if com_bounds else geometry.entire_table_com.top,
        waste_in=measurement.waste_in,
        com_text_bottom_in=com_text_bottom_in,
        hl_table_px=hl_table_px,
        ink_bbox_px=ink.ink_bbox_px,
        text_bottom_px=ink.text_bottom_px,
        hl_bottom_px=hl_bottom_px,
        slide_width_in=slide_width_in,
        slide_height_in=slide_height_in,
        geometry=geometry,
        content_row_px=content_row_px,
        content_cell_px=content_cell_px,
        pptx_table_bottom_px=pptx_table_bottom_px,
        scan_roi_px=scan_roi_px,
        scan_roi_in=scan_roi_in,
        pptx_scan_roi_in=pptx_scan_roi_in,
    )


def render_hl_bounds_debug_image(
    image_path: Path | str,
    snapshot: HlBoundsDebugSnapshot,
    output_path: Path | str,
) -> Path:
    """Draw measurement overlays on the slide PNG and save."""
    from PIL import Image, ImageDraw, ImageFont

    image_path = Path(image_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    base = Image.open(image_path).convert("RGBA")
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    hl_x0, hl_y0, hl_x1, hl_y1 = snapshot.hl_table_px
    ink_x0, ink_y0, ink_x1, ink_y1 = snapshot.ink_bbox_px
    text_y = snapshot.text_bottom_px
    hl_bottom_y = snapshot.hl_bottom_px

    # Yellow shaded band — measured internal whitespace (draw first, behind lines)
    waste_top = min(text_y, hl_bottom_y)
    waste_bottom = max(text_y, hl_bottom_y)
    if waste_bottom > waste_top:
        draw.rectangle(
            [hl_x0 + 2, waste_top, hl_x1 - 2, waste_bottom],
            fill=(255, 255, 0, 60),
            outline=(255, 200, 0, 120),
            width=1,
        )

    # Purple rectangle — content row boundary (COM row 3)
    if snapshot.content_row_px:
        cr_x0, cr_y0, cr_x1, cr_y1 = snapshot.content_row_px
        draw.rectangle(
            [cr_x0, cr_y0, cr_x1, cr_y1],
            outline=(160, 32, 240, 255),
            width=3,
        )

    # Orange rectangle — content cell boundary (COM cell 3,1)
    if snapshot.content_cell_px:
        cc_x0, cc_y0, cc_x1, cc_y1 = snapshot.content_cell_px
        draw.rectangle(
            [cc_x0, cc_y0, cc_x1, cc_y1],
            outline=(255, 140, 0, 255),
            width=3,
        )

    # Cyan dashed rectangle — image ink-detection scan ROI
    if snapshot.scan_roi_px:
        roi_x0, roi_y0, roi_x1, roi_y1 = snapshot.scan_roi_px
        for edge_y in (roi_y0, roi_y1):
            for x in range(roi_x0, roi_x1, 12):
                draw.line(
                    [(x, edge_y), (min(x + 6, roi_x1), edge_y)],
                    fill=(0, 220, 220, 255),
                    width=2,
                )
        for edge_x in (roi_x0, roi_x1):
            for y in range(roi_y0, roi_y1, 12):
                draw.line(
                    [(edge_x, y), (edge_x, min(y + 6, roi_y1))],
                    fill=(0, 220, 220, 255),
                    width=2,
                )

    # Red rectangle — complete HL table boundary (COM rendered geometry)
    draw.rectangle(
        [hl_x0, hl_y0, hl_x1, hl_y1],
        outline=(255, 0, 0, 255),
        width=4,
    )

    # Green rectangle — visible text ink bounding box (image scan)
    draw.rectangle(
        [ink_x0, ink_y0, ink_x1, ink_y1],
        outline=(0, 200, 0, 255),
        width=3,
    )

    # Blue horizontal line — bottom-most visible text pixel row
    draw.line(
        [(hl_x0, text_y), (hl_x1, text_y)],
        fill=(0, 120, 255, 255),
        width=3,
    )

    # Red horizontal line — measured HL table bottom (COM entire table)
    draw.line(
        [(hl_x0, hl_bottom_y), (hl_x1, hl_bottom_y)],
        fill=(255, 0, 0, 255),
        width=3,
    )

    # Dashed gray line — python-pptx inflated table bottom (for comparison only)
    if snapshot.pptx_table_bottom_px is not None:
        pptx_y = snapshot.pptx_table_bottom_px
        for x in range(hl_x0, hl_x1, 16):
            draw.line(
                [(x, pptx_y), (min(x + 8, hl_x1), pptx_y)],
                fill=(128, 128, 128, 200),
                width=2,
            )

    composed = Image.alpha_composite(base, overlay)
    draw_final = ImageDraw.Draw(composed)

    try:
        font = ImageFont.truetype("arial.ttf", 22)
        font_sm = ImageFont.truetype("arial.ttf", 18)
    except OSError:
        font = ImageFont.load_default()
        font_sm = font

    legend_x = hl_x1 + 12 if hl_x1 + 480 < base.width else 20
    legend_y = max(20, hl_y0)
    geo = snapshot.geometry
    pptx_bottom = geo.entire_table_pptx.bottom if geo else None
    lines = [
        f"Slide {snapshot.slide_index}: {snapshot.title[:48]}",
        f"method: {snapshot.measurement_method}",
        f"hl_waste_below_text_in = {snapshot.waste_in:.4f} in",
        f"text_bottom (image) = {snapshot.text_bottom_in:.4f} in",
        f"hl_bottom (COM table) = {snapshot.hl_bottom_in:.4f} in",
        f"COM text_bottom = {snapshot.com_text_bottom_in:.4f} in"
        if snapshot.com_text_bottom_in is not None
        else "COM text_bottom = n/a",
        f"pptx table bottom = {pptx_bottom:.4f} in (NOT used)"
        if pptx_bottom is not None
        else "",
        "",
        "RED = complete HL table (COM)",
        "PURPLE = content row (COM row 3)",
        "ORANGE = content cell (COM cell 3,1)",
        "GREEN = visible text ink bbox",
        "BLUE = bottom-most text pixel",
        "YELLOW = measured whitespace",
        "GRAY dashed = pptx inflated bottom",
        "CYAN dashed = image scan ROI",
    ]
    lines = [ln for ln in lines if ln]
    box_h = 24 * len(lines) + 16
    draw_final.rectangle(
        [legend_x - 8, legend_y - 8, legend_x + 500, legend_y + box_h],
        fill=(255, 255, 255, 220),
        outline=(80, 80, 80, 255),
    )
    for i, line in enumerate(lines):
        draw_final.text(
            (legend_x, legend_y + i * 24),
            line,
            fill=(0, 0, 0, 255),
            font=font if i == 0 else font_sm,
        )

    # Waste dimension arrow on left side of HL box
    arrow_x = max(8, hl_x0 - 18)
    mid_y = (text_y + hl_bottom_y) // 2
    draw_final.line([(arrow_x, text_y), (arrow_x, hl_bottom_y)], fill=(255, 140, 0, 255), width=2)
    draw_final.line([(arrow_x - 5, text_y), (arrow_x + 5, text_y)], fill=(255, 140, 0, 255), width=2)
    draw_final.line(
        [(arrow_x - 5, hl_bottom_y), (arrow_x + 5, hl_bottom_y)],
        fill=(255, 140, 0, 255),
        width=2,
    )
    draw_final.text(
        (arrow_x + 8, mid_y - 10),
        f"{snapshot.waste_in:.3f} in",
        fill=(200, 80, 0, 255),
        font=font_sm,
    )

    composed.convert("RGB").save(output_path)
    return output_path


def generate_hl_bounds_debug_image(
    ppt_path: str | Path,
    slide_index: int,
    output_path: str | Path,
    *,
    image_path: Path | str | None = None,
) -> dict[str, Any]:
    """Build snapshot, render overlay, return summary dict."""
    snapshot = build_debug_snapshot(ppt_path, slide_index, image_path=image_path)
    if image_path is None:
        image_path = Path(ppt_path).parent / "debug" / f"slide_{slide_index:02d}.png"
    out = render_hl_bounds_debug_image(image_path, snapshot, output_path)
    report = format_geometry_report(snapshot.geometry) if snapshot.geometry else ""
    roi_report = ""
    if snapshot.scan_roi_in is not None:
        roi_report = format_scan_roi_report(
            scan_roi_in=snapshot.scan_roi_in,
            pptx_roi_in=snapshot.pptx_scan_roi_in,
        )
    return {
        "output_path": str(out),
        "slide_index": snapshot.slide_index,
        "title": snapshot.title,
        "measurement_method": snapshot.measurement_method,
        "hl_waste_below_text_in": snapshot.waste_in,
        "text_bottom_in": snapshot.text_bottom_in,
        "hl_bottom_in": snapshot.hl_bottom_in,
        "com_text_bottom_in": snapshot.com_text_bottom_in,
        "geometry_report": report,
        "scan_roi_report": roi_report,
        "snapshot": snapshot,
    }
