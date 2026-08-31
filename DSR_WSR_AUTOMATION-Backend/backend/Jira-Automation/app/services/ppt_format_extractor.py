"""Extract structured format metrics from a delivery-status PPTX for AI evaluation."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.oxml.ns import qn

from app.constants.ppt_bullets import (
    CATEGORY_HEADER_BULLET,
    CATEGORY_HEADER_BULLET_FONT,
    CATEGORY_HEADER_LEVEL,
    FORBIDDEN_CATEGORY_HEADER_BULLETS,
    is_valid_category_header_bullet,
)
from app.services.ppt_layout_metrics import (
    CANONICAL_KA_ITEM_SLOTS,
    CANONICAL_PARA_SLOTS,
    count_visual_lines_in_hl,
    effective_hl_utilization,
    hl_waste_below_text_in,
    ka_rendered_text_bottom_in,
    ka_waste_below_text_in,
    rendered_text_bottom_emu,
    text_ka_clearance_in,
    utilization_ratio,
    count_ka_items,
)
from app.services.ppt_shape_utils import (
    is_highlights_table,
    is_ka_table,
    iter_all_shapes,
)

EMU_PER_INCH = 914400

from app.paths import G10X_TEMPLATE, SCRIPTS_DIR

_g10x_prs_cache = None
_uds_module = None


def _uds_helpers():
    """Lazy-load builder layout helpers (shared text-bottom math with repair)."""
    global _g10x_prs_cache, _uds_module
    if _uds_module is None:
        import sys

        scripts = str(SCRIPTS_DIR)
        if scripts not in sys.path:
            sys.path.insert(0, scripts)
        import update_delivery_status as uds  # noqa: WPS433

        _uds_module = uds
        from pptx import Presentation

        _g10x_prs_cache = Presentation(str(G10X_TEMPLATE))
    return _uds_module, _g10x_prs_cache


def _service_base_title(title: str) -> str:
    base = re.sub(r"^Delivery status\s*[–-]\s*", "", title, flags=re.I)
    return re.sub(r"\s*\(Contd.*\)\s*$", "", base, flags=re.I).strip()


def _hl_waste_below_text_in(hl) -> float:
    """Wrap-aware empty HL area below rendered text (matches on-slide appearance)."""
    ref_r2 = hl.table.rows[2].height
    return hl_waste_below_text_in(hl, ref_r2=ref_r2)


def _emu_in(value: int | None) -> float | None:
    if value is None:
        return None
    return round(value / EMU_PER_INCH, 4)


def _paragraph_text(p_elem) -> str:
    return "".join(n.text or "" for n in p_elem.iter() if n.tag.endswith("}t"))


def _para_metrics(p_elem) -> dict[str, Any]:
    text = _paragraph_text(p_elem).strip()
    pPr = p_elem.find(qn("a:pPr"))
    info: dict[str, Any] = {"text": text}
    if pPr is not None:
        info["level"] = int(pPr.get("lvl") or 0)
        bu_char = pPr.find(qn("a:buChar"))
        if bu_char is not None:
            info["bullet"] = bu_char.get("char")
        elif pPr.find(qn("a:buNone")) is not None:
            info["bullet"] = "none"
        bu_font = pPr.find(qn("a:buFont"))
        if bu_font is not None:
            info["bullet_font"] = bu_font.get("typeface")
        spc_bef = pPr.find(qn("a:spcBef"))
        if spc_bef is not None:
            pt = spc_bef.find(qn("a:spcPts"))
            if pt is not None:
                info["spc_bef_pt"] = int(pt.get("val", 0)) / 100
        ln_spc = pPr.find(qn("a:lnSpc"))
        if ln_spc is not None:
            pt = ln_spc.find(qn("a:spcPts"))
            if pt is not None:
                info["line_spacing_pt"] = int(pt.get("val", 0)) / 100
            pct = ln_spc.find(qn("a:spcPct"))
            if pct is not None:
                info["line_spacing_pct"] = int(pct.get("val", 0)) / 1000
    runs = []
    for r in p_elem.findall(qn("a:r")):
        t_elem = r.find(qn("a:t"))
        run: dict[str, Any] = {"text": ((t_elem.text if t_elem is not None else None) or "")[:60]}
        rPr = r.find(qn("a:rPr"))
        if rPr is not None:
            if rPr.get("b") is not None:
                run["bold"] = rPr.get("b") == "1"
            sz_attr = rPr.get("sz")
            if sz_attr is not None:
                run["size_pt"] = int(sz_attr) / 100
            else:
                sz = rPr.find(qn("a:sz"))
                if sz is not None:
                    run["size_pt"] = int(sz.get("val", 0)) / 100
            latin = rPr.find(qn("a:latin"))
            if latin is not None:
                run["font"] = latin.get("typeface")
        runs.append(run)
    if runs:
        info["runs"] = runs
    return info


def _classify_paragraph(p: dict[str, Any]) -> str:
    text = p.get("text", "")
    if not text:
        return "blank"
    low = text.lower()
    if low == "current week sprint status":
        return "current_week"
    if text.startswith("Sprint") or text.startswith("sprint"):
        return "sprint_line"
    if "stories completed this week" in low:
        return "category_completed"
    if "released for partner review" in low:
        return "category_released"
    if "in-progress this week" in low or "in progress this week" in low:
        return "category_inprogress"
    lvl = p.get("level", 0)
    if lvl == 7:
        return "category_other"
    if lvl == 1:
        return "story_item"
    if lvl == 0 and p.get("bullet") in ("•", None) and text.isupper():
        return "project_name"
    return "other"


def _get_highlights_shape(slide):
    for shape in iter_all_shapes(slide.shapes):
        if is_highlights_table(shape):
            return shape
    return None


def _get_ka_shape(slide):
    for shape in iter_all_shapes(slide.shapes):
        if is_ka_table(shape):
            return shape
    return None


def _title_shape(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_id == 2:
            return shape
    return None


def extract_slide(slide, slide_index: int) -> dict[str, Any] | None:
    """Extract format metrics for one slide. Returns None if not a delivery-status slide."""
    title_shape = _title_shape(slide)
    if not title_shape:
        return None
    title = title_shape.text_frame.text.strip()
    if not re.search(r"delivery status|\(contd", title, re.I):
        return None

    entry: dict[str, Any] = {
        "slide_index": slide_index,
        "title": title,
        "is_contd": bool(re.search(r"\(Contd", title, re.I)),
        "title_metrics": _para_metrics(title_shape.text_frame.paragraphs[0]._p),
        "title_position_in": {
            "top": _emu_in(title_shape.top),
            "left": _emu_in(title_shape.left),
        },
    }

    hl = _get_highlights_shape(slide)
    ka = _get_ka_shape(slide)

    if hl:
        cell = hl.table.cell(2, 0)
        paras = [_para_metrics(p) for p in cell.text_frame._txBody.findall(qn("a:p"))]
        for p in paras:
            p["role"] = _classify_paragraph(p)

        filled = [p for p in paras if p.get("text")]
        story_items = [p for p in paras if p.get("role") == "story_item"]
        blank_between_sprints = 0
        for i, p in enumerate(paras):
            if p.get("role") != "blank":
                continue
            prev_role = paras[i - 1].get("role") if i > 0 else ""
            next_role = paras[i + 1].get("role") if i + 1 < len(paras) else ""
            if prev_role in ("story_item", "category_inprogress", "category_released", "category_completed") and next_role == "sprint_line":
                blank_between_sprints += 1

        category_to_story_gaps = 0
        for i, p in enumerate(paras[:-1]):
            if not p.get("role", "").startswith("category_"):
                continue
            if paras[i + 1].get("role") != "blank":
                continue
            # Only flag when blank sits between category header and a story line.
            for j in range(i + 2, len(paras)):
                nxt = paras[j].get("role", "")
                if nxt == "blank":
                    continue
                if nxt == "story_item":
                    category_to_story_gaps += 1
                break

        category_bullet_violations = []
        for p in paras:
            role = p.get("role", "")
            if not role.startswith("category_"):
                continue
            bullet = p.get("bullet")
            bullet_font = p.get("bullet_font")
            level = p.get("level")
            if not is_valid_category_header_bullet(bullet, bullet_font, level):
                issue = "invalid_category_bullet"
                if bullet in FORBIDDEN_CATEGORY_HEADER_BULLETS:
                    issue = "forbidden_bullet_char"
                elif level != CATEGORY_HEADER_LEVEL:
                    issue = "wrong_level"
                elif bullet != CATEGORY_HEADER_BULLET:
                    issue = "wrong_bullet_char"
                elif not bullet_font or "Wingdings" not in bullet_font:
                    issue = "missing_wingdings_font"
                category_bullet_violations.append(
                    {
                        "role": role,
                        "issue": issue,
                        "level": level,
                        "bullet": bullet,
                        "bullet_font": bullet_font,
                        "required": {
                            "level": CATEGORY_HEADER_LEVEL,
                            "bullet": CATEGORY_HEADER_BULLET,
                            "bullet_font": CATEGORY_HEADER_BULLET_FONT,
                        },
                    }
                )

        header_cell = hl.table.cell(0, 0)
        header_paras = [
            _para_metrics(p) for p in header_cell.text_frame._txBody.findall(qn("a:p"))
        ]
        header_metrics = header_paras[0] if header_paras else {}

        visual_lines = count_visual_lines_in_hl(hl)
        filled_count = len(filled)
        entry["highlights"] = {
            "header_metrics": header_metrics,
            "position_in": {
                "top": _emu_in(hl.top),
                "height": _emu_in(hl.height),
                "bottom": _emu_in(hl.top + hl.height),
                "width": _emu_in(hl.width),
            },
            "paragraph_count": len(paras),
            "filled_paragraph_count": filled_count,
            "visual_line_count": visual_lines,
            "story_item_count": len(story_items),
            "blank_between_sprints": blank_between_sprints,
            "category_to_story_blank_gaps": category_to_story_gaps,
            "category_bullet_violations": category_bullet_violations,
            "utilization_ratio": utilization_ratio(filled_count, CANONICAL_PARA_SLOTS),
            "effective_utilization_ratio": effective_hl_utilization(
                filled_count, visual_lines, CANONICAL_PARA_SLOTS
            ),
            "paragraphs": paras,
            "level_distribution": _level_distribution(paras),
            "bullet_distribution": _bullet_distribution(paras),
        }

    if ka:
        content_cell = ka.table.cell(1, 0)
        item_count = count_ka_items(ka)
        entry["key_activities"] = {
            "position_in": {
                "top": _emu_in(ka.top),
                "height": _emu_in(ka.height),
                "bottom": _emu_in(ka.top + ka.height),
            },
            "item_count": item_count,
            "utilization_ratio": utilization_ratio(item_count, CANONICAL_KA_ITEM_SLOTS),
            "header": _para_metrics(ka.table.cell(0, 0).text_frame.paragraphs[0]._p),
            "content_paragraphs": [
                _para_metrics(p._p) for p in content_cell.text_frame.paragraphs if p.text.strip()
            ],
        }
        entry["ka_waste_below_text_in"] = ka_waste_below_text_in(ka)
        entry["ka_rendered_text_bottom_in"] = ka_rendered_text_bottom_in(ka)

    if hl and ka:
        hl_bottom = hl.top + hl.height
        entry["hl_ka_gap_in"] = round((ka.top - hl_bottom) / EMU_PER_INCH, 4)
        try:
            uds, g10x_prs = _uds_helpers()
            service = _service_base_title(title)
            layout = uds.get_g10x_layout_slide(g10x_prs, service)
            if layout is not None:
                profile = uds.build_layout_profile(layout)
                profile["canonical_fill_para_count"] = uds.get_canonical_fill_para_count(
                    g10x_prs
                )
                profile["canonical_line_height_emu"] = uds.get_canonical_line_height_emu(
                    g10x_prs
                )
                text_bottom = uds._hl_rendered_text_bottom(hl, profile)
                entry["estimated_text_bottom_in"] = _emu_in(text_bottom)
                entry["hl_text_bottom_for_fit_in"] = _emu_in(text_bottom)
                clearance = round((ka.top - text_bottom) / EMU_PER_INCH, 4)
                entry["text_ka_clearance_in"] = clearance
                hl_bottom = hl.top + hl.height
                entry["hl_text_overflow_in"] = round(
                    max((text_bottom - hl_bottom) / EMU_PER_INCH, 0), 4
                )
        except Exception:  # noqa: BLE001
            ref_r2 = hl.table.rows[2].height
            est_bottom = rendered_text_bottom_emu(hl, ref_r2=ref_r2)
            entry["estimated_text_bottom_in"] = _emu_in(est_bottom)
            clearance = text_ka_clearance_in(hl, ka, ref_r2=ref_r2)
            if clearance is not None:
                entry["text_ka_clearance_in"] = clearance

    hl_metrics = entry.get("highlights")
    if hl and hl_metrics and not entry.get("_rendered_bounds_deferred"):
        entry["hl_waste_below_text_in"] = _hl_waste_below_text_in(hl)
        entry["hl_text_bounds_method"] = "estimated"

    ka_has_items = ka is not None and count_ka_items(ka) > 0
    if hl and ka is not None:
        entry["layout_type"] = (
            "hl_and_ka" if not entry.get("is_contd") else "hl_only_contd"
        )
    elif hl:
        entry["layout_type"] = (
            "hl_only_contd" if entry.get("is_contd") else "hl_only_main"
        )
    elif ka is not None:
        entry["layout_type"] = (
            "ka_only_contd" if entry.get("is_contd") else "ka_only_main"
        )
    else:
        entry["layout_type"] = "unknown"

    return entry


def _level_distribution(paras: list[dict]) -> dict[str, int]:
    dist: dict[str, int] = {}
    for p in paras:
        if not p.get("text"):
            continue
        key = str(p.get("level", 0))
        dist[key] = dist.get(key, 0) + 1
    return dist


def _bullet_distribution(paras: list[dict]) -> dict[str, int]:
    dist: dict[str, int] = {}
    for p in paras:
        if not p.get("text"):
            continue
        key = str(p.get("bullet", "unknown"))
        dist[key] = dist.get(key, 0) + 1
    return dist


def extract_deck(
    ppt_path: str | Path,
    *,
    use_rendered_bounds: bool = True,
) -> dict[str, Any]:
    """Extract all delivery-status slides from the entire deck."""
    ppt_path = Path(ppt_path)
    prs = Presentation(str(ppt_path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        data = extract_slide(slide, i)
        if data:
            if use_rendered_bounds:
                data["_rendered_bounds_deferred"] = True
            slides.append(data)
    deck = {
        "file": str(ppt_path.name),
        "slide_count": len(slides),
        "slides": slides,
    }
    if use_rendered_bounds and slides:
        from app.services.ppt_rendered_bounds import enrich_deck_rendered_hl_metrics

        enrich_deck_rendered_hl_metrics(deck, ppt_path)
        for slide_data in slides:
            slide_data.pop("_rendered_bounds_deferred", None)
    return deck
