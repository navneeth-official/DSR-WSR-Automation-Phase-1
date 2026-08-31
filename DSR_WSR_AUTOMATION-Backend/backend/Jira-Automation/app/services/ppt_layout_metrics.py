"""Shared HL/KA layout metrics (wrap-aware line estimates) for extract + repair."""

from __future__ import annotations

from pptx.oxml.ns import qn

EMU_PER_INCH = 914400
MIN_TEXT_KA_CLEARANCE_IN = 0.15
HL_KA_TAB_GAP_LINES = 2  # white space between HL table bottom and KA header (user reference)
CANONICAL_LINE_HEIGHT_EMU = 142875  # legacy supplier reference (~0.156 in)
# Fixed story line height: 2 lines = 0.472 in between HL table bottom and KA tab top.
FIXED_STORY_LINE_HEIGHT_EMU = int(0.472 * EMU_PER_INCH / HL_KA_TAB_GAP_LINES)
HL_KA_TARGET_BORDER_GAP_IN = round(
    HL_KA_TAB_GAP_LINES * FIXED_STORY_LINE_HEIGHT_EMU / EMU_PER_INCH, 4
)
HL_VISUAL_LINE_BUFFER = 1.12
CANONICAL_PARA_SLOTS = 20
UTILIZATION_THRESHOLD = 0.85
SPARSE_CONTD_MAX_FILLED = 3
MIN_KA_BLOCK_IN = 1.0  # min vertical room to place KA below sparse HL on main
DEFAULT_EMPTY_KA_HEIGHT_IN = 0.73  # typical empty KA table height on main slide
SPARSE_HL_MAX_WASTE_IN = 2.4  # template allows large intentional slack (Location ~1.14 in)
SPARSE_KA_MAX_WASTE_IN = 1.2
# Border gap between HL/KA tables is NOT enforced — template uses overlap; use text clearance.
HL_KA_MIN_BORDER_GAP_IN = -0.15
HL_KA_MAX_BORDER_GAP_IN = 3.0
FOOTER_MAX_BOTTOM_IN = 6.29
FOOTER_MAX_HL_DENSE_BOTTOM_IN = 6.55  # dense HL-only main may fill toward dotted line
FOOTER_MAX_BOTTOM_EMU = int(FOOTER_MAX_BOTTOM_IN * EMU_PER_INCH)
DEFAULT_KA_TABLE_HEIGHT_EMU = int(DEFAULT_EMPTY_KA_HEIGHT_IN * EMU_PER_INCH)
CANONICAL_KA_ITEM_SLOTS = 6  # reference KA capacity for utilization

_CHARS_PER_VISUAL_LINE = {0: 92, 1: 72, 7: 80}
_DEFAULT_CHARS_PER_LINE = 85
_VISUAL_LINE_HEIGHT_BUFFER = 1.12


def hl_ka_tab_gap_emu(canonical_line_height_emu: int | None = None) -> int:
    """EMU gap between Highlights table bottom and Key Activities table top."""
    if canonical_line_height_emu is not None:
        return HL_KA_TAB_GAP_LINES * canonical_line_height_emu
    return HL_KA_TAB_GAP_LINES * FIXED_STORY_LINE_HEIGHT_EMU


def apply_fixed_story_line_metrics(profile: dict) -> dict:
    """Use one line height everywhere for HL sizing and HL→KA tab gap."""
    merged = dict(profile)
    merged["canonical_line_height_emu"] = FIXED_STORY_LINE_HEIGHT_EMU
    merged["canonical_per_line_emu"] = FIXED_STORY_LINE_HEIGHT_EMU
    return merged


def _profile_hl_top(profile: dict) -> int:
    hl_top = profile.get("ref_hl_top")
    if hl_top is not None:
        return int(hl_top)
    ref_hl = profile.get("ref_hl")
    if ref_hl is not None:
        return int(ref_hl.top)
    return 0


def estimated_hl_table_height_emu(profile: dict, line_count: int) -> int:
    """Minimum HL table height for a wrap-aware line count using fixed line height."""
    r0, r1 = profile["r0"], profile["r1"]
    content_h = int(FIXED_STORY_LINE_HEIGHT_EMU * max(line_count, 1) * HL_VISUAL_LINE_BUFFER)
    pad = int(0.1 * EMU_PER_INCH)
    return int(r0 + r1 + content_h + pad)


def max_hl_height_with_ka_emu(profile: dict, ka_height_emu: int) -> int:
    """Vertical room for HL table when KA shares the slide (footer-safe)."""
    gap = hl_ka_tab_gap_emu()
    return int(FOOTER_MAX_BOTTOM_EMU - ka_height_emu - gap - _profile_hl_top(profile))


def hl_ka_fits_on_main_slide(profile: dict, line_count: int, ka_height_emu: int | None = None) -> bool:
    """True when estimated HL height + KA + tab gap fits above the footer."""
    ka_h = ka_height_emu or profile.get("ref_ka_height")
    if ka_h is None:
        ref_ka = profile.get("ref_ka")
        if ref_ka is not None and hasattr(ref_ka, "height"):
            ka_h = int(ref_ka.height)
    if ka_h is None:
        ka_h = DEFAULT_KA_TABLE_HEIGHT_EMU
    return estimated_hl_table_height_emu(profile, line_count) <= max_hl_height_with_ka_emu(profile, int(ka_h))


def paragraph_text(p_elem) -> str:
    return "".join(n.text or "" for n in p_elem.iter() if n.tag.endswith("}t"))


def paragraph_bullet_level(p_elem) -> int:
    p_pr = p_elem.find(qn("a:pPr"))
    if p_pr is None or p_pr.get("lvl") is None:
        return 0
    return int(p_pr.get("lvl"))


def estimate_para_visual_lines(p_elem) -> int:
    text = paragraph_text(p_elem).strip()
    if not text:
        return 0
    lvl = paragraph_bullet_level(p_elem)
    chars = _CHARS_PER_VISUAL_LINE.get(lvl, _DEFAULT_CHARS_PER_LINE)
    return max(1, (len(text) + chars - 1) // chars)


def count_hl_paragraphs(hl_shape) -> int:
    tx_body = hl_shape.table.cell(2, 0).text_frame._txBody
    return sum(
        1 for p in tx_body.findall(qn("a:p")) if paragraph_text(p).strip()
    )


def count_visual_lines_in_hl(hl_shape) -> int:
    tx_body = hl_shape.table.cell(2, 0).text_frame._txBody
    total = 0
    for p in tx_body.findall(qn("a:p")):
        if paragraph_text(p).strip():
            total += estimate_para_visual_lines(p)
    return total


def cell_margins(cell) -> tuple[int, int]:
    body_pr = cell._tc.find(qn("a:txBody")).find(qn("a:bodyPr"))
    return (
        int(body_pr.get("marT") or 45720),
        int(body_pr.get("marB") or 45720),
    )


def inner_content_bottom_emu(hl_shape) -> int:
    r0, r1, r2 = (hl_shape.table.rows[i].height for i in range(3))
    _, mar_b = cell_margins(hl_shape.table.cell(2, 0))
    return hl_shape.top + r0 + r1 + r2 - mar_b


def estimated_text_bottom_emu(
    hl_shape,
    *,
    ref_para_count: int = 15,
    ref_r2: int | None = None,
    per_line_emu: int | None = None,
) -> int:
    """Wrap-aware text end position inside Highlights content cell."""
    r0, r1 = hl_shape.table.rows[0].height, hl_shape.table.rows[1].height
    mar_t, _ = cell_margins(hl_shape.table.cell(2, 0))
    para_count = count_hl_paragraphs(hl_shape)
    visual = count_visual_lines_in_hl(hl_shape)
    line_count = max(para_count, visual)
    if per_line_emu is not None:
        per_line = per_line_emu
    elif ref_r2 is not None:
        ref_para_count = max(ref_para_count, 1)
        per_line = ref_r2 / max(ref_para_count - 2, 1)
    else:
        per_line = CANONICAL_LINE_HEIGHT_EMU
    content_h = int(per_line * line_count * _VISUAL_LINE_HEIGHT_BUFFER)
    return hl_shape.top + r0 + r1 + mar_t + content_h


def rendered_text_bottom_emu(
    hl_shape,
    *,
    ref_para_count: int = 15,
    ref_r2: int | None = None,
    per_line_emu: int | None = None,
) -> int:
    """Wrap-aware text end for G10X text-to-KA clearance (never underestimate)."""
    est = estimated_text_bottom_emu(
        hl_shape,
        ref_para_count=ref_para_count,
        ref_r2=ref_r2,
        per_line_emu=per_line_emu or CANONICAL_LINE_HEIGHT_EMU,
    )
    inner = inner_content_bottom_emu(hl_shape)
    return max(est, inner)


def text_ka_clearance_in(hl_shape, ka_shape, *, ref_para_count: int = 15, ref_r2: int | None = None) -> float | None:
    if hl_shape is None or ka_shape is None:
        return None
    text_bottom = rendered_text_bottom_emu(
        hl_shape, ref_para_count=ref_para_count, ref_r2=ref_r2
    )
    return round((ka_shape.top - text_bottom) / EMU_PER_INCH, 4)


def hl_waste_below_text_in_estimated(
    hl_shape,
    *,
    ref_para_count: int = 15,
    ref_r2: int | None = None,
) -> float:
    """Legacy estimate: empty HL area below text using canonical line-height model."""
    if ref_r2 is None:
        ref_r2 = hl_shape.table.rows[2].height
    hl_bottom = hl_shape.top + hl_shape.height
    text_bottom = rendered_text_bottom_emu(
        hl_shape, ref_para_count=ref_para_count, ref_r2=ref_r2
    )
    return round(max((hl_bottom - text_bottom) / EMU_PER_INCH, 0), 4)


def hl_waste_below_text_in(
    hl_shape,
    *,
    ref_para_count: int = 15,
    ref_r2: int | None = None,
    measured_text_bottom_in: float | None = None,
    measured_hl_bottom_in: float | None = None,
) -> float:
    """
    Empty HL area below visible text.

    Prefer ``measured_*`` from COM/image bounds; falls back to estimation.
    """
    if measured_text_bottom_in is not None:
        if measured_hl_bottom_in is not None:
            hl_bottom_in = measured_hl_bottom_in
        else:
            hl_bottom_in = (hl_shape.top + hl_shape.height) / EMU_PER_INCH
        return round(max(hl_bottom_in - measured_text_bottom_in, 0.0), 4)
    return hl_waste_below_text_in_estimated(
        hl_shape, ref_para_count=ref_para_count, ref_r2=ref_r2
    )


def ka_rendered_text_bottom_emu(ka_shape) -> int:
    """Bottom of visible KA content (header when empty, last bullet otherwise)."""
    r0 = ka_shape.table.rows[0].height
    items = count_ka_items(ka_shape)
    if items == 0:
        return ka_shape.top + r0
    mar_t, _ = cell_margins(ka_shape.table.cell(1, 0))
    per_line = 50292
    content_h = int(items * per_line * _VISUAL_LINE_HEIGHT_BUFFER)
    return ka_shape.top + r0 + mar_t + content_h


def ka_rendered_text_bottom_in(ka_shape) -> float:
    return round(ka_rendered_text_bottom_emu(ka_shape) / EMU_PER_INCH, 4)


def count_ka_items(ka_shape) -> int:
    content_cell = ka_shape.table.cell(1, 0)
    return sum(1 for p in content_cell.text_frame.paragraphs if p.text.strip())


def ka_waste_below_text_in(ka_shape) -> float:
    """Empty KA area below visible items (content-relative, not row-height coupled)."""
    items = count_ka_items(ka_shape)
    if items == 0:
        return 0.0
    ka_bottom = ka_shape.top + ka_shape.height
    r0 = ka_shape.table.rows[0].height
    mar_t, _ = cell_margins(ka_shape.table.cell(1, 0))
    per_line = 50292
    text_bottom = ka_shape.top + r0 + mar_t + int(items * per_line * _VISUAL_LINE_HEIGHT_BUFFER)
    return round(max((ka_bottom - text_bottom) / EMU_PER_INCH, 0), 4)


def utilization_ratio(filled_count: int, slot_count: int = CANONICAL_PARA_SLOTS) -> float:
    if slot_count <= 0:
        return 0.0
    return round(filled_count / slot_count, 4)


def effective_hl_utilization(
    filled_paras: int,
    visual_lines: int,
    slot_count: int = CANONICAL_PARA_SLOTS,
) -> float:
    """Density using max(paragraph count, wrap-aware visual lines) vs canonical slots."""
    return utilization_ratio(max(filled_paras, visual_lines), slot_count)


def hl_is_dense(hl: dict) -> bool:
    """True when HL is dense by paragraph or visual-line fill (>= 85%)."""
    util = hl.get("effective_utilization_ratio")
    if util is None:
        util = hl.get("utilization_ratio")
    return util is not None and util >= UTILIZATION_THRESHOLD
