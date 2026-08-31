"""Parse and export WSR PowerPoint decks for the in-browser editor."""

from __future__ import annotations

import base64
import io
import json
import uuid
from datetime import date
from pathlib import Path
from typing import Any, Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from app.paths import OUTPUT_DIR, wsr_output_paths

CANVAS_WIDTH = 1280
CANVAS_HEIGHT = 720


def _safe_rgb_hex(color_format) -> str | None:
    """Return #RRGGBB from a python-pptx ColorFormat, or None when color is unset."""
    if color_format is None:
        return None
    try:
        rgb = color_format.rgb
        if rgb is not None:
            return f"#{rgb}"
    except AttributeError:
        pass
    return None


def _emu_to_px(value: int, slide_extent_emu: int, canvas_extent_px: int) -> float:
    if not slide_extent_emu:
        return float(value)
    return round(value * canvas_extent_px / slide_extent_emu, 2)


def _px_to_emu(value: float, slide_extent_emu: int, canvas_extent_px: int) -> int:
    if not canvas_extent_px:
        return int(value)
    return int(value * slide_extent_emu / canvas_extent_px)


def _text_frame_font(text_frame) -> dict[str, Any]:
    font_size = 11
    font_family = "Calibri"
    color = "#000000"
    bold = False
    italic = False
    align = "left"
    if text_frame.paragraphs:
        para = text_frame.paragraphs[0]
        if para.alignment is not None:
            align_map = {1: "left", 2: "center", 3: "right"}
            align = align_map.get(int(para.alignment), "left")
        if para.runs:
            run = para.runs[0]
            if run.font.size:
                font_size = int(run.font.size.pt)
            if run.font.name:
                font_family = run.font.name
            if run.font.bold:
                bold = True
            if run.font.italic:
                italic = True
            hex_color = _safe_rgb_hex(run.font.color)
            if hex_color:
                color = hex_color
    return {
        "fontSize": font_size,
        "fontFamily": font_family,
        "color": color,
        "bold": bold,
        "italic": italic,
        "align": align,
    }


def _cell_fill(cell) -> str:
    try:
        fill = cell.fill
        if fill.type is not None:
            hex_fill = _safe_rgb_hex(fill.fore_color)
            if hex_fill:
                return hex_fill
    except Exception:
        pass
    return "#ffffff"


def _shape_fill_hex(shape) -> str | None:
    try:
        if shape.fill.type is not None:
            return _safe_rgb_hex(shape.fill.fore_color)
    except Exception:
        pass
    return None


def _make_text_element(
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    rotation: float,
    text: str,
    style: dict[str, Any],
    source_shape_id: int,
    table_row: int | None = None,
    table_col: int | None = None,
    position_locked: bool = False,
    cell_fill: str | None = None,
    cell_border: str | None = None,
) -> dict[str, Any]:
    element: dict[str, Any] = {
        "id": str(uuid.uuid4()),
        "type": "text",
        "x": x,
        "y": y,
        "width": max(8.0, width),
        "height": max(8.0, height),
        "rotation": rotation,
        "text": text,
        "style": style,
        "sourceShapeId": source_shape_id,
        "locked": False,
        "positionLocked": position_locked,
    }
    if cell_fill:
        element["cellFill"] = cell_fill
    if cell_border:
        element["cellBorder"] = cell_border
    if table_row is not None:
        element["tableRow"] = table_row
    if table_col is not None:
        element["tableCol"] = table_col
    return element


def _parse_table_cells(
    shape,
    abs_left: int,
    abs_top: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> list[dict[str, Any]]:
    table = shape.table
    elements: list[dict[str, Any]] = []
    row_top = abs_top
    rotation = float(getattr(shape, "rotation", 0) or 0)

    for row_idx, row in enumerate(table.rows):
        col_left = abs_left
        row_height = row.height
        for col_idx in range(len(table.columns)):
            col_width = table.columns[col_idx].width
            cell = table.cell(row_idx, col_idx)
            x = _emu_to_px(col_left, slide_width_emu, CANVAS_WIDTH)
            y = _emu_to_px(row_top, slide_height_emu, CANVAS_HEIGHT)
            w = _emu_to_px(col_width, slide_width_emu, CANVAS_WIDTH)
            h = _emu_to_px(row_height, slide_height_emu, CANVAS_HEIGHT)
            cell_fill = _cell_fill(cell)

            elements.append(
                {
                    "id": str(uuid.uuid4()),
                    "type": "shape",
                    "shapeKind": "rect",
                    "x": x,
                    "y": y,
                    "width": max(8.0, w),
                    "height": max(8.0, h),
                    "rotation": rotation,
                    "fill": cell_fill,
                    "stroke": "#cbd5e1",
                    "strokeWidth": 1,
                    "sourceShapeId": shape.shape_id,
                    "tableRow": row_idx,
                    "tableCol": col_idx,
                    "isTableCell": True,
                    "locked": True,
                }
            )

            text = (cell.text_frame.text or "").strip()
            if text:
                elements.append(
                    _make_text_element(
                        x=x,
                        y=y,
                        width=w,
                        height=h,
                        rotation=rotation,
                        text=text,
                        style=_text_frame_font(cell.text_frame),
                        source_shape_id=shape.shape_id,
                        table_row=row_idx,
                        table_col=col_idx,
                        position_locked=True,
                        cell_fill=cell_fill,
                        cell_border="1px solid #cbd5e1",
                    )
                )
            col_left += col_width
        row_top += row_height

    return elements


def _parse_shape(
    shape,
    abs_left: int,
    abs_top: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> list[dict[str, Any]]:
    rotation = float(getattr(shape, "rotation", 0) or 0)
    left = _emu_to_px(abs_left, slide_width_emu, CANVAS_WIDTH)
    top = _emu_to_px(abs_top, slide_height_emu, CANVAS_HEIGHT)
    width = max(8.0, _emu_to_px(shape.width, slide_width_emu, CANVAS_WIDTH))
    height = max(8.0, _emu_to_px(shape.height, slide_height_emu, CANVAS_HEIGHT))

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        elements: list[dict[str, Any]] = []
        for child in shape.shapes:
            elements.extend(
                _collect_shape_elements(
                    child,
                    abs_left + child.left,
                    abs_top + child.top,
                    slide_width_emu,
                    slide_height_emu,
                )
            )
        return elements

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return _parse_table_cells(shape, abs_left, abs_top, slide_width_emu, slide_height_emu)

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        blob = shape.image.blob
        mime = shape.image.content_type or "image/png"
        encoded = base64.b64encode(blob).decode("ascii")
        return [
            {
                "id": str(uuid.uuid4()),
                "type": "image",
                "x": left,
                "y": top,
                "width": width,
                "height": height,
                "rotation": rotation,
                "src": f"data:{mime};base64,{encoded}",
                "sourceShapeId": shape.shape_id,
                "locked": False,
            }
        ]

    if getattr(shape, "has_text_frame", False):
        text = (shape.text_frame.text or "").strip()
        if text:
            return [
                _make_text_element(
                    x=left,
                    y=top,
                    width=width,
                    height=height,
                    rotation=rotation,
                    text=text,
                    style=_text_frame_font(shape.text_frame),
                    source_shape_id=shape.shape_id,
                    cell_fill=_shape_fill_hex(shape),
                )
            ]

    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        fill = "#e5e7eb"
        hex_fill = _safe_rgb_hex(getattr(shape.fill, "fore_color", None))
        if hex_fill:
            fill = hex_fill
        return [
            {
                "id": str(uuid.uuid4()),
                "type": "shape",
                "shapeKind": "rect",
                "x": left,
                "y": top,
                "width": width,
                "height": height,
                "rotation": rotation,
                "fill": fill,
                "stroke": "#9ca3af",
                "strokeWidth": 1,
                "sourceShapeId": shape.shape_id,
                "locked": False,
            }
        ]

    return []


def _collect_shape_elements(
    shape,
    abs_left: int,
    abs_top: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> list[dict[str, Any]]:
    try:
        return _parse_shape(shape, abs_left, abs_top, slide_width_emu, slide_height_emu)
    except Exception:
        return []


def _iter_slide_shapes(slide) -> Iterator[tuple[Any, int, int]]:
    for shape in slide.shapes:
        yield shape, shape.left, shape.top


def parse_pptx_to_editor_document(ppt_path: Path) -> dict[str, Any]:
    """Convert a .pptx file into the editor JSON model."""
    prs = Presentation(str(ppt_path))
    slide_width_emu = int(prs.slide_width)
    slide_height_emu = int(prs.slide_height)

    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        elements: list[dict[str, Any]] = []
        title = f"Slide {index}"
        for shape, abs_left, abs_top in _iter_slide_shapes(slide):
            elements.extend(
                _collect_shape_elements(shape, abs_left, abs_top, slide_width_emu, slide_height_emu)
            )
            if getattr(shape, "has_text_frame", False) and shape.shape_id == 2:
                title = (shape.text_frame.text or "").strip() or title

        slides.append(
            {
                "id": str(uuid.uuid4()),
                "index": index,
                "title": title,
                "background": "#ffffff",
                "backgroundImage": None,
                "elements": elements,
            }
        )

    return {
        "id": str(uuid.uuid4()),
        "filename": ppt_path.name,
        "canvasWidth": CANVAS_WIDTH,
        "canvasHeight": CANVAS_HEIGHT,
        "slideWidthEmu": slide_width_emu,
        "slideHeightEmu": slide_height_emu,
        "sourcePptPath": str(ppt_path),
        "slides": slides,
    }


def load_wsr_editor_deck(
    start_date: date,
    end_date: date,
    *,
    preview_slides: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    """Load editor document for a WSR week, attaching slide preview images."""
    paths = wsr_output_paths(start_date, end_date)
    if not paths.ppt_path.is_file():
        raise FileNotFoundError(f"No deck found for {start_date} to {end_date}")

    document = parse_pptx_to_editor_document(paths.ppt_path)
    document["sourcePptPath"] = str(paths.ppt_path)

    preview_by_index = {int(item["slide_index"]): item for item in (preview_slides or [])}
    for slide in document["slides"]:
        preview = preview_by_index.get(slide["index"])
        if preview:
            slide["title"] = preview.get("title") or slide["title"]
            slide["backgroundImage"] = preview.get("image_url")
        for element in slide.get("elements", []):
            if element.get("type") == "text":
                element["originalText"] = element.get("text", "")
                element["isDirty"] = False
            element["originalX"] = element.get("x", 0)
            element["originalY"] = element.get("y", 0)
            element["originalWidth"] = element.get("width", 0)
            element["originalHeight"] = element.get("height", 0)

    editor_json_path = OUTPUT_DIR / f"WSR_{start_date}_{end_date}_editor.json"
    editor_json_path.write_text(json.dumps(document, indent=2), encoding="utf-8")
    return document


def save_wsr_editor_deck(start_date: date, end_date: date, document: dict[str, Any]) -> Path:
    path = OUTPUT_DIR / f"WSR_{start_date}_{end_date}_editor.json"
    path.write_text(json.dumps(document, indent=2), encoding="utf-8")
    return path


def _apply_geometry(shape, element: dict[str, Any], slide_width_emu: int, slide_height_emu: int) -> None:
    shape.left = _px_to_emu(float(element.get("x", 0)), slide_width_emu, CANVAS_WIDTH)
    shape.top = _px_to_emu(float(element.get("y", 0)), slide_height_emu, CANVAS_HEIGHT)
    shape.width = _px_to_emu(float(element.get("width", 100)), slide_width_emu, CANVAS_WIDTH)
    shape.height = _px_to_emu(float(element.get("height", 40)), slide_height_emu, CANVAS_HEIGHT)


def _apply_updates_to_slide(
    slide,
    elements: list[dict[str, Any]],
    slide_width_emu: int,
    slide_height_emu: int,
) -> None:
    by_shape: dict[int, list[dict[str, Any]]] = {}
    for element in elements:
        shape_id = element.get("sourceShapeId")
        if shape_id is None:
            continue
        by_shape.setdefault(int(shape_id), []).append(element)

    for shape in slide.shapes:
        updates = by_shape.get(shape.shape_id, [])
        if not updates:
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for element in updates:
                if element.get("type") != "text":
                    continue
                row = element.get("tableRow")
                col = element.get("tableCol")
                if row is None or col is None:
                    continue
                try:
                    shape.table.cell(int(row), int(col)).text_frame.text = element.get("text", "")
                except Exception:
                    continue
            continue

        text_el = next(
            (
                e
                for e in updates
                if e.get("type") == "text" and e.get("tableRow") is None
            ),
            None,
        )
        if text_el and getattr(shape, "has_text_frame", False):
            shape.text_frame.text = text_el.get("text", "")
            if not text_el.get("positionLocked"):
                _apply_geometry(shape, text_el, slide_width_emu, slide_height_emu)

        image_el = next((e for e in updates if e.get("type") == "image"), None)
        if image_el and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _apply_geometry(shape, image_el, slide_width_emu, slide_height_emu)

        shape_el = next(
            (
                e
                for e in updates
                if e.get("type") == "shape" and not e.get("isTableCell")
            ),
            None,
        )
        if shape_el and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not text_el:
            _apply_geometry(shape, shape_el, slide_width_emu, slide_height_emu)


def export_editor_document_to_pptx(document: dict[str, Any], output_path: Path) -> Path:
    """Patch text in the original .pptx so layout and formatting are preserved."""
    source_path = document.get("sourcePptPath")
    if source_path and Path(source_path).is_file():
        prs = Presentation(str(source_path))
    elif output_path.is_file():
        prs = Presentation(str(output_path))
    else:
        slide_width_emu = int(document.get("slideWidthEmu") or 12192000)
        slide_height_emu = int(document.get("slideHeightEmu") or 6858000)
        prs = Presentation()
        prs.slide_width = slide_width_emu
        prs.slide_height = slide_height_emu
        blank_layout = prs.slide_layouts[6]
        for _slide_data in document.get("slides", []):
            prs.slides.add_slide(blank_layout)

    slides_by_index = {int(s.get("index", i + 1)): s for i, s in enumerate(document.get("slides", []))}
    slide_width_emu = int(document.get("slideWidthEmu") or int(prs.slide_width))
    slide_height_emu = int(document.get("slideHeightEmu") or int(prs.slide_height))
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = slides_by_index.get(slide_idx)
        if slide_data:
            _apply_updates_to_slide(
                slide,
                slide_data.get("elements", []),
                slide_width_emu,
                slide_height_emu,
            )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
