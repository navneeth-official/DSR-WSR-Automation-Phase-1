"""Export delivery-status slides from a .pptx to PNG for vision review."""

from __future__ import annotations

import re
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.rendering.com_backend import ComSlideRendererBackend


def _slide_title(slide) -> str:
    for sh in slide.shapes:
        if sh.has_text_frame and getattr(sh, "shape_id", None) == 2:
            return sh.text_frame.text.strip()
    return ""


def list_all_slide_indices(ppt_path: str | Path) -> list[dict[str, Any]]:
    """Return every slide in the deck (1-based index + best-effort title)."""
    prs = Presentation(str(ppt_path))
    out: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        if not title:
            title = _first_text_on_slide(slide) or f"Slide {i}"
        out.append({"slide_index": i, "title": title})
    return out


def _first_text_on_slide(slide) -> str:
    for sh in slide.shapes:
        if sh.has_text_frame:
            text = sh.text_frame.text.strip()
            if text:
                first_line = text.splitlines()[0].strip()
                if first_line:
                    return first_line[:120]
    return ""


def export_all_slides_to_png(
    ppt_path: str | Path,
    output_dir: str | Path,
    *,
    width_px: int = 1280,
) -> list[dict[str, Any]]:
    """Export every slide in the deck to PNG (Windows COM)."""
    ppt_path = Path(ppt_path).resolve()
    output_dir = Path(output_dir)
    indices = [s["slide_index"] for s in list_all_slide_indices(ppt_path)]
    return export_slides_to_png(
        ppt_path,
        output_dir,
        slide_indices=indices,
        width_px=width_px,
    )


def list_delivery_slide_indices(ppt_path: str | Path) -> list[dict[str, Any]]:
    """Return physical slide indices (1-based) for Delivery status slides."""
    prs = Presentation(str(ppt_path))
    out: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        if "Delivery status" in title or re.search(r"\(Contd", title, re.I):
            out.append({"slide_index": i, "title": title})
    return out


def _titles_by_index(ppt_path: Path, indices: list[int]) -> dict[int, str]:
    prs = Presentation(str(ppt_path))
    titles: dict[int, str] = {}
    for idx in indices:
        if 1 <= idx <= len(prs.slides):
            titles[idx] = _slide_title(prs.slides[idx - 1])
    return titles


def export_slides_to_png(
    ppt_path: str | Path,
    output_dir: str | Path,
    *,
    slide_indices: list[int] | None = None,
    width_px: int = 1920,
) -> list[dict[str, Any]]:
    """
    Export selected slides to PNG using PowerPoint COM (Windows only).

    Returns list of {slide_index, title, image_path}.
    """
    ppt_path = Path(ppt_path).resolve()
    output_dir = Path(output_dir)

    indices = slide_indices
    if indices is None:
        indices = [s["slide_index"] for s in list_delivery_slide_indices(ppt_path)]

    backend = ComSlideRendererBackend()
    image_paths = backend.render_slides(
        ppt_path,
        output_dir,
        slide_indices=indices,
        width_px=width_px,
    )

    titles = _titles_by_index(ppt_path, indices)
    exported: list[dict[str, Any]] = []
    for path in image_paths:
        stem = path.stem
        idx = int(stem.split("_", 1)[1]) if "_" in stem else 0
        exported.append(
            {
                "slide_index": idx,
                "title": titles.get(idx, ""),
                "image_path": str(path),
            }
        )
    return exported


def export_deck_pngs(
    ppt_path: str | Path,
    output_dir: str | Path | None = None,
) -> tuple[Path, list[dict[str, Any]]]:
    """Export all delivery-status slides; uses a temp dir when output_dir omitted."""
    if output_dir is None:
        output_dir = Path(tempfile.mkdtemp(prefix="ppt_vision_"))
    out_dir = Path(output_dir)
    return out_dir, export_slides_to_png(ppt_path, out_dir)
