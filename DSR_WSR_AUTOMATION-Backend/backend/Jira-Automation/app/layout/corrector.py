"""Deterministic PowerPoint layout correction from vision measurements."""



from __future__ import annotations



from pathlib import Path

from typing import Sequence



from pptx import Presentation



from app.layout.config import LayoutCorrectorConfig

from app.layout.exceptions import LayoutFailureError

from app.layout.geometry_debug import get_geometry_logger

from app.layout.measurement_mapper import VisionMeasurementMapper

from app.layout.shape_ops import (

    PixelScale,

    expand_table_height,

    maintain_alignment,

    maintain_gap,

    move_shape,

    restore_template_position,

    shrink_table_height,

)

from app.layout.template_geometry import TemplateGeometry, TemplateGeometryProvider

from app.layout.types import (

    CorrectionAction,

    CorrectionActionType,

    LayoutCorrectionResult,

    SlideCorrectionInput,

    SlideCorrectionResult,

)

from app.vision.types import RecommendedAction, SlideEvaluationResult, SlideMeasurements





class LayoutCorrector:

    """

    Apply deterministic layout fixes to a presentation using parsed vision measurements.



    This module does not call AI and does not re-estimate text layout. Corrections

    are driven by vision pixel measurements converted to slide EMU coordinates and

    template geometry from G10X reference slides.

    """



    def __init__(

        self,

        *,

        config: LayoutCorrectorConfig | None = None,

        geometry_provider: TemplateGeometryProvider | None = None,

    ) -> None:

        self._config = config or LayoutCorrectorConfig()

        self._geometry = geometry_provider or TemplateGeometryProvider()



    def correct(

        self,

        presentation: Presentation | Path | str,

        slides: Sequence[SlideCorrectionInput],

        *,

        save: bool = True,

    ) -> LayoutCorrectionResult:

        """

        Modify the presentation in response to vision evaluation results.



        Args:

            presentation: Open ``Presentation`` or path to a ``.pptx`` file.

            slides: Per-slide vision results (optionally with rendered image paths).

            save: When ``presentation`` is a path, persist changes back to disk.

        """

        ppt_path: Path | None = None

        if isinstance(presentation, (str, Path)):

            ppt_path = Path(presentation).resolve()

            prs = Presentation(str(ppt_path))

        else:

            prs = presentation



        uds, _g10x_prs = self._geometry.uds_helpers()

        scale = PixelScale.from_paths(

            slide_width_emu=prs.slide_width,

            slide_height_emu=prs.slide_height,

            image_path=None,

            default_width_px=self._config.default_image_width_px,

            default_height_px=self._config.default_image_height_px,

        )



        slide_results: list[SlideCorrectionResult] = []

        deck_modified = False

        failures: list[str] = []



        for item in slides:

            evaluation = item.evaluation

            slide_number = evaluation.slide_number

            if slide_number is None or slide_number < 1 or slide_number > len(prs.slides):

                slide_results.append(

                    SlideCorrectionResult(

                        slide_number=slide_number,

                        failures=[f"Invalid slide number: {slide_number}"],

                    )

                )

                continue



            item_scale = PixelScale.from_paths(

                slide_width_emu=prs.slide_width,

                slide_height_emu=prs.slide_height,

                image_path=item.image_path,

                default_width_px=self._config.default_image_width_px,

                default_height_px=self._config.default_image_height_px,

            )



            slide = prs.slides[slide_number - 1]

            geometry = self._geometry.for_slide(slide, slide_number=slide_number)

            if geometry is None:

                slide_results.append(

                    SlideCorrectionResult(

                        slide_number=slide_number,

                        failures=["Not a delivery-status slide"],

                    )

                )

                continue



            result = self._correct_slide(

                slide,

                evaluation,

                geometry=geometry,

                scale=item_scale,

                uds=uds,

            )

            slide_results.append(result)

            deck_modified = deck_modified or result.modified

            failures.extend(

                f"Slide {slide_number}: {msg}" for msg in result.failures

            )



        if save and ppt_path is not None and deck_modified:

            prs.save(str(ppt_path))



        return LayoutCorrectionResult(

            ppt_path=ppt_path or Path("presentation.pptx"),

            modified=deck_modified,

            slides=slide_results,

            failures=failures,

        )



    def _correct_slide(

        self,

        slide,

        evaluation: SlideEvaluationResult,

        *,

        geometry: TemplateGeometry,

        scale: PixelScale,

        uds,

    ) -> SlideCorrectionResult:

        actions: list[CorrectionAction] = []

        failures: list[str] = []

        modified = False

        mapper = VisionMeasurementMapper(scale, self._config)



        hl = self._safe_highlights(slide, uds)

        ka = self._safe_key_activities(slide, uds)



        measurements = evaluation.measurements

        if self._apply_measurement_rules(

            hl,

            ka,

            geometry=geometry,

            mapper=mapper,

            uds=uds,

            measurements=measurements,

            actions=actions,

            failures=failures,

        ):

            modified = True



        for issue in evaluation.issues:

            try:

                if self._apply_issue_action(

                    issue.recommended_action,

                    hl,

                    ka,

                    geometry=geometry,

                    mapper=mapper,

                    uds=uds,

                    slide_measurements=measurements,

                    issue_measurement=issue.measurement,

                    actions=actions,

                ):

                    modified = True

            except LayoutFailureError as exc:

                failures.append(str(exc))



        return SlideCorrectionResult(

            slide_number=evaluation.slide_number,

            modified=modified,

            actions=actions,

            failures=failures,

        )



    def _apply_measurement_rules(

        self,

        hl,

        ka,

        *,

        geometry: TemplateGeometry,

        mapper: VisionMeasurementMapper,

        uds,

        measurements: SlideMeasurements,

        actions: list[CorrectionAction],

        failures: list[str],

    ) -> bool:

        changed = False

        if hl and ka:

            gap_deficit = mapper.gap_deficit_emu(measurements)

            if gap_deficit is not None and gap_deficit > 0:

                gap_px = mapper.gap_px(measurements)

                upper_bottom = mapper.upper_content_bottom_emu(measurements)

                if upper_bottom is None:

                    upper_bottom = hl.top + hl.height

                new_top = maintain_gap(

                    upper_bottom,

                    ka,

                    min_gap_emu=mapper.min_gap_emu,

                )

                actions.append(

                    CorrectionAction(

                        action_type=CorrectionActionType.MAINTAIN_GAP,

                        target="key_activities",

                        detail=(

                            f"Moved Key Activities down to maintain "

                            f"{self._config.min_text_ka_clearance_in:.2f} in gap "

                            f"(vision gap {gap_px}px → +{gap_deficit} EMU)"

                        ),

                        delta_emu=new_top - ka.top,

                    )

                )

                self._try_maintain_ka_alignment(

                    ka,

                    geometry,

                    failures=failures,

                    context="maintain_gap after vision gap correction",

                )

                changed = True



        if hl:

            shrink_emu = mapper.hl_shrink_excess_emu(measurements)

            if shrink_emu is not None and shrink_emu > 0:

                waste_px = mapper.unused_space_px(measurements)

                shrink_table_height(

                    hl,

                    delta_emu=shrink_emu,

                    profile=geometry.profile,

                    uds_module=uds,

                )

                actions.append(

                    CorrectionAction(

                        action_type=CorrectionActionType.RESIZE_SHAPE,

                        target="highlights",

                        detail=(

                            f"Reduced Highlights height; unused space "

                            f"{waste_px}px exceeds tolerance "

                            f"(shrink {shrink_emu} EMU)"

                        ),

                        delta_emu=-shrink_emu,

                    )

                )

                changed = True

                if ka:

                    upper_bottom = hl.top + hl.height

                    maintain_gap(

                        upper_bottom,

                        ka,

                        min_gap_emu=mapper.min_gap_emu,

                    )

        return changed



    def _apply_issue_action(

        self,

        action: RecommendedAction,

        hl,

        ka,

        *,

        geometry: TemplateGeometry,

        mapper: VisionMeasurementMapper,

        uds,

        slide_measurements: SlideMeasurements,

        issue_measurement: dict,

        actions: list[CorrectionAction],

    ) -> bool:

        if action == RecommendedAction.NO_ACTION:

            return False



        if action == RecommendedAction.RESTORE_TEMPLATE_POSITION:

            ref_hl = geometry.profile["ref_hl"]

            if hl:

                restore_template_position(hl, ref_hl)

                actions.append(

                    CorrectionAction(

                        action_type=CorrectionActionType.RESTORE_TEMPLATE,

                        target="highlights",

                        detail="Restored Highlights to template position",

                    )

                )

            ref_ka = geometry.profile.get("ref_ka")

            if ka and ref_ka:

                restore_template_position(ka, ref_ka)

                actions.append(

                    CorrectionAction(

                        action_type=CorrectionActionType.RESTORE_TEMPLATE,

                        target="key_activities",

                        detail="Restored Key Activities to template position",

                    )

                )

            return bool(hl or ka)



        if action == RecommendedAction.MOVE_SECTION_DOWN:

            if not ka:

                return False

            delta_emu = mapper.overlap_emu(slide_measurements, issue_measurement)

            if delta_emu is None or delta_emu <= 0:

                return False

            new_top = max(0, ka.top + delta_emu)

            move_shape(ka, top=new_top)

            self._try_maintain_ka_alignment(

                ka,

                geometry,

                failures=[],

                context="move_section_down",

            )

            overlap_px = mapper.overlap_px(slide_measurements, issue_measurement)

            actions.append(

                CorrectionAction(

                    action_type=CorrectionActionType.MOVE_SHAPE,

                    target="key_activities",

                    detail=(

                        f"Moved Key Activities down by {delta_emu} EMU "

                        f"(overlap {overlap_px}px)"

                    ),

                    delta_emu=delta_emu,

                )

            )

            return True



        if action == RecommendedAction.MOVE_SECTION_UP:

            if not ka:

                return False

            excess_emu = mapper.gap_excess_emu(slide_measurements, issue_measurement)

            if excess_emu is None or excess_emu <= 0:

                return False

            delta_emu = -excess_emu

            new_top = max(0, ka.top + delta_emu)

            move_shape(ka, top=new_top)

            self._try_maintain_ka_alignment(

                ka,

                geometry,

                failures=[],

                context="move_section_up",

            )

            gap_px = mapper.gap_px(slide_measurements, issue_measurement)

            actions.append(

                CorrectionAction(

                    action_type=CorrectionActionType.MOVE_SHAPE,

                    target="key_activities",

                    detail=(

                        f"Moved Key Activities up by {excess_emu} EMU "

                        f"(gap {gap_px}px exceeds min clearance)"

                    ),

                    delta_emu=delta_emu,

                )

            )

            return True



        if action in (

            RecommendedAction.REDUCE_UNUSED_SPACE,

            RecommendedAction.DECREASE_TEXTBOX_HEIGHT,

        ):

            if not hl:

                return False

            delta_emu = mapper.hl_shrink_to_remove_emu(

                slide_measurements,

                issue_measurement,

            )

            if delta_emu is None or delta_emu <= 0:

                return False

            shrink_table_height(

                hl,

                delta_emu=delta_emu,

                profile=geometry.profile,

                uds_module=uds,

            )

            waste_px = mapper.unused_space_px(slide_measurements, issue_measurement)

            actions.append(

                CorrectionAction(

                    action_type=CorrectionActionType.RESIZE_SHAPE,

                    target="highlights",

                    detail=(

                        f"Reduced Highlights textbox height by {delta_emu} EMU "

                        f"(unused space {waste_px}px)"

                    ),

                    delta_emu=-delta_emu,

                )

            )

            if ka:

                maintain_gap(

                    hl.top + hl.height,

                    ka,

                    min_gap_emu=mapper.min_gap_emu,

                )

            return True



        if action in (

            RecommendedAction.INCREASE_TEXTBOX_HEIGHT,

            RecommendedAction.EXPAND_PLACEHOLDER,

            RecommendedAction.OVERFLOW_DETECTED,

        ):

            if not hl:

                return False

            delta_emu = mapper.hl_expand_emu(slide_measurements, issue_measurement)

            if delta_emu is None or delta_emu <= 0:

                return False

            expand_table_height(

                hl,

                delta_emu=delta_emu,

                profile=geometry.profile,

                uds_module=uds,

                max_bottom_emu=geometry.footer_max_bottom_emu,

            )

            actions.append(

                CorrectionAction(

                    action_type=CorrectionActionType.RESIZE_SHAPE,

                    target="highlights",

                    detail=f"Expanded Highlights textbox height by {delta_emu} EMU",

                    delta_emu=delta_emu,

                )

            )

            if ka:

                upper_bottom = mapper.upper_content_bottom_emu(slide_measurements)

                if upper_bottom is None:

                    upper_bottom = hl.top + hl.height

                maintain_gap(

                    max(upper_bottom, hl.top + hl.height),

                    ka,

                    min_gap_emu=mapper.min_gap_emu,

                )

            return True



        return False



    def _try_maintain_ka_alignment(

        self,

        ka,

        geometry: TemplateGeometry,

        *,

        failures: list[str],

        context: str,

    ) -> bool:

        """

        Align KA horizontally to template geometry when a reference exists.



        Skips gracefully (with logging) when no reference shape is available.

        """

        ref_ka = geometry.ka_alignment_reference()

        if ka is None:

            return False

        if ref_ka is None:

            msg = (

                f"Skipped KA alignment ({context}) on slide "

                f"{geometry.slide_number} {geometry.slide_title!r}: "

                f"{geometry.ref_ka_none_reason or 'no ref_ka available'}"

            )

            get_geometry_logger().warning(msg)

            failures.append(msg)

            return False

        maintain_alignment(ka, ref_ka)

        get_geometry_logger().debug(

            "Aligned KA on slide %s via ref_ka source=%s context=%s",

            geometry.slide_number,

            geometry.ref_ka_source,

            context,

        )

        return True



    @staticmethod

    def _safe_highlights(slide, uds):

        try:

            return uds.get_highlights_shape(slide)

        except ValueError:

            return None



    @staticmethod

    def _safe_key_activities(slide, uds):

        return uds.get_key_activities_shape(slide)


