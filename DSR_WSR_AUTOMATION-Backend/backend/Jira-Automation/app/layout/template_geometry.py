"""Template geometry resolved from G10X reference slides (no hardcoded coordinates)."""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.slide import Slide

from app.layout.geometry_debug import (
    configure_geometry_logging,
    log_g10x_shape_inventory,
    log_ref_ka_resolution,
    summarize_ref_shape,
)

from app.paths import G10X_TEMPLATE, GEOMETRY_DEBUG_LOG, SCRIPTS_DIR


@dataclass(frozen=True)
class TemplateGeometry:
    """Placeholder positions and spacing from the matching G10X reference slide."""

    service_title: str
    slide_title: str
    slide_number: int | None
    is_contd: bool
    g10x_slide_index: int
    ref_ka_source: str
    ref_ka_none_reason: str
    ref_hl_top: int
    ref_hl_left: int
    ref_hl_width: int
    ref_hl_height: int
    ref_ka_top: int | None
    ref_ka_left: int | None
    ref_ka_width: int | None
    ref_ka_height: int | None
    standard_gap_emu: int | None
    footer_max_bottom_emu: int
    on_slide_ka_in_g10x: bool
    profile: dict[str, Any]

    @property
    def has_key_activities_reference(self) -> bool:
        return self.profile.get("ref_ka") is not None

    @property
    def has_key_activities(self) -> bool:
        return self.ref_ka_top is not None

    def ka_alignment_reference(self):
        """Shape used for horizontal KA alignment (may be canonical fallback)."""
        return self.profile.get("ref_ka")


def _service_title(slide_title: str) -> str:
    base = re.sub(r"^Delivery status\s*[–-]\s*", "", slide_title, flags=re.I)
    return re.sub(r"\s*\(Contd.*\)\s*$", "", base, flags=re.I).strip()


def _is_contd_title(slide_title: str) -> bool:
    return bool(re.search(r"\(Contd", slide_title, re.I))


def _slide_title(slide: Slide) -> str:
    for shape in slide.shapes:
        if shape.has_text_frame and getattr(shape, "shape_id", None) == 2:
            return shape.text_frame.text.strip()
    return ""


def _g10x_slide_index(uds, g10x_prs, service: str) -> int:
    return uds.G10X_LAYOUT_BY_TITLE.get(service, 2)


def _g10x_slide_title(g10x_slide) -> str:
    return _slide_title(g10x_slide)


@lru_cache(maxsize=1)
def _uds_bundle() -> tuple[Any, Presentation]:
    scripts = str(SCRIPTS_DIR)
    if scripts not in sys.path:
        sys.path.insert(0, scripts)
    import update_delivery_status as uds  # noqa: WPS433

    g10x_prs = Presentation(str(G10X_TEMPLATE))
    return uds, g10x_prs


class TemplateGeometryProvider:
    """Resolve template geometry for a delivery-status slide from G10X metadata."""

    def __init__(
        self,
        *,
        template_path: Path | None = None,
        log_path: Path | None = None,
    ) -> None:
        self._template_path = template_path or G10X_TEMPLATE
        if log_path is None:
            log_path = GEOMETRY_DEBUG_LOG
        configure_geometry_logging(log_path=log_path)

    def for_slide(
        self,
        slide: Slide,
        *,
        slide_number: int | None = None,
    ) -> TemplateGeometry | None:
        title = _slide_title(slide)
        if "Delivery status" not in title and "Contd" not in title:
            return None

        service = _service_title(title)
        is_contd = _is_contd_title(title)
        uds, g10x_prs = _uds_bundle()
        g10x_index = _g10x_slide_index(uds, g10x_prs, service)
        g10x_layout = g10x_prs.slides[g10x_index]
        g10x_title = _g10x_slide_title(g10x_layout)

        log_g10x_shape_inventory(
            slide_number=slide_number,
            slide_title=title,
            service_title=service,
            g10x_slide_index=g10x_index,
            g10x_slide_title=g10x_title,
            shapes=list(g10x_layout.shapes),
        )

        base_profile = uds.build_layout_profile(g10x_layout)
        build_has_ka = base_profile.get("ref_ka") is not None
        on_slide_ka = uds.g10x_main_has_on_slide_ka(g10x_layout)

        # Match the layout engine: canonical KA fallback when the service G10X
        # main slide is HL-only (e.g. Supplier) but the deck slide still has KA.
        ka_profile = uds.get_ka_layout_profile(g10x_prs, g10x_layout)
        ref_ka = ka_profile.get("ref_ka")

        if build_has_ka:
            ref_ka_source = "g10x_on_slide"
            reason = (
                f"G10X slide {g10x_index} ({g10x_title!r}) contains a Key Activities "
                "table in cell (0,0) header text."
            )
        elif ref_ka is not None:
            ref_ka_source = "canonical_fallback"
            reason = (
                f"G10X slide {g10x_index} ({g10x_title!r}) is HL-only on the main "
                "reference; using Cost Core canonical KA for alignment geometry."
            )
        else:
            ref_ka_source = "none"
            reason = (
                f"No Key Activities table on G10X slide {g10x_index} and no canonical "
                "KA reference available."
            )

        if is_contd and not build_has_ka:
            reason += (
                " Note: contd deck slide maps to the service main G10X slide, not the "
                "contd+KA G10X slide; canonical fallback is used for horizontal alignment."
            )

        log_ref_ka_resolution(
            slide_number=slide_number,
            slide_title=title,
            service_title=service,
            is_contd=is_contd,
            g10x_slide_index=g10x_index,
            on_slide_ka_in_g10x=on_slide_ka,
            ref_ka_from_build_profile=build_has_ka,
            ref_ka_after_ka_profile=ref_ka is not None,
            ref_ka_source=ref_ka_source,
            reason=reason,
            ref_ka_summary=summarize_ref_shape(ref_ka),
        )

        ref_hl = ka_profile["ref_hl"]

        return TemplateGeometry(
            service_title=service,
            slide_title=title,
            slide_number=slide_number,
            is_contd=is_contd,
            g10x_slide_index=g10x_index,
            ref_ka_source=ref_ka_source,
            ref_ka_none_reason=reason if ref_ka is None else "",
            ref_hl_top=ref_hl.top,
            ref_hl_left=ref_hl.left,
            ref_hl_width=ref_hl.width,
            ref_hl_height=ref_hl.height,
            ref_ka_top=ref_ka.top if ref_ka else None,
            ref_ka_left=ref_ka.left if ref_ka else None,
            ref_ka_width=ref_ka.width if ref_ka else None,
            ref_ka_height=ref_ka.height if ref_ka else None,
            standard_gap_emu=ka_profile.get("standard_gap"),
            footer_max_bottom_emu=uds.MAX_KA_BOTTOM_EMU,
            on_slide_ka_in_g10x=on_slide_ka,
            profile=ka_profile,
        )

    def uds_helpers(self) -> tuple[Any, Presentation]:
        return _uds_bundle()
